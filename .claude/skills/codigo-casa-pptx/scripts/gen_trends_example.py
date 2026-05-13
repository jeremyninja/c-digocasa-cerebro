#!/usr/bin:env python3
"""Trends Identidad — 3 macro tendencias en formato Trend Development."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

ASSETS = "/Users/jeremyrodriguez/Downloads/Código Casa/Gráfica"
OUTPUT = "/Users/jeremyrodriguez/Claude/codigo-casa/Trends-Familia-Identidad.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_M = RGBColor(0x88, 0x88, 0x88)
GRAY_D = RGBColor(0x55, 0x55, 0x55)

TF = "Instrument Serif"
BF = "Poppins"

DUST = os.path.join(ASSETS, "Dust texture", "04.png")
UPPER = os.path.join(ASSETS, "Upper bar.png")
MASK = os.path.join(ASSETS, "Logo Final", "Mascara CC.png")
MASK_RATIO = 1.17

CAP = "FAMILIA E IDENTIDAD"
SRC = "Source: Trends Hunting — Código Casa 2026 | Ninja Thinking"


def make_dust(op=0.15):
    cache = f"/tmp/dust_04_op{int(op*100)}.png"
    if os.path.exists(cache):
        return cache
    d = Image.open(DUST).convert("RGBA")
    b = Image.new("RGBA", d.size, (0, 0, 0, 255))
    a = d.split()[-1].point(lambda p: int(p * op))
    d.putalpha(a)
    out = Image.alpha_composite(b, d)
    out.convert("RGB").save(cache)
    return cache


DUST_BG = make_dust(0.15)


def bg(s):
    s.shapes.add_picture(DUST_BG, 0, 0, SLIDE_W, SLIDE_H)


def upper(s):
    s.shapes.add_picture(UPPER, 0, 0, SLIDE_W, Cm(0.35))


def mask_corner(s):
    h = Cm(1.0)
    w = Cm(1.0 * MASK_RATIO)
    s.shapes.add_picture(MASK, SLIDE_W - w - Cm(0.8), SLIDE_H - h - Cm(0.8), w, h)


def t(s, text, x, y, w, h, font=BF, size=11, color=WHITE,
      bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for old in list(p.runs):
        old._r.getparent().remove(old._r)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def header(s):
    y = Cm(0.85)
    t(s, f"CÓDIGO CASA® - {CAP}", Cm(0.8), y, Cm(15), Cm(0.4), size=7, color=WHITE)
    t(s, "NINJA THINKING", SLIDE_W - Cm(15.8), y, Cm(15), Cm(0.4),
      size=7, color=WHITE, align=PP_ALIGN.RIGHT)


def chrome(s):
    bg(s)
    upper(s)
    header(s)
    mask_corner(s)


def source(s, txt=None):
    t(s, txt or SRC, Cm(1.5), Cm(17.5), Cm(28), Cm(0.4),
      size=8, color=GRAY_L, italic=True)


def build_trend(prs, headline, context, contrast, transformation, stats):
    """
    Build one trend development slide per spec.

    headline: str (UPPERCASE, 2 líneas ideal)
    context, contrast, transformation: str (paragraphs)
    stats: list of (number, description) — 3 items
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    # =============== COLUMNA 1 — NARRATIVA (x=1.5 → 10.5, w=9cm) ===============
    col1_x = Cm(1.5)
    col1_w = Cm(9)

    # Headline y=3, height 2.5cm
    t(s, headline, col1_x, Cm(3.0), col1_w, Cm(2.5),
      font=BF, size=20, color=WHITE, bold=True)

    # 3 bloques narrativa (y=7 → 16.5, c/u ~3cm + gap 0.5cm)
    blocks_y = Cm(7.0)
    block_h = Cm(2.8)
    block_gap = Cm(0.5)

    narrative = [
        ("The context", context),
        ("The contrast", contrast),
        ("The transformation", transformation),
    ]
    for i, (label, text_content) in enumerate(narrative):
        by = blocks_y + i * (block_h + block_gap)
        # Label italic 11pt gris claro
        t(s, label, col1_x, by, col1_w, Cm(0.5),
          font=BF, size=11, color=GRAY_L, italic=True)
        # Gap 0.3cm
        # Paragraph Poppins Regular 9pt blanco
        t(s, text_content, col1_x, by + Cm(0.5), col1_w, block_h - Cm(0.5),
          font=BF, size=9, color=WHITE)

    # =============== COLUMNA 2 — STATS (x=12 → 21, w=9cm) ===============
    col2_x = Cm(12)
    col2_w = Cm(9)

    # 3 stats stacked (y=3 → 16.5)
    # Cada stat: número 1.9cm + gap 0.6cm + descripción 2.5cm = 5.0cm
    # Gap entre stats: 0.8cm
    # Total: 5.0*3 + 0.8*2 = 16.6cm → empieza en y=3 termina y=19.6 (se sale)
    # Ajustar: altura stat = 4.2cm, total = 4.2*3 + 0.8*2 = 14.2 → y=3 a 17.2 ✓
    stat_num_h = Cm(1.9)   # 48pt × 0.039
    stat_gap_num_desc = Cm(0.6)
    stat_desc_h = Cm(1.7)
    stat_block_h = stat_num_h + stat_gap_num_desc + stat_desc_h  # ~4.2cm
    stat_vgap = Cm(0.8)

    stats_y = Cm(3.0)

    for i, (num, desc) in enumerate(stats):
        sy = stats_y + i * (stat_block_h + stat_vgap)
        # Número italic 48pt LEFT
        t(s, num, col2_x, sy, col2_w, stat_num_h,
          font=TF, size=48, color=WHITE, italic=True,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        # Descripción
        t(s, desc, col2_x, sy + stat_num_h + stat_gap_num_desc, col2_w, stat_desc_h,
          font=BF, size=9, color=WHITE)

    # =============== COLUMNA 3 — THE PROOF (x=22.5 → 33, w=10.5cm) ===============
    col3_x = Cm(22.5)
    col3_w = Cm(10.5)
    col3_y = Cm(2.5)
    col3_h = Cm(13)

    proof = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                col3_x, col3_y, col3_w, col3_h)
    proof.fill.solid()
    proof.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    proof.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    proof.line.width = Pt(1.0)
    proof.adjustments[0] = 0.08

    # Label "The proof" top-left interior (+0.6cm, +0.6cm)
    t(s, "The proof", col3_x + Cm(0.6), col3_y + Cm(0.6), Cm(5), Cm(0.5),
      font=BF, size=12, color=WHITE)

    source(s)
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    print("Generando Trends Identidad (3 macro tendencias)...")

    # ========== MACRO 1 — INVENTOLOGÍA DE LA ADULTEZ ==========
    build_trend(
        prs,
        headline="UN MUNDO EN CRISIS ESTÁ REESCRIBIENDO QUÉ SIGNIFICA SER ADULTO Y FORMAR FAMILIA",
        context="El mundo entró en crisis múltiple: económica, climática, política, demográfica. Se rompieron los marcadores que durante un siglo definieron qué era ser adulto: casarse, parir, comprar casa. Una generación entera está inventando otras formas de llegar a adultez.",
        contrast="En 1971, el CIAS describía al padre dominicano como figura cuya autoridad era razón última. Cincuenta y tres años después, la adultez ya no se hereda como trono — se inventa como acuerdo.",
        transformation="En cinco años, la familia biparental dejará de ser el default mental del dominicano. Las marcas que sigan vendiendo el modelo de los 90 estarán pagando para hablarle a un país que ya no existe.",
        stats=[
            ("2.3", "Tasa de fecundidad global en 2023, cayendo desde 4.9 en los años 50. ONU World Population Prospects 2024."),
            ("38%", "De los hogares dominicanos son biparentales clásicos. El 62% restante vive en estructuras distintas al ideal. Código Casa 2025."),
            ("18%", "De Gen Z global ve subir la escalera corporativa como inteligente. El rol de proveedor único murió. Fiverr 2025."),
        ]
    )
    print("  ✓ 1 Inventología de la Adultez")

    # ========== MACRO 2 — LOS HERNÁNDEZ ARE PROMPTED ==========
    build_trend(
        prs,
        headline="EL DOMINICANO YA ENTRÓ AL MUNDO PROMPTEADO. SOLO NO LO NOMBRA ASÍ.",
        context="El dominicano entró a la era post-GenAI digitalmente activo: 72% declara usar IA, y la casa ya tiene Alexa antes que biblioteca. Pero el uso todavía es instrumental: se delega tarea, no criterio.",
        contrast="La IA entra al hogar dominicano por tres puertas: productividad, vigilancia y delegación generacional. No todavía por la puerta de la intimidad emocional.",
        transformation="En 3-5 años la IA dejará de ser herramienta del muchacho y se volverá consejera de la casa. La infraestructura de confianza ya está montada.",
        stats=[
            ("72%", "De dominicanos afirma haber usado herramientas de IA. Global Public Confidence Study 2025, IRIS Network."),
            ("47%", "De usuarios fuertes de IA se sienten más confiados al usarla en momentos de inseguridad para comprar. Accenture 2025."),
            ("44.96", "Score dominicano en el ILIA: noveno de 19 países latinoamericanos en adopción activa de inteligencia artificial. CEPAL 2025."),
        ]
    )
    print("  ✓ 2 Los Hernández Are Prompted")

    # ========== MACRO 3 — ALGORITMO DEL HOGAR ==========
    build_trend(
        prs,
        headline="EL FEED SE SENTÓ EN LA MESA Y NADIE LE OFRECIÓ SILLA",
        context="Percibimos el mundo desde lo digital antes que desde lo físico. El feed gobierna qué se cocina, qué se compra, cómo se cría, qué se discute en la mesa. Las redes son el sistema operativo del hogar dominicano.",
        contrast="El algoritmo entra al hogar dominicano con dos caras: ritual nocturno compartido y fuente primaria de criterio de crianza y pareja. La abuela, el pastor y el padre compiten por atención con él.",
        transformation="En la próxima década, el algoritmo dejará de ser infraestructura invisible para volverse miembro reconocido del hogar: con nombre, con voz, con opinión pedida. La pregunta será qué asiento le damos en la mesa.",
        stats=[
            ("15%", "De consumidores globales compra productos puramente porque son tendencia en TikTok. SAP Emarsys 2025."),
            ("21B", "Tamaño del mercado de family/parenting influencers en 2024. Proyección de US$32B para 2027. Influencer Marketing Hub 2025."),
            ("71%", "De Gen Z y millennials dice haber comprado un producto de hogar después de verlo en TikTok. Adweek/TikTok 2024."),
        ]
    )
    print("  ✓ 3 Algoritmo del Hogar")

    prs.save(OUTPUT)
    print(f"\n✅ {OUTPUT}")


if __name__ == "__main__":
    main()
