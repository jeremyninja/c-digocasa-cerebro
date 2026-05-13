#!/usr/bin/env python3
"""T01 v3 — UI pulido: opacity 15%, mask ratio, 30pt titles, 45pt connectors, separator lines, spacing ≥4cm"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

ASSETS = "/Users/jeremyrodriguez/Downloads/Código Casa/Gráfica"
OUTPUT = "/Users/jeremyrodriguez/Claude/codigo-casa/T01-Familia-Identidad-v3.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_M = RGBColor(0x88, 0x88, 0x88)
GRAY_D = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0xFF, 0xFF, 0xFF)  # deprecated: ahora siempre blanco

TF = "Instrument Serif"
BF = "Poppins"

DUST = os.path.join(ASSETS, "Dust texture", "04.png")
CONN = os.path.join(ASSETS, "Gradients finales", "DESIGN-SYNDROME-HOLOGRAPHIX-TEXTURE-GENERATOR-V2-23.jpg")
UPPER = os.path.join(ASSETS, "Upper bar.png")
MASK = os.path.join(ASSETS, "Logo Final", "Mascara CC.png")
MASK_RATIO = 1.17  # width/height

CAP = "FAMILIA E IDENTIDAD"
SRC = "Source: Código Casa: Estudio cuantitativo y cualitativo familias dominicanas 2025 | 500 personas | +70 horas de grupos focales"
SRC_L = "SOURCE: NINJA SOCIAL LISTENING IMPACT® — DATA RECOPILADA DESDE EL 01 DE SEPTIEMBRE AL 31 DE DICIEMBRE DEL 2025."


def make_dust(opacity=0.15):
    """Dust texture at very low opacity over black."""
    cache = f"/tmp/dust_04_op{int(opacity*100)}.png"
    if os.path.exists(cache): return cache
    dust = Image.open(DUST).convert("RGBA")
    black = Image.new("RGBA", dust.size, (0, 0, 0, 255))
    alpha = dust.split()[-1].point(lambda p: int(p * opacity))
    dust.putalpha(alpha)
    out = Image.alpha_composite(black, dust)
    out.convert("RGB").save(cache)
    return cache

DUST_BG = make_dust(0.15)


def bg(s, conn=False):
    s.shapes.add_picture(CONN if conn else DUST_BG, 0, 0, SLIDE_W, SLIDE_H)


def upper(s):
    s.shapes.add_picture(UPPER, 0, 0, SLIDE_W, Cm(0.35))


def mask_corner(s):
    """Bottom-right mask, aspect ratio preserved."""
    h = Cm(1.0)
    w = Cm(1.0 * MASK_RATIO)
    s.shapes.add_picture(MASK, SLIDE_W - w - Cm(0.8), SLIDE_H - h - Cm(0.8), w, h)


def mask_center_bottom(s):
    """Mask centered at bottom (for connectors)."""
    h = Cm(1.2)
    w = Cm(1.2 * MASK_RATIO)
    s.shapes.add_picture(MASK, (SLIDE_W - w) / 2, SLIDE_H - h - Cm(1.5), w, h)


def header(s):
    """7pt double-sided header."""
    y = Cm(0.85)
    t(s, f"CÓDIGO CASA® - {CAP}", Cm(0.8), y, Cm(15), Cm(0.4),
      size=7, color=WHITE)
    t(s, "NINJA THINKING", SLIDE_W - Cm(15.8), y, Cm(15), Cm(0.4),
      size=7, color=WHITE, align=PP_ALIGN.RIGHT)


def chrome(s):
    bg(s)
    upper(s)
    header(s)
    mask_corner(s)


def t(s, text, x, y, w, h, font=BF, size=11, color=WHITE,
      bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if anchor: tf.vertical_anchor = anchor
    # Clear default paragraph, add fresh run with explicit formatting
    p = tf.paragraphs[0]
    p.alignment = align
    # Remove any existing runs to avoid inherited formatting
    for old in list(p.runs):
        old._r.getparent().remove(old._r)
    p.text = ""  # reset
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def rich(s, runs, x, y, w, h, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    # Clear any inherited runs
    for old in list(p.runs):
        old._r.getparent().remove(old._r)
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", BF)
        run.font.size = Pt(r.get("size", 11))
        run.font.color.rgb = r.get("color", WHITE)
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
    return tf


def vline(s, x, y, height_cm=6, width_pt=0.75):
    """Thin vertical separator — use thin rectangle for guaranteed visibility."""
    w = Cm(0.02)  # ~0.5pt thin
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - w/2, y, w, Cm(height_cm))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
    rect.line.fill.background()


def card(s, x, y, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x15)
    sh.line.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
    sh.line.width = Pt(0.5)
    sh.adjustments[0] = 0.08
    return sh


def pill(s, text, x, y, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.background()
    sh.line.color.rgb = GRAY_M
    sh.line.width = Pt(0.75)
    sh.adjustments[0] = 0.5
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Cm(0.3)
    tf.margin_top = tf.margin_bottom = Cm(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = BF
    r.font.size = Pt(8)
    r.font.color.rgb = WHITE
    r.font.bold = True


def source(s, txt=None):
    t(s, txt or SRC, Cm(0.8), SLIDE_H - Cm(1.0), Cm(25), Cm(0.4),
      size=6, color=GRAY_M)


# ============ SLIDES ============

def s1_connector(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, conn=True)
    header(s)  # Firmas top: CÓDIGO CASA® - CAP / NINJA THINKING

    # Text container: 25cm × 8cm centered horizontally
    tx = (SLIDE_W - Cm(25)) / 2
    ty = Cm(5)

    t(s, "Tensión 01", tx, ty + Cm(0.2), Cm(25), Cm(0.8),
      font=TF, size=20, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    rich(s, [
        {"text": "En RD ", "font": TF, "size": 42, "color": WHITE, "bold": True},
        {"text": '"mamá y papá"', "font": TF, "size": 42, "color": WHITE, "bold": True, "italic": True},
        {"text": " es el ideal", "font": TF, "size": 42, "color": WHITE, "bold": True},
    ], tx, ty + Cm(1.7), Cm(25), Cm(1.8), align=PP_ALIGN.CENTER)

    t(s, '"-Abuela, tía y nana" es la realidad-',
      tx, ty + Cm(3.8), Cm(25), Cm(1.8),
      font=TF, size=42, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    t(s, "Aquí no cría la familia nuclear, cría la red",
      tx, ty + Cm(6.5), Cm(25), Cm(0.8),
      font=BF, size=12, color=GRAY_L, align=PP_ALIGN.CENTER)

    mask_center_bottom(s)


def s2_hallazgo_3data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    rich(s, [
        {"text": "PARA EL DOMINICANO LA CRIANZA\n", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "ES ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "COLECTIVA", "font": TF, "size": 32, "color": WHITE, "italic": True},
        {"text": ", NO NUCLEAR", "font": BF, "size": 30, "color": WHITE, "bold": True},
    ], Cm(0), Cm(2.5), SLIDE_W, Cm(3.0), align=PP_ALIGN.CENTER)

    data = [
        ("48%", "Dice que sus padres participaron activamente en su crianza, pero sumando abuelas, tíos y hermanos mayores, la red extendida alcanza al 33% adicional."),
        ("89%", "Califica a sus figuras cuidadoras como 'muy influyentes' en su formación emocional."),
        ("38%", "Vive en hogares biparentales clásicos con hijos menores."),
    ]
    col_w = Cm(9.5)
    total = col_w * 3
    start = (SLIDE_W - total) / 2
    # Number at 90pt → glyph ~3.2cm tall, box Cm(4.5) safe
    num_y = Cm(7.5)
    num_h = Cm(3.5)
    ctx_y = num_y + num_h + Cm(0.85)

    ctx_offset = (col_w - Cm(7)) / 2
    for i, (num, ctx) in enumerate(data):
        x = start + i * col_w
        t(s, num, x, num_y, col_w, num_h,
          font=TF, size=90, color=WHITE, italic=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        t(s, ctx, x + ctx_offset, ctx_y, Cm(7), Cm(4),
          font=BF, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Separator lines: 5cm, centered vertically on the number block
    line_y = num_y + (num_h - Cm(5)) / 2
    vline(s, start + col_w, line_y, 5)
    vline(s, start + col_w * 2, line_y, 5)

    source(s)


def s3_hallazgo_2data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    rich(s, [
        {"text": "LA FAMILIA ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": '"IDEAL"', "font": TF, "size": 32, "color": WHITE, "italic": True},
        {"text": " EXISTE.\n", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "SOLO QUE EN EL ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "38%", "font": TF, "size": 32, "color": WHITE, "italic": True},
        {"text": " DE LOS HOGARES.", "font": BF, "size": 30, "color": WHITE, "bold": True},
    ], Cm(0), Cm(2.5), SLIDE_W, Cm(3.0), align=PP_ALIGN.CENTER)

    col_w = Cm(13)
    start = (SLIDE_W - col_w * 2) / 2
    num_y = Cm(7.0)
    num_h = Cm(4.3)
    ctx_y = num_y + num_h + Cm(0.85)

    ctx_offset = (col_w - Cm(7)) / 2
    t(s, "38%", start, num_y, col_w, num_h,
      font=TF, size=110, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    t(s, "de los hogares son biparentales clásicos. El 62% restante vive en estructuras distintas al ideal.",
      start + ctx_offset, ctx_y, Cm(7), Cm(4),
      font=BF, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    t(s, "24%", start + col_w, num_y, col_w, num_h,
      font=TF, size=110, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    t(s, "son monoparentales — donde tías, abuelas y hermanas mayores suplen los vacíos.",
      start + col_w + ctx_offset, ctx_y, Cm(7), Cm(4),
      font=BF, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    line_y = num_y + (num_h - Cm(5)) / 2
    vline(s, start + col_w, line_y, 5)
    source(s)


def s4_hallazgo_apellido(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    rich(s, [
        {"text": "EL QUE ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "CRÍA", "font": TF, "size": 32, "color": WHITE, "italic": True},
        {"text": " ES FAMILIA.\n", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "TENGA O NO EL APELLIDO.", "font": BF, "size": 30, "color": WHITE, "bold": True},
    ], Cm(0), Cm(2.5), SLIDE_W, Cm(3.0), align=PP_ALIGN.CENTER)

    col_w = Cm(13)
    start = (SLIDE_W - col_w * 2) / 2
    num_y = Cm(7.0)
    num_h = Cm(4.3)
    ctx_y = num_y + num_h + Cm(0.85)

    ctx_offset = (col_w - Cm(7)) / 2
    t(s, "89%", start, num_y, col_w, num_h,
      font=TF, size=110, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    t(s, "Dice que sus figuras de crianza fueron muy influyentes en su formación emocional.",
      start + ctx_offset, ctx_y, Cm(7), Cm(4),
      font=BF, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    t(s, "79%", start + col_w, num_y, col_w, num_h,
      font=TF, size=110, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    t(s, 'Recuerda la crianza a través de "me cuidaba cuando estaba enfermo/a."',
      start + col_w + ctx_offset, ctx_y, Cm(7), Cm(4),
      font=BF, size=11, color=WHITE, align=PP_ALIGN.CENTER)

    line_y = num_y + (num_h - Cm(5)) / 2
    vline(s, start + col_w, line_y, 5)
    source(s)


def s5_digital(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    pill(s, "CONVERSACIÓN DIGITAL", (SLIDE_W - Cm(5.5)) / 2, Cm(2.2), Cm(5.5), Cm(0.7))

    # 30pt title with italic % emphasis
    rich(s, [
        {"text": "53% ", "font": TF, "size": 36, "color": WHITE, "italic": True},
        {"text": "DE LAS MENCIONES SOBRE CRIANZA NO\n", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "BIOLÓGICA ES NEGATIVA POR DENUNCIAS A LA\n", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "FIGURA QUE NO ASUMIÓ SU ROL", "font": BF, "size": 30, "color": WHITE, "bold": True},
    ], Cm(0), Cm(3.7), SLIDE_W, Cm(4), align=PP_ALIGN.CENTER)

    # Wide rounded rectangle container: 4cm margin from slide edges
    rect_x = Cm(4)
    rect_w = SLIDE_W - Cm(8)
    rect_y = Cm(8.2)
    rect_h = Cm(4.2)
    card(s, rect_x, rect_y, rect_w, rect_h)

    # Inside: MENCIONES | line | ALCANCE | line | description (breathing 1cm inside)
    inner_pad = Cm(1.2)
    content_y = rect_y + Cm(0.9)
    num_h = Cm(1.9)  # 48pt × 0.039 ≈ 1.87cm

    t(s, "MENCIONES", rect_x + inner_pad, content_y, Cm(4), Cm(0.4),
      font=BF, size=8, color=GRAY_M, bold=True, align=PP_ALIGN.CENTER)
    t(s, "258", rect_x + inner_pad, content_y + Cm(0.6), Cm(4), num_h,
      font=TF, size=48, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    vline(s, rect_x + inner_pad + Cm(4.4), content_y, 2.5)

    t(s, "ALCANCE", rect_x + inner_pad + Cm(4.8), content_y, Cm(4), Cm(0.4),
      font=BF, size=8, color=GRAY_M, bold=True, align=PP_ALIGN.CENTER)
    t(s, "1.3M", rect_x + inner_pad + Cm(4.8), content_y + Cm(0.6), Cm(4), num_h,
      font=TF, size=48, color=WHITE, italic=True,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    vline(s, rect_x + inner_pad + Cm(9.2), content_y, 2.5)

    t(s, "La conversación digital no castiga a la familia no tradicional por existir, pero sí expone el peso simbólico, social y económico que todavía recae sobre ella. Más que mostrar qué estructuras familiares predominan, el debate revela cuáles siguen siendo más cuestionadas, más discutidas y más difíciles de validar fuera del modelo tradicional.",
      rect_x + inner_pad + Cm(9.6), content_y, rect_w - inner_pad*2 - Cm(9.6), Cm(3),
      font=BF, size=9, color=WHITE)

    # 3 verbatim cards — rounded edges
    card_w = Cm(10)
    gap = Cm(0.8)
    total = card_w * 3 + gap * 2
    start = (SLIDE_W - total) / 2
    verbatims = [
        ("Ismael_", "Tik Tok", '"crecí sin amor de padre y madre y a pesar de que me siento bien, siempre quise sentir ese amor…"'),
        ("Jaydennoficial", "Tik Tok", '"Quieren regalo por el día de la madre y al muchacho lo crió la abuela"'),
        ("cdn37", "Tik Tok", '"Niño de 12 años conmueve al defender a su padrastro frente a su padre biológico"'),
    ]
    cy = Cm(13.2)
    for i, (name, plat, q) in enumerate(verbatims):
        x = start + i * (card_w + gap)
        card(s, x, cy, card_w, Cm(3.8))
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Cm(0.6), cy + Cm(0.6), Cm(0.3), Cm(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()
        t(s, name, x + Cm(1.1), cy + Cm(0.45), card_w - Cm(1.3), Cm(0.5),
          font=BF, size=11, color=WHITE, bold=True)
        t(s, plat, x + Cm(1.1), cy + Cm(1.0), card_w - Cm(1.3), Cm(0.4),
          font=BF, size=8, color=GRAY_M)
        t(s, q, x + Cm(0.6), cy + Cm(1.8), card_w - Cm(1.2), Cm(2),
          font=BF, size=9, color=WHITE, italic=True)

    source(s, SRC_L)


def s6_verdades(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    rich(s, [
        {"text": "VERDADES ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "INCÓMODAS", "font": TF, "size": 32, "color": WHITE, "italic": True},
    ], Cm(0), Cm(3.5), SLIDE_W, Cm(1.5), align=PP_ALIGN.CENTER)

    t(s, "QUE NADIE TE QUIERE DECIR",
      Cm(0), Cm(5.3), SLIDE_W, Cm(0.6),
      font=BF, size=11, color=GRAY_M, align=PP_ALIGN.CENTER)

    truths = [
        "Las marcas hablan del momento en que las familias están compartiendo, pero eso no sucede por más de 7 horas a la semana.",
        "La crianza en la familia dominicana es colectiva: pero tu marca solo habla del día de las madres.",
        "El mercado sigue hablándole a la familia nuclear que aparece en los comerciales.",
    ]
    card_w = Cm(9.5)
    gap = Cm(0.8)
    start = (SLIDE_W - (card_w * 3 + gap * 2)) / 2
    for i, tr in enumerate(truths):
        x = start + i * (card_w + gap)
        y = Cm(8.5)
        card(s, x, y, card_w, Cm(6.0))
        t(s, f"0{i+1}", x + Cm(0.8), y + Cm(0.8), Cm(2), Cm(1.5),
          font=TF, size=40, color=GRAY_M, italic=True)
        t(s, tr, x + Cm(0.8), y + Cm(2.8), card_w - Cm(1.6), Cm(3),
          font=BF, size=12, color=WHITE)


def s7_oportunidades(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    rich(s, [
        {"text": "MAPEO DE ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "OPORTUNIDADES", "font": TF, "size": 32, "color": WHITE, "italic": True},
    ], Cm(0), Cm(2.8), SLIDE_W, Cm(1.5), align=PP_ALIGN.CENTER)

    opps = [
        ("La crianza en la familia dominicana es colectiva: pero tu marca solo habla del día de las madres.",
         "Diseñar productos / plataformas para la crianza compartida.",
         "¿Qué pasaría si tu marca fuera la primera en reconocer que la familia dominicana cría en tribu?"),
        ("El mercado sigue hablándole a la familia nuclear que aparece en los comerciales.",
         "Reconfigurar el concepto de familia: Familia es a donde sientes que llegaste.",
         "¿Y si tu marca comienza a verse tan real como la verdadera familia dominicana?"),
        ("Las marcas hablan del momento en que las familias están compartiendo, pero eso no sucede más de 7 horas a la semana.",
         "Construir productos y comunicación para la conexión asincrónica de las familias.",
         "¿Qué pasaría si tu marca existiera en las horas en que la familia no está junta?"),
    ]

    card_w = Cm(9.5)
    gap = Cm(0.8)
    start = (SLIDE_W - (card_w * 3 + gap * 2)) / 2
    for i, (v, o, d) in enumerate(opps):
        x = start + i * (card_w + gap)
        y = Cm(5.8)
        card(s, x, y, card_w, Cm(9))
        t(s, "VERDAD INCÓMODA", x + Cm(0.8), y + Cm(0.8), card_w - Cm(1.6), Cm(0.4),
          font=BF, size=8, color=GREEN, bold=True)
        t(s, v, x + Cm(0.8), y + Cm(1.3), card_w - Cm(1.6), Cm(2.4),
          font=BF, size=10, color=WHITE, italic=True)
        t(s, "OPORTUNIDAD", x + Cm(0.8), y + Cm(4.0), card_w - Cm(1.6), Cm(0.4),
          font=BF, size=8, color=GREEN, bold=True)
        t(s, o, x + Cm(0.8), y + Cm(4.5), card_w - Cm(1.6), Cm(2.2),
          font=BF, size=11, color=WHITE, bold=True)
        t(s, "DETONADOR", x + Cm(0.8), y + Cm(6.9), card_w - Cm(1.6), Cm(0.4),
          font=BF, size=8, color=GREEN, bold=True)
        t(s, d, x + Cm(0.8), y + Cm(7.4), card_w - Cm(1.6), Cm(1.5),
          font=BF, size=10, color=WHITE, italic=True)

    source(s, "Source: Ninja Thinking 2026")


def s8_caso(prs):
    """Caso referencia: foto izquierda + texto derecha."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    # Label arriba
    t(s, "CASO REFERENCIA", Cm(0), Cm(2.0), SLIDE_W, Cm(0.5),
      font=BF, size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Layout 2 columnas
    margin = Cm(3)
    gap = Cm(1.2)
    col_w = (SLIDE_W - margin * 2 - gap) / 2  # ~13.7cm cada columna
    col_y = Cm(3.5)
    col_h = Cm(11)

    # LEFT: Image frame placeholder (card con label "[FOTO DEL CASO]")
    img_frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    margin, col_y, col_w, col_h)
    img_frame.fill.solid()
    img_frame.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    img_frame.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    img_frame.line.width = Pt(1.0)
    img_frame.line.dash_style = 7  # dashed line
    img_frame.adjustments[0] = 0.08

    t(s, "[ INSERTAR FOTO DEL CASO ]",
      margin, col_y + col_h/2 - Cm(0.4), col_w, Cm(0.8),
      font=BF, size=11, color=GRAY_M, bold=True, align=PP_ALIGN.CENTER)

    # RIGHT: Text column
    rx = margin + col_w + gap
    ty = col_y + Cm(0.3)

    # Headline 30pt
    t(s, "Ayuda para ser padres cuando no tengo tiempo.",
      rx, ty, col_w, Cm(3),
      font=TF, size=30, color=WHITE, italic=True, align=PP_ALIGN.LEFT)

    # Brand label
    t(s, "WHATSAPP · MARZO 2026",
      rx, ty + Cm(3.5), col_w, Cm(0.5),
      font=BF, size=9, color=GRAY_M, bold=True)

    # Description
    t(s, "WhatsApp lanzó en marzo de 2026 cuentas para preadolescentes gestionadas por padres. Permiten supervisar contactos, grupos y privacidad mediante un panel con PIN, manteniendo el cifrado de extremo a extremo en los mensajes.",
      rx, ty + Cm(4.3), col_w, Cm(6),
      font=BF, size=13, color=WHITE)

    source(s, "Source: Ninja Thinking 2026")


def s9_connector_kit(prs):
    """Kit pregunta connector — gradient bg, 45pt headline, firmas top, text 25x8."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, conn=True)
    header(s)

    tx = (SLIDE_W - Cm(25)) / 2
    ty = Cm(5)

    rich(s, [
        {"text": "¿Cómo puede tu marca\n", "font": TF, "size": 45, "color": WHITE, "bold": True},
        {"text": "transformar ", "font": TF, "size": 45, "color": WHITE, "italic": True},
        {"text": "su narrativa\npara responder a esta\n", "font": TF, "size": 45, "color": WHITE, "bold": True},
        {"text": "tensión?", "font": TF, "size": 45, "color": WHITE, "bold": True, "italic": True},
    ], tx, ty, Cm(25), Cm(8), align=PP_ALIGN.CENTER)

    mask_center_bottom(s)


def s10_framework_kit(prs):
    """Framework diagram: Situación → 3 ingredientes → Nueva Narrativa."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    # Title
    rich(s, [
        {"text": "KIT DE ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "RESPUESTA", "font": TF, "size": 32, "color": WHITE, "italic": True},
    ], Cm(0), Cm(1.8), SLIDE_W, Cm(1.5), align=PP_ALIGN.CENTER)

    t(s, "Marco para construir tu nueva narrativa de marca",
      Cm(0), Cm(3.4), SLIDE_W, Cm(0.6),
      font=BF, size=11, color=GRAY_L, italic=True, align=PP_ALIGN.CENTER)

    # SITUACIÓN (top band, full-width)
    band_x = Cm(4)
    band_w = SLIDE_W - Cm(8)
    card(s, band_x, Cm(4.6), band_w, Cm(1.8))
    t(s, "SITUACIÓN", band_x, Cm(4.9), band_w, Cm(0.5),
      font=BF, size=9, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    t(s, "Desde la tensión identificada — planteamiento a resolver",
      band_x, Cm(5.5), band_w, Cm(0.7),
      font=BF, size=13, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    # Arrow down
    t(s, "↓", Cm(0), Cm(6.7), SLIDE_W, Cm(0.6),
      font=BF, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 3 ingredient cards
    card_w = Cm(8.5)
    gap = Cm(0.6)
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    ing_y = Cm(7.6)
    card_h = Cm(4.5)

    ingredients = [
        ("01", "OPORTUNIDAD CLAVE", "Usa el cuadro de oportunidades para detectar dónde tu marca puede actuar."),
        ("02", "POSTURA / ROL DE MARCA", "Determina la postura o rol ante el hallazgo."),
        ("03", "VERBO TRANSFORMADOR", "Establece el verbo de acción de tu marca."),
    ]
    for i, (num, label, desc) in enumerate(ingredients):
        x = start_x + i * (card_w + gap)
        card(s, x, ing_y, card_w, card_h)
        t(s, num, x + Cm(0.8), ing_y + Cm(0.7), card_w - Cm(1.6), Cm(1.2),
          font=TF, size=32, color=GREEN, italic=True)
        t(s, label, x + Cm(0.8), ing_y + Cm(2.0), card_w - Cm(1.6), Cm(0.6),
          font=BF, size=10, color=WHITE, bold=True)
        t(s, desc, x + Cm(0.8), ing_y + Cm(2.8), card_w - Cm(1.6), Cm(1.5),
          font=BF, size=9, color=GRAY_L)

    # Arrow down
    t(s, "↓", Cm(0), Cm(12.4), SLIDE_W, Cm(0.6),
      font=BF, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Output: Fórmula
    card(s, band_x, Cm(13.3), band_w, Cm(3.2))
    t(s, "NUEVA NARRATIVA DE MARCA", band_x, Cm(13.6), band_w, Cm(0.5),
      font=BF, size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    rich(s, [
        {"text": "[Tu marca]", "font": TF, "size": 16, "color": WHITE, "italic": True},
        {"text": "  quiere  ", "font": BF, "size": 13, "color": GRAY_L},
        {"text": "[rol de marca]", "font": TF, "size": 16, "color": WHITE, "italic": True},
        {"text": "  pero sabe que  ", "font": BF, "size": 13, "color": GRAY_L},
        {"text": "[oportunidad clave]", "font": TF, "size": 16, "color": WHITE, "italic": True},
        {"text": "  por eso  ", "font": BF, "size": 13, "color": GRAY_L},
        {"text": "[transformación]", "font": TF, "size": 16, "color": WHITE, "italic": True},
    ], band_x + Cm(1), Cm(14.4), band_w - Cm(2), Cm(2.5), align=PP_ALIGN.CENTER)

    source(s, "Source: Ninja Thinking 2026")


def s11_kit_respuesta(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    t(s, "KIT DE RESPUESTA — TENSIÓN 01",
      Cm(0), Cm(2.2), SLIDE_W, Cm(0.5),
      font=BF, size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    rich(s, [
        {"text": "Nueva ", "font": TF, "size": 30, "color": WHITE},
        {"text": "Narrativa", "font": TF, "size": 30, "color": WHITE, "italic": True},
        {"text": " de Marca", "font": TF, "size": 30, "color": WHITE},
    ], Cm(0), Cm(3.2), SLIDE_W, Cm(1.5), align=PP_ALIGN.CENTER)

    # Narrative card — with breathing room
    card(s, Cm(4), Cm(6), SLIDE_W - Cm(8), Cm(4.5))
    t(s, "[Tu marca] Quiere hacer visible a la familia dominicana",
      Cm(5), Cm(6.7), SLIDE_W - Cm(10), Cm(1),
      font=BF, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    t(s, "Pero sabe que tu verdadera familia va más allá del núcleo tradicional,\npor eso celebra la familia real sin importar cómo esté compuesta.",
      Cm(5), Cm(8.2), SLIDE_W - Cm(10), Cm(2),
      font=BF, size=12, color=GRAY_L, italic=True, align=PP_ALIGN.CENTER)

    # Meta grid — centered with 4cm breathing
    meta = [
        ("POSTURA/ROL", 'La marca que "hace visible" a mi familia'),
        ("OPORTUNIDAD CLAVE", "Familia extendida ignorada como estructura funcional"),
        ("VERBO TRANSFORMADOR", "Celebrar"),
        ("DESDE LA TENSIÓN", 'En RD "mamá y papá" es el ideal. "Abuela, tía y nana" es la realidad.'),
    ]
    y = Cm(11.5)
    for label, value in meta:
        t(s, label, Cm(5), y, Cm(8), Cm(0.4),
          font=BF, size=8, color=GREEN, bold=True)
        t(s, value, Cm(13), y, Cm(17), Cm(0.4),
          font=BF, size=11, color=WHITE)
        y += Cm(0.65)

    source(s, "Source: Ninja Thinking 2026")


def s12_espacios_vacios(prs, slide_num, slots):
    """Slide de Espacios Vacíos — 3 cards por slide.
    slots = [(num, titulo, desc), ...] — 3 items
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(s)

    # Title 30pt mix
    rich(s, [
        {"text": "ESPACIOS ", "font": BF, "size": 30, "color": WHITE, "bold": True},
        {"text": "VACÍOS", "font": TF, "size": 32, "color": WHITE, "italic": True},
    ], Cm(0), Cm(2.0), SLIDE_W, Cm(1.5), align=PP_ALIGN.CENTER)

    # Subtitle
    t(s, f"Que ninguna marca está habitando  —  {slide_num}/4",
      Cm(0), Cm(3.7), SLIDE_W, Cm(0.6),
      font=BF, size=12, color=GRAY_L, italic=True, align=PP_ALIGN.CENTER)

    # 3 cards
    margin = Cm(3)
    gap = Cm(0.8)
    card_w = (SLIDE_W - margin * 2 - gap * 2) / 3  # ~8.9cm cada
    card_y = Cm(5.5)
    card_h = Cm(9.5)

    for i, (num, titulo, desc) in enumerate(slots):
        x = margin + i * (card_w + gap)
        card(s, x, card_y, card_w, card_h)
        # Número grande
        t(s, num, x + Cm(0.8), card_y + Cm(0.8), card_w - Cm(1.6), Cm(2),
          font=TF, size=48, color=WHITE, italic=True,
          anchor=MSO_ANCHOR.TOP)
        # Línea sutil debajo del número (opcional, gris fino)
        sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Cm(0.8), card_y + Cm(3.3), Cm(1.5), Cm(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
        sep.line.fill.background()
        # Título del espacio
        t(s, titulo, x + Cm(0.8), card_y + Cm(3.8), card_w - Cm(1.6), Cm(2),
          font=BF, size=13, color=WHITE, bold=True)
        # Descripción
        t(s, desc, x + Cm(0.8), card_y + Cm(5.8), card_w - Cm(1.6), Cm(3.5),
          font=BF, size=10, color=GRAY_L)

    source(s, "Source: Ninja Thinking 2026")


ESPACIOS = [
    # Slide 1
    [
        ("01", "CRIANZA EN TRIBU",
         "Producto o plataforma que reconoce que la crianza dominicana es colectiva — no nuclear."),
        ("02", "FAMILIA REAL VISIBLE",
         "Comunicar la familia real — no la ideal del comercial de televisión."),
        ("03", "CONEXIÓN ASINCRÓNICA",
         "Existir en las horas que la familia no está junta. Más allá del 'momento en familia'."),
    ],
    # Slide 2
    [
        ("04", "ABUELAS COMO ALIADAS",
         "Reconocer el rol central de la figura abuela en la crianza dominicana."),
        ("05", "PADRES NO BIOLÓGICOS",
         "Visibilizar al padrastro, al tío, al vecino que crió — los que la sociedad no nombra."),
        ("06", "HOGAR MONOPARENTAL",
         "Diseñar para el 24% de hogares monoparentales que el mercado ignora."),
    ],
    # Slide 3
    [
        ("07", "FAMILIA EXTENDIDA ECONÓMICA",
         "Productos y servicios para estructuras multi-generacionales bajo un mismo techo."),
        ("08", "RED DE CUIDADO",
         "Apoyar a quien cuida — no solo a quien recibe el cuidado. El cuidador invisible."),
        ("09", "IDENTIDAD DOMINICANA MODERNA",
         "Qué significa ser orgullosamente dominicano hoy — más allá del estereotipo."),
    ],
    # Slide 4
    [
        ("10", "FAMILIA HOMOPARENTAL",
         "Reconocer al 0.6% que nunca se nombra ni se le habla desde las marcas."),
        ("11", "HIJOS ADULTOS EN CASA",
         "Hogares con adultos que siguen viviendo con sus padres — la estructura silenciada."),
        ("12", "REDEFINICIÓN DEL HOGAR",
         "El hogar como red de afecto, no como estructura física o consanguínea."),
    ],
]


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    print("Generando T01 v3...")
    builders = [
        (s1_connector, "Connector tensión"),
        (s2_hallazgo_3data, "Hallazgo + 3 datos (48/89/38) + líneas separadoras"),
        (s3_hallazgo_2data, "Familia ideal 38/24"),
        (s4_hallazgo_apellido, "El que cría 89/79"),
        (s5_digital, "Social listening + cards rounded"),
        (s6_verdades, "Verdades incómodas"),
        (s7_oportunidades, "Mapeo oportunidades"),
        (s8_caso, "Caso referencia"),
        (s9_connector_kit, "Connector: ¿Cómo transformar? (45pt)"),
        (s10_framework_kit, "Framework Kit (Situación → 3 ingredientes → Fórmula)"),
        (s11_kit_respuesta, "Kit de respuesta T01"),
    ]
    for i, (fn, name) in enumerate(builders, 1):
        fn(prs)
        print(f"  ✓ {i:2d} {name}")

    # 4 slides Espacios Vacíos
    for idx, slots in enumerate(ESPACIOS, 1):
        s12_espacios_vacios(prs, idx, slots)
        print(f"  ✓ {11+idx:2d} Espacios Vacíos {idx}/4")

    prs.save(OUTPUT)
    print(f"\n✅ {OUTPUT}")


if __name__ == "__main__":
    main()
