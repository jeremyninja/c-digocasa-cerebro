#!/usr/bin/env python3
"""
build_alimentacion_deck.py
==========================
Deck alimentacion-deck-flat.pptx desde cero con python-pptx.
21 slides — set editorial cerrado alimentacion-hallazgos-editados.md (12 hallazgos).

Distribución por hallazgo:
  H01 (T1): cuanti 2 stats + Consumer Voice     → 2 slides
  H02 (T2): cuanti 2 stats + Consumer Voice     → 2 slides
  H03 (T3): cuanti 2 stats + Consumer Voice     → 2 slides
  H04 (T4): cuanti 1 stat  + Consumer Voice     → 2 slides
  H05 (T5): cuanti 2 stats + Consumer Voice     → 2 slides
  H06 (T6): cuanti 2 stats + Consumer Voice     → 2 slides
  H07 (T7): cuanti 2 stats + Consumer Voice     → 2 slides
  H08 (T7): cuanti 1 stat  + Consumer Voice     → 2 slides
  H09 (T3): cuanti 2 stats + Consumer Voice     → 2 slides
  H10 (T8): solo-cuali 2 verbatims cards s/s    → 1 slide
  H11 (T9): solo-cuali 2 verbatims cards s/s    → 1 slide
  H12 (T10): solo-cuali 2 verbatims cards s/s   → 1 slide

TOTAL: 9×2 + 3×1 = 21 slides

Specs visuales (aprendizajes-montador-cc.md § v6 — helpers copiados de build_sistema_creencias_deck.py):
  Slide: 1920×1080 px = 18,288,000 × 10,287,000 EMU
  Fonts: Instrument Serif (headlines, stats, verbatims CV) + Poppins (cuerpo, atribución)
  Fondo: #000000 negro plano (NO masterslide)
  Consumer Voice (Layout 4): verbatim SUELTO sobre fondo negro — SIN card
  Cards cuali (Layout 5): 567×227pt side by side con fill rgba(0,0,0,0.45)
  Stat number: 135pt, Instrument Serif italic
  Stat desc: 13pt Poppins
  Headline: 50pt fijo Instrument Serif MAYÚSCULAS
  NO Source al pie
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

def pt_from_px(px_val):
    return round(px_val * 0.75, 1)

# ── Fuentes canónicas (aprendizajes v6) ────────────────────────────────────────
HEADLINE_PT       = 50.0
STAT_NUMBER_PT    = pt_from_px(180)  # 135pt
STAT_DESC_PT      = 13.0
CV_HEADER_PT      = 14.0
CV_VERBATIM_PT    = 50.0
CV_ATTRIB_PT      = 20.0
CARD_VERBATIM_PT  = 15.0
CARD_ATTRIB_PT    = 12.0

# ── Dimensiones de stat boxes (v6: 10cm × 3cm = 378×113pt) ────────────────────
STAT_BOX_W_1STAT  = 378
STAT_BOX_W_2STATS = 378
STAT_BOX_W_3STATS = 378
STAT_BOX_H        = 113
STAT_TOP_PX       = 480

# ── Dimensiones de cards cuali (v6: 567×227pt side by side) ───────────────────
CARD_W  = 567
CARD_H  = 227
CARD_GAP_2 = 80
CARD_GAP_3 = 36
CARD_TOP    = 429

# ── Consumer Voice dimensions ─────────────────────────────────────────────────
CV_TEXT_W   = 1637
CV_TEXT_LEFT = (SLIDE_W_PX - CV_TEXT_W) // 2  # = 141px


# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en el run XML."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """Headline 50pt fijo Instrument Serif MAYUSCULAS centrado."""
    txBox = slide.shapes.add_textbox(
        px(100), px(80), px(1720), px(320)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(HEADLINE_PT)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(HEADLINE_PT)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    """Linea vertical fina #2E2E2E entre stats."""
    line = slide.shapes.add_shape(
        1,
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left_px, top_px, width_px, height_px):
    """Card rounded rectangle: fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt."""
    shape = slide.shapes.add_shape(
        5,
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass

    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_el.set('val', '45000')  # 45% opacidad

    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def stat_box_width(n_stats):
    if n_stats == 1:
        return STAT_BOX_W_1STAT
    elif n_stats == 2:
        return STAT_BOX_W_2STATS
    else:
        return STAT_BOX_W_3STATS


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w = stat_box_width(n_stats)
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=STAT_BOX_H,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


# ── Layouts de slides ──────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic, stats):
    """
    Slide de Hallazgo cuanti.
    NO source (eliminado per aprendizajes v5/v6).
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice (Layout 4 — verbatim SUELTO sobre fondo negro, SIN card).
    Header "CONSUMER VOICE" Poppins 14pt gris arriba.
    Verbatim Instrument Serif 50pt blanco, comillas «».
    Atribucion Poppins italic 20pt blanco.
    """
    slide = blank_slide(prs)

    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    full_quote = f"«{quote}»"
    char_count = len(quote)
    if char_count > 200:
        verb_top = 280
    elif char_count > 120:
        verb_top = 330
    else:
        verb_top = 380

    verb_box = slide.shapes.add_textbox(
        px(CV_TEXT_LEFT), px(verb_top),
        px(CV_TEXT_W), px(450)
    )
    tf = verb_box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_q = p.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(CV_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    attrib_top = verb_top + 300
    add_textbox(
        slide,
        left_px=0, top_px=attrib_top,
        width_px=1920, height_px=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIB_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_cuali_slide_side_by_side(prs, headline_plain, headline_italic, verbatims):
    """
    Slide Card cualitativa (Layout 5 — cards 567×227pt SIDE BY SIDE).
    Headline arriba 50pt Instrument Serif MAYUSCULAS.
    Cards con fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt.
    Verbatim Poppins 15pt blanco, atribucion Poppins italic 12pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)

    if n == 1:
        card_left = (SLIDE_W_PX - CARD_W) // 2
        positions = [(card_left, CARD_TOP)]
    elif n == 2:
        total_w = CARD_W * 2 + CARD_GAP_2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_2, CARD_TOP)
        ]
    else:
        total_w = CARD_W * 3 + CARD_GAP_3 * 2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_3, CARD_TOP),
            (start_x + CARD_W * 2 + CARD_GAP_3 * 2, CARD_TOP)
        ]

    for i, v in enumerate(verbatims):
        left, top = positions[i]
        make_card_rgba_45_with_outline(slide, left, top, CARD_W, CARD_H)

        full_quote = f"«{v['quote']}»"
        pad_h = 30
        pad_v = 25
        text_w = CARD_W - (pad_h * 2)
        text_h = CARD_H - (pad_v * 2)

        txBox = slide.shapes.add_textbox(
            px(left + pad_h), px(top + pad_v),
            px(text_w), px(text_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left   = Emu(0)
        tf.margin_right  = Emu(0)
        tf.margin_top    = Emu(0)
        tf.margin_bottom = Emu(0)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run_q = p1.add_run()
        run_q.text = full_quote
        run_q.font.name   = FONT_BODY
        run_q.font.size   = Pt(CARD_VERBATIM_PT)
        run_q.font.italic = False
        run_q.font.bold   = False
        run_q.font.color.rgb = WHITE
        _set_kern(run_q._r)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run_a = p2.add_run()
        run_a.text = f"— {v['attribution']}"
        run_a.font.name   = FONT_BODY
        run_a.font.size   = Pt(CARD_ATTRIB_PT)
        run_a.font.italic = True
        run_a.font.bold   = False
        run_a.font.color.rgb = WHITE
        _set_kern(run_a._r)

    return slide


# ── BUILD DECK ─────────────────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — T1: El dominicano sabe que es comer bien. Lo sabe y no lo hace.
    # Headline: "46% define la buena alimentacion como dieta equilibrada.
    #  57% no sigue ninguna dieta — y 24% solo 'intenta comer mejor', sin metodo."
    # Stats: 46.4%, 57.0% + 23.8%
    # Verbatim: Familia Mixta
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 01 — H01 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="46% DEFINE LA BUENA ALIMENTACIÓN COMO DIETA EQUILIBRADA. ",
        headline_italic="57% NO SIGUE NINGUNA DIETA — Y 24% SOLO “INTENTA COMER MEJOR”, SIN MÉTODO.",
        stats=[
            {
                'value': '46.4%',
                'desc_runs': [
                    make_run('define buena alimentación como “'),
                    make_run('dieta equilibrada con proteínas, carbohidratos y vegetales', bold=True),
                    make_run('”. Es la respuesta dominante de la batería P20. Base 500.')
                ]
            },
            {
                'value': '57.0%',
                'desc_runs': [
                    make_run('declara no seguir ninguna dieta; '),
                    make_run('23.8% responde “intento comer mejor, pero no sigo una dieta específica”', bold=True),
                    make_run('. Sumados, 80.8% del hogar dominicano opera sin pauta alimenticia. P23, Base 500.')
                ]
            }
        ]
    )

    # Slide 02 — H01 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Bueno, a mí me ha tocado de como alrededor de dos años para acá aprender a comer la comida saludable. Yo antes, a pesar de que ustedes me ven con estas libritas, antes yo comía de una forma muy grotesca. Yo comía de todo.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — T2: La cocina es el centro del hogar dominicano. Aunque cada vez este mas sola.
    # Headline: "85% del hogar dominicano cocina en casa siempre.
    #  La calle compite, el delivery crece — la cocina no cede."
    # Stats: 85.0%, 1.8%
    # Verbatim: Familia Biparental con Hijos Adultos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 03 — H02 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="85% DEL HOGAR DOMINICANO COCINA EN CASA SIEMPRE. ",
        headline_italic="LA CALLE COMPITE, EL DELIVERY CRECE — LA COCINA NO CEDE.",
        stats=[
            {
                'value': '85.0%',
                'desc_runs': [
                    make_run('come comida hecha en casa con frecuencia TOP (5 = siempre). '),
                    make_run('Sumando frecuencias altas y medias, 98.2% come en casa con regularidad', bold=True),
                    make_run('. P21, Base 500.')
                ]
            },
            {
                'value': '1.8%',
                'desc_runs': [
                    make_run('come comida hecha en casa con '),
                    make_run('frecuencia baja o nunca', bold=True),
                    make_run('. La cifra tan baja confirma que el hallazgo real está en el 98.2% que sí lo hace. P21, Base 500.')
                ]
            }
        ]
    )

    # Slide 04 — H02 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Me identifico con las 7, muchísimo, porque imíginate, yo tengo cocina diaria, y tengo que hacer arroz, habichuelas, todos los días.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — T3: El estrato define quien cocina y quien compra.
    # Headline: "Cocinar en casa es habito popular, no burgues. En el estrato D, 89% siempre cocina.
    #  Quien tiene dinero compra el tiempo que otros no tienen."
    # Stats: 89.0% D + 83.9% C, 52.2% AB (tendencia)
    # Verbatim: Familia Biparental con Hijos Pequenos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 05 — H03 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="COCINAR EN CASA ES HÁBITO POPULAR, NO BURGUES. EN EL ESTRATO D, 89% SIEMPRE COCINA. ",
        headline_italic="QUIEN TIENE DINERO COMPRA EL TIEMPO QUE OTROS NO TIENEN.",
        stats=[
            {
                'value': '89.0%',
                'desc_runs': [
                    make_run('del estrato D come en casa con frecuencia TOP (siempre); '),
                    make_run('83.9% del estrato C', bold=True),
                    make_run('. La cocina doméstica es práctica de mayoría en los estratos de mayor peso muestral. P21 × NSE, Base D n=172 / C n=218.')
                ]
            },
            {
                'value': '52.2%',
                'desc_runs': [
                    make_run('del estrato AB declara “siempre” cocinar en casa — la más baja del espectro. '),
                    make_run('Tendencia direccional', bold=True),
                    make_run(', base AB n≤23 no publicable como dato aislado. P21 × NSE.')
                ]
            }
        ]
    )

    # Slide 06 — H03 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo gastaba mucho dinero en mi parámetro, mucho hasta que me tocó decir no esperate vamos a cocinar más en la casa vamos a pedir menos.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — T4: La cocina sigue siendo territorio de ella.
    # Headline: "89.5% femenino come comida hecha en casa siempre, vs 77.9% masculino.
    #  El hombre que cocina es la excepcion. Ella, la regla invisible."
    # Stats: 89.5% vs 77.9% — brecha 11.6 pp (1 stat, 2 cifras en desc)
    # Verbatim: Familia Sin Hijos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 07 — H04 Hallazgo cuanti (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="89.5% FEMENINO COME COMIDA HECHA EN CASA SIEMPRE, VS 77.9% MASCULINO. ",
        headline_italic="EL HOMBRE QUE COCINA ES LA EXCEPCIÓN. ELLA, LA REGLA INVISIBLE.",
        stats=[
            {
                'value': '11.6 pp',
                'desc_runs': [
                    make_run('es la '),
                    make_run('brecha de género', bold=True),
                    make_run(': 89.5% femenino vs 77.9% masculino con frecuencia TOP “siempre”. P21 × Sexo, Base 305 femenino / 195 masculino.')
                ]
            }
        ]
    )

    # Slide 08 — H04 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Bueno, yo tengo el rol de hacer la lista de compra. Nosotros el 100% del tiempo vamos al supermercado juntos. (…) Como ya casi siempre digamos el 90% del tiempo ella es quien cocina.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — T5: El precio manda, la salud espera.
    # Headline: "54% elige lo que compra por precio, no por salud.
    #  Lo nutricional vive en tercer lugar, despues del costo y la calidad."
    # Stats: 53.8% precio, 45.8% calidad + 34.4% salud
    # Verbatim: Familia Monoparental
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 09 — H05 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="54% ELIGE LO QUE COMPRA POR PRECIO, NO POR SALUD. ",
        headline_italic="LO NUTRICIONAL VIVE EN TERCER LUGAR, DESPUÉS DEL COSTO Y LA CALIDAD.",
        stats=[
            {
                'value': '53.8%',
                'desc_runs': [
                    make_run('menciona el '),
                    make_run('precio', bold=True),
                    make_run(' como factor que más influye en la compra de alimentos. TOP-1 de la batería P22 (multirespuesta). Calidad: 45.8%. Salud: 34.4%. Base 500.')
                ]
            },
            {
                'value': '34.4%',
                'desc_runs': [
                    make_run('menciona '),
                    make_run('salud o valor nutricional', bold=True),
                    make_run(' — tercer lugar, detrás de precio y calidad. Marca (14.4%), facilidad (10.4%) y rapidez (7.4%) quedan por debajo. P22, Base 500 (multirespuesta).')
                ]
            }
        ]
    )

    # Slide 10 — H05 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo siempre vivo como le dije la cosa del reciclaje en los supermercados, oferta dos por uno, marca propia, a veces no podemos comprar la que queremos por los intereses, en manera propia que es casi igual de bueno.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — T6: El regimen es proyecto de minoria.
    # Headline: "Solo 17% sigue una dieta especifica.
    #  El otro 83% come — sin regla, sin pauta, sin metodo."
    # Stats: 17.2% con dieta, dietas premium < 3%
    # Verbatim: Familia Biparental con Hijos Pequenos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 11 — H06 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="SOLO 17% SIGUE UNA DIETA ESPECÍFICA. ",
        headline_italic="EL OTRO 83% COME — SIN REGLA, SIN PAUTA, SIN MÉTODO.",
        stats=[
            {
                'value': '17.2%',
                'desc_runs': [
                    make_run('declara seguir una dieta específica: '),
                    make_run('sin azúcar 8.6%, baja en carbohidratos 7.8%', bold=True),
                    make_run(', fasting 1.6%, libre de gluten 1.0%, keto 0.2%. P23, Base 500.')
                ]
            },
            {
                'value': '<3%',
                'desc_runs': [
                    make_run('agregan las dietas “premium” combinadas '),
                    make_run('(keto + gluten-free + fasting)', bold=True),
                    make_run('. La opción vegana no aparece en tabulación (n=0 reportable). P23, Base 500.')
                ]
            }
        ]
    )

    # Slide 12 — H06 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo no es que yo como tanto, ya yo cené. Pero es un reto, llevar una dieta, porque ni por mi cuerpo ni por mi salud lo he logrado yo.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — T7: La mesa familiar no es ritual, es el unico momento que queda.
    # Headline: "34% se sienta a la mesa para compartir.
    #  La comida es la ultima excusa que queda — antes era la mesa, ahora es el rato."
    # Stats: 34.2% compartir, 16.6% + 14.6% + 15.2%
    # Verbatim: Familia Sin Hijos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 13 — H07 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="34% SE SIENTA A LA MESA PARA COMPARTIR. ",
        headline_italic="LA COMIDA ES LA ÚLTIMA EXCUSA QUE QUEDA — ANTES ERA LA MESA, AHORA ES EL RATO.",
        stats=[
            {
                'value': '34.2%',
                'desc_runs': [
                    make_run('declara que su familia comparte horas de comida “'),
                    make_run('para poder compartir un poco juntos', bold=True),
                    make_run('”. TOP-1 de la batería P25. Segunda razón: “seguir la tradición de comer en la mesa” (16.6%). Base 500.')
                ]
            },
            {
                'value': '15.2%',
                'desc_runs': [
                    make_run('usa la mesa “'),
                    make_run('para relajarnos', bold=True),
                    make_run('”; 14.6% “para hablar sobre las novedades del día”. Solo 2.4% la usa para hablar de noticias y política. P25, Base 500.')
                ]
            }
        ]
    )

    # Slide 14 — H07 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Sí, nosotros cada vez que estamos en casa, ya sea para la cena o fines de semana, desayuno, comida, nosotros tenemos que sentarnos los dos a comer. Eso es algo, no es de que tú en la sala y yo en la habitación, no. Es los dos en el comedor, hasta tenemos una regla de no comer en las habitaciones, ni en la cama ni nada.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — T7: La mesa familiar no es ritual, es el unico momento que queda.
    # Headline: "17% de los hogares no comparte ninguna hora de comida en familia.
    #  Una de cada seis casas come — pero no se sienta."
    # Stats: 17.0% no comparte (1 stat con contexto vs 16.6% tradicion)
    # Verbatim: Familia Monoparental
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 15 — H08 Hallazgo cuanti (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="17% DE LOS HOGARES NO COMPARTE NINGUNA HORA DE COMIDA EN FAMILIA. ",
        headline_italic="UNA DE CADA SEIS CASAS COME — PERO NO SE SIENTA.",
        stats=[
            {
                'value': '17.0%',
                'desc_runs': [
                    make_run('responde “'),
                    make_run('mi familia no comparte en estas horas', bold=True),
                    make_run('”. Segunda respuesta más alta de P25, por encima del 16.6% que cita “seguir la tradición de comer en la mesa”. La grieta y la tradición son del mismo tamaño. Base 500.')
                ]
            }
        ]
    )

    # Slide 16 — H08 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Me identifico con la 1, la 5 y la 7. (…) La 1 eso pasa siempre en los fines de semana porque en los días de semana no comemos, yo como en mi trabajo. Y ella cuando llega del colegio, entonces come sola, porque ya ha comido mi hermanito y la señora.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — T3: El estrato define quien cocina y quien compra.
    # Headline: "Compartir la mesa es mas frecuente donde hay mas ingreso.
    #  En el estrato D, la mesa familiar se disputa con horarios que el bolsillo no controla."
    # Stats: C 32.6% + D 30.8% + C+ 47.2%, AB 52.2% (tendencia)
    # Verbatim: Familia Sin Hijos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 17 — H09 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="COMPARTIR LA MESA ES MÁS FRECUENTE DONDE HAY MÁS INGRESO. ",
        headline_italic="EN EL ESTRATO D, LA MESA FAMILIAR SE DISPUTA CON HORARIOS QUE EL BOLSILLO NO CONTROLA.",
        stats=[
            {
                'value': '47.2%',
                'desc_runs': [
                    make_run('del estrato C+ comparte la mesa “para compartir un poco juntos”. '),
                    make_run('En C: 32.6%; en D: 30.8%', bold=True),
                    make_run('. La brecha entre estratos medios y populares supera los 15 puntos. P25 × NSE, Base C n=218 / D n=172.')
                ]
            },
            {
                'value': '52.2%',
                'desc_runs': [
                    make_run('del estrato AB declara compartir mesa como fin en sí mismo. '),
                    make_run('Tendencia direccional', bold=True),
                    make_run(', base AB n≤23 no publicable como dato aislado. La dirección —más ingreso, más mesa compartida— es robusta desde C hasta D. P25 × NSE.')
                ]
            }
        ]
    )

    # Slide 18 — H09 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='En el almuerzo nosotros siempre compramos comida porque no hacemos comida para comer regularmente. (…) Comemos, damos un poco, luego que terminamos de comer, ya la limpieza y hacer otra cosa.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — T8: Saludable es aspiracion. El entorno no ayuda. (FUSIONADO H10+H12 crudo)
    # Solo-cuali — 2 verbatims en cards side by side
    # Headline: "El hogar dominicano quiere comer saludable.
    #  El bolsillo, el tiempo y la calle se lo cobran."
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 19 — H10 Card cualitativa (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="EL HOGAR DOMINICANO QUIERE COMER SALUDABLE. ",
        headline_italic="EL BOLSILLO, EL TIEMPO Y LA CALLE SE LO COBRAN.",
        verbatims=[
            {
                'quote': 'Comer cosas saludables salen más baratos que comer saludables. Si tú quieres un buen huevo (…) la logística para tú y buscar algo. (…) A mí me gustaría comer mucho más saludable pero la realidad es que el bolsillo.',
                'attribution': 'Familia Biparental con Hijos Pequeños'
            },
            {
                'quote': 'El problema también con la comida saludable es el exceso porque como tú dices en la calle hay mucha comida rápida, hay muchas pizzas, mucho pollo, mucha cosa empanizada, mucha cosa frita pero a la hora de tu buscar opciones saludables son escasas y las que hay son muy caros.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — T9: El nino dicta el menu.
    # Solo-cuali — 2 verbatims en cards side by side
    # Headline: "En hogares con hijos, la canasta se disena primero para ellos.
    #  El adulto elige despues — o consume lo mismo."
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 20 — H11 Card cualitativa (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="EN HOGARES CON HIJOS, LA CANASTA SE DISEÑA PRIMERO PARA ELLOS. ",
        headline_italic="EL ADULTO ELIGE DESPUÉS — O CONSUME LO MISMO.",
        verbatims=[
            {
                'quote': 'Yo creo que un 80% de lo que tienen todos, yo pienso en qué ellos van a comer. Mi alimentación depende de lo que ellos puedan comer. Yo no voy a ir al supermercado a comprar algo pensando en mi suegra. Yo voy al supermercado y después…',
                'attribution': 'Familia Extendida'
            },
            {
                'quote': 'La hija mía costumbra comer mucha harina (…) ella es lo de espagueti, pizza, donuts, ella es amante de esa línea, entonces ahí la cosa se complica mucho en cuanto a la compra. (…) Uno se ve obligado a tener que consumir lo demás porque no se puede desperdiciar la otra parte.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H12 — T10: El fin de semana tiene su propio menu.
    # Solo-cuali — 2 verbatims en cards side by side
    # Headline: "La semana es eficiencia; el fin de semana, ritual.
    #  El hogar dominicano no come lo mismo de lunes a viernes que el sabado."
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 21 — H12 Card cualitativa (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="LA SEMANA ES EFICIENCIA; EL FIN DE SEMANA, RITUAL. ",
        headline_italic="EL HOGAR DOMINICANO NO COME LO MISMO DE LUNES A VIERNES QUE EL SÁBADO.",
        verbatims=[
            {
                'quote': 'El número 7 suelo más los sábados y domingos. Me toca de cocinar en la casa porque yo por ejemplo yo me levanto a las 5 de la mañana todos los días, yo hago comida porque me la llevo a mi trabajo. (…) El sábado es el único día que yo tengo para yo realmente me dedico a cocinar.',
                'attribution': 'Familia Mixta'
            },
            {
                'quote': 'A fines de semana hay que comer yaroa, pizza y cosas así.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "alimentacion-deck-flat.pptx")
    prs.save(output_path)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    print("\n-- Verificacion specs -------------------------------------------")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"  = {int(prs.slide_width / PX)}px x {int(prs.slide_height / PX)}px")
    ok = prs.slide_width == Emu(SLIDE_W_PX * PX) and prs.slide_height == Emu(SLIDE_H_PX * PX)
    print(f"  Slide 1920x1080 OK: {ok}")
    print(f"Stat number: {STAT_NUMBER_PT}pt Instrument Serif italic")
    print(f"Stat desc:   {STAT_DESC_PT}pt Poppins")
    print(f"Headline:    {HEADLINE_PT}pt fijo (no auto-size)")
    print(f"CV Verbatim: {CV_VERBATIM_PT}pt Instrument Serif regular")
    print(f"CV Attrib:   {CV_ATTRIB_PT}pt Poppins italic")
    print(f"Card Verb:   {CARD_VERBATIM_PT}pt Poppins")
    print(f"Card Attrib: {CARD_ATTRIB_PT}pt Poppins italic")
    print(f"Card size:   {CARD_W}x{CARD_H}pt side by side")
    print("Cards cuali: rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt")
    print("Consumer Voice: verbatim SUELTO sobre fondo negro (SIN card)")
    print("Comillas espanolas en todos los verbatims")
    print("Kerning 0 en todo el deck")
    print("NO masterslide - fondo negro plano")
    print("NO source en ningun slide (v6)")

    print("\n-- QA FLAGS --------------------------------------------------")
    print("H01 (slides 1-2):  cuanti 2 stats. 46.4%, 57.0% — decimales exactos")
    print("H02 (slides 3-4):  cuanti 2 stats. 85.0%, 1.8% — decimales exactos")
    print("H03 (slides 5-6):  cuanti 2 stats. 89.0% (D), 52.2% (AB tendencia) — caveat inline")
    print("H04 (slides 7-8):  cuanti 1 stat. 11.6pp brecha de genero")
    print("H05 (slides 9-10): cuanti 2 stats. 53.8%, 34.4% — decimales exactos")
    print("H06 (slides 11-12): cuanti 2 stats. 17.2%, <3%")
    print("H07 (slides 13-14): cuanti 2 stats. 34.2%, 15.2%")
    print("H08 (slides 15-16): cuanti 1 stat. 17.0%")
    print("H09 (slides 17-18): cuanti 2 stats. 47.2% (C+), 52.2% (AB tendencia) — caveat inline")
    print("H10 (slide 19): solo-cuali 2 verbatims side by side. Fusion H10+H12 crudo.")
    print("H11 (slide 20): solo-cuali 2 verbatims side by side. Caveat 80% en verbatim 1.")
    print("H12 (slide 21): solo-cuali 2 verbatims side by side.")
    print("Total slides verificados: 21")
    print("Verbatims en blanco (WHITE) — ningun verde lima")
    print("NO lineas decorativas entre headline y stats")
    print("NO conclusiones italic separadas")

    return output_path


if __name__ == "__main__":
    build_deck()
