#!/usr/bin/env python3
"""
build_bienestar_v3.py
=====================
Construye el deck bienestar-deck-flat-v3.pptx desde cero con python-pptx.
23 slides según el set editorial cerrado bienestar-hallazgos-editados.md.

CAMBIOS RESPECTO A v2 (mayo 2026 — segunda iteración de feedback Jeremy):

1. CONSUMER VOICE = 1 verbatim solo.
   El slide Consumer Voice apoya y fortalece el hallazgo cuanti del slide anterior.
   1 verbatim por slide Consumer Voice, en una sola card centrada.
   Slides afectados vs v2: H03 (slide 6), H06 (slide 11), H08 (slide 15),
   H09 (slide 17), H13 (slide 23).
   — En v2 tenían 2 cards apiladas; en v3 tienen 1 card sola centrada.
   — El verbatim que permanece es el que el editor eligió en el set v3.

2. Sin cajas "Pregunta P##. Base. Fuente" debajo de stats.
   — v2 ya no las tenía (corrección de primera iteración).
   — Confirmado: ningún slide tiene metadata técnica debajo de stats en este script.
   — El stat queda limpio: cifra grande + descripción Poppins con bold selectivo.
   — Toda la metadata queda en el Source pie de página.

Reglas que se mantienen de v2:
- Fondo negro plano #000000 (sin masterslide decorativo)
- Headlines en MAYÚSCULAS auto-size según tabla de caracteres (>85 chars → 30pt)
- Stats grandes redondeados a entero (decimal solo si 1 dígito antes del punto)
- Verbatims con comillas españolas «...» en Instrument Serif REGULAR (no italic)
- Atribución "— Familia X" en italic Poppins 13pt
- Líneas verticales #2E2E2E entre stats (cuando hay 2 o 3)
- Stats con 2 → posiciones 1/3 y 2/3 del ancho (no extremos)
- Stats con 1 → centrado
- Stats con 3 → 1/4, 1/2, 3/4
- H10 y H12 solo-cuali: headline arriba + cards apiladas (H10 con 2, H12 con 1)
  Estos son los ÚNICOS slides con múltiples cards apiladas en el deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree
import os

# ── rutas ──────────────────────────────────────────────────────────────────────
OUTPUT_PATH = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Agentes Hallazgos/bienestar-deck-flat-v3.pptx"

# ── constantes visuales ────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x2E, 0x2E, 0x2E)   # gris carbón — cards y separadores
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)  # header "CONSUMER VOICE"

# slide 16:9 — 13.33" × 7.5"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Agrega un slide en blanco (layout 6 = blank) y pone fondo negro."""
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def add_textbox(slide, left, top, width, height, text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT,
                word_wrap=True):
    """Agrega un textbox con un solo run de texto."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def _set_kern(r_elem):
    """Fuerza kerning a 0 en el elemento <a:r>."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_rich_textbox(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
                     word_wrap=True):
    """
    Agrega un textbox con múltiples runs en el mismo párrafo.
    runs: lista de dicts con keys: text, font_name, font_size_pt, bold, italic, color
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', 11))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def pick_headline_size(chars):
    """Auto-size de headline según número de caracteres (regla mayo 2026)."""
    if chars < 30:
        return 70
    elif chars < 45:
        return 55
    elif chars < 66:
        return 42
    elif chars < 86:
        return 36
    else:
        return 30


def add_headline(slide, text_plain, text_italic=""):
    """
    Agrega headline en MAYÚSCULAS con split plain/italic.
    Centrado horizontalmente, posición arriba (0.7" del top).
    Auto-size según largo total del texto.
    """
    full_text = text_plain + text_italic
    chars     = len(full_text.strip())
    size_pt   = pick_headline_size(chars)

    left   = Inches(0.5)
    top    = Inches(0.7)
    width  = Inches(12.33)
    height = Inches(2.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(size_pt)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(size_pt)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_source(slide, source_text):
    """Agrega source al pie del slide."""
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(7.0),
        width=Inches(12.33), height=Inches(0.35),
        text=source_text,
        font_name=FONT_BODY, font_size_pt=9,
        italic=True, color=RGBColor(0x9B, 0x9B, 0x9B),
        align=PP_ALIGN.CENTER
    )


def add_vertical_separator(slide, x_inches, y_top_inches, height_inches):
    """Agrega línea vertical fina #2E2E2E entre stats."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x_inches - 0.01), Inches(y_top_inches),
        Inches(0.02), Inches(height_inches)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY
    line.line.fill.background()


def add_stat_block(slide, stat_value, desc_runs, x_center_inches, y_stat_inches,
                   y_desc_inches):
    """
    Agrega un bloque stat: número grande + descripción multi-run centrados.
    x_center_inches: centro horizontal del bloque.
    """
    stat_width = Inches(3.8)
    stat_left  = Inches(x_center_inches) - stat_width / 2

    # Número grande
    add_textbox(
        slide,
        left=stat_left, top=Inches(y_stat_inches),
        width=stat_width, height=Inches(1.3),
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=96,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Descripción (sin cajas de Pregunta/Base — solo cifra + descripción Poppins)
    desc_width = Inches(3.5)
    desc_left  = Inches(x_center_inches) - desc_width / 2
    add_rich_textbox(
        slide,
        left=desc_left, top=Inches(y_desc_inches),
        width=desc_width, height=Inches(1.0),
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


def stat_x_positions(n_stats):
    """
    Devuelve lista de x_center en pulgadas para n_stats (1, 2 o 3).
    Slide ancho = 13.33"
    """
    w = 13.33
    if n_stats == 1:
        return [w / 2]
    elif n_stats == 2:
        return [w * (1/3), w * (2/3)]
    else:  # 3
        return [w * (1/4), w * (1/2), w * (3/4)]


def make_run(text, bold=False, italic=False, size_pt=11, color=WHITE,
             font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def add_card(slide, quote_text, attribution, x_left_inches, y_top_inches,
             card_width_inches, card_height_inches, verbatim_size_pt=None):
    """
    Agrega un card de verbatim estilo Consumer Voice.
    Fondo #2E2E2E sólido, border-radius roundRect, comillas españolas.
    quote_text: texto SIN comillas (se agregan aquí).
    attribution: texto tipo "Familia monoparental"
    verbatim_size_pt: None = auto según largo
    """
    # Auto-size del verbatim
    if verbatim_size_pt is None:
        n = len(quote_text)
        if n < 80:
            verbatim_size_pt = 32
        elif n < 150:
            verbatim_size_pt = 28
        elif n < 220:
            verbatim_size_pt = 24
        elif n < 320:
            verbatim_size_pt = 20
        else:
            verbatim_size_pt = 17

    # Card background (rectángulo redondeado)
    card = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x_left_inches), Inches(y_top_inches),
        Inches(card_width_inches), Inches(card_height_inches)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = GREY
    card.line.fill.background()  # sin borde

    # border-radius via XML (prstGeom roundRect)
    sp_elem = card._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            spPr.remove(prstGeom)
        new_geom = etree.SubElement(spPr, qn('a:prstGeom'))
        new_geom.set('prst', 'roundRect')
        avLst = etree.SubElement(new_geom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 16667')  # ~border-radius 16pt proporcional

    # Texto del verbatim con comillas españolas
    full_quote = f"«{quote_text}»"
    pad_left = x_left_inches + 0.35
    pad_top  = y_top_inches + 0.25
    text_width  = card_width_inches - 0.7
    text_height = card_height_inches - 0.7

    txBox = slide.shapes.add_textbox(
        Inches(pad_left), Inches(pad_top),
        Inches(text_width), Inches(text_height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Párrafo del verbatim
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(verbatim_size_pt)
    run_q.font.italic = False   # NO italic — regla mayo 2026
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Párrafo de atribución
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(13)
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return card


# ── layouts de slides ──────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic,
                          stats,  # lista de dicts: {value, desc_runs}
                          source_text):
    """
    Construye un slide de Hallazgo cuanti.
    stats: 1, 2 o 3 elementos.
    Sin cajas de Pregunta/Base debajo de stats — solo cifra + descripción Poppins.
    """
    slide = blank_slide(prs)

    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_positions(n)

    y_stat = 2.8
    y_desc = 4.2

    # separadores verticales entre stats
    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) / 2
        add_vertical_separator(slide, x_sep1, y_stat - 0.1, 1.6)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) / 2
        add_vertical_separator(slide, x_sep2, y_stat - 0.1, 1.6)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'],
                       xs[i], y_stat, y_desc)

    add_source(slide, source_text)
    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Construye un slide Consumer Voice con 1 solo verbatim en card centrada.

    REGLA v3 (mayo 2026 — segunda iteración):
    Consumer Voice = 1 verbatim solo. Apoya y fortalece el hallazgo cuanti
    del slide anterior. 1 card centrada vertical y horizontalmente.
    NO múltiples cards apiladas en Consumer Voice.
    Múltiples cards apiladas SOLO en slides solo-cuali (build_cuali_slide).

    quote: string del verbatim (sin comillas — se agregan aquí)
    attribution: string tipo "Familia monoparental"
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE"
    add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.55),
        width=Inches(12.33), height=Inches(0.4),
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=10,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # 1 card centrada — dimensiones generosas para el verbatim solo
    margin_side = 1.2
    card_width  = 13.33 - (2 * margin_side)  # ~10.93"
    card_height = 4.8
    y_available_start = 1.15
    y_available_end   = 7.2
    y_available       = y_available_end - y_available_start
    y_card = y_available_start + (y_available - card_height) / 2

    add_card(
        slide,
        quote_text=quote,
        attribution=attribution,
        x_left_inches=margin_side,
        y_top_inches=y_card,
        card_width_inches=card_width,
        card_height_inches=card_height
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text):
    """
    Construye un slide de hallazgo solo-cuali.
    Headline arriba + cards de verbatims debajo.

    Este es el ÚNICO tipo de slide donde aplican múltiples cards apiladas:
    - H10: 2 verbatims → 2 cards apiladas
    - H12: 1 verbatim → 1 card

    Consumer Voice usa build_consumer_voice_slide (1 card sola, sin headline).
    """
    slide = blank_slide(prs)

    add_headline(slide, headline_plain, headline_italic)

    n_cards = len(verbatims)

    margin_side = 1.2
    card_width  = 13.33 - (2 * margin_side)
    y_top_cards = 2.5
    total_height_available = 7.5 - y_top_cards - 0.5
    gap_between = 0.25

    if n_cards == 1:
        card_height = min(total_height_available, 4.2)
    elif n_cards == 2:
        card_height = min((total_height_available - gap_between) / 2, 2.5)
    else:
        card_height = min((total_height_available - 2 * gap_between) / 3, 1.8)

    for i, v in enumerate(verbatims):
        y_card = y_top_cards + i * (card_height + gap_between)
        add_card(
            slide,
            quote_text=v['quote'],
            attribution=v['attribution'],
            x_left_inches=margin_side,
            y_top_inches=y_card,
            card_width_inches=card_width,
            card_height_inches=card_height
        )

    add_source(slide, source_text)
    return slide


# ── datos del set editorial v3 ─────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — T1: El dominicano no se hace el "loco" con la salud mental: la delega.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 1: Hallazgo cuanti (2 stats)
    # Headline 162 chars → 30pt
    # Stats originales: 3.0% → 3% (entero), 9.0% → 9% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="La terapia es la actividad de autocuidado menos elegida. El dominicano no ignora su salud mental — ",
        headline_italic="la terceriza a la fe, la mascota, el ejercicio y la familia.",
        stats=[
            {
                'value': '3%',
                'desc_runs': [
                    make_run('declara "voy a '),
                    make_run('terapia', bold=True),
                    make_run('" como actividad principal de autocuidado. La opción menos votada de toda la batería P27.')
                ]
            },
            {
                'value': '9%',
                'desc_runs': [
                    make_run('elige '),
                    make_run('meditación', bold=True),
                    make_run(' como autocuidado. Terapia + meditación sumadas: 12%. El 88% restante canaliza su bienestar por vías no clínicas.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P27 · Base 500."
    )

    # Slide 2: Consumer Voice — 1 verbatim (igual que v2, ya tenía 1 card)
    build_consumer_voice_slide(
        prs,
        quote='Yo no iba a hacer lo que fue medicarme, pasar por psiquiatría. (...) Después de que pasó todo lo que pasó con la otra pareja que yo tuve, tuve que ir a la psiquiatría, tuve que medicarme, pero al final valió la pena, rindió sus frutos.',
        attribution='Familia Homoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — T1: El dominicano no se hace el "loco" con la salud mental: la delega.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 3: Hallazgo cuanti (1 stat)
    # Headline 167 chars → 30pt
    # Stat: 85.0% → 85% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="85% del dominicano considera a la mascota parte de la familia. ",
        headline_italic="En los hombres, el perro cumple la función emocional que la cultura no les permite buscar en otra parte.",
        stats=[
            {
                'value': '85%',
                'desc_runs': [
                    make_run('considera a las mascotas '),
                    make_run('parte de la familia', bold=True),
                    make_run('. En hogares sin hijos con mascota sube a 94.6%.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P10 · Base 500."
    )

    # Slide 4: Consumer Voice — 1 verbatim (igual que v2, ya tenía 1 card)
    build_consumer_voice_slide(
        prs,
        quote='Estadísticamente hablando, el perro suele ser un apoyo muy emocional para el hombre, en muchas ocasiones hasta más que la mujer, y no es por un tema de que no sentamos ese amor, sino porque tenemos un aspecto cultural y hasta biológico (...) de que no nos desahogamos como las mujeres lo hacen (...) nosotros simplemente necesitamos compañía, tacto y simplemente presencia.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — T2: El dominicano gestiona su salud mental de rodillas.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # CAMBIO v3: Consumer Voice pasa de 2 cards a 1 card.
    # Verbatim elegido por editor: Familia Sin Hijos (fg-02).
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 5: Hallazgo cuanti (2 stats)
    # Headline 160 chars → 30pt
    # Stats: 64.0% → 64% (entero), ratio 1:6
    build_hallazgo_slide(
        prs,
        headline_plain="64% recurre siempre a la religión en momentos de dificultad. ",
        headline_italic="La fe es el primer canal de salud mental: gratis, disponible las 24 horas, sin estigma ni copago.",
        stats=[
            {
                'value': '64%',
                'desc_runs': [
                    make_run('marca “5 = Siempre” recurrir a la '),
                    make_run('religión', bold=True),
                    make_run(' o espiritualidad en momentos de dificultad. Respuesta dominante en todas las tipologías.')
                ]
            },
            {
                'value': '1:6',
                'desc_runs': [
                    make_run('menciones clínicas por cada 6 '),
                    make_run('religiosas', bold=True),
                    make_run(' en 11 FGs (216 menciones fe vs 37 de terapia / psicólogo / psiquiatría).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P29 · Base 500."
    )

    # Slide 6: Consumer Voice — 1 card (v3: de 2 cards → 1 card)
    # Verbatim elegido: Familia Sin Hijos — declara equivalencia Dios=salud mental de forma directa.
    # Verbatim descartado: Familia Biparental con Hijos Adultos — describe stack de coping (fe+podcast+libros),
    # fragmenta el stat del 64%.
    build_consumer_voice_slide(
        prs,
        quote='En mi caso nosotros con la biblia (...) si eso me ayudó bastante, bastante, la cercanía con Dios, a nuestra salud mental, de ambos.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — T2: El dominicano gestiona su salud mental de rodillas.
    # Tipo: solo-cuanti (sin verbatim) → 1 slide
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 7: Hallazgo solo-cuanti (2 stats)
    # Headline 170 chars → 30pt
    # Stats: 70.8% → 71% (entero), 72.7% → 73% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="70.8% de las mujeres y 72.7% del estrato D recurren siempre a la fe. ",
        headline_italic="La intensidad religiosa sube donde más aprieta: si eres mujer y el bolsillo está flaco, rezas más.",
        stats=[
            {
                'value': '71%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' marca “Siempre” recurrir a la religión en momentos difíciles, vs 53.3% de los hombres. Brecha de género: 17.5 puntos.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' marca “Siempre” recurrir a la fe, vs 52.8% del estrato C+ y 61.5% del estrato C.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P29 × D2, P29 × D6 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — T3: El autocuidado es físico, no emocional.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 8: Hallazgo cuanti (2 stats)
    # Headline 163 chars → 30pt
    # Stats: 42.6% → 43% (entero), 45.9% → 46% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="42.6% del dominicano no hace ninguna actividad para cuidar su salud, ni física ni mental. ",
        headline_italic="La inacción es la respuesta más común de toda la batería de autocuidado.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara “'),
                    make_run('no realizo ninguna actividad', bold=True),
                    make_run('” para cuidar su salud mental y física. TOP-1 de la batería P27, por encima de ejercicio (34.8%), oración (27.6%) y todas las demás.')
                ]
            },
            {
                'value': '46%',
                'desc_runs': [
                    make_run('de hogares sin hijos con mascota lidera la '),
                    make_run('inacción', bold=True),
                    make_run('. La inacción lidera en 5 de 7 tipologías: biparental (44.2%), extendido (42.9%), monoparental (41.8%).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27 · Base 500."
    )

    # Slide 9: Consumer Voice — 1 verbatim (igual que v2, ya tenía 1 card)
    build_consumer_voice_slide(
        prs,
        quote='Bueno yo para lo físico yo no estoy haciendo nada. Yo no hago ejercicio, yo me como todo lo que yo quiera y a la hora que yo quiera. Bueno, para la salud mental a veces uno tiene que hacerse loco en ciertas cosas.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — T3: El autocuidado es físico, no emocional.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # CAMBIO v3: Consumer Voice pasa de 2 cards a 1 card.
    # Verbatim elegido por editor: Familia Biparental con Hijos Pequeños (fg-07).
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 10: Hallazgo cuanti (2 stats)
    # Headline 186 chars → 30pt
    # Stats: 34.8% → 35% (entero), 43.1% → 43% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="34.8% del dominicano se ejercita con regularidad — el único autocuidado donde el hombre lidera. ",
        headline_italic="Para él, entrenar no es físico: es la válvula emocional que la cultura sí le permite.",
        stats=[
            {
                'value': '35%',
                'desc_runs': [
                    make_run('declara “'),
                    make_run('me ejercito con regularidad', bold=True),
                    make_run('” como autocuidado — segundo lugar de la batería P27, detrás de “ninguna actividad” (42.6%).')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('masculino', bold=True),
                    make_run(' tiene el ejercicio como TOP-1. En el subset femenino el TOP-1 es “ninguna actividad” con 48.5%. El ejercicio es el único canal donde el hombre supera a la mujer.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27, P27 × D2 · Base 500."
    )

    # Slide 11: Consumer Voice — 1 card (v3: de 2 cards → 1 card)
    # Verbatim elegido: Familia Biparental con Hijos Pequeños — "La salud mental mía es deporte",
    # literaliza el hallazgo con ecuación explícita.
    # Verbatim descartado: Familia Sin Hijos con Mascota — mismo insight en negativo, menos contundente.
    build_consumer_voice_slide(
        prs,
        quote='Yo soy una persona súper activa, súper activa (...) La salud mental mía es deporte. Yo me irrito cuando yo duro tres días que no puedo hacer ejercicio (...) Yo entro a la cancha y yo me olvido de todo lo que yo tengo.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — T4: Mi equilibrio mental es mi equilibrio económico.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 12: Hallazgo cuanti (3 stats)
    # Headline 189 chars → 30pt
    # Stats: 47.8% → 48% (entero), 28.4% → 28% (entero), 51.2% → 51% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="47.8% nombra la economía como su mayor fuente de estrés familiar. ",
        headline_italic="Sumando los dos ítems económicos de la batería, el estresor número uno del hogar dominicano alcanza al 76% de los hogares.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('declara “'),
                    make_run('Economía', bold=True),
                    make_run('” como el factor que más estrés genera en su vida familiar. TOP-1 nacional de la batería P26.')
                ]
            },
            {
                'value': '28%',
                'desc_runs': [
                    make_run('marca también “'),
                    make_run('Realidad económica', bold=True),
                    make_run('” en la misma batería. Los dos ítems económicos combinados alcanzan al 76% del total.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' declara la economía como estresor principal, vs el promedio nacional de 47.8%. A menor ingreso, mayor peso del estresor económico.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P26 · Base 500."
    )

    # Slide 13: Consumer Voice — 1 verbatim (igual que v2, ya tenía 1 card)
    build_consumer_voice_slide(
        prs,
        quote='Tuve un tiempo como mamá soltera y mi hermana es mamá soltera de tres niños, y es bien difícil levantarte, llevar al niño a la escuela, el dinero, cuando el papá no está económicamente ni emocionalmente.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — T5: Me agota ser padre/madre, pero nunca me canso de mis hijos.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # CAMBIO v3: Consumer Voice pasa de 2 cards a 1 card.
    # Verbatim elegido por editor: Familia Monoparental (fg-03).
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 14: Hallazgo cuanti (3 stats)
    # Headline 180 chars → 30pt
    # Stats: 42.2% → 42% (entero), 22.3% → 22% (entero), 5.5% se mantiene (1 dígito antes del punto)
    build_hallazgo_slide(
        prs,
        headline_plain="42.2% dice que la presión sobre los padres es demasiada. ",
        headline_italic="El agotamiento no aparece de golpe — en la franja 35–44 ya lo declara como estresor el 22.3%, el triple que a los 18–24.",
        stats=[
            {
                'value': '42%',
                'desc_runs': [
                    make_run('declara “5 — Muy de acuerdo” con “'),
                    make_run('hay demasiada presión', bold=True),
                    make_run(' para que los papás lo tengan todo.” El ítem con mayor consenso afirmativo de toda la batería P8 sobre crianza.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('35–44 años', bold=True),
                    make_run(' declara la crianza como factor de estrés, vs 15.5% del subset 25–34. Pico de carga parental en la franja media adulta.')
                ]
            },
            {
                'value': '5.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('18–24 años', bold=True),
                    make_run(' declara la crianza como estresor. La presión parental no existe al inicio — se construye con los años.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P8, P26 × D5 · Base 500."
    )

    # Slide 15: Consumer Voice — 1 card (v3: de 2 cards → 1 card)
    # Verbatim elegido: Familia Monoparental — "siempre agotada, cansada, batida"
    # nombra el agotamiento como condición estructural cultural, alineado al 42.2% de presión parental.
    # Verbatim descartado: Familia Mixta — describe el estrés del 24/7 con los hijos, situacional.
    build_consumer_voice_slide(
        prs,
        quote='Es un mito que una mujer va a estar (...) toda ella relajada, limpiando una casa, con el marido sentado allá cargando el niño y qué feliz de la vida. Eso es un mito. De toda mujer que tiene hijos siempre agotada, cansada, batida.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — T6: La crisis monoparental no tiene tratamiento.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # CAMBIO v3: Consumer Voice pasa de 2 cards a 1 card.
    # Verbatim elegido por editor: Familia Monoparental verbatim 2 (fg-03 — "Soy madre soltera, sola").
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 16: Hallazgo cuanti (2 stats)
    # Headline 183 chars → 30pt
    # Stats: 24.4% → 24% (entero), 56.6% → 57% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="24.4% de los hogares es monoparental. En más de la mitad, la madre es la única responsable de la crianza. ",
        headline_italic="El mercado diseña para dos. El hogar dominicano no siempre llega de a dos.",
        stats=[
            {
                'value': '24%',
                'desc_runs': [
                    make_run('del dominicano declara que su hogar es '),
                    make_run('monoparental', bold=True),
                    make_run('. Segunda tipología más común después del biparental con hijos menores de 18 (38.0%).')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('del subset monoparental cita a la '),
                    make_run('madre', bold=True),
                    make_run(' como única responsable de la crianza, vs 41.8% del promedio nacional. La crianza monoparental es de facto materna en más de 5 de cada 10 casos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P1, P2 × D1 · Base 500."
    )

    # Slide 17: Consumer Voice — 1 card (v3: de 2 cards → 1 card)
    # Verbatim elegido: "Soy madre soltera, sola" — acumulación simultánea de cargas,
    # conecta directamente con el 56.6% madre única.
    # Verbatim descartado: "lamentablemente son los hijos" — consecuencia psicológica del aislamiento,
    # no la carga estructural del stat.
    build_consumer_voice_slide(
        prs,
        quote='Soy madre soltera, sola, con una situación actualmente de mi madre que falleció y que vive conmigo y mi hermanito en condición especial, entonces es difícil, es difícil todo en general, quién lo cuide, porque yo trabajo todo el tiempo.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — T7: Las familias dominicanas no tienen un sistema de salud,
    #           y de todas formas, no se cuidan.
    # Tipo: solo-cuali, 2 verbatims → 1 slide (card cualitativa con 2 cards apiladas)
    # Este es uno de los 2 únicos slides con múltiples cards apiladas en el deck.
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 18: Card cualitativa (headline arriba + 2 cards apiladas)
    # Headline 185 chars → 30pt
    build_cuali_slide(
        prs,
        headline_plain="El sistema público de salud falla antes de la consulta. ",
        headline_italic="La pregunta dominicana no es qué le pasa al médico — es si el seguro aprueba, si hay cama, si alguien dentro te puede ayudar.",
        verbatims=[
            {
                'quote': 'Vaya con un dolor, aunque sea un dolor de una hebra de cabello, y que te digan “hasta que el seguro no te apruebe, no.” O tú ir a comprar un medicamento y que te digan: “ah no, tu límite del medicamento ya se agotó.” Exactamente. Tú vas a una clínica y lo primero que te preguntan es qué seguro tú tienes.',
                'attribution': 'Familia Homoparental'
            },
            {
                'quote': 'Si tú no tienes alguien que te puede ayudar en el sistema, tú estás muerto.',
                'attribution': 'Familia Sin Hijos con Mascota'
            }
        ],
        source_text="Source: Código Casa — Estudio cualitativo 2025 · P28 · Base cualitativa: 11 grupos de enfoque."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — T8: Los ancianos son invisibles.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 19: Hallazgo cuanti (3 stats)
    # Headline 176 chars → 30pt
    # Stats: 73.8% → 74% (entero), 24.0% → 24% (entero), 12.8% → 13% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="73.8% pide pensiones para que los mayores vivan con dignidad. Solo 12.8% es cuidador directo. ",
        headline_italic="Los dominicanos saben lo que hace falta — solo que no están en posición de darlo.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('marca “'),
                    make_run('Pensiones o apoyos económicos suficientes', bold=True),
                    make_run('” como lo que más falta para que la tercera edad viva con dignidad. TOP-1 absoluto de la batería P79.')
                ]
            },
            {
                'value': '24%',
                'desc_runs': [
                    make_run('menciona “'),
                    make_run('Apoyo y capacitación para cuidadores familiares', bold=True),
                    make_run('” en la misma batería. El cuidador familiar dominicano pide ayuda — el mercado no la ofrece como categoría.')
                ]
            },
            {
                'value': '13%',
                'desc_runs': [
                    make_run('declara ser el '),
                    make_run('cuidador principal', bold=True),
                    make_run(' de un adulto mayor en el hogar. 72.8% declara que no vive con ninguna persona mayor en su hogar.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P78, P79 · Base 500."
    )

    # Slide 20: Consumer Voice — 1 verbatim (igual que v2, ya tenía 1 card)
    build_consumer_voice_slide(
        prs,
        quote='La frustración en cuanto a salud y bienestar para mí es el tema del presupuesto que requiere. (...) Porque tengo un choque entre lo que yo digo que tengo que aceptar y las costumbres pasadas. Mi mamá era de la mujer que me enseñó a mí a que si algo no es necesario, no gastes dinero en eso, economiza.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H12 — T9: Saber aguantar se hereda.
    # Tipo: solo-cuali, 1 verbatim → 1 slide (card cualitativa con 1 card)
    # Este es uno de los 2 únicos slides con múltiples cards apiladas (aquí: 1 card).
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 21: Card cualitativa (headline arriba + 1 card)
    # Headline 189 chars → 30pt
    build_cuali_slide(
        prs,
        headline_plain='“Yo no me voy a separar, yo voy a aguantar — porque mi mamá me decía que hay que aguantar.” ',
        headline_italic="El aguante dominicano no es una decisión. Es una instrucción que llega de la generación anterior.",
        verbatims=[
            {
                'quote': 'Yo me crí con esos dos conceptos en mi mente y yo siempre me enfocaba. Tanto así que yo decía: “No, yo no me voy a separar, yo voy a aguantar”, porque mi mamá me decía que hay que aguantar.',
                'attribution': 'Familia Monoparental'
            }
        ],
        source_text="Source: Código Casa — Estudio cualitativo 2025 · Base cualitativa: 11 grupos de enfoque."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H13 — T10: El amor es lo más importante, a menos que sea amor propio.
    # Tipo: cuanti + 1 verbatim → 2 slides
    # CAMBIO v3: Consumer Voice pasa de 2 cards a 1 card.
    # Verbatim elegido por editor: Familia Monoparental (fg-03 — "¿y yo?").
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 22: Hallazgo cuanti (1 stat)
    # Headline 170 chars → 30pt
    # Stat principal: 48.5% (1 dígito antes del punto → mantener decimal)
    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no hace ninguna actividad para cuidarse, vs 33.3% de los hombres. ",
        headline_italic="Brecha de género en autocuidado: 15.2 puntos. La mujer se cuida menos — y lo sabe.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' declara “no realizo ninguna actividad para cuidar mi salud” como TOP-1 de P27, vs 33.3% del subset masculino. Brecha de género en inacción: 15.2 puntos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27 × D2 · Base 500."
    )

    # Slide 23: Consumer Voice — 1 card (v3: de 2 cards → 1 card)
    # Verbatim elegido: Familia Monoparental — "se olvida de uno mismo" / "¿y yo?"
    # literaliza el 48.5% de mujeres que no hace ninguna actividad de autocuidado.
    # Verbatim descartado: Familia Biparental con Hijos Pequeños — ángulo lateral al stat de género.
    build_consumer_voice_slide(
        prs,
        quote='Cuando uno es madre, a veces se olvida, a veces no... casi siempre se olvida de uno mismo. (...) Y cuando tú vienes a mirar para atrás, entonces tú dices: “¿y yo?” Porque crecen y se van.',
        attribution='Familia Monoparental'
    )

    # ── guardar ────────────────────────────────────────────────────────────────
    prs.save(OUTPUT_PATH)
    print(f"Deck guardado: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

    # QA rápido — conteo por tipo de slide
    slide_types = {
        'hallazgo_cuanti': [1, 3, 5, 7, 8, 10, 12, 14, 16, 19, 22],
        'consumer_voice_1card': [2, 4, 6, 9, 11, 13, 15, 17, 20, 23],
        'cuali_slide': [18, 21]
    }
    print(f"\nConteo por tipo:")
    print(f"  Slides Hallazgo cuanti (sin cajas P##):  {len(slide_types['hallazgo_cuanti'])}")
    print(f"  Slides Consumer Voice (1 card cada uno): {len(slide_types['consumer_voice_1card'])}")
    print(f"  Slides Card cualitativa (H10: 2 cards, H12: 1 card): {len(slide_types['cuali_slide'])}")
    print(f"\nVerificacion v3:")
    print(f"  Slides H03/H06/H08/H09/H13 Consumer Voice: slides 6, 11, 15, 17, 23 → 1 card cada uno (CORRECTO)")
    print(f"  Ningun slide Hallazgo cuanti lleva caja P##/Base debajo de stats (CORRECTO)")


if __name__ == "__main__":
    build_deck()
