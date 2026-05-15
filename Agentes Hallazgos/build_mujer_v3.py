#!/usr/bin/env python3
"""
build_mujer_v3.py
=================
Deck mujer-v3-deck-flat.pptx desde cero con python-pptx.
Pilar 11 — MUJER v3 — 25 hallazgos, fuentes duales CC + Finance Women (FW).
Dos bloques: BLOQUE A (Código Casa, A01-A08) + BLOQUE B (Finance Women, B01-B17).

Specs v6 activas (aprendizajes-montador-cc.md § 3.9 v6):
- Slide: 1920×1080 px → 18_288_000 × 10_287_000 EMU
- HEADLINE_PT = 50        fijo, no auto-size
- STAT_DESC_PT = 13       13pt fijo (v5)
- STAT_BOX_W = 378        10cm fijo
- STAT_BOX_H = 113        3cm fijo
- CV_VERBATIM_PT = 50     v6: era 60
- CARD_W = 567            v6: 15cm — era 549
- CARD_H = 227            v6: 6cm — era 221
- Cifras: Instrument Serif italic 180pt. Redondeo a entero por defecto;
  decimal solo si 1 dígito antes (8.7%) o brecha matemática (15.2 pts).
- Consumer Voice: verbatim SUELTO sin card. Header "CONSUMER VOICE" Poppins 14pt gris.
- Cards cualitativas: fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt.
- Source: Poppins italic 12pt, gris GREY_SOFT, centrado, top 1010pt.
- NUNCA mencionar "Banreservas" en ningún elemento.
- Fondo negro plano — sin masterslide decorativo.
- Kerning 0 en todo el deck.

Convención de unidades:
  1 pt Keynote = 1 px @ 96 dpi en python-pptx
  1 px @ 96 dpi = 9525 EMU
  Se define PX = 9525 y se usa como multiplicador.

Bloques:
  BLOQUE A — Código Casa (CC): portada + A01-A08 → 15 slides
  BLOQUE B — Finance Women (FW): portada + B01-B17 → 33 slides
  TOTAL: 48 slides
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Conversión px/pt → EMU ─────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU (también 1 pt Keynote = 9525 EMU)

def px(n):
    """Convierte píxeles/puntos Keynote a EMU."""
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W = 1920   # pt (equiv px @ 96dpi)
SLIDE_H = 1080   # pt

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)   # separadores verticales entre stats
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)   # header CONSUMER VOICE + source

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── Tamaños en pt (v6 activa — specs exactas de aprendizajes-montador-cc.md) ──
HEADLINE_PT       = 50    # headlines — fijo, no auto-size
STAT_NUMBER_PT    = 180   # cifra dominante Instrument Serif italic
STAT_DESC_PT      = 13    # descripción stat Poppins (v5: era 16pt, ahora 13pt)
CV_VERBATIM_PT    = 50    # v6: Consumer Voice verbatim Instrument Serif (era 60, ahora 50)
CV_HEADER_PT      = 14    # "CONSUMER VOICE" Poppins regular
CV_ATTRIBUTION_PT = 20    # atribución Consumer Voice Poppins italic
CARD_VERBATIM_PT  = 15    # verbatim dentro de card cualitativa Poppins
CARD_ATTRIBUTION_PT = 12  # atribución card cualitativa Poppins italic
SOURCE_PT         = 12    # source line: Poppins italic, GREY_SOFT, centrado

# ── Dimensiones de stat box (v5: 378×113pt — 10×3cm) ─────────────────────────
STAT_BOX_W = 378   # pt = 10 cm
STAT_BOX_H = 113   # pt = 3 cm

# ── Dimensiones de card cualitativa (v6: 567×227pt — 15×6cm) ──────────────────
CARD_W = 567   # pt = 15 cm
CARD_H = 227   # pt = 6 cm

# ── Posiciones verticales clave ────────────────────────────────────────────────
HEADLINE_TOP   = 100   # pt — headline arriba
STAT_NUM_TOP   = 420   # pt — cifra grande
STAT_DESC_TOP  = 660   # pt — caja de descripción debajo de cifra
CV_HEADER_TOP  = 80    # pt — header "CONSUMER VOICE"
CV_VERBATIM_TOP = 350  # pt — inicio del bloque de verbatim Consumer Voice


# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    """Presentation con slide 1920×1080 pt."""
    prs = Presentation()
    prs.slide_width  = px(SLIDE_W)
    prs.slide_height = px(SLIDE_H)
    return prs


def blank_slide(prs):
    """Slide en blanco con fondo negro plano."""
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en <a:r>."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left, top, width, height,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox simple de un solo run, posicionado en pt."""
    txBox = slide.shapes.add_textbox(
        px(left), px(top), px(width), px(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_source(slide, source_text):
    """
    Línea de source en la parte inferior del slide.
    Poppins italic SOURCE_PT (12pt), GREY_SOFT, centrado, top 1010pt.
    Solo se llama en slides de hallazgo — NO en portadas de bloque.
    NUNCA incluir "Banreservas".
    """
    add_textbox(
        slide,
        left=0, top=1010,
        width=1920, height=50,
        text=source_text,
        font_name=FONT_BODY, font_size_pt=SOURCE_PT,
        italic=True, color=GREY_SOFT,
        align=PP_ALIGN.CENTER
    )


def add_rich_textbox(slide, left, top, width, height,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox con múltiples runs en el mismo párrafo."""
    txBox = slide.shapes.add_textbox(
        px(left), px(top), px(width), px(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """
    Headline MAYÚSCULAS, 50pt fijo, centrado.
    left 100, top 100, width 1720pt.
    Split plain/italic para la frase del giro.
    """
    txBox = slide.shapes.add_textbox(
        px(100), px(HEADLINE_TOP), px(1720), px(300)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(HEADLINE_PT)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(HEADLINE_PT)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_pt, y_top_pt, height_pt):
    """Línea vertical fina #2E2E2E (2pt de ancho) entre stats."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(x_pt - 1), px(y_top_pt),
        px(2), px(height_pt)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left, top, width, height):
    """
    Card ROUNDED_RECTANGLE con fill rgba(0,0,0,0.45) + outline blanco 1pt.
    border-radius: adjustment 0.06 ≈ 15pt en card de ~227pt de alto.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(left), px(top), px(width), px(height)
    )
    shape.adjustments[0] = 0.06  # ~15pt border-radius

    # Outline blanco 1pt
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    # Fill negro sólido base
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    # Opacidad 45% vía XML
    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')  # 45% en notación OOXML

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_positions(n_stats):
    """
    Posiciones X del lado izquierdo de cada caja de stat.
    v5: cajas 378pt fijas, gaps según número de stats.
    3 stats: total 378*3 + 100*2 = 1334pt; left start = (1920-1334)/2 = 293pt
    2 stats: total 378*2 + 200 = 956pt; left start = (1920-956)/2 = 482pt
    1 stat: left = (1920-378)/2 = 771pt
    """
    if n_stats == 1:
        return [(SLIDE_W - STAT_BOX_W) // 2]
    elif n_stats == 2:
        gap = 200
        total = STAT_BOX_W * 2 + gap
        start = (SLIDE_W - total) // 2
        return [start, start + STAT_BOX_W + gap]
    else:  # 3
        gap = 100
        total = STAT_BOX_W * 3 + gap * 2
        start = (SLIDE_W - total) // 2
        return [start, start + STAT_BOX_W + gap, start + (STAT_BOX_W + gap) * 2]


def add_stat_block(slide, stat_value, desc_runs, left_pt, n_stats):
    """
    Bloque stat: cifra 180pt + caja descripción 378×113pt (13pt).
    left_pt: coordenada izquierda de la caja de descripción.
    La cifra se centra sobre la caja.
    Si hay 3 stats, reduce cifra a 150pt para evitar solapamiento.
    """
    num_pt = STAT_NUMBER_PT if n_stats < 3 else 150

    # Cifra grande — centrada horizontalmente sobre la caja
    cifra_w = 500   # pt — ancho suficiente para cifra centrada
    cifra_left = left_pt + STAT_BOX_W // 2 - cifra_w // 2

    add_textbox(
        slide,
        left=cifra_left, top=STAT_NUM_TOP,
        width=cifra_w, height=200,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=num_pt,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Descripción — 378×113pt, Poppins 13pt
    add_rich_textbox(
        slide,
        left=left_pt, top=STAT_DESC_TOP,
        width=STAT_BOX_W, height=STAT_BOX_H,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


def build_hallazgo_slide(prs, headline_plain, headline_italic, stats, source_text=""):
    """
    Slide Hallazgo cuanti.
    stats: lista de dicts {value, desc_runs} — 1, 2 o 3 elementos.
    source_text: si se pasa, se agrega en top 1010pt (Poppins italic 12pt gris).
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_positions(n)

    # Separadores verticales finos entre stats
    sep_top    = STAT_NUM_TOP - 20   # 400pt
    sep_height = 380                  # pt — cubre cifra + descripción

    if n >= 2:
        x_sep1 = xs[0] + STAT_BOX_W + (xs[1] - xs[0] - STAT_BOX_W) // 2
        add_vertical_separator(slide, x_sep1, sep_top, sep_height)
    if n == 3:
        x_sep2 = xs[1] + STAT_BOX_W + (xs[2] - xs[1] - STAT_BOX_W) // 2
        add_vertical_separator(slide, x_sep2, sep_top, sep_height)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    if source_text:
        add_source(slide, source_text)

    return slide


def build_consumer_voice_slide(prs, quote, attribution, source_text=""):
    """
    Slide Consumer Voice — verbatim SUELTO sobre fondo negro (sin card).
    v6: NO card. Solo texto suelto centrado.
    Header "CONSUMER VOICE" arriba, Poppins 14pt gris.
    Verbatim: Instrument Serif regular 50pt, comillas «...».
    Atribución: Poppins italic 20pt.
    source_text: si se pasa, se agrega en top 1010pt.
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE"
    add_textbox(
        slide,
        left=0, top=CV_HEADER_TOP,
        width=SLIDE_W, height=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Verbatim suelto — ancho 1637pt (left 141, width 1637)
    full_quote = f"«{quote}»"
    verbatim_left  = 141
    verbatim_width = 1637
    verbatim_height = 500  # suficiente para multi-línea

    add_textbox(
        slide,
        left=verbatim_left, top=CV_VERBATIM_TOP,
        width=verbatim_width, height=verbatim_height,
        text=full_quote,
        font_name=FONT_HEADLINE, font_size_pt=CV_VERBATIM_PT,
        italic=False, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Atribución — debajo del verbatim (~870pt)
    add_textbox(
        slide,
        left=0, top=870,
        width=SLIDE_W, height=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIBUTION_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    if source_text:
        add_source(slide, source_text)

    return slide


def add_cuali_card(slide, quote_text, attribution, left, top, card_w, card_h):
    """
    Card cualitativa: ROUNDED_RECTANGLE rgba(0,0,0,0.45) + outline blanco 1pt.
    Texto interno: Poppins 15pt + atribución Poppins italic 12pt.
    Padding: 30pt lados, 25pt arriba/abajo.
    """
    make_card_rgba_45_with_outline(slide, left, top, card_w, card_h)

    # Padding interno
    pad_h = 30
    pad_v = 25
    text_w = card_w - (pad_h * 2)
    text_h = card_h - (pad_v * 2)

    full_quote = f"«{quote_text}»"

    txBox = slide.shapes.add_textbox(
        px(left + pad_h), px(top + pad_v),
        px(text_w), px(text_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    # Verbatim — Poppins 15pt
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_BODY
    run_q.font.size   = Pt(CARD_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Atribución — Poppins italic 12pt
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(CARD_ATTRIBUTION_PT)
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text=""):
    """
    Slide card cualitativa — headline arriba + 1-3 cards side by side.
    Layout 5: posiciones canónicas según número de verbatims.
    source_text: si se pasa, se agrega en top 1010pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)
    card_top = 429   # pt — posición Y canónica

    if n == 1:
        left = (SLIDE_W - CARD_W) // 2
        add_cuali_card(slide, verbatims[0]['quote'], verbatims[0]['attribution'],
                       left, card_top, CARD_W, CARD_H)
    elif n == 2:
        gap = 80
        total = CARD_W * 2 + gap
        start = (SLIDE_W - total) // 2
        for i, v in enumerate(verbatims):
            x = start + i * (CARD_W + gap)
            add_cuali_card(slide, v['quote'], v['attribution'],
                           x, card_top, CARD_W, CARD_H)
    else:  # 3
        gap = 36
        total = CARD_W * 3 + gap * 2
        start = (SLIDE_W - total) // 2
        for i, v in enumerate(verbatims):
            x = start + i * (CARD_W + gap)
            add_cuali_card(slide, v['quote'], v['attribution'],
                           x, card_top, CARD_W, CARD_H)

    if source_text:
        add_source(slide, source_text)

    return slide


def build_block_cover(prs, block_label, block_title, block_subtitle=""):
    """
    Slide portada de bloque — SIN source.
    block_label: etiqueta arriba (ej. "BLOQUE A", "BLOQUE B")
    block_title: título principal centrado
    block_subtitle: subtítulo opcional en gris
    """
    slide = blank_slide(prs)

    # Etiqueta del bloque arriba
    add_textbox(
        slide,
        left=0, top=200,
        width=SLIDE_W, height=60,
        text=block_label,
        font_name=FONT_BODY, font_size_pt=16,
        bold=False, italic=False,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Título principal del bloque
    add_textbox(
        slide,
        left=100, top=300,
        width=1720, height=300,
        text=block_title.upper(),
        font_name=FONT_HEADLINE, font_size_pt=HEADLINE_PT,
        bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.CENTER,
        word_wrap=True
    )

    # Subtítulo opcional
    if block_subtitle:
        add_textbox(
            slide,
            left=200, top=680,
            width=1520, height=120,
            text=block_subtitle,
            font_name=FONT_BODY, font_size_pt=18,
            italic=True, color=GREY_SOFT,
            align=PP_ALIGN.CENTER,
            word_wrap=True
        )

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SOURCE STRINGS — mapeados por hallazgo
# CC: "Source: Código Casa — Pilar Mujer · QXXX · Mujeres n=305."
# FW: "Source: Finance Women — Estudio cuantitativo 2025 · MXX · Mujeres n=240."
# NUNCA "Banreservas" en ningún source.
# ═══════════════════════════════════════════════════════════════════════════════

# BLOQUE A — Código Casa
SRC_A01 = "Source: Código Casa — Pilar Mujer · Q054 · Mujeres n=305."
SRC_A02 = "Source: Código Casa — Pilar Mujer · Q056 · Mujeres n=305."
SRC_A03 = "Source: Código Casa — Pilar Mujer · Q057, Q056 · Mujeres n=305."
SRC_A04 = "Source: Código Casa — Pilar Mujer · Q058 · Mujeres n=305."
SRC_A05 = "Source: Código Casa — Pilar Mujer · Q059 · Mujeres n=305."
SRC_A06 = "Source: Código Casa — Pilar Mujer · Q060 · Mujeres n=305."
SRC_A07 = "Source: Código Casa — Pilar Mujer · Q055 · Mujeres n=305."
SRC_A08 = "Source: Código Casa — Pilar Mujer · Q053 · Mujeres n=305."

# BLOQUE B — Finance Women
SRC_B01 = "Source: Finance Women — Estudio cuantitativo 2025 · M20, M21 · Mujeres n=240."
SRC_B02 = "Source: Finance Women — Estudio cuantitativo 2025 · M23, M24, M54 · Mujeres n=240."
SRC_B03 = "Source: Finance Women — Estudio cuantitativo 2025 · M16, M28, M32, M60 · Mujeres n=240."
SRC_B04 = "Source: Finance Women — Estudio cuantitativo 2025 · M45, M45b · Mujeres n=240."
SRC_B05 = "Source: Finance Women — Estudio cuantitativo 2025 · M19, M43, M44 · Mujeres n=240."
SRC_B06 = "Source: Finance Women — Estudio cuantitativo 2025 · M33, M34, M36, M37 · Mujeres n=240."
SRC_B07 = "Source: Finance Women — Estudio cuantitativo 2025 · M38, M39, M40, M61 · Mujeres n=240."
SRC_B08 = "Source: Finance Women — Estudio cuantitativo 2025 · M5, M8, M17 · Mujeres n=240."
SRC_B09 = "Source: Finance Women — Estudio cuantitativo 2025 · M11, M12, M14 · Mujeres n=240."
SRC_B10 = "Source: Finance Women — Estudio cuantitativo 2025 · M15, M18, M29 · Mujeres n=240."
SRC_B11 = "Source: Finance Women — Estudio cuantitativo 2025 · M22, M39 · Mujeres n=240."
SRC_B12 = "Source: Finance Women — Estudio cuantitativo 2025 · M30, M31 · Mujeres n=240."
SRC_B13 = "Source: Finance Women — Estudio cuantitativo 2025 · M17 · Mujeres n=240."
SRC_B14 = "Source: Finance Women — Estudio cuantitativo 2025 · M11, M46 · Mujeres n=240."
SRC_B15 = "Source: Finance Women — Estudio cuantitativo 2025 · M52, M56 · Mujeres n=240."
SRC_B16 = "Source: Finance Women — Estudio cuantitativo 2025 · M57 · Mujeres n=240."
SRC_B17 = "Source: Finance Women — Estudio cuantitativo 2025 · M62, M63 · Mujeres n=240."


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD DECK
# ═══════════════════════════════════════════════════════════════════════════════

def build_deck():
    prs = new_prs()

    # ═══════════════════════════════════════════════════════════════════════════
    # PORTADA — BLOQUE A · Código Casa
    # ═══════════════════════════════════════════════════════════════════════════

    build_block_cover(
        prs,
        block_label="BLOQUE A",
        block_title="Pilar Mujer — Código Casa",
        block_subtitle="Roles de género, crianza, trabajo doméstico y autocuidado. Subsample mujer n≈305 sobre n=500."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A01 — CRIANZA ASIGNADA
    # Headline: "49.5% de las mujeres dice que la madre es la principal responsable
    #            de la crianza. Solo el hombre distribuye la carga — ella ya la tiene asignada."
    # 3 stats + 1 verbatim (cuanti + Consumer Voice)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="49.5% de las mujeres dice que la madre es la principal responsable de la crianza. ",
        headline_italic="Solo el hombre distribuye la carga — ella ya la tiene asignada.",
        stats=[
            {
                'value': '42%',
                'desc_runs': [
                    make_run('del total cree que la mayor responsabilidad de la crianza recae en la '),
                    make_run('madre', bold=True),
                    make_run('; 41.2% dice "ambos padres" y solo 15.8% se\xf1ala al padre.')
                ]
            },
            {
                'value': '49.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "Madre" como principal responsable, frente a 45.6% del masculino que responde "Ambos padres".')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('de los hogares '),
                    make_run('monoparentales', bold=True),
                    make_run(' y 59.4% del estrato E responden "Madre". La carga se concentra donde no hay con qui\xe9n repartirla.')
                ]
            }
        ],
        source_text=SRC_A01
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo tengo un pap\xe1 que es el pap\xe1 bueno, el pap\xe1 consentidor, el pap\xe1 que dice vete suave, vete en paz. Pero no es el que est\xe1 con ella todos los d\xedas. No es el que tiene que reg\xe1\xf1arla por las tareas, porque tiene que ir a impulsar los deberes de la casa.',
        attribution='Familia Monoparental',
        source_text=SRC_A01
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A02 — TERRITORIO DEL TRABAJO DOMÉSTICO
    # Headline: "Cocinar, lavar y limpiar siguen siendo territorio de ella.
    #            El único reparto equitativo llega cuando hay que pagar o ir al supermercado."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="Cocinar, lavar y limpiar siguen siendo territorio de ella. ",
        headline_italic="El \xfanico reparto equitativo llega cuando hay que pagar o ir al supermercado.",
        stats=[
            {
                'value': '67%',
                'desc_runs': [
                    make_run('reporta que '),
                    make_run('cocinar', bold=True),
                    make_run(' lo hace "m\xe1s la mujer"; lavar la ropa 63.0%; limpiar y organizar la casa 59.6%.')
                ]
            },
            {
                'value': '48%',
                'desc_runs': [
                    make_run('dice que el '),
                    make_run('supermercado', bold=True),
                    make_run(' lo hacen "ambos por igual". Manejo del dinero del hogar: 51.0%. Las tareas de dinero y compras s\xed se reparten.')
                ]
            },
            {
                'value': '45%',
                'desc_runs': [
                    make_run('acompa\xf1ar a '),
                    make_run('citas m\xe9dicas o escolares', bold=True),
                    make_run(' lo hace "m\xe1s la mujer". El cuidado log\xedstico de los hijos tampoco se reparte.')
                ]
            }
        ],
        source_text=SRC_A02
    )

    build_consumer_voice_slide(
        prs,
        quote='Como mujer, bueno, no sab\xeda cocinar y tuve que hacerlo. Obligada, aprend\xed a mam\xe1, pero... no me gusta, pero tengo que hacerlo.',
        attribution='Familia Monoparental',
        source_text=SRC_A02
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A03 — PERCEPCIÓN DE CAMBIO VS. REPARTO REAL
    # Headline: "74.8% dice que el rol de género cambió mucho. La cocina, el lavado y
    #            la limpieza siguieron siendo de ella. El discurso avanzó; el reparto, no."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="74.8% dice que el rol de g\xe9nero cambi\xf3 mucho. La cocina, el lavado y la limpieza siguieron siendo de ella. ",
        headline_italic="El discurso avanz\xf3; el reparto, no.",
        stats=[
            {
                'value': '75%',
                'desc_runs': [
                    make_run('responde "5 = Mucho" a cu\xe1nto cambi\xf3 el rol del hombre y la mujer en la familia. '),
                    make_run('85%', bold=True),
                    make_run(' en el top-2 (4+5).')
                ]
            },
            {
                'value': '77%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' percibe el cambio como "Mucho" vs 71.8% del masculino. Las mujeres ven m\xe1s el cambio que no cambia.')
                ]
            },
            {
                'value': '67%',
                'desc_runs': [
                    make_run('reporta cocinar '),
                    make_run('"m\xe1s la mujer"', bold=True),
                    make_run(' (Q056). Convive con la percepci\xf3n de cambio (Q057). El reparto material no acompa\xf1a la percepci\xf3n.')
                ]
            }
        ],
        source_text=SRC_A03
    )

    build_consumer_voice_slide(
        prs,
        quote='Las que salimos a echar para adelante, a trabajar, a luchar, la que ya no dependemos de un hombre para salir adelante. Y yo en mi caso, que aunque est\xe9 casada, trabajo igualito como si no, hago y lucho y no tengo que depender como tiempo atr\xe1s que las mujeres dominicanas nos creaban como que era para depender del hombre.',
        attribution='Familia Preferente Roles de Liderazgo',
        source_text=SRC_A03
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A04 — AUTORIDAD EN LA COMPRA DE ALIMENTOS
    # Headline: "89.5% de las mujeres decide siempre qué se compra de comer.
    #            La nutrición familiar es un poder casi exclusivo de ella, no compartido."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="89.5% de las mujeres decide siempre qu\xe9 se compra de comer. ",
        headline_italic="La nutrici\xf3n familiar es un poder casi exclusivo de ella, no compartido.",
        stats=[
            {
                'value': '90%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "5 = Siempre" a si toma la decisi\xf3n principal sobre la compra de alimentos, frente a 77.9% del masculino.')
                ]
            },
            {
                'value': '85%',
                'desc_runs': [
                    make_run('del '),
                    make_run('total', bold=True),
                    make_run(' responde "5 = Siempre". Es la decisi\xf3n menos disputada de toda la matriz de roles.')
                ]
            },
            {
                'value': '89%',
                'desc_runs': [
                    make_run('en '),
                    make_run('NSE D', bold=True),
                    make_run(' y 87.5% en NSE E responden "Siempre", frente a 52.2% en NSE AB. [AB n≈23: contraste, no cifra robusta.]')
                ]
            }
        ],
        source_text=SRC_A04
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo considero saludable la comida que yo cocino, ese tipo de ingredientes que le estoy echando, y tambi\xe9n las porciones. Lo que t\xfa cocinas en tu casa.',
        attribution='Familia Monoparental',
        source_text=SRC_A04
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A05 — EL ESTRÉS ECONÓMICO COMO PRIORIDAD
    # Headline: "47.8% de las mujeres nombra la economía como su mayor fuente de estrés.
    #            La crianza y el hogar no se acercan: el miedo más grande es la plata."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="47.8% de las mujeres nombra la econom\xeda como su mayor fuente de estr\xe9s. ",
        headline_italic="La crianza y el hogar no se acercan: el miedo m\xe1s grande es la plata.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('menciona "'),
                    make_run('Econom\xeda', bold=True),
                    make_run('" como factor que m\xe1s estr\xe9s genera; sumada a "Realidad econ\xf3mica" (28.4%), lo econ\xf3mico domina el ranking.')
                ]
            },
            {
                'value': '21%',
                'desc_runs': [
                    make_run('nombra expl\xedcitamente "'),
                    make_run('Falta de tiempo personal', bold=True),
                    make_run('" como fuente de estr\xe9s. Un quinto de las mujeres pone nombre a la ausencia de tiempo propio.')
                ]
            },
            {
                'value': '16%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Crianza de los hijos', bold=True),
                    make_run('" 16.2% y "Responsabilidades del hogar" 16.2%. El trabajo de cuidado estresa menos que la incertidumbre del dinero.')
                ]
            }
        ],
        source_text=SRC_A05
    )

    build_consumer_voice_slide(
        prs,
        quote='No tengo paciencia cuando no tengo dinero. Cuando no tengo dinero, que no s\xe9 qu\xe9 voy a hacer con los compromisos, entonces yo pienso, le doy vuelta a la cosa, un estr\xe9s. La mayor parte del problema hoy en d\xeda es el dinero, todo lo resuelve el dinero aunque la gente no lo quiera admitir.',
        attribution='Familia Monoparental',
        source_text=SRC_A05
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A06 — AUTOCUIDADO SACRIFICADO
    # Headline: "48.5% de las mujeres no hace nada por su salud mental ni física.
    #            El hombre se ejercita. Ella no tiene espacio ni para eso."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no hace nada por su salud mental ni f\xedsica. ",
        headline_italic="El hombre se ejercita. Ella no tiene espacio ni para eso.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "No realizo ninguna actividad para cuidar mi salud mental y f\xedsica" — frente al masculino cuya respuesta top es "Me ejercito" (43.1%).')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('del '),
                    make_run('total', bold=True),
                    make_run(' no hace ninguna actividad de autocuidado; solo 3.0% va a terapia.')
                ]
            },
            {
                'value': '60%',
                'desc_runs': [
                    make_run('de las mujeres '),
                    make_run('18-24 a\xf1os', bold=True),
                    make_run(' no realiza ninguna actividad de autocuidado. El abandono es mayor en las m\xe1s j\xf3venes.')
                ]
            }
        ],
        source_text=SRC_A06
    )

    build_consumer_voice_slide(
        prs,
        quote='Ese tema de que dependa de m\xed 24-7 lo que se vaya a hacer te da un poco de estr\xe9s.',
        attribution='Familia Mixta',
        source_text=SRC_A06
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A07 — MATERNIDAD COMO CASI-DEFAULT BIOGRÁFICO (solo-cuanti)
    # Headline: "85.4% de las mujeres ya tiene hijos. El no tener hijos no se vive como
    #            elección — se vive como espera, postergación o falta de pareja."
    # 3 stats, sin verbatim (hallazgo solo-cuanti declarado)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="85.4% de las mujeres ya tiene hijos. ",
        headline_italic="El no tener hijos no se vive como elecci\xf3n — se vive como espera, postergaci\xf3n o falta de pareja.",
        stats=[
            {
                'value': '85%',
                'desc_runs': [
                    make_run('del subset reporta "'),
                    make_run('Tengo hijos actualmente', bold=True),
                    make_run('"; entre el 14.6% sin hijos, las razones top: "No es una prioridad" (3.8%), "No me siento lista/o" (3.6%), "Prioridad profesional" (3.4%).')
                ]
            },
            {
                'value': '1%',
                'desc_runs': [
                    make_run('menciona "'),
                    make_run('No deseo tener hijos en absoluto', bold=True),
                    make_run('" (n=5) y solo 0.4% "Mi pareja no quiere hijos". La elecci\xf3n de no maternidad por convicci\xf3n aparece en menos del 1.5%.')
                ]
            },
            {
                'value': '2.2%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Razones econ\xf3micas', bold=True),
                    make_run('" iguala "No tengo pareja o condiciones emocionales adecuadas" (2.2%). Ambas superan el "no quiero" rotundo.')
                ]
            }
        ],
        source_text=SRC_A07
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # A08 — COMPOSICIÓN DEL HOGAR (solo-cuanti)
    # Headline: "1 de cada 4 hogares dominicanos es monoparental, y la mujer ocupa la cabeza.
    #            La familia nuclear no es la norma — es una de varias formas."
    # 3 stats, sin verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="1 de cada 4 hogares dominicanos es monoparental, y la mujer ocupa la cabeza. ",
        headline_italic="La familia nuclear no es la norma — es una de varias formas.",
        stats=[
            {
                'value': '24%',
                'desc_runs': [
                    make_run('del total responde "'),
                    make_run('Hogar monoparental', bold=True),
                    make_run('" (n=122 del n=500). Un cuarto de los hogares est\xe1 encabezado por un solo padre o madre.')
                ]
            },
            {
                'value': '25%',
                'desc_runs': [
                    make_run('en el '),
                    make_run('estrato E', bold=True),
                    make_run(' encabeza monoparental como respuesta TOP — la \xfanica tipolog\xeda donde supera a la biparental en el ranking. [NSE E n≈32: tendencia direccional.]')
                ]
            },
            {
                'value': '31%',
                'desc_runs': [
                    make_run('de hogares ya no encaja en biparental cl\xe1sico ni monoparental: '),
                    make_run('sin hijos + mascotas', bold=True),
                    make_run(' 7.4%, extendido 12.6%, otro 11.4%. La diversidad de formas es la norma.')
                ]
            }
        ],
        source_text=SRC_A08
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # PORTADA — BLOQUE B · Finance Women
    # ═══════════════════════════════════════════════════════════════════════════

    build_block_cover(
        prs,
        block_label="BLOQUE B",
        block_title="Pilar Mujer — Estudio Mujer y Finanzas",
        block_subtitle="Comportamiento financiero, sue\xf1os, barreras y autonom\xeda econ\xf3mica. Mujeres n=240."
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B01 — META FINANCIERA: UN TECHO PROPIO
    # Headline: "La meta financiera #1 de la mujer del estudio FW es un techo propio.
    #            El sueño es no depender del próximo sueldo."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="La meta financiera #1 de la mujer del estudio FW es un techo propio. ",
        headline_italic="El sue\xf1o es no depender del pr\xf3ximo sueldo.",
        stats=[
            {
                'value': '60%',
                'desc_runs': [
                    make_run('elige "'),
                    make_run('Comprar vivienda', bold=True),
                    make_run('" como meta financiera actual (1\xaa selecci\xf3n); 50.8% "Lograr estabilidad para no depender del pr\xf3ximo sueldo"; 48.3% "Emprender / Invertir".')
                ]
            },
            {
                'value': '48%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Comprar vivienda', bold=True),
                    make_run('" se repite como #1 en la 2\xaa selecci\xf3n (47.9%) y "Salir de deudas" sube a #2 (32.1%). La prioridad no cambia.')
                ]
            },
            {
                'value': '14%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Mejorar imagen personal', bold=True),
                    make_run('" queda \xfaltimo o pen\xfaltimo en ambas selecciones (13.8% y 2.1%). El prejuicio de la mujer que sue\xf1a con consumo est\xe9tico no aguanta la data.')
                ]
            }
        ],
        source_text=SRC_B01
    )

    build_consumer_voice_slide(
        prs,
        quote='Siempre ha sido mi sue\xf1o, siempre ha sido como que pagar mi cosa yo, si yo pudiera pagar todo, o sea, todo lo m\xedo yo.',
        attribution='Familia Preferente J\xf3venes',
        source_text=SRC_B01
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B02 — BLOQUEOS AL SUEÑO FINANCIERO
    # Headline: "74.9% siente que algo le bloquea el sueño financiero.
    #            87.1% apunta a lo mismo: falta plata."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="74.9% siente que algo le bloquea el sue\xf1o financiero. ",
        headline_italic="87.1% apunta a lo mismo: falta plata.",
        stats=[
            {
                'value': '75%',
                'desc_runs': [
                    make_run('responde "S\xed" a si siente que hay algo que le impide lograr su '),
                    make_run('sue\xf1o financiero', bold=True),
                    make_run('.')
                ]
            },
            {
                'value': '87%',
                'desc_runs': [
                    make_run('de las que sienten un bloqueo se\xf1ala "'),
                    make_run('Falta de ingresos', bold=True),
                    make_run('" como impedimento; muy por encima de "Falta de h\xe1bitos" (44.9%) o "Falta de educaci\xf3n financiera" (35.4%).')
                ]
            },
            {
                'value': '60%',
                'desc_runs': [
                    make_run('de todas las mujeres nombra "'),
                    make_run('Falta de ingresos', bold=True),
                    make_run('" como lo que le impide tomar mejores decisiones financieras hoy.')
                ]
            }
        ],
        source_text=SRC_B02
    )

    build_consumer_voice_slide(
        prs,
        quote='Comparado con un a\xf1o atr\xe1s que yo no ten\xeda libertad financiera, ten\xeda que depender de otro y es dif\xedcil, porque aunque el otro quiera, a veces uno quiere hacer algo y no puede. Ahora como yo tengo mi negocio, si quiero comprar algo lo compro, si le quiero comprar algo a mis hijos no tengo que esperar que otro me d\xe9.',
        attribution='Familia Masivo Emprendedoras',
        source_text=SRC_B02
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B03 — LA ADMINISTRADORA ES LA NORMA
    # Headline: "74.2% tiene presupuesto, 74.1% planifica siempre o casi siempre,
    #            78.7% revisa sus cuentas cada semana o cada día.
    #            La administradora es la norma, no la excepción."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="74.2% tiene presupuesto, 74.1% planifica siempre o casi siempre, 78.7% revisa sus cuentas cada semana o cada d\xeda. ",
        headline_italic="La administradora es la norma, no la excepci\xf3n.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('tiene un '),
                    make_run('presupuesto mensual', bold=True),
                    make_run('. Y 74.1% planifica sus gastos "Siempre" (42.1%) o "Casi siempre" (32.1%).')
                ]
            },
            {
                'value': '79%',
                'desc_runs': [
                    make_run('revisa sus cuentas a '),
                    make_run('diario o semanalmente', bold=True),
                    make_run(' (44.2% diario, 34.6% semanal). Control financiero activo como pr\xe1ctica cotidiana.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('se describe como "'),
                    make_run('Cuidadosa', bold=True),
                    make_run('" con el dinero y 20% como "Pr\xe1ctica"; solo 6.7% se define "Libre". La autopercepci\xf3n es de control, no de derroche.')
                ]
            }
        ],
        source_text=SRC_B03
    )

    build_consumer_voice_slide(
        prs,
        quote='Hay mujeres que no, que son ocultas, que si se presenta un caso, ellas siempre est\xe1n ah\xed tapando la falta del esposo. Hay mujeres que son muy honradas con el dinero.',
        attribution='Familia Masivo Jefas de Hogar',
        source_text=SRC_B03
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B04 — PREJUICIO DE GÉNERO EN EL DINERO
    # Headline: "39.2% siente que se espera algo distinto de ella en dinero por ser mujer.
    #            La esperan derrochadora y a la vez le exigen ser la administradora perfecta del hogar."
    # 2 stats + 2 verbatims en cuali (hallazgo cuanti + cuali integrado)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="39.2% siente que se espera algo distinto de ella en dinero por ser mujer. ",
        headline_italic="La esperan derrochadora y a la vez le exigen ser la administradora perfecta del hogar.",
        stats=[
            {
                'value': '39%',
                'desc_runs': [
                    make_run('responde "S\xed" a si se espera algo distinto de ella en temas de '),
                    make_run('dinero por ser mujer', bold=True),
                    make_run('; solo 31.7% dice "No" y 29.2% "No estoy segura". La mayor\xeda no descarta la presi\xf3n.')
                ]
            },
            {
                'value': '94',
                'desc_runs': [
                    make_run('respuestas abiertas (M45b) muestran dos polos: "que sea '),
                    make_run('buena administradora', bold=True),
                    make_run(', que siempre tenga guardado" vs "siempre piensan que las mujeres malgastan". El doble est\xe1ndar en una sola pregunta.')
                ]
            }
        ],
        source_text=SRC_B04
    )

    build_cuali_slide(
        prs,
        headline_plain="Dos prejuicios, una mujer: ",
        headline_italic="derrochadora y mala administradora al mismo tiempo.",
        verbatims=[
            {
                'quote': 'Algunas personas tienen la mentalidad que las mujeres son derrochadoras. El otro extremo es que la mujer siempre siempre debe guardar. Ambos extremos son malos.',
                'attribution': 'Respuesta abierta M45b · FW'
            },
            {
                'quote': 'Que al no ser proveedoras no somos buenas administradoras.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            }
        ],
        source_text=SRC_B04
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B05 — EL DINERO SE VIVE COMO AMENAZA
    # Headline: "La emoción que más siente la mujer al pensar en su dinero es estrés (44.3%).
    #            El orgullo solo llega al 5.5%: el dinero se vive como amenaza, no como logro."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="La emoci\xf3n que m\xe1s siente la mujer al pensar en su dinero es estr\xe9s (44.3%). ",
        headline_italic="El orgullo solo llega al 5.5%: el dinero se vive como amenaza, no como logro.",
        stats=[
            {
                'value': '44%',
                'desc_runs': [
                    make_run('responde "'),
                    make_run('Estr\xe9s', bold=True),
                    make_run('" como la emoci\xf3n que m\xe1s siente al pensar en sus finanzas; "Tranquilidad" 29.8%, "Miedo" 12.3%, "Culpa" 8.1%, "Orgullo" apenas 5.5%.')
                ]
            },
            {
                'value': '41%',
                'desc_runs': [
                    make_run('siente '),
                    make_run('ansiedad o estr\xe9s financiero', bold=True),
                    make_run(' con frecuencia (top-2, 4+5); el promedio en escala 1-5 es 3.48.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('est\xe1 satisfecha con su '),
                    make_run('situaci\xf3n financiera actual', bold=True),
                    make_run(' (top-2 22.5%); 30.8% insatisfecha (bottom-2). Promedio 2.75 sobre 5.')
                ]
            }
        ],
        source_text=SRC_B05
    )

    build_consumer_voice_slide(
        prs,
        quote='Tengo miedo a no ser lo que yo espero ser, o a que las cosas no se den como yo quiero. Siento que el tiempo est\xe1 pasando muy r\xe1pido, que no tengo una casa m\xeda o no tengo un auto propio.',
        attribution='Familia Masivo J\xf3venes',
        source_text=SRC_B05
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B06 — LA COMPRA IMPULSIVA COMO DESAHOGO
    # Headline: "47.9% hizo una compra impulsiva el último mes.
    #            El motivo #1 es consentirse — el único espacio que el rol le deja solo para ella."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="47.9% hizo una compra impulsiva el \xfaltimo mes. ",
        headline_italic="El motivo #1 es consentirse — el \xfanico espacio que el rol le deja solo para ella.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('hizo una '),
                    make_run('compra impulsiva', bold=True),
                    make_run(' en el \xfaltimo mes; 42.5% admite haber comprado algo recientemente solo para sentirse mejor emocionalmente.')
                ]
            },
            {
                'value': '44%',
                'desc_runs': [
                    make_run('El motivo #1: "'),
                    make_run('Para consentirme o darme un gusto', bold=True),
                    make_run('" (44.2%), por encima de "buena oportunidad" (39.2%) o "promoci\xf3n" (33.8%). 25.8% compra cuando est\xe1 estresada.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Comidas/Bebidas', bold=True),
                    make_run('" 72.6% y "Ropa/Accesorios" 59.8% lideran las categor\xedas. "Electr\xf3nica" apenas 5.0%. Desahogo cotidiano, no derroche aspiracional.')
                ]
            }
        ],
        source_text=SRC_B06
    )

    build_consumer_voice_slide(
        prs,
        quote='Uno se siente como empoderada y se\xf1or me lo merezco, por eso yo trabajo. Tiene como su momento, le llega de vez en vez, no siempre. D\xe9jame comprar ese porchecito, porque se me suba el \xe1nimo.',
        attribution='Familia Masivo Jefas de Hogar',
        source_text=SRC_B06
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B07 — EL AHORRO QUE NO LLEGA
    # Headline: "50% no logra ahorrar con facilidad. La estrategia que más funciona:
    #            esconderse el dinero de sí mismas antes de que llegue a la cuenta."
    # 3 stats + 2 verbatims en cuali (convergencia del mecanismo)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="50% no logra ahorrar con facilidad. La estrategia que m\xe1s funciona: ",
        headline_italic="esconderse el dinero de s\xed mismas antes de que llegue a la cuenta.",
        stats=[
            {
                'value': '50%',
                'desc_runs': [
                    make_run('responde que NO le resulta f\xe1cil '),
                    make_run('ahorrar', bold=True),
                    make_run('. En escala 1-5, el bottom-2 (35.8%) supera al top-2 (35%).')
                ]
            },
            {
                'value': '24%',
                'desc_runs': [
                    make_run('no destina ning\xfan porcentaje de su dinero al '),
                    make_run('ahorro', bold=True),
                    make_run('; otro 28.3% destina menos del 5%. M\xe1s de la mitad ahorra poco o nada.')
                ]
            },
            {
                'value': '66%',
                'desc_runs': [
                    make_run('"'),
                    make_run('No me alcanza el dinero', bold=True),
                    make_run('" es el obst\xe1culo #1 (65.5%), por encima de "Falta de educaci\xf3n financiera" (58.4%). La brecha de ahorro es de ingresos, no de conocimiento.')
                ]
            }
        ],
        source_text=SRC_B07
    )

    build_cuali_slide(
        prs,
        headline_plain="El mecanismo del ahorro: ",
        headline_italic="interceptar el dinero antes de que llegue a la mano.",
        verbatims=[
            {
                'quote': 'Me entr\xe9 a la cooperativa de la empresa, no me lo van a extraer de la cuenta. Lo hice con ese fin porque una de mis metas al principio del a\xf1o fue esa: ahorrar, porque crear un fondo de emergencia.',
                'attribution': 'Familia Preferente Jefas de Hogar'
            },
            {
                'quote': 'Me entr\xe9 en la cooperativa del trabajo, que es la \xfanica forma que ese dinero no entre a mi cuenta y salga.',
                'attribution': 'Familia Preferente Jefas de Hogar'
            }
        ],
        source_text=SRC_B07
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B08 — LA PROVEEDORA INVISIBLE
    # Headline: "40% de las mujeres del estudio FW es la principal proveedora de su hogar.
    #            El rol que se le atribuye al hombre lo ocupa ella."
    # 3 stats + 2 verbatims en cuali (complementarios)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="40% de las mujeres del estudio FW es la principal proveedora de su hogar. ",
        headline_italic="El rol que se le atribuye al hombre lo ocupa ella.",
        stats=[
            {
                'value': '40%',
                'desc_runs': [
                    make_run('responde que s\xed es la '),
                    make_run('principal proveedora', bold=True),
                    make_run(' de su hogar. [FW — perfil: j\xf3venes, jefas de hogar, l\xedderes, emprendedoras. No proyectable a poblaci\xf3n general.]')
                ]
            },
            {
                'value': '68%',
                'desc_runs': [
                    make_run('tiene '),
                    make_run('trabajo formal', bold=True),
                    make_run('; 13.9% informal o emprendimiento; solo 4.2% se declara ama de casa. La mujer del estudio FW est\xe1 mayoritariamente vinculada al ingreso.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Las tomo yo sola', bold=True),
                    make_run('" es la frase con mayor acuerdo top-2 en decisiones financieras (42.5%, μ=3.01). "No las tomo yo, lo hace mi esposo" tiene el mayor desacuerdo (bottom-2 79.6%).')
                ]
            }
        ],
        source_text=SRC_B08
    )

    build_cuali_slide(
        prs,
        headline_plain="El discurso del proveedor no cedi\xf3. ",
        headline_italic="Ellas proveen, el rol sigue siendo del hombre en el cuento social.",
        verbatims=[
            {
                'quote': 'No, por tu trabajo t\xfa eres proveedor. Pero se supone que el rol del proveedor siempre ha sido el hombre.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            },
            {
                'quote': 'Yo no resuelvo de nadie. Mi esposo es la madre. [Modismo: ella ocupa el rol de proveedor que culturalmente se le asigna al hombre.]',
                'attribution': 'Familia Masivo Jefas de Hogar'
            }
        ],
        source_text=SRC_B08
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B09 — EDUCACIÓN FINANCIERA AUTODIDACTA
    # Headline: "60% aprendió de finanzas por su cuenta y solo 11.7% en la escuela
    #            o universidad. El sistema formal no le enseñó; ella se enseñó."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="60% aprendi\xf3 de finanzas por su cuenta y solo 11.7% en la escuela o universidad. ",
        headline_italic="El sistema formal no le ense\xf1\xf3; ella se ense\xf1\xf3.",
        stats=[
            {
                'value': '60%',
                'desc_runs': [
                    make_run('reporta haber aprendido sobre finanzas "'),
                    make_run('Por tu cuenta', bold=True),
                    make_run('", seguido de "Familia" (45%); "Escuela/universidad" apenas 11.7%, "Trabajo" 15%; 9.6% responde "No he aprendido".')
                ]
            },
            {
                'value': '55%',
                'desc_runs': [
                    make_run('cree que la responsable de la educaci\xf3n financiera deber\xeda ser "'),
                    make_run('Las escuelas/colegios', bold=True),
                    make_run('" (54.6%); solo 15.1% nombra "Los bancos" y 8.8% "Yo misma". El Estado es el responsable esperado.')
                ]
            },
            {
                'value': '30%',
                'desc_runs': [
                    make_run('top-2 (4+5) en conocimiento sobre '),
                    make_run('productos financieros', bold=True),
                    make_run(' (29.6%); promedio 3.12 sobre 5. La mayor\xeda se ubica en neutro.')
                ]
            }
        ],
        source_text=SRC_B09
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo cuando chiquita, porque mi mam\xe1 me ense\xf1\xf3 desde peque\xf1a c\xf3mo yo ten\xeda que manejarme y guardar mi dinero desde peque\xf1a, pero me ense\xf1\xf3 que primero tengo que ganarme. Me compraba art\xedculos para que yo los vendiera en la escuela. Me ense\xf1\xf3 lo que era la venta desde peque\xf1a.',
        attribution='Familia Masivo Emprendedoras',
        source_text=SRC_B09
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B10 — AUTOCONOCIMIENTO FINANCIERO EN CONSTRUCCIÓN (solo-cuanti)
    # Headline: "46.7% se reconoce aprendiendo poco a poco por su cuenta y solo 3.3%
    #            tiene un alto nivel de conocimiento.
    #            La seguridad financiera está en construcción, no asumida."
    # 3 stats, sin verbatim (hallazgo solo-cuanti declarado)
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="46.7% se reconoce aprendiendo poco a poco por su cuenta y solo 3.3% tiene un alto nivel de conocimiento. ",
        headline_italic="La seguridad financiera est\xe1 en construcci\xf3n, no asumida.",
        stats=[
            {
                'value': '47%',
                'desc_runs': [
                    make_run('se identifica con "'),
                    make_run('Estoy aprendiendo poco a poco por mi cuenta', bold=True),
                    make_run('"; 24.2% con "S\xe9 lo b\xe1sico, pero no me siento segura"; solo 3.3% con "Tengo un alto nivel de conocimiento".')
                ]
            },
            {
                'value': '38%',
                'desc_runs': [
                    make_run('top-2 en '),
                    make_run('seguridad tomando decisiones financieras', bold=True),
                    make_run(' (μ=3.34); 50.8% se sit\xfaa en neutro. La confianza es moderada.')
                ]
            },
            {
                'value': '93%',
                'desc_runs': [
                    make_run('(top-2) considera muy importante el '),
                    make_run('conocimiento sobre el dinero', bold=True),
                    make_run(' en su vida (μ=4.58). La importancia es absoluta; el conocimiento real, modesto.')
                ]
            }
        ],
        source_text=SRC_B10
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B11 — LA JERARQUÍA DEL GASTO (solo-cuanti)
    # Headline: "El gasto #1 son alimentos. El ahorro queda en posición 8 de 10.
    #            El cuidado personal y el entretenimiento van últimos porque no caben."
    # 3 stats, sin verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="El gasto #1 son alimentos. El ahorro queda en posici\xf3n 8 de 10. ",
        headline_italic="El cuidado personal y el entretenimiento van \xfaltimos porque no caben.",
        stats=[
            {
                'value': '#1',
                'desc_runs': [
                    make_run('"'),
                    make_run('Alimentos y productos b\xe1sicos', bold=True),
                    make_run('" (μ=3.07) ocupa la posici\xf3n 1 en el ranking de gastos; "Pago de servicios" en posici\xf3n 2 (μ=4.10).')
                ]
            },
            {
                'value': '#8',
                'desc_runs': [
                    make_run('"'),
                    make_run('Ahorro', bold=True),
                    make_run('" queda en posici\xf3n 8 (μ=6.04); solo 7.6% lo coloca como #1. "Cuidado personal" en posici\xf3n 9 y "Entretenimiento" en posici\xf3n 10.')
                ]
            },
            {
                'value': '36%',
                'desc_runs': [
                    make_run('bottom-2 en facilidad para '),
                    make_run('ahorrar', bold=True),
                    make_run(' (35.8%). La mitad ahorra poco o nada mensualmente.')
                ]
            }
        ],
        source_text=SRC_B11
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B12 — DEUDA COMO PALANCA (no como recompensa)
    # Headline: "49.2% tomaría un préstamo solo si es para comprar una casa.
    #            La deuda se acepta como palanca, no como recompensa."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="49.2% tomar\xeda un pr\xe9stamo solo si es para comprar una casa. ",
        headline_italic="La deuda se acepta como palanca, no como recompensa.",
        stats=[
            {
                'value': '49%',
                'desc_runs': [
                    make_run('responde "'),
                    make_run('Comprar casa', bold=True),
                    make_run('" como \xfanica raz\xf3n para tomar un pr\xe9stamo; 21.7% "Iniciar/expandir un negocio"; 13.8% "Comprar un veh\xedculo"; 8.3% "Emergencia m\xe9dica".')
                ]
            },
            {
                'value': '1%',
                'desc_runs': [
                    make_run('solo '),
                    make_run('0.8% tomar\xeda un pr\xe9stamo para "Viajes o experiencias"', bold=True),
                    make_run(' y 2.9% para "Educaci\xf3n"; 3.3% responde "Ninguna". La deuda no se justifica como consumo.')
                ]
            },
            {
                'value': '38%',
                'desc_runs': [
                    make_run('top-2 en comodidad '),
                    make_run('manejando deudas', bold=True),
                    make_run(' (μ=3.30); 18.3% bottom-2. M\xe1s c\xf3moda que inc\xf3moda, pero la mitad se ubica en neutro.')
                ]
            }
        ],
        source_text=SRC_B12
    )

    build_consumer_voice_slide(
        prs,
        quote='Se me da\xf1\xf3 la lavadora y dije no, no voy a fu\xf1ir con lavadora peque\xf1a, voy a comprarme mi torre. Y en ese momento, como gracias a Dios tengo un buen historial crediticio, pues nada, la saqu\xe9 financiada. Eso s\xed, me he cuidado de manejarme con mis tarjetas y con mis pr\xe9stamos al d\xeda, gracias a Dios, sin faltar.',
        attribution='Familia Preferente Jefas de Hogar',
        source_text=SRC_B12
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B13 — ELLA DECIDE (no le entrega la decisión)
    # Headline: "79.6% rechaza la frase 'no las tomo yo, las toma mi esposo'.
    #            Ella decide sola o lo discute — pero no le entrega la decisión a él."
    # 3 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="79.6% rechaza la frase 'no las tomo yo, las toma mi esposo'. ",
        headline_italic="Ella decide sola o lo discute — pero no le entrega la decisi\xf3n a \xe9l.",
        stats=[
            {
                'value': '80%',
                'desc_runs': [
                    make_run('bottom-2 (en desacuerdo) frente a "'),
                    make_run('No las tomo yo, lo hace mi esposo', bold=True),
                    make_run('"; solo 7.9% top-2; μ=1.67. Rechazo masivo de la dependencia financiera.')
                ]
            },
            {
                'value': '37%',
                'desc_runs': [
                    make_run('top-2 frente a "Mi esposo y yo tomamos las decisiones '),
                    make_run('equitativamente', bold=True),
                    make_run('"; pero 44.6% bottom-2; μ=2.77. M\xe1s rechazo que acuerdo al equitativo declarado.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Las tomo yo sola', bold=True),
                    make_run('" — top-2 42.5%, μ=3.01. La frase con mayor adhesi\xf3n en toda la bater\xeda de decisiones financieras.')
                ]
            }
        ],
        source_text=SRC_B13
    )

    build_consumer_voice_slide(
        prs,
        quote='Mi pareja y yo, parecido, trabajamos juntos, entonces como que es dif\xedcil que me digan, no, no he cobrado, cobramos juntos. Pero s\xed somos muy claros, yo desde un principio, y le dije, yo soy muy desorganizada, estoy aprendiendo, estoy pidiendo ayuda por aqu\xed, mis amigos que trabajan en bancas, les pregunto, \xbfqu\xe9 tarjeta t\xfa me recomiendas?',
        attribution='Familia Preferente Jefas de Hogar',
        source_text=SRC_B13
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B14 — EDUCACIÓN FINANCIERA TRANSMITIDA
    # Headline: "Las mujeres del estudio enseñan finanzas a sus hijos como les enseñaron
    #            a ellas. Alcancía de leche evaporada, monedas de 5 y 10, y
    #            'guarda pan pa' mayo' antes que cualquier app."
    # 2 stats + 2 verbatims en cuali
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="Las mujeres del estudio ense\xf1an finanzas a sus hijos como les ense\xf1aron a ellas. ",
        headline_italic="Alcanc\xeda de leche evaporada, monedas de 5 y 10, y 'guarda pan pa' mayo' antes que cualquier app.",
        stats=[
            {
                'value': '19%',
                'desc_runs': [
                    make_run('se identifica con "'),
                    make_run("Guarda pan pa' mayo", bold=True),
                    make_run('" como filosof\xeda de dinero (18.8%), empatando con "Dios proveer\xe1" (18.8%). "El dinero da libertad" lidera con 23.3%.')
                ]
            },
            {
                'value': '45%',
                'desc_runs': [
                    make_run('report\xf3 haber aprendido finanzas de su '),
                    make_run('Familia', bold=True),
                    make_run(' — la segunda fuente despu\xe9s de "Por tu cuenta" (60%). La transmisi\xf3n informal es la escuela real.')
                ]
            }
        ],
        source_text=SRC_B14
    )

    build_cuali_slide(
        prs,
        headline_plain="La alcanc\xeda artesanal como primer curso de finanzas: ",
        headline_italic="el m\xe9todo de la madre, transmitido a los hijos.",
        verbatims=[
            {
                'quote': 'Yo soy de pinter\xf3n. Las alcanc\xedas a tres meses — en latas de leche rica. Es una alcanc\xeda natural, yo la lavo, la seco, y esa la uso de 5 y de 10, porque me di cuenta que de los 5 y los 10 yo boto mucho dinero.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            },
            {
                'quote': 'Con los muchachos, con los ni\xf1os hay que — yo quiero para diciembre tal, este a\xf1o es bueno, dame ya una alcanc\xeda, lo voy a empezar a echar.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            }
        ],
        source_text=SRC_B14
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B15 — LO QUE PIDE: RESPALDO PARA SOSTENER LO QUE YA HACE
    # Headline: "22.3% querría tener un fondo de emergencia y 13.4% una herramienta
    #            para llevar su presupuesto.
    #            Lo que pide es respaldo para sostener lo que ya hace."
    # 2 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="22.3% querr\xeda tener un fondo de emergencia y 13.4% una herramienta para llevar su presupuesto. ",
        headline_italic="Lo que pide es respaldo para sostener lo que ya hace.",
        stats=[
            {
                'value': '22%',
                'desc_runs': [
                    make_run('responde "'),
                    make_run('Un fondo de emergencia', bold=True),
                    make_run('" como producto financiero que le gustar\xeda tener; 20.6% "Cuenta de ahorro con beneficios"; 18.1% "Fondo de inversi\xf3n l\xedquido".')
                ]
            },
            {
                'value': '35%',
                'desc_runs': [
                    make_run('elige "Necesito opciones de ahorro que se ajusten a mis '),
                    make_run('ingresos irregulares', bold=True),
                    make_run('" (34.6%); 33.8% "Quisiera beneficios para emprender"; 31.3% "Herramientas para manejar mejor mi dinero".')
                ]
            }
        ],
        source_text=SRC_B15
    )

    build_consumer_voice_slide(
        prs,
        quote='Ahorro y salir de deudas. Es lo que pude, porque lo mir\xe9 m\xe1s de cerca, yo no puedo estar sin un dinero para una emergencia, entonces hice lo mismo que Andy, me entr\xe9 en la cooperativa del trabajo, que es la \xfanica forma que ese dinero no entre a mi cuenta y salga.',
        attribution='Familia Preferente Jefas de Hogar',
        source_text=SRC_B15
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B16 — LA BARRERA ES EL MIEDO, NO LA OFERTA
    # Headline: "39.2% no accede a productos financieros porque le da miedo endeudarse.
    #            Dos de cada diez agregan: nadie me ha orientado bien."
    # 2 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="39.2% no accede a productos financieros porque le da miedo endeudarse. ",
        headline_italic="Dos de cada diez agregan: nadie me ha orientado bien.",
        stats=[
            {
                'value': '39%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Me da miedo endeudarme', bold=True),
                    make_run('" es la barrera principal (39.2%); 30.2% "No tengo ingresos fijos"; 27.2% "No califico por mi historial crediticio"; 22.8% "Nadie me ha orientado bien".')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('reconoce "'),
                    make_run('No s\xe9 c\xf3mo funcionan', bold=True),
                    make_run('" (22.4%); 21.1% "Mi falta de h\xe1bitos"; 19.0% "No he tenido tiempo para averiguar". Solo 6.9% nombra "No me siento c\xf3moda con los bancos".')
                ]
            }
        ],
        source_text=SRC_B16
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo estoy averiguando en diferentes instrumentos de inversi\xf3n para ver cu\xe1l se acomoda mejor a m\xed, a mi objetivo, a mi estilo de ser porque no quiero correr muchos riesgos. Entonces estoy como que ahorrando para invertir.',
        attribution='Familia Preferente Jefas de Hogar',
        source_text=SRC_B16
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # B17 — LA VIDA FINANCIERA PLENA: POSIBLE PERO PROPIA
    # Headline: "42.9% ve posible una vida financiera plena. En la abierta, el patrón
    #            dominante es uno: me toca a mí, con organización y más ingresos."
    # 2 stats + 1 verbatim
    # ═══════════════════════════════════════════════════════════════════════════

    build_hallazgo_slide(
        prs,
        headline_plain="42.9% ve posible una vida financiera plena. En la abierta, el patr\xf3n dominante es uno: ",
        headline_italic="me toca a m\xed, con organizaci\xf3n y m\xe1s ingresos.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('top-2 (4+5) ve "muy posible" o "posible" tener una '),
                    make_run('vida financiera plena', bold=True),
                    make_run(' en su situaci\xf3n actual; 18.3% bottom-2 la ve imposible; μ=3.38.')
                ]
            },
            {
                'value': '238',
                'desc_runs': [
                    make_run('respuestas abiertas (M63) muestran el patr\xf3n: "'),
                    make_run('falta de ingresos', bold=True),
                    make_run('" + "organizaci\xf3n" + "estoy trabajando para ello". La lectura es activa individual, no estructural pasiva.')
                ]
            }
        ],
        source_text=SRC_B17
    )

    build_consumer_voice_slide(
        prs,
        quote='Al no tener ingresos fijos, no te puedes planificar. Eso limita todo.',
        attribution='Respuesta abierta M63 · FW',
        source_text=SRC_B17
    )

    return prs


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    script_dir   = os.path.dirname(os.path.abspath(__file__))
    output_path  = os.path.join(script_dir, "mujer-v3-deck-flat.pptx")

    print("Construyendo deck Mujer v3...")
    print("  Specs activas: v6 (HEADLINE_PT=50, STAT_BOX_W=378, STAT_BOX_H=113,")
    print("                  CV_VERBATIM_PT=50, CARD_W=567, CARD_H=227)")
    print("  Bloques: A (CC, A01-A08) + B (FW, B01-B17)")
    print("  Sin Banreservas en ningún elemento.")

    prs = build_deck()
    prs.save(output_path)

    total = len(prs.slides)
    print(f"\nDeck guardado: {output_path}")
    print(f"Total slides: {total}")
    print("\nInventario de slides:")
    for i, _slide in enumerate(prs.slides, 1):
        print(f"  Slide {i:02d}")


if __name__ == "__main__":
    main()
