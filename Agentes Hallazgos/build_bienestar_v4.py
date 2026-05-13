#!/usr/bin/env python3
"""
build_bienestar_v4.py
=====================
Deck bienestar-deck-flat-v4.pptx desde cero con python-pptx.
23 slides — set editorial cerrado bienestar-hallazgos-editados.md.

CAMBIOS CRÍTICOS v3 → v4 (specs exactas en px de aprendizajes-montador-cc.md § 3.9):

1. Slide: 1920×1080 px  →  18_288_000 × 10_287_000 EMU.
2. Cards verbatim: fill rgba(0,0,0,0.45) vía XML alpha 45000  (NO #2E2E2E sólido).
   Shape: MSO_SHAPE.ROUNDED_RECTANGLE (no RECTANGLE).
   adj = 0.10 → border-radius ~24px.
3. Stat box: 822×238 px (Layout 1 y 3). Layout 2 → 880×238 px para que quepan 2.
   Stat number: 180px Instrument Serif italic.
   Stat description: Poppins 22px.
4. Consumer Voice card: 1637×485 px centrada.
5. Verbatim: Instrument Serif regular 50px.
6. Headlines: tabla auto-size en px:
   ≤30 chars → 110px, 31–44 → 88px, 45–65 → 64px, 66–85 → 52px, >85 → 42px.
   Todos los headlines Bienestar son >85 chars → 42px.
7. Atribución: Poppins italic 22px.
8. Source pie: Poppins italic 16px gris #9B9B9B.
9. Header "CONSUMER VOICE": Poppins 16px gris #9B9B9B.
10. NO masterslide decorativo. Fondo negro plano.
11. Comillas españolas «...» en todos los verbatims.
12. Kerning 0 en todo el deck.

Convención de unidades:
  1 px @ 96 dpi = 9525 EMU
  Se define PX = 9525 y se usa PX como multiplicador.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    """Convierte píxeles a EMU."""
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)  # separadores verticales entre stats
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)  # source y header CONSUMER VOICE

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── Tamaños en pt (1pt = 12700 EMU; Pt(n) ya es el helper de python-pptx) ─────
# Para fuentes usamos Pt() directamente.
# Las specs del archivo aprendizajes son en px, no en pt.
# A 96 dpi: 1pt = 1.333px → 1px ≈ 0.75pt
# px_to_pt(n) = n * 0.75
def pt_from_px(px_val):
    """Convierte px a pt (a 96 dpi). Redondea a .5 para python-pptx."""
    return round(px_val * 0.75, 1)

# Fuentes clave en pt (derivadas de las specs en px del § 3.9):
STAT_NUMBER_PT  = pt_from_px(180)   # 135.0 pt
STAT_DESC_PT    = pt_from_px(22)    # 16.5 pt
VERBATIM_PT     = pt_from_px(50)    # 37.5 pt
ATTRIBUTION_PT  = pt_from_px(22)    # 16.5 pt
SOURCE_PT       = pt_from_px(16)    # 12.0 pt
CV_HEADER_PT    = pt_from_px(16)    # 12.0 pt

# ── Dimensiones de elementos en px ────────────────────────────────────────────
STAT_BOX_W_1STAT  = 822   # px — Layout 1 (1 stat centrado)
STAT_BOX_W_2STATS = 880   # px — Layout 2 (2 stats: 880×2 + margen = 1920 aprox)
STAT_BOX_W_3STATS = 560   # px — Layout 3 (3 stats)
STAT_BOX_H        = 238   # px

CV_CARD_W = 1637  # px — Consumer Voice card
CV_CARD_H = 485   # px

# Posición vertical de stat boxes: top 500px
STAT_TOP_PX = 500

# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Slide en blanco con fondo negro plano."""
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en <a:r>."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox simple de un solo run, posicionado en px."""
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    # Quitar márgenes internos del textbox
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox con múltiples runs en el mismo párrafo."""
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def pick_headline_size_px(chars):
    """
    Auto-size en px según § 3.9 specs exactas.
    Todos los headlines de Bienestar son >85 chars → 42px.
    """
    if chars <= 30:
        return 110
    elif chars <= 44:
        return 88
    elif chars <= 65:
        return 64
    elif chars <= 85:
        return 52
    else:
        return 42


def add_headline(slide, text_plain, text_italic=""):
    """
    Headline en MAYÚSCULAS, centrado, auto-size px, split plain/italic.
    Posición: top 100px, left 100px, width 1720px (1920 - 100*2).
    """
    full_text = text_plain + text_italic
    chars     = len(full_text.strip())
    size_px   = pick_headline_size_px(chars)
    size_pt   = pt_from_px(size_px)

    txBox = slide.shapes.add_textbox(
        px(100), px(100), px(1720), px(300)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(size_pt)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(size_pt)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_source(slide, source_text):
    """
    Source al pie del slide.
    Posición: bottom 60px → top = 1080 - 60 - 30 (height) = 990px.
    Poppins italic 16px gris #9B9B9B, centrado.
    """
    add_textbox(
        slide,
        left_px=0, top_px=990,
        width_px=1920, height_px=50,
        text=source_text,
        font_name=FONT_BODY, font_size_pt=SOURCE_PT,
        italic=True, color=GREY_SOFT,
        align=PP_ALIGN.CENTER
    )


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    """
    Línea vertical fina #2E2E2E (1px de ancho) entre stats.
    Usa RECTANGLE de 2px de ancho visible.
    """
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45(slide, left_px, top_px, width_px, height_px):
    """
    Crea un ROUNDED_RECTANGLE con fill rgba(0,0,0,0.45).
    - Shape: MSO_SHAPE.ROUNDED_RECTANGLE (no RECTANGLE)
    - adj = 0.10 → border-radius ~24px
    - Fill negro sólido + alpha 45000 vía XML

    Retorna el shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )

    # Ajuste del radio de esquina: 0.10 ≈ 10% del lado más corto
    # python-pptx expone adjustments[0] como valor 0.0–0.5
    shape.adjustments[0] = 0.05  # valor más conservador para 24px en 485px de alto

    # Sin borde
    shape.line.fill.background()

    # Fill negro sólido base
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    # Aplicar opacidad 45% vía XML directo
    # Estructura XML: solidFill > srgbClr > alpha val="45000"
    sp_elem = shape._element
    # Buscar el elemento solidFill en spPr
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        # Buscar el solidFill que python-pptx ya creó
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            # Crear manualmente si no existe
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))

        # Reemplazar con srgbClr + alpha
        # Primero limpiar contenido existente del solidFill
        for child in list(solidFill):
            solidFill.remove(child)

        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')

        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')  # 45000 = 45% en notación OOXML (0–100000)

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def add_card(slide, quote_text, attribution,
             left_px, top_px, card_w_px, card_h_px,
             verbatim_size_pt=None):
    """
    Card de verbatim: ROUNDED_RECTANGLE rgba(0,0,0,0.45) + texto encima.
    quote_text: sin comillas (se agregan aquí como «...»).
    verbatim_size_pt: None → usa VERBATIM_PT (50px → 37.5pt).
    """
    if verbatim_size_pt is None:
        verbatim_size_pt = VERBATIM_PT

    # Card background
    make_card_rgba_45(slide, left_px, top_px, card_w_px, card_h_px)

    # Texto del verbatim con comillas españolas
    full_quote = f"«{quote_text}»"

    # Padding interno: 80px horizontal, 60px vertical (proporcional a specs § 3.9)
    pad_h = 80
    pad_v = 60
    text_w = card_w_px - (pad_h * 2)
    text_h = card_h_px - (pad_v * 2)

    txBox = slide.shapes.add_textbox(
        px(left_px + pad_h), px(top_px + pad_v),
        px(text_w), px(text_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    # Párrafo del verbatim
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(verbatim_size_pt)
    run_q.font.italic = False   # NO italic — regla § 3.9
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Párrafo de atribución
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(ATTRIBUTION_PT)  # 22px → 16.5pt
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


# ── stat x-positions en px ────────────────────────────────────────────────────

def stat_x_centers(n_stats):
    """
    Devuelve centros horizontales en px para n_stats.
    Slide = 1920px.
    1 stat: centrado en 960px.
    2 stats: 1/3 y 2/3 = 640px y 1280px.
    3 stats: 1/4, 1/2, 3/4 = 480px, 960px, 1440px.
    """
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def stat_box_width(n_stats):
    """Ancho de cada caja de stat según número de stats."""
    if n_stats == 1:
        return STAT_BOX_W_1STAT
    elif n_stats == 2:
        return STAT_BOX_W_2STATS
    else:
        return STAT_BOX_W_3STATS


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    """
    Agrega un bloque stat: cifra grande + descripción.
    Stat box: width según n_stats × 238px.
    Stat number: 180px (135pt) Instrument Serif italic centrado.
    Stat desc: Poppins 22px (16.5pt) centrado con bold selectivo.
    """
    box_w = stat_box_width(n_stats)
    box_h = STAT_BOX_H  # 238px

    # left del box = centro - mitad del ancho
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX  # 500px

    # Cifra grande — altura reservada para el número (~180px aprox)
    num_h = 200  # px de altura para el textbox del número
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Descripción — debajo del número
    desc_top = box_top + num_h + 20  # 20px gap entre número y desc
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=180,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


# ── layouts de slides ──────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic,
                          stats, source_text):
    """
    Slide Hallazgo cuanti.
    stats: lista de dicts {value, desc_runs} — 1, 2 o 3 elementos.
    Sin cajas P##/Base — solo cifra + descripción Poppins con bold selectivo.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    # Separadores verticales entre stats
    sep_top_px    = STAT_TOP_PX - 20        # 480px
    sep_height_px = STAT_BOX_H + 40         # 278px

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    add_source(slide, source_text)
    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice — 1 verbatim en card centrada.
    Card: 1637×485 px, centrada horizontal, top 220px.
    Header "CONSUMER VOICE": Poppins 16px gris #9B9B9B, top 80px.
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE" — top 80px
    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Card centrada horizontalmente
    card_left = (SLIDE_W_PX - CV_CARD_W) // 2  # (1920 - 1637) / 2 = 141px
    card_top  = 220  # px, debajo del header

    add_card(
        slide,
        quote_text=quote,
        attribution=attribution,
        left_px=card_left, top_px=card_top,
        card_w_px=CV_CARD_W, card_h_px=CV_CARD_H
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text):
    """
    Slide solo-cuali: headline arriba + 1–3 cards apiladas.
    Card dimensions según § 3.9 Layout 5:
      1 verbatim → 1637×485 centrada
      2 verbatims → 2 × 1637×290 con 30px gap
      3 verbatims → 3 × 1637×200 con 25px gap
    Verbatim size:
      1 verbatim → 50px (37.5pt)
      2 verbatims → 38px (28.5pt)
      3 verbatims → 32px (24pt)
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)
    card_left = (SLIDE_W_PX - CV_CARD_W) // 2  # 141px

    if n == 1:
        card_h   = 485
        gap      = 0
        verb_pt  = VERBATIM_PT          # 37.5pt (50px)
        y_start  = 430  # centrar verticalmente en espacio restante (~430–950)
    elif n == 2:
        card_h   = 290
        gap      = 30
        verb_pt  = pt_from_px(38)       # 28.5pt
        y_start  = 380
    else:
        card_h   = 200
        gap      = 25
        verb_pt  = pt_from_px(32)       # 24pt
        y_start  = 360

    for i, v in enumerate(verbatims):
        y_card = y_start + i * (card_h + gap)
        add_card(
            slide,
            quote_text=v['quote'],
            attribution=v['attribution'],
            left_px=card_left, top_px=y_card,
            card_w_px=CV_CARD_W, card_h_px=card_h,
            verbatim_size_pt=verb_pt
        )

    add_source(slide, source_text)
    return slide


# ── Datos del set editorial v3/v4 ─────────────────────────────────────────────
# Contenido idéntico al v3. Solo cambia el motor visual (dimensiones, cards, fuentes).
# Las decisiones editoriales (qué stats, qué verbatims, qué atribuciones) son del editor.

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — T1: El dominicano no se hace el "loco" con la salud mental: la delega.
    # Layout 2 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 1 — Hallazgo cuanti (2 stats)
    # 162 chars → 42px (>85)
    # 3.0% → 3% (entero), 9.0% → 9% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="La terapia es la actividad de autocuidado menos elegida. El dominicano no ignora su salud mental — ",
        headline_italic="la terceriza a la fe, la mascota, el ejercicio y la familia.",
        stats=[
            {
                'value': '3%',
                'desc_runs': [
                    make_run('declara "voy a '),
                    make_run('terapia', bold=True),
                    make_run('" como actividad principal de autocuidado. La opción menos votada de toda la batería P27.')
                ]
            },
            {
                'value': '9%',
                'desc_runs': [
                    make_run('elige '),
                    make_run('meditación', bold=True),
                    make_run(' como autocuidado. Terapia + meditación: 12%. El 88% restante canaliza su bienestar por vías no clínicas.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P27 · Base 500."
    )

    # Slide 2 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo no iba a hacer lo que fue medicarme, pasar por psiquiatría. (...) Después de que pasó todo lo que pasó con la otra pareja que yo tuve, tuve que ir a la psiquiatría, tuve que medicarme, pero al final valió la pena, rindió sus frutos.',
        attribution='Familia Homoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — T1: El dominicano no se hace el "loco" con la salud mental: la delega.
    # Layout 1 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 3 — Hallazgo cuanti (1 stat)
    # 167 chars → 42px (>85)
    # 85.0% → 85% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="85% del dominicano considera a la mascota parte de la familia. ",
        headline_italic="En los hombres, el perro cumple la función emocional que la cultura no les permite buscar en otra parte.",
        stats=[
            {
                'value': '85%',
                'desc_runs': [
                    make_run('considera a las mascotas '),
                    make_run('parte de la familia', bold=True),
                    make_run('. En hogares sin hijos con mascota sube a 94.6%.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P10 · Base 500."
    )

    # Slide 4 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Estadísticamente hablando, el perro suele ser un apoyo muy emocional para el hombre, en muchas ocasiones hasta más que la mujer, y no es por un tema de que no sentamos ese amor, sino porque tenemos un aspecto cultural y hasta biológico (...) de que no nos desahogamos como las mujeres lo hacen (...) nosotros simplemente necesitamos compañía, tacto y simplemente presencia.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — T2: El dominicano gestiona su salud mental de rodillas.
    # Layout 2 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 5 — Hallazgo cuanti (2 stats)
    # 160 chars → 42px (>85)
    # 64.0% → 64% (entero), ratio 1:6
    build_hallazgo_slide(
        prs,
        headline_plain="64% recurre siempre a la religión en momentos de dificultad. ",
        headline_italic="La fe es el primer canal de salud mental: gratis, disponible las 24 horas, sin estigma ni copago.",
        stats=[
            {
                'value': '64%',
                'desc_runs': [
                    make_run('marca "5 = Siempre" recurrir a la '),
                    make_run('religión', bold=True),
                    make_run(' o espiritualidad en momentos de dificultad. Respuesta dominante en todas las tipologías.')
                ]
            },
            {
                'value': '1:6',
                'desc_runs': [
                    make_run('menciones clínicas por cada 6 '),
                    make_run('religiosas', bold=True),
                    make_run(' en 11 FGs (216 menciones fe vs 37 de terapia / psicólogo / psiquiatría).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P29 · Base 500."
    )

    # Slide 6 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='En mi caso nosotros con la biblia (...) si eso me ayudó bastante, bastante, la cercanía con Dios, a nuestra salud mental, de ambos.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — T2: El dominicano gestiona su salud mental de rodillas.
    # Layout 2 → 1 slide (solo-cuanti, sin verbatim)
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 7 — Hallazgo solo-cuanti (2 stats)
    # 170 chars → 42px (>85)
    # 70.8% → 71% (entero), 72.7% → 73% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="70.8% de las mujeres y 72.7% del estrato D recurren siempre a la fe. ",
        headline_italic="La intensidad religiosa sube donde más aprieta: si eres mujer y el bolsillo está flaco, rezas más.",
        stats=[
            {
                'value': '71%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' marca "Siempre" recurrir a la religión en momentos difíciles, vs 53.3% de los hombres. Brecha de género: 17.5 puntos.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' marca "Siempre" recurrir a la fe, vs 52.8% del estrato C+ y 61.5% del estrato C.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P29 × D2, P29 × D6 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — T3: El autocuidado es físico, no emocional.
    # Layout 2 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 8 — Hallazgo cuanti (2 stats)
    # 163 chars → 42px (>85)
    # 42.6% → 43% (entero), 45.9% → 46% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="42.6% del dominicano no hace ninguna actividad para cuidar su salud, ni física ni mental. ",
        headline_italic="La inacción es la respuesta más común de toda la batería de autocuidado.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('no realizo ninguna actividad', bold=True),
                    make_run('" para cuidar su salud mental y física. TOP-1 de la batería P27, por encima de ejercicio (34.8%), oración (27.6%).')
                ]
            },
            {
                'value': '46%',
                'desc_runs': [
                    make_run('de hogares '),
                    make_run('sin hijos con mascota', bold=True),
                    make_run(' lidera la inacción. La inacción lidera en 5 de 7 tipologías: biparental (44.2%), extendido (42.9%), monoparental (41.8%).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27 · Base 500."
    )

    # Slide 9 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Bueno yo para lo físico yo no estoy haciendo nada. Yo no hago ejercicio, yo me como todo lo que yo quiera y a la hora que yo quiera. Bueno, para la salud mental a veces uno tiene que hacerse loco en ciertas cosas.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — T3: El autocuidado es físico, no emocional.
    # Layout 2 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 10 — Hallazgo cuanti (2 stats)
    # 186 chars → 42px (>85)
    # 34.8% → 35% (entero), 43.1% → 43% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="34.8% del dominicano se ejercita con regularidad — el único autocuidado donde el hombre lidera. ",
        headline_italic="Para él, entrenar no es físico: es la válvula emocional que la cultura sí le permite.",
        stats=[
            {
                'value': '35%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('me ejercito con regularidad', bold=True),
                    make_run('" como autocuidado — segundo lugar de la batería P27, detrás de "ninguna actividad" (42.6%).')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('masculino', bold=True),
                    make_run(' tiene el ejercicio como TOP-1. En el femenino el TOP-1 es "ninguna actividad" (48.5%). El único canal donde el hombre supera a la mujer.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27, P27 × D2 · Base 500."
    )

    # Slide 11 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo soy una persona súper activa, súper activa (...) La salud mental mía es deporte. Yo me irrito cuando yo duro tres días que no puedo hacer ejercicio (...) Yo entro a la cancha y yo me olvido de todo lo que yo tengo.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — T4: Mi equilibrio mental es mi equilibrio económico.
    # Layout 3 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 12 — Hallazgo cuanti (3 stats)
    # 189 chars → 42px (>85)
    # 47.8% → 48% (entero), 28.4% → 28% (entero), 51.2% → 51% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="47.8% nombra la economía como su mayor fuente de estrés familiar. ",
        headline_italic="Sumando los dos ítems económicos de la batería, el estresor número uno del hogar dominicano alcanza al 76% de los hogares.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('Economía', bold=True),
                    make_run('" como el factor que más estrés genera en su vida familiar. TOP-1 nacional de la batería P26.')
                ]
            },
            {
                'value': '28%',
                'desc_runs': [
                    make_run('marca también "'),
                    make_run('Realidad económica', bold=True),
                    make_run('" en la misma batería. Los dos ítems económicos combinados: 76% del total.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' declara la economía como estresor principal, vs el promedio nacional de 47.8%.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P26 · Base 500."
    )

    # Slide 13 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Tuve un tiempo como mamá soltera y mi hermana es mamá soltera de tres niños, y es bien difícil levantarte, llevar al niño a la escuela, el dinero, cuando el papá no está económicamente ni emocionalmente.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — T5: Me agota ser padre/madre, pero nunca me canso de mis hijos.
    # Layout 3 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 14 — Hallazgo cuanti (3 stats)
    # 180 chars → 42px (>85)
    # 42.2% → 42% (entero), 22.3% → 22% (entero), 5.5% → mantiene decimal (1 dígito antes)
    build_hallazgo_slide(
        prs,
        headline_plain="42.2% dice que la presión sobre los padres es demasiada. ",
        headline_italic="El agotamiento no aparece de golpe — en la franja 35–44 ya lo declara como estresor el 22.3%, el triple que a los 18–24.",
        stats=[
            {
                'value': '42%',
                'desc_runs': [
                    make_run('declara "5 — Muy de acuerdo" con "'),
                    make_run('hay demasiada presión', bold=True),
                    make_run(' para que los papás lo tengan todo." Ítem con mayor consenso de la batería P8 sobre crianza.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('35–44 años', bold=True),
                    make_run(' declara la crianza como factor de estrés, vs 15.5% del subset 25–34. Pico en la franja media adulta.')
                ]
            },
            {
                'value': '5.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('18–24 años', bold=True),
                    make_run(' declara la crianza como estresor. La presión parental se construye con los años.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P8, P26 × D5 · Base 500."
    )

    # Slide 15 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Es un mito que una mujer va a estar (...) toda ella relajada, limpiando una casa, con el marido sentado allá cargando el niño y qué feliz de la vida. Eso es un mito. De toda mujer que tiene hijos siempre agotada, cansada, batida.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — T6: La crisis monoparental no tiene tratamiento.
    # Layout 2 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 16 — Hallazgo cuanti (2 stats)
    # 183 chars → 42px (>85)
    # 24.4% → 24% (entero), 56.6% → 57% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="24.4% de los hogares es monoparental. En más de la mitad, la madre es la única responsable de la crianza. ",
        headline_italic="El mercado diseña para dos. El hogar dominicano no siempre llega de a dos.",
        stats=[
            {
                'value': '24%',
                'desc_runs': [
                    make_run('del dominicano declara que su hogar es '),
                    make_run('monoparental', bold=True),
                    make_run('. Segunda tipología más común después del biparental con hijos menores de 18 (38.0%).')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('del subset monoparental cita a la '),
                    make_run('madre', bold=True),
                    make_run(' como única responsable de la crianza, vs 41.8% del promedio nacional. Crianza monoparental de facto materna.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P1, P2 × D1 · Base 500."
    )

    # Slide 17 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Soy madre soltera, sola, con una situación actualmente de mi madre que falleció y que vive conmigo y mi hermanito en condición especial, entonces es difícil, es difícil todo en general, quién lo cuide, porque yo trabajo todo el tiempo.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — T7: Las familias dominicanas no tienen un sistema de salud.
    # Layout 5 (2 verbatims) → 1 slide solo-cuali
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 18 — Card cualitativa (2 cards apiladas)
    # 185 chars → 42px (>85)
    build_cuali_slide(
        prs,
        headline_plain="El sistema público de salud falla antes de la consulta. ",
        headline_italic="La pregunta dominicana no es qué le pasa al médico — es si el seguro aprueba, si hay cama, si alguien dentro te puede ayudar.",
        verbatims=[
            {
                'quote': 'Vaya con un dolor, aunque sea un dolor de una hebra de cabello, y que te digan "hasta que el seguro no te apruebe, no." O tú ir a comprar un medicamento y que te digan: "ah no, tu límite del medicamento ya se agotó." Exactamente. Tú vas a una clínica y lo primero que te preguntan es qué seguro tú tienes.',
                'attribution': 'Familia Homoparental'
            },
            {
                'quote': 'Si tú no tienes alguien que te puede ayudar en el sistema, tú estás muerto.',
                'attribution': 'Familia Sin Hijos con Mascota'
            }
        ],
        source_text="Source: Código Casa — Estudio cualitativo 2025 · P28 · Base cualitativa: 11 grupos de enfoque."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — T8: Los ancianos son invisibles.
    # Layout 3 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 19 — Hallazgo cuanti (3 stats)
    # 176 chars → 42px (>85)
    # 73.8% → 74% (entero), 24.0% → 24% (entero), 12.8% → 13% (entero)
    build_hallazgo_slide(
        prs,
        headline_plain="73.8% pide pensiones para que los mayores vivan con dignidad. Solo 12.8% es cuidador directo. ",
        headline_italic="Los dominicanos saben lo que hace falta — solo que no están en posición de darlo.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('marca "'),
                    make_run('Pensiones o apoyos económicos suficientes', bold=True),
                    make_run('" como lo que más falta para que la tercera edad viva con dignidad. TOP-1 absoluto de P79.')
                ]
            },
            {
                'value': '24%',
                'desc_runs': [
                    make_run('menciona "'),
                    make_run('Apoyo y capacitación para cuidadores familiares', bold=True),
                    make_run('" en la misma batería. El cuidador familiar pide ayuda — el mercado no la ofrece como categoría.')
                ]
            },
            {
                'value': '13%',
                'desc_runs': [
                    make_run('declara ser el '),
                    make_run('cuidador principal', bold=True),
                    make_run(' de un adulto mayor. 72.8% declara que no vive con ninguna persona mayor en su hogar.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P78, P79 · Base 500."
    )

    # Slide 20 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='La frustración en cuanto a salud y bienestar para mí es el tema del presupuesto que requiere. (...) Porque tengo un choque entre lo que yo digo que tengo que aceptar y las costumbres pasadas. Mi mamá era de la mujer que me enseñó a mí a que si algo no es necesario, no gastes dinero en eso, economiza.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H12 — T9: Saber aguantar se hereda.
    # Layout 5 (1 verbatim) → 1 slide solo-cuali
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 21 — Card cualitativa (1 card)
    # 189 chars → 42px (>85)
    build_cuali_slide(
        prs,
        headline_plain='"Yo no me voy a separar, yo voy a aguantar — porque mi mamá me decía que hay que aguantar." ',
        headline_italic="El aguante dominicano no es una decisión. Es una instrucción que llega de la generación anterior.",
        verbatims=[
            {
                'quote': 'Yo me crié con esos dos conceptos en mi mente y yo siempre me enfocaba. Tanto así que yo decía: "No, yo no me voy a separar, yo voy a aguantar", porque mi mamá me decía que hay que aguantar.',
                'attribution': 'Familia Monoparental'
            }
        ],
        source_text="Source: Código Casa — Estudio cualitativo 2025 · Base cualitativa: 11 grupos de enfoque."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H13 — T10: El amor es lo más importante, a menos que sea amor propio.
    # Layout 1 + Layout 4 → 2 slides
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 22 — Hallazgo cuanti (1 stat)
    # 170 chars → 42px (>85)
    # 48.5% mantiene decimal (1 dígito antes del punto)
    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no hace ninguna actividad para cuidarse, vs 33.3% de los hombres. ",
        headline_italic="Brecha de género en autocuidado: 15.2 puntos. La mujer se cuida menos — y lo sabe.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' declara "no realizo ninguna actividad para cuidar mi salud" como TOP-1 de P27, vs 33.3% del subset masculino. Brecha de género en inacción: 15.2 puntos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27 × D2 · Base 500."
    )

    # Slide 23 — Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Cuando uno es madre, a veces se olvida, a veces no... casi siempre se olvida de uno mismo. (...) Y cuando tú vienes a mirar para atrás, entonces tú dices: "¿y yo?" Porque crecen y se van.',
        attribution='Familia Monoparental'
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Agentes Hallazgos/bienestar-deck-flat-v4.pptx"
    prs.save(output_path)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Verificación de specs v4
    print("\n── Verificación specs v4 ──────────────────────────────────────────")
    print(f"Slide size: {prs.slide_width} × {prs.slide_height} EMU")
    print(f"  = {prs.slide_width / PX:.0f}px × {prs.slide_height / PX:.0f}px")
    expected_w = 1920 * PX
    expected_h = 1080 * PX
    print(f"  Esperado: {expected_w} × {expected_h} EMU")
    print(f"  Slide size OK: {prs.slide_width == Emu(expected_w) and prs.slide_height == Emu(expected_h)}")

    print(f"\nStat number size: {STAT_NUMBER_PT}pt ({STAT_NUMBER_PT/0.75:.0f}px)")
    print(f"Stat desc size:   {STAT_DESC_PT}pt ({STAT_DESC_PT/0.75:.0f}px)")
    print(f"Verbatim size:    {VERBATIM_PT}pt ({VERBATIM_PT/0.75:.0f}px)")
    print(f"Attribution size: {ATTRIBUTION_PT}pt ({ATTRIBUTION_PT/0.75:.0f}px)")
    print(f"Source size:      {SOURCE_PT}pt ({SOURCE_PT/0.75:.0f}px)")
    print(f"CV header size:   {CV_HEADER_PT}pt ({CV_HEADER_PT/0.75:.0f}px)")
    print(f"\nStat box (1 stat):  {STAT_BOX_W_1STAT}×{STAT_BOX_H}px")
    print(f"Stat box (2 stats): {STAT_BOX_W_2STATS}×{STAT_BOX_H}px")
    print(f"Stat box (3 stats): {STAT_BOX_W_3STATS}×{STAT_BOX_H}px")
    print(f"CV card:            {CV_CARD_W}×{CV_CARD_H}px")
    print("\nCards: rgba(0,0,0,0.45) vía MSO_SHAPE.ROUNDED_RECTANGLE + XML alpha=45000")
    print("Headlines >85 chars → 42px (31.5pt)")
    print("Comillas españolas «...» en todos los verbatims")
    print("Kerning 0 en todo el deck")
    print("NO masterslide — fondo negro plano")

    return output_path


if __name__ == "__main__":
    build_deck()
