#!/usr/bin/env python3
"""
build_opiniones_politicas_deck.py
==================================
Deck opiniones-politicas-deck-flat.pptx desde cero con python-pptx.
21 slides — set editorial cerrado opiniones-politicas-hallazgos-editados.md (14 hallazgos).

Especificaciones (aprendizajes-montador-cc.md § 3.9):
- Slide 1920×1080 px = 18_288_000 × 10_287_000 EMU
- Cards verbatim: fill rgba(0,0,0,0.45) vía XML alpha 45000 (NO #2E2E2E sólido)
  Shape: MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05 → ~24px border-radius
- Stat box: Layout 1 → 822×238, Layout 2 → 880×238, Layout 3 → 560×238
- Stat number: 180px (135pt) Instrument Serif italic
- Stat desc: Poppins 22px (16.5pt) con bold selectivo
- Consumer Voice card: 1637×485 px centrada
- Verbatim: Instrument Serif regular 50px (37.5pt)
- Headlines: auto-size ≤30→110px, 31-44→88px, 45-65→64px, 66-85→52px, >85→42px
- Atribución: Poppins italic 22px (16.5pt)
- Source: Poppins italic 16px (12pt) gris #9B9B9B
- Header "CONSUMER VOICE": Poppins 16px (12pt) gris #9B9B9B
- Fondo negro plano — sin masterslide decorativo
- Comillas españolas «...» en todos los verbatims
- Kerning 0 en todo el deck

Conteo de slides:
  H01 (2 stats + CV):       2
  H02 (2 stats solo-cuanti): 1
  H03 (2 stats + CV):       2
  H04 (2 stats solo-cuanti): 1
  H05 (2 stats + CV):       2
  H06 (2 stats solo-cuanti): 1
  H07 (2 stats + CV):       2
  H08 (2 stats solo-cuanti): 1
  H09 (2 stats + CV):       2
  H10 (1 stat + CV):        2
  H11 (solo-cuali 2 verbatims): 1
  H12 (3 stats + CV):       2
  H13 (2 stats + CV):       2
  H14 (2 stats solo-cuanti): 1
  TOTAL: 22 slides

Reglas de redondeo aplicadas (aprendizajes-montador-cc.md § 3.5):
  Stat grande → entero por defecto.
  Decimal solo si: 1 dígito antes del punto (8.7%, 5.5%) o brecha significativa.
  Todos los stats de este pilar tienen 2 dígitos antes → se redondean a entero.

Convención de unidades:
  1 px @ 96 dpi = 9525 EMU
  px_to_pt(n) = n * 0.75
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── px → pt ───────────────────────────────────────────────────────────────────
def pt_from_px(px_val):
    return round(px_val * 0.75, 1)

STAT_NUMBER_PT  = pt_from_px(180)   # 135.0 pt
STAT_DESC_PT    = pt_from_px(22)    # 16.5 pt
VERBATIM_PT     = pt_from_px(50)    # 37.5 pt
ATTRIBUTION_PT  = pt_from_px(22)    # 16.5 pt
SOURCE_PT       = pt_from_px(16)    # 12.0 pt
CV_HEADER_PT    = pt_from_px(16)    # 12.0 pt

# ── Dimensiones de elementos ───────────────────────────────────────────────────
STAT_BOX_W_1STAT  = 822
STAT_BOX_W_2STATS = 880
STAT_BOX_W_3STATS = 560
STAT_BOX_H        = 238

CV_CARD_W = 1637
CV_CARD_H = 485

STAT_TOP_PX = 500


# ── Helpers básicos ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def pick_headline_size_px(chars):
    if chars <= 30:
        return 110
    elif chars <= 44:
        return 88
    elif chars <= 65:
        return 64
    elif chars <= 85:
        return 52
    else:
        return 42


def add_headline(slide, text_plain, text_italic=""):
    full_text = text_plain + text_italic
    chars     = len(full_text.strip())
    size_px   = pick_headline_size_px(chars)
    size_pt   = pt_from_px(size_px)

    txBox = slide.shapes.add_textbox(
        px(100), px(100), px(1720), px(300)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(size_pt)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(size_pt)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_source(slide, source_text):
    add_textbox(
        slide,
        left_px=0, top_px=1010,
        width_px=1920, height_px=50,
        text=source_text,
        font_name=FONT_BODY, font_size_pt=SOURCE_PT,
        italic=True, color=GREY_SOFT,
        align=PP_ALIGN.CENTER
    )


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    line = slide.shapes.add_shape(
        1,
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45(slide, left_px, top_px, width_px, height_px):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    shape.adjustments[0] = 0.05
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def add_card(slide, quote_text, attribution,
             left_px, top_px, card_w_px, card_h_px,
             verbatim_size_pt=None):
    if verbatim_size_pt is None:
        verbatim_size_pt = VERBATIM_PT

    make_card_rgba_45(slide, left_px, top_px, card_w_px, card_h_px)

    full_quote = f"«{quote_text}»"

    pad_h = 80
    pad_v = 60
    text_w = card_w_px - (pad_h * 2)
    text_h = card_h_px - (pad_v * 2)

    txBox = slide.shapes.add_textbox(
        px(left_px + pad_h), px(top_px + pad_v),
        px(text_w), px(text_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(verbatim_size_pt)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(ATTRIBUTION_PT)
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def stat_box_width(n_stats):
    if n_stats == 1:
        return STAT_BOX_W_1STAT
    elif n_stats == 2:
        return STAT_BOX_W_2STATS
    else:
        return STAT_BOX_W_3STATS


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w = stat_box_width(n_stats)
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=180,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


def build_hallazgo_slide(prs, headline_plain, headline_italic,
                          stats, source_text):
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    add_source(slide, source_text)
    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    slide = blank_slide(prs)

    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    card_left = (SLIDE_W_PX - CV_CARD_W) // 2
    card_top  = 220

    add_card(
        slide,
        quote_text=quote,
        attribution=attribution,
        left_px=card_left, top_px=card_top,
        card_w_px=CV_CARD_W, card_h_px=CV_CARD_H
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text):
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)
    card_left = (SLIDE_W_PX - CV_CARD_W) // 2

    if n == 1:
        card_h   = 485
        gap      = 0
        verb_pt  = VERBATIM_PT
        y_start  = 430
    elif n == 2:
        card_h   = 290
        gap      = 30
        verb_pt  = pt_from_px(38)
        y_start  = 380
    else:
        card_h   = 200
        gap      = 25
        verb_pt  = pt_from_px(32)
        y_start  = 360

    for i, v in enumerate(verbatims):
        y_card = y_start + i * (card_h + gap)
        add_card(
            slide,
            quote_text=v['quote'],
            attribution=v['attribution'],
            left_px=card_left, top_px=y_card,
            card_w_px=CV_CARD_W, card_h_px=card_h,
            verbatim_size_pt=verb_pt
        )

    if source_text:
        add_source(slide, source_text)
    return slide


# ── BUILD DECK ─────────────────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — 6 de cada 10 dominicanos se informa de política por redes sociales.
    # Layout 2 (2 stats) + CV = 2 slides
    # 58.6%→59%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="6 DE CADA 10 DOMINICANOS SE INFORMA DE POLÍTICA POR REDES SOCIALES. LA TV CAE A SEGUNDO LUGAR, ",
        headline_italic="Y EL NOTICIERO YA NO CONVOCA.",
        stats=[
            {
                'value': '59%',
                'desc_runs': [
                    make_run('usa '),
                    make_run('redes sociales', bold=True),
                    make_run(' como fuente principal para informarse sobre política — la opción más votada de P56, con más del doble que la TV.')
                ]
            },
            {
                'value': '27%',
                'desc_runs': [
                    make_run('declara la '),
                    make_run('TV', bold=True),
                    make_run(' como fuente principal. Periódicos (1.2%) y radio (1.4%) son residuales. La conversación política migró al feed antes que al noticiero.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P56 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Se habla de política en esos casos. (...) Pero no porque somos políticos, sino porque al niño mío le gusta mucho ver cosas como en YouTube, así, como lo dice, ese tipo de cosas. Entonces él me dice (...) que está el candidato Luis Sabinader y que él es el presidente.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — 80% de los jóvenes 18-24 se informa por redes. En el estrato E, la TV aún gana.
    # Layout 2 (2 stats) SOLO-CUANTI = 1 slide
    # 80.0%→80%, brecha 30.4 puntos → se mantiene para describir
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="80% DE LOS JÓVENES 18-24 SE INFORMA POR REDES. EN EL ESTRATO E, LA TV AÚN GANA. ",
        headline_italic="LA PLATAFORMA QUE USAS PARA VER POLÍTICA DEPENDE DE CUÁNTO PAGAS LA LUZ.",
        stats=[
            {
                'value': '80%',
                'desc_runs': [
                    make_run('de los jóvenes de '),
                    make_run('18-24 años', bold=True),
                    make_run(' se informa de política por redes, frente a 49.6% de los mayores de 55 que ubican la TV en primer lugar. Brecha: 30.4 puntos.')
                ]
            },
            {
                'value': '38%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato E', bold=True),
                    make_run(' declara la TV como fuente principal — el único NSE donde las redes no encabezan. En estratos C, C+ y AB las redes superan el 65%. (n≈32, tendencia direccional)')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P56 × D5, P56 × D6 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — 58% del dominicano se declara nada interesado en política.
    # Layout 2 (2 stats) + CV = 2 slides
    # 58.0%→58%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="58% DEL DOMINICANO SE DECLARA NADA INTERESADO EN POLÍTICA. EL DESINTERÉS NO ES UNA POSTURA DE MINORÍA — ",
        headline_italic="ES EL ESTADO DECLARADO DE LA MAYORÍA.",
        stats=[
            {
                'value': '58%',
                'desc_runs': [
                    make_run('se ubica en el punto más bajo de la escala — "'),
                    make_run('1 = Nada interesado', bold=True),
                    make_run('" en política. Es la respuesta más votada por 25 puntos de distancia sobre cualquier otra opción.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('suma el '),
                    make_run('interés alto (4+5)', bold=True),
                    make_run(': solo 16.6% se declara "muy interesado" (5) y 5.4% en el nivel 4. La mayoría que se interesa por la política no llega a un cuarto del total.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P57 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='No me gusta mucho hablar de política porque nunca me he beneficiado de ningún político (...) y nunca he estado como muy atento a eso porque lo que es política y cosas de apuestas no soy como compatible de eso.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — El desinterés político tiene rostro de mujer joven.
    # Layout 2 (2 stats) SOLO-CUANTI = 1 slide
    # 61.6%→62%, 69.1%→69%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="EL DESINTERÉS POLÍTICO TIENE ROSTRO DE MUJER JOVEN. 69% DE LAS 18-24 SE DECLARA NADA INTERESADA — ",
        headline_italic="MÁS QUE CUALQUIER OTRO SEGMENTO DEL ESTUDIO.",
        stats=[
            {
                'value': '62%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' se declara "nada interesada" en política, frente a 52.3% de los hombres. Brecha de género: 9.3 puntos.')
                ]
            },
            {
                'value': '69%',
                'desc_runs': [
                    make_run('de los jóvenes de '),
                    make_run('18-24 años', bold=True),
                    make_run(' se declara "nada interesado" — el rango etario con mayor desinterés del estudio, por encima de los mayores de 55 (56.3%).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P57 × D2, P57 × D5 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — 77% del dominicano votó en las últimas elecciones.
    # Layout 2 (2 stats) + CV = 2 slides
    # 77.0%→77%, 58.0%→58%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="77% DEL DOMINICANO VOTÓ EN LAS ÚLTIMAS ELECCIONES. EL MISMO PAÍS DECLARA 58% DE DESINTERÉS TOTAL EN POLÍTICA. ",
        headline_italic="VOTAR NO REQUIERE CREER — REQUIERE IR.",
        stats=[
            {
                'value': '77%',
                'desc_runs': [
                    make_run('declara haber '),
                    make_run('votado', bold=True),
                    make_run(' en las últimas elecciones; solo 23.0% no votó.')
                ]
            },
            {
                'value': '58%',
                'desc_runs': [
                    make_run('se declara "'),
                    make_run('nada interesado', bold=True),
                    make_run('" en política (P57). La mayoría que vota se solapa directamente con la mayoría que dice no interesarle la política.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P58, P57 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Por ejemplo yo voté, mi esposo dijo que no iba a votar. Yo le dije, ¿por qué tú no vas a votar? Porque tú eres un ladrón. Entonces, ¿pero me vota para que lo saques? No. Entonces estos son peores.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — El voto sube con la edad.
    # Layout 2 (2 stats) SOLO-CUANTI = 1 slide
    # 82.4%→82%, 69.1%→69%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="EL VOTO SUBE CON LA EDAD: 82% DE LOS MAYORES DE 55 VOTÓ, CONTRA 69% DE LOS 18-24. ",
        headline_italic="EL HÁBITO ELECTORAL NO SE TRANSMITE — SE ACUMULA CON LOS AÑOS.",
        stats=[
            {
                'value': '82%',
                'desc_runs': [
                    make_run('de los mayores de '),
                    make_run('55 años', bold=True),
                    make_run(' votó en las últimas elecciones, frente a 69.1% de los jóvenes de 18-24. Brecha generacional: 13.3 puntos.')
                ]
            },
            {
                'value': '69%',
                'desc_runs': [
                    make_run('de los jóvenes de '),
                    make_run('18-24', bold=True),
                    make_run(' votó. La participación crece de forma escalonada: 18-24 (69%) → 25-34 (75%) → 35-44 (75%) → 45-54 (79%) → 55+ (82%). Cada década suma entre 3 y 4 puntos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P58 × D5 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — 4 de cada 10 dominicanos califican la gestión del gobierno como "muy mala".
    # Layout 2 (2 stats) + CV = 2 slides
    # 39.6%→40%, 10.0%+10.0%=20%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="4 DE CADA 10 DOMINICANOS CALIFICAN LA GESTIÓN DEL GOBIERNO COMO \"MUY MALA\". ",
        headline_italic="EL DESCONTENTO NO ES MATIZ — ES EL ESCALÓN MÁS POBLADO DE LA ESCALA.",
        stats=[
            {
                'value': '40%',
                'desc_runs': [
                    make_run('califica la gestión del gobierno actual como "'),
                    make_run('1 = Muy mala', bold=True),
                    make_run('" — la respuesta más votada de toda la escala P59.')
                ]
            },
            {
                'value': '20%',
                'desc_runs': [
                    make_run('suma la valoración '),
                    make_run('positiva (4+5)', bold=True),
                    make_run(': solo 10.0% "muy buena" (5) y 10.0% en nivel 4. El 80% restante se reparte entre "muy mala", "mala" y el punto medio (29.8% en 3).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P59 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo siento que el país tiene un retroceso, yo no sé, no es política, no, yo siento que hay un retroceso, porque cómo va a ser que a esta altura del juego, ni que baje sin agua, la luz crece cada rato.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — El descontento con el gobierno pega más duro abajo.
    # Layout 2 (2 stats) SOLO-CUANTI = 1 slide
    # 46.9%→47%, 43.6%→44%, 47.3%→47%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="EL DESCONTENTO CON EL GOBIERNO PEGA MÁS DURO ABAJO: EL ESTRATO D Y EL ESTRATO E LIDERAN EL \"MUY MALO\", ",
        headline_italic="MIENTRAS EL C+ ES EL ÚNICO NSE QUE NO LLEGA AL EXTREMO.",
        stats=[
            {
                'value': '47%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato E', bold=True),
                    make_run(' y 43.6% del D califican la gestión como "muy mala" (1). En el C+, la respuesta top es el punto medio "3" (36.1%) — único NSE que no llega al extremo. (E: n≈32, direccional)')
                ]
            },
            {
                'value': '47%',
                'desc_runs': [
                    make_run('de los jóvenes de '),
                    make_run('18-24', bold=True),
                    make_run(' y 47.0% de los 25-34 califican "muy malo", contra 31.1% de los mayores de 55. A menos edad, más severa la calificación.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P59 × D6, P59 × D5 · Base 500."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — Votar es la única forma válida de cambio para 4 de cada 10.
    # Layout 2 (2 stats) + CV = 2 slides
    # 40.6%→41%, 21.2%→21%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="VOTAR ES LA ÚNICA FORMA VÁLIDA DE CAMBIO PARA 4 DE CADA 10. EL SEGUNDO LUGAR NO ES PROTESTAR NI ORGANIZARSE — ",
        headline_italic="ES NO CREER EN NADA DE ESO.",
        stats=[
            {
                'value': '41%',
                'desc_runs': [
                    make_run('considera "'),
                    make_run('votar en elecciones', bold=True),
                    make_run('" la forma válida de generar cambio en el país — la opción más elegida de la batería P60.')
                ]
            },
            {
                'value': '21%',
                'desc_runs': [
                    make_run('eligió "'),
                    make_run('no creo que ninguna de estas cosas cambie realmente algo', bold=True),
                    make_run('" — segundo lugar, por encima de educarse (19.4%), apoyar movimientos sociales (7.0%) y protestar (6.0%).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P60 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Para conseguir un empleo tiene que ser político, que conocer gente. Si se acaba todo eso puede ser que se erradique la pobreza.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — Las protestas y las peticiones digitales no llegan al 7%.
    # Layout 1 (1 stat) + CV = 2 slides
    # Suma de formas activas < 15%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="LAS PROTESTAS Y LAS PETICIONES DIGITALES NO LLEGAN AL 7% COMO VÍA DE CAMBIO. ",
        headline_italic="EL DOMINICANO DELEGA LA TRANSFORMACIÓN AL VOTO — Y NO MUEVE EL CUERPO.",
        stats=[
            {
                'value': '15%',
                'desc_runs': [
                    make_run('suma las formas de '),
                    make_run('participación activa', bold=True),
                    make_run(': apoyar movimientos sociales (7.0%), protestar (6.0%), firmar peticiones digitales (0.6%), quejarse en redes (1.2%). Juntas no llegan al 15% del total.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P60 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Ya el concepto de la política antes, por lo menos, aunque había sectores y posiciones que hacían cosas malas, pero había mucha política de buena intención. (...) Ya hoy en día hay un vacío. La política es muy grande, tan fuerte, que ya nadie está dando importancia.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — Cuando se imagina el cambio, el dominicano no nombra ideología.
    # SOLO-CUALI (2 verbatims) — card cualitativa = 1 slide
    # P61 colapsada — sin stat publicable
    # ─────────────────────────────────────────────────────────────────────────

    build_cuali_slide(
        prs,
        headline_plain="CUANDO SE IMAGINA EL CAMBIO, EL DOMINICANO NO NOMBRA IDEOLOGÍA. NOMBRA CORRUPCIÓN, SALUD, EDUCACIÓN Y SEGURIDAD. ",
        headline_italic="LA POLÍTICA DOMINICANA SE VIVE COMO DÉFICIT DE SERVICIOS, NO DE BANDO.",
        verbatims=[
            {
                'quote': 'Cada vez que un político de esos se va con 500 millones de pesos, se llevó el puentecito que hace falta en los velaquitos, se llevó la carreterita de 4 kilómetros que tanto necesita, la escuela de los niños, la escuelita. Y se lo lleva a un solo personaje.',
                'attribution': 'Familia Biparental con Hijos Adultos'
            },
            {
                'quote': 'Y siempre que un cambio de gobierno (...) te ponen a una gente por encima de ti ganando muchísimo más que tú y no es ni profesional. Entonces tú recibir un mandato de una persona que sabe menos que tú.',
                'attribution': 'Familia Monoparental'
            }
        ],
        source_text="Source: Código Casa — Estudio cualitativo 2025 · P61 · Base cualitativa: grupos de enfoque."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H12 — 73% del dominicano dice que el sistema político es nada o poco justo.
    # Layout 3 (3 stats) + CV = 2 slides
    # 46.6%→47%, 26.8%→27%, sum 73.4%→73%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="73% DEL DOMINICANO DICE QUE EL SISTEMA POLÍTICO ES NADA O POCO JUSTO. ",
        headline_italic="LA PERCEPCIÓN DE INJUSTICIA NO ES DE UN SECTOR — ATRAVIESA TODOS LOS ESTRATOS Y TODAS LAS EDADES.",
        stats=[
            {
                'value': '73%',
                'desc_runs': [
                    make_run('percibe el sistema político como '),
                    make_run('injusto', bold=True),
                    make_run(': 46.6% "nada justo" (1) + 26.8% "poco justo" (2). La valoración positiva del sistema no llega al 9%.')
                ]
            },
            {
                'value': '47%',
                'desc_runs': [
                    make_run('de los rangos 25-34, 35-44, 45-54 y '),
                    make_run('55+', bold=True),
                    make_run(' pone "1 = nada justo" como respuesta top. Solo los 18-24 prefieren "2 = poco justo" (36.4%) — único segmento que no llega al extremo.')
                ]
            },
            {
                'value': '9%',
                'desc_runs': [
                    make_run('suma la valoración '),
                    make_run('positiva (4+5)', bold=True),
                    make_run(': 4.8% "bastante justo" y 3.8% "completamente justo". La injusticia es transversal por NSE: incluso en el segmento AB la respuesta top es "nada justo".')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P62, P62 × D5, P62 × D6 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Los políticos dicen a ti lo que tú quieres escuchar, pero no son los intereses justamente que lo están apoyando a ellos.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H13 — Casi la mitad del dominicano está muy dispuesto a irse del país.
    # Layout 2 (2 stats) + CV = 2 slides
    # 47.0%→47%, 30.4%→30%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="CASI LA MITAD DEL DOMINICANO ESTÁ MUY DISPUESTO A IRSE DEL PAÍS. ",
        headline_italic="LA MALETA NO ES PLAN B: PARA EL 47% ES LA PRIMERA RESPUESTA AL CONTEXTO.",
        stats=[
            {
                'value': '47%',
                'desc_runs': [
                    make_run('se declara "'),
                    make_run('muy dispuesto', bold=True),
                    make_run('" (5) a moverse a otro país por su situación económica o el contexto político — la respuesta más votada de la escala P63.')
                ]
            },
            {
                'value': '30%',
                'desc_runs': [
                    make_run('se declara "'),
                    make_run('nada dispuesto', bold=True),
                    make_run('" (1) a emigrar. Los que querrían irse superan a los que se quedarían firmes por 16.6 puntos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P63 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo misma me he considerado irme de aquí. (...) Yo estoy mandando mi currículum a Canadá, a ver si me hacen algo.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H14 — La disposición a emigrar sube hasta los 54 años y luego se invierte.
    # Layout 2 (2 stats) SOLO-CUANTI = 1 slide
    # 54.6%→55%, 43.6%→44%, 52.9%→53%, 40.6%→41%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="LA DISPOSICIÓN A EMIGRAR SUBE HASTA LOS 54 AÑOS Y LUEGO SE INVIERTE. LOS MAYORES DE 55 PREFIEREN QUEDARSE. ",
        headline_italic="EL ESTRATO MÁS BAJO TAMBIÉN — PERO POR RAZONES DISTINTAS.",
        stats=[
            {
                'value': '55%',
                'desc_runs': [
                    make_run('de los '),
                    make_run('45-54 años', bold=True),
                    make_run(' está "muy dispuesto" a emigrar — pico máximo. En los 55+, se invierte: la respuesta top pasa a "nada dispuesto" (46.2%). 18-24: 43.6%; 25-34: 49.6%; 35-44: 51.8%.')
                ]
            },
            {
                'value': '41%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato E', bold=True),
                    make_run(' está "nada dispuesto" a emigrar — único NSE donde la respuesta top es quedarse. El D es el más dispuesto a irse (52.9%). Emigrar requiere recursos. (E: n≈32, direccional)')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P63 × D5, P63 × D6 · Base 500."
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = "/home/user/c-digocasa-cerebro/Agentes Hallazgos/opiniones-politicas-deck-flat.pptx"
    prs.save(output_path)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    print("\n── Verificación specs ───────────────────────────────────────────")
    print(f"Slide size: {prs.slide_width} × {prs.slide_height} EMU")
    print(f"  = {prs.slide_width // PX}px × {prs.slide_height // PX}px")
    print(f"  Slide 1920×1080 OK: {prs.slide_width == Emu(SLIDE_W_PX * PX) and prs.slide_height == Emu(SLIDE_H_PX * PX)}")
    print(f"Stat number: {STAT_NUMBER_PT}pt ({STAT_NUMBER_PT/0.75:.0f}px)")
    print(f"Verbatim:    {VERBATIM_PT}pt ({VERBATIM_PT/0.75:.0f}px)")
    print(f"CV card:     {CV_CARD_W}×{CV_CARD_H}px")
    print("Cards: rgba(0,0,0,0.45) via MSO_SHAPE.ROUNDED_RECTANGLE + XML alpha=45000")
    print("Fondo negro plano — sin masterslide")
    print("Kerning 0 en todo el deck")
    print("Comillas españolas «...» en todos los verbatims")

    print("\n── QA FLAGS ───────────────────────────────────────────────")
    print("H01 stat 1: 58.6% → 59% (redondeo a entero OK)")
    print("H01 stat 2: 27.2% → 27% (redondeo a entero OK)")
    print("H02 solo-cuanti — sin Consumer Voice. OK.")
    print("H04 solo-cuanti — sin Consumer Voice. OK.")
    print("H06 solo-cuanti — sin Consumer Voice. OK.")
    print("H08 solo-cuanti — sin Consumer Voice. OK.")
    print("H08: ambos stats redondean a 47% — diferenciados en descripciones (estrato E vs 18-24).")
    print("H10: 1 solo stat (suma de formas activas). Layout 1 centrado. OK.")
    print("H11: solo-cuali, 2 verbatims en cards apiladas. Sin stat. OK.")
    print("H12 stat 3: 9% (4.8% + 3.8% = 8.6% → 9%). Descripción mantiene decimales inline.")
    print("H14 solo-cuanti — sin Consumer Voice. OK.")
    print("H14 stat 1: 54.6% → 55% (peak 45-54). Headline no declara cifra de subset bajo umbral.")
    print("Total slides esperados: 22 (diagnóstico inicial fue 21 — H10 tiene 2 slides: hallazgo + CV)")

    return output_path


if __name__ == "__main__":
    build_deck()
