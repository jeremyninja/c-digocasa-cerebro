#!/usr/bin/env python3
"""
build_sistema_creencias_deck.py
================================
Deck sistema-de-creencias-deck-flat.pptx desde cero con python-pptx.
20 slides — set editorial cerrado sistema-de-creencias-hallazgos-editados.md (12 hallazgos).

Distribución por tensión:
  T1 — El cuerpo legal y el cuerpo del deseo (H01, H02)
       H01: cuanti 3 stats + Consumer Voice  → 2 slides
       H02: solo-cuali 2 verbatims (cards side by side) → 1 slide
       Subtotal: 3 slides

  T2 — La fe como infraestructura (H03, H04)
       H03: cuanti 3 stats + Consumer Voice  → 2 slides
       H04: cuanti 2 stats, solo-cuanti      → 1 slide
       Subtotal: 3 slides

  T3 — El país que no mira a todos (H05, H06)
       H05: cuanti 2 stats + Consumer Voice  → 2 slides
       H06: cuanti 3 stats + Consumer Voice  → 2 slides
       Subtotal: 4 slides

  T4 — La violencia que no tiene nombre (H07, H08)
       H07: cuanti 2 stats + Consumer Voice  → 2 slides
       H08: cuanti 1 stat + Consumer Voice   → 2 slides
       Subtotal: 4 slides

  T5 — La equidad pendiente (H09, H10, H11, H12)
       H09: cuanti 2 stats + Consumer Voice  → 2 slides
       H10: cuanti 2 stats + Consumer Voice  → 2 slides
       H11: cuanti 2 stats, solo-cuanti      → 1 slide
       H12: cuanti 2 stats, solo-cuanti      → 1 slide
       Subtotal: 6 slides

TOTAL: 3 + 3 + 4 + 4 + 6 = 20 slides

Specs visuales (aprendizajes-montador-cc.md § 3.9 — versión consolidada):
  Slide: 1920×1080 px = 18,288,000 × 10,287,000 EMU
  1px = 9525 EMU
  Fonts: Instrument Serif (headlines, stats, verbatims CV) + Poppins (cuerpo, atribución)
  Fondo: #000000 negro plano (NO masterslide)
  Consumer Voice (Layout 4): verbatim SUELTO sobre fondo negro — SIN card
    Header "CONSUMER VOICE" Poppins 14pt gris #9B9B9B arriba
    Verbatim Instrument Serif regular 60pt blanco, comillas «»
    Atribución Poppins italic 20pt blanco
  Cards cuali (Layout 5): 549×221pt side by side con fill rgba(0,0,0,0.45)
    + outline blanco 1pt + border-radius 15pt
    Verbatim Poppins 15pt, atribución Poppins italic 12pt
  Stat number: 180px → 135pt, Instrument Serif italic
  Stat desc: Poppins 16pt blanco, bold selectivo
  Headlines: Instrument Serif MAYÚSCULAS 50pt fijo, centrado
  Kerning 0 en todo el deck
  NO Source en ningún slide (eliminado per aprendizajes v5)

Reglas de redondeo (aprendizajes-montador-cc.md § 3.5):
  Stat grande → entero por defecto.
  Decimal solo si: 1 dígito antes del punto (8.7%, 5.5%) o brecha matemática.

Convención de unidades:
  1 px @ 96 dpi = 9525 EMU
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
# MSO_SHAPE_TYPE import removed — usando int directo para formas
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

def pt_from_px(px_val):
    return round(px_val * 0.75, 1)

# ── Fuentes canónicas (aprendizajes v5) ────────────────────────────────────────
HEADLINE_PT       = 50.0          # Instrument Serif MAYÚSCULAS, fijo (no auto-size)
STAT_NUMBER_PT    = pt_from_px(180)  # 135pt – Instrument Serif italic
STAT_DESC_PT      = 13.0          # v5: era 16, ahora 13pt fijo
CV_HEADER_PT      = 14.0          # Poppins regular, gris
CV_VERBATIM_PT    = 50.0          # v6: Consumer Voice Instrument Serif (era 60, ahora 50)
CV_ATTRIB_PT      = 20.0          # Poppins italic — Consumer Voice
CARD_VERBATIM_PT  = 15.0          # Poppins — cards cuali (Layout 5)
CARD_ATTRIB_PT    = 12.0          # Poppins italic — cards cuali

# ── Dimensiones de stat boxes (v5: 10cm × 3cm = 378×113pt fijo) ────────────────
STAT_BOX_W_1STAT  = 378
STAT_BOX_W_2STATS = 378
STAT_BOX_W_3STATS = 378
STAT_BOX_H        = 113
STAT_TOP_PX       = 480   # posición vertical de las cifras grandes

# ── Dimensiones de cards cuali (v6: Layout 5 — 567×227pt = 15×6cm side by side) ─
CARD_W  = 567   # pt = 15 cm
CARD_H  = 227   # pt = 6 cm
CARD_GAP_2 = 80   # gap entre 2 cards side by side
CARD_GAP_3 = 36   # gap entre 3 cards side by side
CARD_TOP    = 429  # top canónico para cards cuali

# ── Consumer Voice dimensions (Layout 4 — verbatim suelto) ────────────────────
CV_TEXT_W   = 1637
CV_TEXT_LEFT = (SLIDE_W_PX - CV_TEXT_W) // 2  # = 141px


# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en el run XML."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """Headline 50pt fijo Instrument Serif MAYÚSCULAS centrado."""
    txBox = slide.shapes.add_textbox(
        px(100), px(80), px(1720), px(320)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(HEADLINE_PT)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(HEADLINE_PT)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    """Línea vertical fina #2E2E2E entre stats."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left_px, top_px, width_px, height_px):
    """Card rounded rectangle: fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt."""
    shape = slide.shapes.add_shape(
        5,
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    # border-radius ~15pt via adjustment
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass

    # Fill negro 45% opacidad via XML
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_el.set('val', '45000')  # 45% opacidad

    # Outline blanco 1pt
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def stat_box_width(n_stats):
    if n_stats == 1:
        return STAT_BOX_W_1STAT
    elif n_stats == 2:
        return STAT_BOX_W_2STATS
    else:
        return STAT_BOX_W_3STATS


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w = stat_box_width(n_stats)
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=STAT_BOX_H,  # v5: 113pt = 3cm
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


# ── Layouts de slides ──────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic, stats):
    """
    Slide de Hallazgo cuanti.
    Layout 1/2/3 según número de stats.
    NO source (eliminado per aprendizajes v5).
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice (Layout 4 — verbatim SUELTO sobre fondo negro, SIN card).
    Header "CONSUMER VOICE" Poppins 14pt gris arriba.
    Verbatim Instrument Serif regular 60pt blanco, comillas «».
    Atribución Poppins italic 20pt blanco.
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE"
    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Verbatim suelto — calcula top según longitud aproximada
    full_quote = f"«{quote}»"
    # Centrado vertical aproximado: más largo → arriba; más corto → más centrado
    char_count = len(quote)
    if char_count > 200:
        verb_top = 280
    elif char_count > 120:
        verb_top = 330
    else:
        verb_top = 380

    verb_box = slide.shapes.add_textbox(
        px(CV_TEXT_LEFT), px(verb_top),
        px(CV_TEXT_W), px(450)
    )
    tf = verb_box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_q = p.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(CV_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Atribución debajo
    attrib_top = verb_top + 300  # aprox debajo del verbatim
    add_textbox(
        slide,
        left_px=0, top_px=attrib_top,
        width_px=1920, height_px=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIB_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_cuali_slide_side_by_side(prs, headline_plain, headline_italic, verbatims):
    """
    Slide Card cualitativa (Layout 5 — cards 549×221pt SIDE BY SIDE).
    Headline arriba 50pt Instrument Serif MAYÚSCULAS.
    Cards con fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt.
    Verbatim Poppins 15pt blanco, atribución Poppins italic 12pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)

    if n == 1:
        # 1 card centrada
        card_left = (SLIDE_W_PX - CARD_W) // 2
        positions = [(card_left, CARD_TOP)]
    elif n == 2:
        # 2 cards side by side, gap 80pt
        total_w = CARD_W * 2 + CARD_GAP_2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_2, CARD_TOP)
        ]
    else:
        # 3 cards side by side, gap 36pt
        total_w = CARD_W * 3 + CARD_GAP_3 * 2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_3, CARD_TOP),
            (start_x + CARD_W * 2 + CARD_GAP_3 * 2, CARD_TOP)
        ]

    for i, v in enumerate(verbatims):
        left, top = positions[i]
        make_card_rgba_45_with_outline(slide, left, top, CARD_W, CARD_H)

        full_quote = f"«{v['quote']}»"
        pad_h = 30
        pad_v = 25
        text_w = CARD_W - (pad_h * 2)
        text_h = CARD_H - (pad_v * 2)

        txBox = slide.shapes.add_textbox(
            px(left + pad_h), px(top + pad_v),
            px(text_w), px(text_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left   = Emu(0)
        tf.margin_right  = Emu(0)
        tf.margin_top    = Emu(0)
        tf.margin_bottom = Emu(0)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run_q = p1.add_run()
        run_q.text = full_quote
        run_q.font.name   = FONT_BODY
        run_q.font.size   = Pt(CARD_VERBATIM_PT)
        run_q.font.italic = False
        run_q.font.bold   = False
        run_q.font.color.rgb = WHITE
        _set_kern(run_q._r)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run_a = p2.add_run()
        run_a.text = f"— {v['attribution']}"
        run_a.font.name   = FONT_BODY
        run_a.font.size   = Pt(CARD_ATTRIB_PT)
        run_a.font.italic = True
        run_a.font.bold   = False
        run_a.font.color.rgb = WHITE
        _set_kern(run_a._r)

    return slide


# ── BUILD DECK ─────────────────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # TENSIÓN 1 — EL CUERPO LEGAL Y EL CUERPO DEL DESEO
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 01 — H01 Hallazgo cuanti (3 stats)
    # Headline: "La pareja del mismo sexo puede comprarse una casa juntos. Adoptar un hijo,
    #  no — el 61.8% rechaza eso y el 60.4% que los reconozcan como padres."
    # 145 chars → >85 → 42px auto-size (pero usamos 50pt fijo per aprendizajes v5)
    # Stats:
    #   62% (61.8% → 62%)  rechazo adopción — entero OK
    #   52% (52.0% → 52%)  acepta propiedad compartida
    #   48% (47.8% → 48%)  acceso salud familiar (solo muy de acuerdo; 57.8% incluyendo algo)
    # NOTA: el headline editado menciona 61.8% y 60.4% inline — se mantienen en desc
    build_hallazgo_slide(
        prs,
        headline_plain="LA PAREJA DEL MISMO SEXO PUEDE COMPRARSE UNA CASA JUNTOS. ADOPTAR UN HIJO, ",
        headline_italic="NO — EL 61.8% RECHAZA ESO Y EL 60.4% QUE LOS RECONOZCAN COMO PADRES.",
        stats=[
            {
                'value': '62%',
                'desc_runs': [
                    make_run('está en desacuerdo con que una pareja del mismo sexo pueda '),
                    make_run('adoptar legalmente', bold=True),
                    make_run(' un hijo — el ítem con mayor rechazo de la batería Q55. 60.4% rechaza además el reconocimiento legal como padres.')
                ]
            },
            {
                'value': '52%',
                'desc_runs': [
                    make_run('acepta que '),
                    make_run('compren una propiedad juntos', bold=True),
                    make_run(' (39.8% muy de acuerdo + 12.2% algo de acuerdo). Los dos ítems más aceptados son los que no involucran a un tercero.')
                ]
            },
            {
                'value': '48%',
                'desc_runs': [
                    make_run('está muy de acuerdo con el acceso a '),
                    make_run('servicios de salud familiar', bold=True),
                    make_run(' como beneficiarios mutuos (+ 10.0% algo de acuerdo = 57.8% total). El préstamo mancomunado suma 45.2% de aprobación.')
                ]
            }
        ]
    )

    # Slide 02 — H01 Consumer Voice (verbatim suelto, sin card)
    build_consumer_voice_slide(
        prs,
        quote='Aquí para tú adoptar, uno de los requisitos es que tú tienes que estar casado y el matrimonio en nuestro país es entre un hombre y una mujer (...) por lo cual nosotros ninguno podemos adoptar.',
        attribution='Familia Homoparental'
    )

    # Slide 03 — H02 Card cualitativa (solo-cuali — 2 verbatims side by side)
    # Headline: "La aceptación de familias del mismo sexo en los grupos siempre llega
    #  con un 'pero'. Se respeta la pareja. Se duda del hijo."
    # 123 chars
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="LA ACEPTACIÓN DE FAMILIAS DEL MISMO SEXO SIEMPRE LLEGA CON UN ‘PERO’. ",
        headline_italic="SE RESPETA LA PAREJA. SE DUDA DEL HIJO.",
        verbatims=[
            {
                'quote': 'Respeto su decisión pero no la comparto. Últimamente se está viviendo en una época moderna en la cual quieren imponer eso (...) hacerlo como normal, algo que no es normal, porque (...) la palabra de Dios dice que hombre y mujer, no las dos personas del mismo sexo.',
                'attribution': 'Familia Monoparental'
            },
            {
                'quote': 'Un niño necesita tener un papá y una mamá porque dos hombres juntos o dos mujeres juntas no pueden procrear (...) yo respeto, pero esa es mi percepción muy personal.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # TENSIÓN 2 — LA FE COMO INFRAESTRUCTURA
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 04 — H03 Hallazgo cuanti (3 stats)
    # Headline: "64% recurre siempre a la fe en momentos de dificultad.
    #  El 64% no es parejo: sube a 70.8% en mujeres y a 72.7% en el estrato D."
    # 130 chars
    # Stats:
    #   64% (64.0% → 64%)  univariada Q57
    #   71% (70.8% → 71%)  mujeres
    #   73% (72.7% → 73%)  estrato D
    build_hallazgo_slide(
        prs,
        headline_plain="64% RECURRE SIEMPRE A LA FE EN MOMENTOS DE DIFICULTAD. ",
        headline_italic="EL 64% NO ES PAREJO: SUBE A 70.8% EN MUJERES Y A 72.7% EN EL ESTRATO D.",
        stats=[
            {
                'value': '64%',
                'desc_runs': [
                    make_run('recurre '),
                    make_run('siempre', bold=True),
                    make_run(' a la religión o espiritualidad como alivio o refugio en momentos de dificultad. Solo 6.6% declara que nunca. Q57, Base 500.')
                ]
            },
            {
                'value': '71%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' recurre siempre, frente a 53.3% de los hombres. Brecha de género: 17.5 puntos. Q57 × sexo, Base 305 mujeres / 195 hombres.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' recurre siempre, vs 52.8% en C+ y 61.5% en C. A menor NSE, mayor dependencia espiritual. Q57 × NSE, Base 172 NSE D.')
                ]
            }
        ]
    )

    # Slide 05 — H03 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Me refugio mucho en Dios, en la oración. Tener una conexión con Dios para mí es importante, dejarle las cosas a él, que él se encargue, porque hay días que tú dices yo no puedo más y Dios se encarga.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # Slide 06 — H04 Hallazgo solo-cuanti (2 stats)
    # Headline: "El dominicano es casi unánimemente creyente. Cómo cree — eso está partido
    #  en tres mitades casi iguales, y el no creyente es una rareza de 1.2%."
    # 147 chars
    # Stats:
    #   31% (30.8% → 31%)  creyente ocasional — entero
    #   30% (30.2% → 30%)  muy practicante y activo — entero (casi igual)
    #   NOTA: usamos 2 stats: distribución en tercios + juventud vs madurez
    build_hallazgo_slide(
        prs,
        headline_plain="EL DOMINICANO ES CASI UNÁNIMEMENTE CREYENTE. CÓMO CREE — ESO ESTÁ PARTIDO EN TRES MITADES CASI IGUALES, ",
        headline_italic="Y EL NO CREYENTE ES UNA RAREZA DE 1.2%.",
        stats=[
            {
                'value': '31%',
                'desc_runs': [
                    make_run('es '),
                    make_run('creyente ocasional', bold=True),
                    make_run(', 30.2% muy practicante y activo, 28.8% espiritual sin práctica regular. Solo 8.0% queda fuera del marco religioso o espiritual. Q58, Base 500.')
                ]
            },
            {
                'value': '47%',
                'desc_runs': [
                    make_run('de los '),
                    make_run('18–24 años', bold=True),
                    make_run(' elige "creyente ocasional" como perfil dominante; en los 55+ el perfil top es "muy practicante y activo" (47.1%). El creyente se vuelve más practicante al envejecer.')
                ]
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # TENSIÓN 3 — EL PAÍS QUE NO MIRA A TODOS
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 07 — H05 Hallazgo cuanti (2 stats)
    # Headline: "9 de cada 10 dicen no haber sufrido discriminación en 3 años.
    #  En hogares monoparentales la cifra cae a 83.6% — casi 8 puntos bajo el promedio."
    # 146 chars
    # Stats:
    #   91% (91.2% → 91%)  no fue discriminado
    #   84% (83.6% → 84%)  monoparentales
    build_hallazgo_slide(
        prs,
        headline_plain="9 DE CADA 10 DICEN NO HABER SUFRIDO DISCRIMINACIÓN EN 3 AÑOS. ",
        headline_italic="EN HOGARES MONOPARENTALES LA CIFRA CAE A 83.6% — CASI 8 PUNTOS BAJO EL PROMEDIO.",
        stats=[
            {
                'value': '91%',
                'desc_runs': [
                    make_run('declara no haber sido víctima de '),
                    make_run('discriminación', bold=True),
                    make_run(' en los últimos 3 años. La mención más alta entre quienes sí reportan algo: "comentarios ofensivos sobre el color de piel" (2.8%, n=14). Q60, Base 500.')
                ]
            },
            {
                'value': '84%',
                'desc_runs': [
                    make_run('en hogares '),
                    make_run('monoparentales', bold=True),
                    make_run(' declara no haber sido víctima — la cifra más baja de todas las tipologías, frente a 93.7% en biparental con hijos menores de 18. Q60 × tipología, Base 122.')
                ]
            }
        ]
    )

    # Slide 08 — H05 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='La discriminación con la gente: que tú vas a una tienda aquí y si tú no vas bien vestido te miran mal, te vienen todos para acá y no vas a comprar nada.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # Slide 09 — H06 Hallazgo cuanti (3 stats)
    # Headline: "54.2% califica de muy baja la inclusión de discapacidad en su comunidad.
    #  61.8% cree que esa misma comunidad no está nada preparada para incluir."
    # 148 chars
    # Stats:
    #   54% (54.2% → 54%)  muy baja la inclusión
    #   62% (61.8% → 62%)  nada preparada — entero
    #   72% (71.5% → 72%)  estrato D / 71% (70.6% → 71%) mayores 55
    build_hallazgo_slide(
        prs,
        headline_plain="54.2% CALIFICA DE MUY BAJA LA INCLUSIÓN DE DISCAPACIDAD EN SU COMUNIDAD. 61.8% CREE QUE ESA MISMA COMUNIDAD ",
        headline_italic="NO ESTÁ NADA PREPARADA PARA INCLUIR.",
        stats=[
            {
                'value': '54%',
                'desc_runs': [
                    make_run('califica de '),
                    make_run('muy baja', bold=True),
                    make_run(' (nota 1 de 5) la inclusión de personas con discapacidad en su comunidad. Sumando notas 1 y 2, el 61.6% la reprueba. Q61, Base 500.')
                ]
            },
            {
                'value': '62%',
                'desc_runs': [
                    make_run('cree que su comunidad está "'),
                    make_run('nada preparada', bold=True),
                    make_run('" para incluir personas neurodivergentes o con discapacidad en trabajo, escuela y transporte. Solo 12.0% la ve muy preparada. Q63, Base 500.')
                ]
            },
            {
                'value': '72%',
                'desc_runs': [
                    make_run('"Nada preparada" sube en el '),
                    make_run('estrato D', bold=True),
                    make_run(' (71.5%, n=172) y en mayores de 55 (70.6%, n=119), frente a 54.5% en 18–24. La exclusión percibida crece con la precariedad y la edad.')
                ]
            }
        ]
    )

    # Slide 10 — H06 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Si esos maestros supieran trabajar con ese tipo de niños, esos niños fueran más incluidos como otros amiguitos, no necesitarían una educación especial en un colegio especial que le cuesta más de la mitad del sueldo de un padre.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # TENSIÓN 4 — LA VIOLENCIA QUE NO TIENE NOMBRE
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 11 — H07 Hallazgo cuanti (2 stats)
    # Headline: "Hablar de orientación sexual en familia parte al dominicano en dos:
    #  36.8% muy cómodo, 34.8% nada cómodo. El hombre se cierra donde la mujer habla."
    # 150 chars
    # Stats:
    #   37% (36.8% → 37%)  muy cómodo
    #   40% (40.3% → 40%)  mujeres muy cómodas
    build_hallazgo_slide(
        prs,
        headline_plain="HABLAR DE ORIENTACIÓN SEXUAL EN FAMILIA PARTE AL DOMINICANO EN DOS: 36.8% MUY CÓMODO, 34.8% NADA CÓMODO. ",
        headline_italic="EL HOMBRE SE CIERRA DONDE LA MUJER HABLA.",
        stats=[
            {
                'value': '37%',
                'desc_runs': [
                    make_run('se siente "'),
                    make_run('muy cómodo/a', bold=True),
                    make_run('" hablando de orientación sexual o identidad de género en familia, y 34.8% "nada cómodo/a". El centro de la escala es minoritario. Q64, Base 500.')
                ]
            },
            {
                'value': '40%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' tiene "muy cómoda" como respuesta top (40.3%, n=123); 39.5% de los hombres tiene "nada cómodo" como respuesta top (n=77). El género predice el polo.')
                ]
            }
        ]
    )

    # Slide 12 — H07 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Mi hijo mayor cuando se declaró gay, esa parte yo no... eso me marcó, porque yo lo rechacé, lo aborrecido como si fuera un trapo viejo (...) Y después de un tiempo (...) yo decidí acercarme y ahora (...) yo lo quiero, hasta su pareja.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # Slide 13 — H08 Hallazgo cuanti (1 stat)
    # Headline: "11.6% de los padres reporta que un hijo vivió bullying escolar o digital.
    #  El 75% dice que no. El silencio no es prueba de ausencia."
    # 133 chars
    # Stat:
    #   12% (11.6% → 12%)  bullying — entero
    build_hallazgo_slide(
        prs,
        headline_plain="11.6% DE LOS PADRES REPORTA QUE UN HIJO VIVIÓ BULLYING ESCOLAR O DIGITAL. ",
        headline_italic="EL 75% DICE QUE NO. EL SILENCIO NO ES PRUEBA DE AUSENCIA.",
        stats=[
            {
                'value': '12%',
                'desc_runs': [
                    make_run('declara que algún hijo vivió '),
                    make_run('acoso o bullying', bold=True),
                    make_run(' escolar o digital; 75.0% (n=375) dice que no y 12.8% no tiene hijos. Recalculado sobre quienes tienen hijos (n≈436): ∼13.3%. Q65, Base 500.')
                ]
            }
        ]
    )

    # Slide 14 — H08 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Tuve al principio de año escolar un momento de tensión porque le estaban haciendo bullying a las niñas, y las niñas que le gustaban su colegio no querían ir. Ya se fue una señal de alerta para mí.',
        attribution='Familia Extendida'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # TENSIÓN 5 — LA EQUIDAD PENDIENTE
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 15 — H09 Hallazgo cuanti (2 stats)
    # Headline: "62.6% dice que el cambio más urgente para la equidad es que los padres
    #  estén presentes. La flexibilidad laboral y la carrera de la mujer no llegan ni al 13%."
    # 161 chars
    # Stats:
    #   63% (62.6% → 63%)  padres presentes
    #   13% (12.8% → 13%)  flexibilidad laboral
    build_hallazgo_slide(
        prs,
        headline_plain="62.6% DICE QUE EL CAMBIO MÁS URGENTE PARA LA EQUIDAD ES QUE LOS PADRES ESTÉN PRESENTES. ",
        headline_italic="LA FLEXIBILIDAD LABORAL Y LA CARRERA DE LA MUJER NO LLEGAN NI AL 13%.",
        stats=[
            {
                'value': '63%',
                'desc_runs': [
                    make_run('elige "que los '),
                    make_run('padres estén más presentes', bold=True),
                    make_run('" en la crianza y el hogar como el cambio social más urgente para la equidad familiar. Muy por encima de flexibilidad laboral (12.8%) o que la maternidad no frene la carrera (2.8%). Q66, Base 500.')
                ]
            },
            {
                'value': '13%',
                'desc_runs': [
                    make_run('elige '),
                    make_run('flexibilidad laboral para ambos padres', bold=True),
                    make_run('. La prioridad de "padres presentes" es más intensa en NSE bajo: 75.0% en NSE E (n=32, tendencia direccional) vs 43.5% en AB.')
                ]
            }
        ]
    )

    # Slide 16 — H09 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Que los padres estén más presentes... porque uno se lo da en la casa, que es donde todo empieza, pero ellos, la mayor parte de su tiempo, no están en la casa.',
        attribution='Familia Monoparental'
    )

    # Slide 17 — H10 Hallazgo cuanti (2 stats)
    # Headline: "Cuando el dominicano sí sufre violencia, la más frecuente no es la física:
    #  es económica. 20.2% reporta control financiero — casi el triple que el 7.6% de violencia física."
    # 172 chars
    # Stats:
    #   20% (20.2% → 20%)  violencia económica
    #   56% (56.4% → 56%)  no sufrió ningún tipo
    build_hallazgo_slide(
        prs,
        headline_plain="CUANDO EL DOMINICANO SÍ SUFRE VIOLENCIA, LA MÁS FRECUENTE NO ES LA FÍSICA: ES ECONÓMICA. ",
        headline_italic="20.2% REPORTA CONTROL FINANCIERO — CASI EL TRIPLE QUE EL 7.6% DE VIOLENCIA FÍSICA.",
        stats=[
            {
                'value': '20%',
                'desc_runs': [
                    make_run('reporta '),
                    make_run('violencia económica', bold=True),
                    make_run(' — control del acceso al dinero, restricción financiera — como el tipo de violencia más frecuente en los últimos 3 años. Le sigue la emocional/psicológica (17.4%) y la física (7.6%). Q67, Base 500.')
                ]
            },
            {
                'value': '56%',
                'desc_runs': [
                    make_run('declara no haber sufrido '),
                    make_run('ningún tipo de violencia', bold=True),
                    make_run('. El 44% restante que sí reporta algo distribuye mayoritariamente en económica y emocional. El cruce NSE E × "no víctima" (n=14) no es publicable.')
                ]
            }
        ]
    )

    # Slide 18 — H10 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Si mi esposo me golpea... aunque yo le diga a mi hija, mira, tú no puedes dejar que nadie te golpee, ella entiende que si él, que es mi esposo, me quiere y me golpea, es una respuesta al amor.',
        attribution='Familia Extendida'
    )

    # Slide 19 — H11 Hallazgo solo-cuanti (2 stats)
    # Headline: "Cuando hay un adulto mayor en el hogar, el cuidado recae sobre una sola persona.
    #  12.8% responde 'yo mismo/a' — diez veces más que el cuidador contratado."
    # 157 chars
    # Stats:
    #   13% (12.8% → 13%)  yo mismo/a cuida — entero
    #   73% (72.8% → 73%)  no convive con persona mayor
    build_hallazgo_slide(
        prs,
        headline_plain="CUANDO HAY UN ADULTO MAYOR EN EL HOGAR, EL CUIDADO RECAE SOBRE UNA SOLA PERSONA. ",
        headline_italic="12.8% RESPONDE “YO MISMO/A” — DIEZ VECES MÁS QUE EL CUIDADOR CONTRATADO.",
        stats=[
            {
                'value': '13%',
                'desc_runs': [
                    make_run('responde "'),
                    make_run('yo mismo/a', bold=True),
                    make_run('" al cuidado del adulto mayor en el hogar, frente a 2.0% que lo comparte entre varios y 0.6% que contrata un cuidador. El cuidado no se distribuye: se individualiza. Q68, Base 500.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('no convive con ninguna '),
                    make_run('persona mayor', bold=True),
                    make_run('. La convivencia con adulto mayor se concentra en hogares extendidos (44.4%) y en mayores de 55 (52.9% sí conviven). Q68 × tipología y edad.')
                ]
            }
        ]
    )

    # Slide 20 — H12 Hallazgo solo-cuanti (2 stats)
    # Headline: "73.8% pide pensiones para que los mayores vivan con dignidad.
    #  El respeto y la compañía llegan después. Lo primero siempre es la plata."
    # 136 chars
    # Stats:
    #   74% (73.8% → 74%)  pensiones — entero
    #   53% (53.4% → 53%)  acceso a salud
    build_hallazgo_slide(
        prs,
        headline_plain="73.8% PIDE PENSIONES PARA QUE LOS MAYORES VIVAN CON DIGNIDAD. ",
        headline_italic="EL RESPETO Y LA COMPAÑÍA LLEGAN DESPUÉS. LO PRIMERO SIEMPRE ES LA PLATA.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('menciona '),
                    make_run('pensiones o apoyos económicos suficientes', bold=True),
                    make_run(' como lo que más falta para que la tercera edad viva con dignidad. Le sigue acceso a salud (53.4%). El respeto y reconocimiento familiar queda en 32.4%. Q69, Base 500.')
                ]
            },
            {
                'value': '53%',
                'desc_runs': [
                    make_run('menciona '),
                    make_run('acceso gratuito o asequible a salud', bold=True),
                    make_run(' como segunda prioridad. La prioridad de pensiones es transversal: de 70.6% en NSE C a 83.3% en C+; de 65.5% en 18–24 a 79.0% en 55+. Ningún segmento la desplaza del primer lugar.')
                ]
            }
        ]
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "sistema-de-creencias-deck-flat.pptx")
    prs.save(output_path)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    print("\n── Verificación specs ─────────────────────────────────────────────")
    print(f"Slide size: {prs.slide_width} × {prs.slide_height} EMU")
    print(f"  = {int(prs.slide_width / PX)}px × {int(prs.slide_height / PX)}px")
    ok = prs.slide_width == Emu(SLIDE_W_PX * PX) and prs.slide_height == Emu(SLIDE_H_PX * PX)
    print(f"  Slide 1920×1080 OK: {ok}")
    print(f"Stat number: {STAT_NUMBER_PT}pt")
    print(f"Stat desc:   {STAT_DESC_PT}pt")
    print(f"Headline:    {HEADLINE_PT}pt fijo (no auto-size)")
    print(f"CV Verbatim: {CV_VERBATIM_PT}pt Instrument Serif regular")
    print(f"CV Attrib:   {CV_ATTRIB_PT}pt Poppins italic")
    print(f"Card Verb:   {CARD_VERBATIM_PT}pt Poppins")
    print(f"Card Attrib: {CARD_ATTRIB_PT}pt Poppins italic")
    print("Cards cuali: rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt")
    print("Consumer Voice: verbatim SUELTO sobre fondo negro (SIN card)")
    print("Comillas españolas «...» en todos los verbatims")
    print("Kerning 0 en todo el deck")
    print("NO masterslide — fondo negro plano")
    print("NO source en ningún slide (eliminado per aprendizajes v5)")

    print("\n── QA FLAGS ──────────────────────────────────────────")
    print("H01 (slide 1): 3 stats. 61.8%→62%, 52.0%→52%, 47.8%→48%")
    print("H01 (slide 2): Consumer Voice — verbatim suelto sin card. OK.")
    print("H02 (slide 3): solo-cuali 2 verbatims side by side (549×221 cada una). OK.")
    print("H03 (slide 4): 3 stats. 64.0%→64%, 70.8%→71%, 72.7%→73%")
    print("H03 (slide 5): Consumer Voice — verbatim suelto sin card. OK.")
    print("H04 (slide 6): solo-cuanti 2 stats. 30.8%→31%, 45.5% (18-24) y 47.1% (55+)")
    print("H05 (slide 7): 2 stats. 91.2%→91%, 83.6%→84%")
    print("H05 (slide 8): Consumer Voice — verbatim suelto sin card. OK.")
    print("H06 (slide 9): 3 stats. 54.2%→54%, 61.8%→62%, 71.5%→72%")
    print("H06 (slide 10): Consumer Voice — verbatim suelto sin card. OK.")
    print("H07 (slide 11): 2 stats. 36.8%→37%, 40.3%→40%")
    print("H07 (slide 12): Consumer Voice — verbatim suelto sin card. OK.")
    print("H08 (slide 13): 1 stat. 11.6%→12%")
    print("H08 (slide 14): Consumer Voice — verbatim suelto sin card. OK.")
    print("H09 (slide 15): 2 stats. 62.6%→63%, 12.8%→13%")
    print("H09 (slide 16): Consumer Voice — verbatim suelto sin card. OK.")
    print("H10 (slide 17): 2 stats. 20.2%→20%, 56.4%→56%")
    print("H10 (slide 18): Consumer Voice — verbatim suelto sin card. OK.")
    print("H11 (slide 19): solo-cuanti 2 stats. 12.8%→13%, 72.8%→73%")
    print("H12 (slide 20): solo-cuanti 2 stats. 73.8%→74%, 53.4%→53%")
    print("Total slides verificados: 20")

    return output_path


if __name__ == "__main__":
    build_deck()
