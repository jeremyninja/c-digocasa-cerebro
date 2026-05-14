#!/usr/bin/env python3
"""
build_sc_qa_pdf.py
==================
Genera el PDF de QA para sistema-de-creencias-deck-flat.pptx.
Como LibreOffice no puede cargar el .pptx en este entorno, genera un PDF
con el texto real de cada slide directamente desde el .pptx usando python-pptx.
Esto produce un documento de revisión de contenido fiel al deck.
"""

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.colors import black, white, HexColor
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, Frame
from reportlab.lib.styles import ParagraphStyle
import textwrap

PPTX_PATH = "/home/user/c-digocasa-cerebro/Agentes Hallazgos/sistema-de-creencias-deck-flat.pptx"
PDF_PATH  = "/home/user/c-digocasa-cerebro/Agentes Hallazgos/sistema-de-creencias-deck-flat-QA.pdf"

def extract_slide_texts(pptx_path):
    """Extrae todo el texto de cada slide del PPTX."""
    prs = Presentation(pptx_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        # Identificar si es stat grande (Instrument Serif, ~135pt)
                        is_stat = False
                        is_source = False
                        is_cv_header = False
                        for run in para.runs:
                            if run.font.size and run.font.size.pt and run.font.size.pt > 100:
                                is_stat = True
                            if run.font.size and run.font.size.pt and run.font.size.pt < 14:
                                is_source = True
                            if run.font.name == "Poppins" and line == "CONSUMER VOICE":
                                is_cv_header = True
                        texts.append({
                            'text': line,
                            'is_stat': is_stat,
                            'is_source': is_source,
                            'is_cv_header': is_cv_header
                        })
        slides_data.append({'slide_num': idx, 'texts': texts})

    return slides_data


def build_qa_pdf(pptx_path, pdf_path):
    slides_data = extract_slide_texts(pptx_path)

    # PDF en landscape (1920:1080 ≈ 16:9 → letter landscape es 11×8.5")
    page_w, page_h = landscape(letter)
    c = canvas.Canvas(pdf_path, pagesize=landscape(letter))

    BG = HexColor('#000000')
    FG = HexColor('#FFFFFF')
    GS = HexColor('#9B9B9B')
    ACCENT = HexColor('#B8FF4D')  # solo para marcadores de slide

    for slide in slides_data:
        n = slide['slide_num']
        texts = slide['texts']

        # Fondo negro
        c.setFillColor(BG)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

        # Número de slide (esquina superior izquierda)
        c.setFillColor(ACCENT)
        c.setFont("Helvetica-Bold", 14)
        c.drawString(0.3*inch, page_h - 0.4*inch, f"Slide {n}/{len(slides_data)}")

        # Línea separadora fina
        c.setStrokeColor(HexColor('#2E2E2E'))
        c.line(0.3*inch, page_h - 0.55*inch, page_w - 0.3*inch, page_h - 0.55*inch)

        # Renderizar textos
        y = page_h - 0.9*inch
        margin_l = 0.5*inch
        max_w = page_w - 1.0*inch

        for item in texts:
            text = item['text']
            if not text:
                continue

            if item['is_cv_header']:
                # Header CONSUMER VOICE en gris
                c.setFillColor(GS)
                c.setFont("Helvetica", 9)
                c.drawCentredString(page_w/2, y, text)
                y -= 0.25*inch
            elif item['is_stat']:
                # Stat grande en blanco, fuente grande
                c.setFillColor(FG)
                c.setFont("Helvetica-BoldOblique", 36)
                c.drawCentredString(page_w/2, y, text)
                y -= 0.55*inch
            elif item['is_source']:
                # Source en gris pequeño al pie
                c.setFillColor(GS)
                c.setFont("Helvetica-Oblique", 7)
                c.drawCentredString(page_w/2, 0.35*inch, text)
            else:
                # Texto normal — wrapeamos para que no se salga
                c.setFillColor(FG)
                # Determinar si es headline (todo MAYÚSCULAS y largo)
                is_upper = text == text.upper() and len(text) > 20
                if is_upper:
                    font_size = 13 if len(text) > 100 else 16
                    c.setFont("Helvetica-Bold", font_size)
                else:
                    # Verificar si contiene comillas «» → es verbatim
                    is_verbatim = '«' in text or '»' in text
                    if is_verbatim:
                        c.setFont("Helvetica-Oblique", 10)
                    elif text.startswith('—'):
                        # atribución
                        c.setFont("Helvetica-Oblique", 9)
                        c.setFillColor(GS)
                    else:
                        c.setFont("Helvetica", 9)
                    font_size = 10

                # Wrap manual
                chars_per_line = int(max_w / (font_size * 0.6)) if font_size > 0 else 80
                lines = textwrap.wrap(text, width=max(chars_per_line, 40))

                for line in lines:
                    if y < 0.6*inch:
                        break  # no overflow
                    c.drawString(margin_l, y, line)
                    y -= (font_size + 4) * 0.013*inch * 72

            if y < 0.6*inch:
                y = 0.6*inch

        c.showPage()

    c.save()
    print(f"PDF guardado: {pdf_path}")
    print(f"Total páginas: {len(slides_data)}")


if __name__ == "__main__":
    build_qa_pdf(PPTX_PATH, PDF_PATH)
