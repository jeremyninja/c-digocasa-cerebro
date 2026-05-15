#!/usr/bin/env python3
"""
build_mujer_v2.py
=================
Deck mujer-v2-deck-flat.pptx desde cero con python-pptx.
Pilar 11 — MUJER v2 — 14 hallazgos, fuentes duales CC + Finance Women (FW).

Specs v5 activas (aprendizajes-montador-cc.md § 3.9 v5):
- Slide: 1920×1080 px → 18_288_000 × 10_287_000 EMU
- Headline: 50pt fijo, Instrument Serif MAYÚSCULAS, centrado
- Stat number: 180pt Instrument Serif italic
- Stat box descripción: 378×113pt (10×3cm), Poppins 13pt
- Consumer Voice: sin card, verbatim suelto Instrument Serif regular 60pt, header Poppins 14pt gris
- Cards cualitativas: 549×221pt, fill rgba(0,0,0,0.45), outline blanco 1pt, border-radius 15pt
- Verbatim en cards cuali: Poppins 15pt blanco, comillas «...»
- Atribución: Poppins italic 20pt (Consumer Voice) / 12pt (cards cuali)
- Kerning 0 en todo el deck
- Source: Poppins italic 12pt, gris GREY_SOFT, centrado, top 1010pt (solo hallazgos, no portadas)
- NO masterslide decorativo

Convención de unidades:
  1 pt Keynote = 1 px @ 96 dpi en python-pptx
  1 px @ 96 dpi = 9525 EMU
  Se define PX = 9525 y se usa como multiplicador.
  Las specs v5 están en pt — se usan directamente como px en este script.
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Conversión px/pt → EMU ─────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU (también 1 pt Keynote = 9525 EMU)

def px(n):
    """Convierte píxeles/puntos Keynote a EMU."""
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W = 1920   # pt (equiv px @ 96dpi)
SLIDE_H = 1080   # pt

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)   # separadores verticales entre stats
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)   # header CONSUMER VOICE + source

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── Tamaños en pt (v5 activa — specs exactas de aprendizajes-montador-cc.md) ──
HEADLINE_PT       = 50    # headlines — fijo, no auto-size
STAT_NUMBER_PT    = 180   # cifra dominante Instrument Serif italic
STAT_DESC_PT      = 13    # descripción stat Poppins (v5: era 16pt, ahora 13pt)
CV_VERBATIM_PT    = 50    # v6: Consumer Voice verbatim Instrument Serif (era 60, ahora 50)
CV_HEADER_PT      = 14    # "CONSUMER VOICE" Poppins regular
CV_ATTRIBUTION_PT = 20    # atribución Consumer Voice Poppins italic
CARD_VERBATIM_PT  = 15    # verbatim dentro de card cualitativa Poppins
CARD_ATTRIBUTION_PT = 12  # atribución card cualitativa Poppins italic
SOURCE_PT         = 12    # source line: Poppins italic, GREY_SOFT, centrado

# ── Dimensiones de stat box (v5: 378×113pt — 10×3cm) ─────────────────────────
STAT_BOX_W = 378   # pt = 10 cm
STAT_BOX_H = 113   # pt = 3 cm

# ── Dimensiones de card cualitativa (v6: 567×227pt — 15×6cm) ──────────────────
CARD_W = 567   # pt = 15 cm
CARD_H = 227   # pt = 6 cm

# ── Posiciones verticales clave ────────────────────────────────────────────────
HEADLINE_TOP   = 100   # pt — headline arriba
STAT_NUM_TOP   = 420   # pt — cifra grande
STAT_DESC_TOP  = 660   # pt — caja de descripción debajo de cifra
CV_HEADER_TOP  = 80    # pt — header "CONSUMER VOICE"
CV_VERBATIM_TOP = 370  # pt — inicio del bloque de verbatim Consumer Voice


# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    """Presentation con slide 1920×1080 pt."""
    prs = Presentation()
    prs.slide_width  = px(SLIDE_W)
    prs.slide_height = px(SLIDE_H)
    return prs


def blank_slide(prs):
    """Slide en blanco con fondo negro plano."""
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en <a:r>."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left, top, width, height,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True,
                left_px=None, top_px=None, width_px=None, height_px=None):
    """Textbox simple de un solo run, posicionado en pt."""
    # Acepta también parámetros _px para compatibilidad con add_source
    l = left_px if left_px is not None else left
    t = top_px  if top_px  is not None else top
    w = width_px if width_px is not None else width
    h = height_px if height_px is not None else height

    txBox = slide.shapes.add_textbox(
        px(l), px(t), px(w), px(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_source(slide, source_text):
    """
    Line de source en la parte inferior del slide.
    Poppins italic SOURCE_PT (12pt), GREY_SOFT, centrado, top 1010pt.
    Solo se llama en slides de hallazgo — NO en portadas de tensión.
    """
    add_textbox(
        slide,
        left=0, top=1010,
        width=1920, height=50,
        text=source_text,
        font_name=FONT_BODY, font_size_pt=SOURCE_PT,
        italic=True, color=GREY_SOFT,
        align=PP_ALIGN.CENTER
    )


def add_rich_textbox(slide, left, top, width, height,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox con múltiples runs en el mismo párrafo."""
    txBox = slide.shapes.add_textbox(
        px(left), px(top), px(width), px(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """
    Headline MAYÚSCULAS, 50pt fijo, centrado.
    left 100, top 100, width 1720pt.
    Split plain/italic para la frase del giro.
    """
    txBox = slide.shapes.add_textbox(
        px(100), px(HEADLINE_TOP), px(1720), px(300)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(HEADLINE_PT)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(HEADLINE_PT)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_pt, y_top_pt, height_pt):
    """Línea vertical fina #2E2E2E (2pt de ancho) entre stats."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(x_pt - 1), px(y_top_pt),
        px(2), px(height_pt)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left, top, width, height):
    """
    Card ROUNDED_RECTANGLE con fill rgba(0,0,0,0.45) + outline blanco 1pt.
    border-radius: adjustment 0.06 ≈ 15pt en card de ~221pt de alto.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(left), px(top), px(width), px(height)
    )
    shape.adjustments[0] = 0.06  # ~15pt border-radius

    # Outline blanco 1pt
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    # Fill negro sólido base
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    # Opacidad 45% vía XML
    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')  # 45% en notación OOXML

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_positions_3(n_stats):
    """
    Posiciones X del lado izquierdo de cada caja de stat.
    v5: cajas 378pt, gaps según número de stats.
    3 stats: total 378*3 + 100*2 = 1334pt; left start = (1920-1334)/2 = 293pt
    2 stats: total 378*2 + 200 = 956pt; left start = (1920-956)/2 = 482pt
    1 stat: left = (1920-378)/2 = 771pt
    """
    if n_stats == 1:
        return [(SLIDE_W - STAT_BOX_W) // 2]
    elif n_stats == 2:
        gap = 200
        total = STAT_BOX_W * 2 + gap
        start = (SLIDE_W - total) // 2
        return [start, start + STAT_BOX_W + gap]
    else:  # 3
        gap = 100
        total = STAT_BOX_W * 3 + gap * 2
        start = (SLIDE_W - total) // 2
        return [start, start + STAT_BOX_W + gap, start + (STAT_BOX_W + gap) * 2]


def add_stat_block(slide, stat_value, desc_runs, left_pt, n_stats):
    """
    Bloque stat: cifra 180pt + caja descripción 378×113pt (13pt).
    left_pt: coordenada izquierda de la caja de descripción.
    La cifra se centra sobre la caja.
    """
    # Cifra grande — centrada horizontalmente sobre la caja
    cifra_w = 500   # pt — ancho suficiente para 180pt centrado
    cifra_left = left_pt + STAT_BOX_W // 2 - cifra_w // 2

    add_textbox(
        slide,
        left=cifra_left, top=STAT_NUM_TOP,
        width=cifra_w, height=200,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Descripción — 378×113pt, Poppins 13pt
    add_rich_textbox(
        slide,
        left=left_pt, top=STAT_DESC_TOP,
        width=STAT_BOX_W, height=STAT_BOX_H,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


def build_hallazgo_slide(prs, headline_plain, headline_italic, stats, source_text=""):
    """
    Slide Hallazgo cuanti.
    stats: lista de dicts {value, desc_runs} — 1, 2 o 3 elementos.
    source_text: si se pasa, se agrega en top 1010pt (Poppins italic 12pt gris).
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_positions_3(n)

    # Separadores verticales finos entre stats
    sep_top    = STAT_NUM_TOP - 20   # 400pt
    sep_height = 380                  # px/pt — cubre cifra + descripción

    if n >= 2:
        x_sep1 = xs[0] + STAT_BOX_W + (xs[1] - xs[0] - STAT_BOX_W) // 2
        add_vertical_separator(slide, x_sep1, sep_top, sep_height)
    if n == 3:
        x_sep2 = xs[1] + STAT_BOX_W + (xs[2] - xs[1] - STAT_BOX_W) // 2
        add_vertical_separator(slide, x_sep2, sep_top, sep_height)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    if source_text:
        add_source(slide, source_text)

    return slide


def build_consumer_voice_slide(prs, quote, attribution, source_text=""):
    """
    Slide Consumer Voice — verbatim SUELTO sobre fondo negro (sin card).
    v5/v6: NO card. Solo texto suelto centrado.
    Header "CONSUMER VOICE" arriba, Poppins 14pt gris.
    Verbatim: Instrument Serif regular 50pt, comillas «...».
    Atribución: Poppins italic 20pt.
    source_text: si se pasa, se agrega en top 1010pt.
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE"
    add_textbox(
        slide,
        left=0, top=CV_HEADER_TOP,
        width=SLIDE_W, height=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Verbatim suelto — ancho 1637pt (left 141, width 1637)
    full_quote = f"«{quote}»"
    verbatim_left  = 141
    verbatim_width = 1637
    verbatim_height = 500  # suficiente para multi-línea

    add_textbox(
        slide,
        left=verbatim_left, top=CV_VERBATIM_TOP,
        width=verbatim_width, height=verbatim_height,
        text=full_quote,
        font_name=FONT_HEADLINE, font_size_pt=CV_VERBATIM_PT,
        italic=False, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Atribución — debajo del verbatim (~870pt)
    add_textbox(
        slide,
        left=0, top=870,
        width=SLIDE_W, height=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIBUTION_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    if source_text:
        add_source(slide, source_text)

    return slide


def add_cuali_card(slide, quote_text, attribution, left, top, card_w, card_h):
    """
    Card cualitativa: ROUNDED_RECTANGLE rgba(0,0,0,0.45) + outline blanco 1pt.
    Texto interno: Poppins 15pt + atribución Poppins italic 12pt.
    """
    make_card_rgba_45_with_outline(slide, left, top, card_w, card_h)

    # Padding interno: 30pt lados, 25pt arriba/abajo
    pad_h = 30
    pad_v = 25
    text_w = card_w - (pad_h * 2)
    text_h = card_h - (pad_v * 2)

    full_quote = f"«{quote_text}»"

    txBox = slide.shapes.add_textbox(
        px(left + pad_h), px(top + pad_v),
        px(text_w), px(text_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    # Verbatim — Poppins 15pt
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_BODY
    run_q.font.size   = Pt(CARD_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Atribución — Poppins italic 12pt
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(CARD_ATTRIBUTION_PT)
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text=""):
    """
    Slide card cualitativa — headline arriba + 1-3 cards side by side.
    Layout 5: posiciones canónicas según número de verbatims.
    source_text: si se pasa, se agrega en top 1010pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)
    card_top = 429   # pt — posición Y canónica

    if n == 1:
        left = (SLIDE_W - CARD_W) // 2
        add_cuali_card(slide, verbatims[0]['quote'], verbatims[0]['attribution'],
                       left, card_top, CARD_W, CARD_H)
    elif n == 2:
        gap = 80
        total = CARD_W * 2 + gap
        start = (SLIDE_W - total) // 2
        for i, v in enumerate(verbatims):
            x = start + i * (CARD_W + gap)
            add_cuali_card(slide, v['quote'], v['attribution'],
                           x, card_top, CARD_W, CARD_H)
    else:  # 3
        gap = 36
        total = CARD_W * 3 + gap * 2
        start = (SLIDE_W - total) // 2
        for i, v in enumerate(verbatims):
            x = start + i * (CARD_W + gap)
            add_cuali_card(slide, v['quote'], v['attribution'],
                           x, card_top, CARD_W, CARD_H)

    if source_text:
        add_source(slide, source_text)

    return slide


# ── PORTADAS DE TENSIÓN ────────────────────────────────────────────────────────

def build_tension_cover(prs, tension_num, tension_title, tension_subtitle=""):
    """
    Slide portada de tensión — SIN source.
    Título grande centrado en el slide.
    """
    slide = blank_slide(prs)

    # Etiqueta "TENSIÓN XX" arriba
    etiqueta = f"TENSIÓN {tension_num:02d}"
    add_textbox(
        slide,
        left=0, top=200,
        width=SLIDE_W, height=80,
        text=etiqueta,
        font_name=FONT_BODY, font_size_pt=18,
        bold=False, italic=False,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Título principal de la tensión
    add_textbox(
        slide,
        left=100, top=320,
        width=1720, height=300,
        text=tension_title.upper(),
        font_name=FONT_HEADLINE, font_size_pt=HEADLINE_PT,
        bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.CENTER,
        word_wrap=True
    )

    # Subtítulo opcional
    if tension_subtitle:
        add_textbox(
            slide,
            left=200, top=680,
            width=1520, height=120,
            text=tension_subtitle,
            font_name=FONT_BODY, font_size_pt=18,
            italic=True, color=GREY_SOFT,
            align=PP_ALIGN.CENTER,
            word_wrap=True
        )

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SOURCE STRINGS — mapeados por hallazgo
# CC: "Source: Código Casa — Pilar Mujer · QXXX · Base 500."
# FW: "Source: Finance Women — Estudio cuantitativo 2025 · MXX · Mujeres n=240."
# NO Banreservas en ningún source FW.
# ═══════════════════════════════════════════════════════════════════════════════

SRC_H01 = "Source: Código Casa — Pilar Mujer · Q054 · Base 500."
SRC_H02 = "Source: Código Casa — Pilar Mujer · Q056 · Base 500."
SRC_H03 = "Source: Código Casa — Pilar Mujer · Q057, Q056 · Base 500."
SRC_H04 = "Source: Código Casa — Pilar Mujer · Q058 · Base 500."
SRC_H05 = "Source: Código Casa — Pilar Mujer · Q059 · Base 500."
SRC_H06 = "Source: Código Casa — Pilar Mujer · Q060 · Base 500."
SRC_H07 = "Source: Finance Women — Estudio cuantitativo 2025 · M20, M21 · Mujeres n=240."
SRC_H08 = "Source: Finance Women — Estudio cuantitativo 2025 · M23, M24, M54 · Mujeres n=240."
SRC_H09 = "Source: Finance Women — Estudio cuantitativo 2025 · M16, M28, M32, M60 · Mujeres n=240."
SRC_H10 = "Source: Finance Women — Estudio cuantitativo 2025 · M45, M45b · Mujeres n=240."
SRC_H11 = "Source: Finance Women — Estudio cuantitativo 2025 · M19, M43, M44 · Mujeres n=240."
SRC_H12 = "Source: Finance Women — Estudio cuantitativo 2025 · M33, M34, M36, M37 · Mujeres n=240."
SRC_H13 = "Source: Finance Women — Estudio cuantitativo 2025 · M38, M39, M40, M61 · Mujeres n=240."
SRC_H14 = "Source: Finance Women — Estudio cuantitativo 2025 · M8, M5, M17 · Mujeres n=240."


# ═══════════════════════════════════════════════════════════════════════════════
# DATOS DEL SET EDITORIAL — MUJER v2
# Fuentes: CC (n=305/500) y FW (n=240)
# 5 tensiones × hallazgos asignados
# ═══════════════════════════════════════════════════════════════════════════════

def build_deck():
    prs = new_prs()

    # ═══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 01 — LA CARGA INVISIBLE
    # H1 (crianza), H2 (tareas del hogar), H6 (autocuidado sacrificado)
    # ═══════════════════════════════════════════════════════════════════════════

    build_tension_cover(
        prs,
        tension_num=1,
        tension_title="La carga invisible",
        tension_subtitle="Crianza, tareas domésticas y el autocuidado que se pospone para siempre."
    )

    # ── H01 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "tiene nombre de mujer"

    build_hallazgo_slide(
        prs,
        headline_plain="49.5% de las mujeres dice que la madre es la principal responsable de la crianza. ",
        headline_italic="Solo el hombre distribuye la carga — ella ya la tiene asignada.",
        stats=[
            {
                'value': '42%',
                'desc_runs': [
                    make_run('del total cree que la mayor responsabilidad de la crianza recae en la '),
                    make_run('madre', bold=True),
                    make_run('; 41.2% dice "ambos padres" y solo 15.8% señala al padre.')
                ]
            },
            {
                'value': '49.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "Madre" como principal responsable, frente a 45.6% del masculino que responde "Ambos padres". La mujer se adjudica la carga que el hombre distribuye.')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('de los hogares '),
                    make_run('monoparentales', bold=True),
                    make_run(' y 59.4% del estrato E responden "Madre". La carga se concentra donde no hay con quién repartirla.')
                ]
            }
        ],
        source_text=SRC_H01
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo tengo un papá que es el papá bueno, el papá consentidor, el papá que dice vete suave, vete en paz. Pero no es el que está con ella todos los días. No es el que tiene que regañarla por las tareas, porque tiene que ir a impulsar los deberes de la casa.',
        attribution='Familia Monoparental',
        source_text=SRC_H01
    )

    # ── H02 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="Cocinar, lavar y limpiar siguen siendo territorio de ella. ",
        headline_italic="El único reparto equitativo llega cuando hay que pagar o ir al supermercado.",
        stats=[
            {
                'value': '67%',
                'desc_runs': [
                    make_run('reporta que '),
                    make_run('cocinar', bold=True),
                    make_run(' lo hace "más la mujer"; lavar la ropa 63.0%; limpiar y organizar 59.6%.')
                ]
            },
            {
                'value': '48%',
                'desc_runs': [
                    make_run('dice que el '),
                    make_run('supermercado', bold=True),
                    make_run(' lo hacen "ambos por igual". Manejo del dinero del hogar: 51.0%. Las tareas de dinero sí se reparten.')
                ]
            },
            {
                'value': '45%',
                'desc_runs': [
                    make_run('acompañar a '),
                    make_run('citas médicas o escolares', bold=True),
                    make_run(' lo hace "más la mujer". El cuidado logístico de los hijos tampoco se reparte.')
                ]
            }
        ],
        source_text=SRC_H02
    )

    build_consumer_voice_slide(
        prs,
        quote='Como mujer, bueno, no sabía cocinar y tuve que hacerlo. Obligada, aprendí a mamá, pero... no me gusta, pero tengo que hacerlo.',
        attribution='Familia Monoparental',
        source_text=SRC_H02
    )

    # ── H06 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no hace nada por su salud mental ni física. ",
        headline_italic="El hombre se ejercita. Ella no tiene espacio ni para eso.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "No realizo ninguna actividad". El masculino lidera en ejercicio (43.1%).')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('del '),
                    make_run('total', bold=True),
                    make_run(' no hace ninguna actividad de autocuidado; solo 3.0% va a terapia.')
                ]
            },
            {
                'value': '60%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres 18-24', bold=True),
                    make_run(' no realiza ninguna actividad de autocuidado. El abandono es mayor en las más jóvenes.')
                ]
            }
        ],
        source_text=SRC_H06
    )

    build_consumer_voice_slide(
        prs,
        quote='Ese tema de que dependa de mí 24-7 lo que se vaya a hacer te da un poco de estrés.',
        attribution='Familia Mixta',
        source_text=SRC_H06
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 02 — EL DISCURSO DEL CAMBIO VS. EL REPARTO REAL
    # H3 (percepción de cambio), H5 (estrés económico como prioritario)
    # ═══════════════════════════════════════════════════════════════════════════

    build_tension_cover(
        prs,
        tension_num=2,
        tension_title="El discurso del cambio vs. el reparto real",
        tension_subtitle="74.8% dice que los roles cambiaron mucho. Las tareas domésticas siguieron siendo de ella."
    )

    # ── H03 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "no de reparto"

    build_hallazgo_slide(
        prs,
        headline_plain="74.8% dice que el rol de género cambió mucho. La cocina, el lavado y la limpieza siguieron siendo de ella. ",
        headline_italic="El discurso avanzó; el reparto, no.",
        stats=[
            {
                'value': '75%',
                'desc_runs': [
                    make_run('responde "5 = Mucho" a cuánto cambió el rol del hombre y la mujer en la familia. '),
                    make_run('85%', bold=True),
                    make_run(' en el top-2 (4+5).')
                ]
            },
            {
                'value': '77%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' percibe el cambio como "Mucho" vs 71.8% del masculino. Las mujeres ven más el cambio que no cambia.')
                ]
            },
            {
                'value': '67%',
                'desc_runs': [
                    make_run('reporta cocinar '),
                    make_run('"más la mujer"', bold=True),
                    make_run(' (Q056). Convive con la percepción de cambio (Q057). El reparto material no acompaña la percepción.')
                ]
            }
        ],
        source_text=SRC_H03
    )

    build_consumer_voice_slide(
        prs,
        quote='Las que salimos a echar para adelante, a trabajar, a luchar, la que ya no dependemos de un hombre para salir adelante. Y yo en mi caso, que aunque esté casada, trabajo igualito como si no, hago y lucho y no tengo que depender como tiempo atrás que las mujeres dominicanas nos creaban como que era para depender del hombre.',
        attribution='Familia Preferente Roles de Liderazgo',
        source_text=SRC_H03
    )

    # ── H05 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="47.8% de las mujeres nombra la economía como su mayor fuente de estrés. ",
        headline_italic="La crianza y el hogar no se acercan: el miedo más grande es la plata.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('menciona "'),
                    make_run('Economía', bold=True),
                    make_run('" como factor que más estrés genera. Sumada a "Realidad económica" (28.4%), lo económico domina el ranking.')
                ]
            },
            {
                'value': '21%',
                'desc_runs': [
                    make_run('nombra explícitamente "'),
                    make_run('Falta de tiempo personal', bold=True),
                    make_run('" como fuente de estrés. Un quinto de las mujeres pone nombre a la ausencia de tiempo propio.')
                ]
            },
            {
                'value': '16%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Crianza de los hijos', bold=True),
                    make_run('" 16.2% y "Responsabilidades del hogar" 16.2%. El trabajo de cuidado estresa menos que el dinero.')
                ]
            }
        ],
        source_text=SRC_H05
    )

    build_consumer_voice_slide(
        prs,
        quote='No tengo paciencia cuando no tengo dinero. Cuando no tengo dinero, que no sé qué voy a hacer con los compromisos, entonces yo pienso, le doy vuelta a la cosa, un estrés. La mayor parte del problema hoy en día es el dinero, todo lo resuelve el dinero aunque la gente no lo quiera admitir.',
        attribution='Familia Monoparental',
        source_text=SRC_H05
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 03 — AUTORIDAD Y PODER
    # H4 (autoridad en alimentación), H14 (rol económico en el hogar)
    # ═══════════════════════════════════════════════════════════════════════════

    build_tension_cover(
        prs,
        tension_num=3,
        tension_title="Autoridad y poder",
        tension_subtitle="La mujer decide qué se come y también quién provee. El discurso social tarda en reconocerlo."
    )

    # ── H04 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="89.5% de las mujeres decide siempre qué se compra de comer. ",
        headline_italic="La nutrición familiar es un poder casi exclusivo de ella, no compartido.",
        stats=[
            {
                'value': '90%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' responde "5 = Siempre" a si toma la decisión principal sobre la compra de alimentos, frente a 77.9% del masculino.')
                ]
            },
            {
                'value': '85%',
                'desc_runs': [
                    make_run('del '),
                    make_run('total', bold=True),
                    make_run(' responde "5 = Siempre". Es la decisión menos disputada de toda la matriz de roles.')
                ]
            },
            {
                'value': '89%',
                'desc_runs': [
                    make_run('en '),
                    make_run('NSE D', bold=True),
                    make_run(' y 87.5% en NSE E responden "Siempre", frente a 52.2% en NSE AB. [AB n≈23: contraste, no cifra robusta.]')
                ]
            }
        ],
        source_text=SRC_H04
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo considero saludable la comida que yo cocino, ese tipo de ingredientes que le estoy echando, y también las porciones. Lo que tú cocinas en tu casa.',
        attribution='Familia Monoparental',
        source_text=SRC_H04
    )

    # ── H14 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "imaginario" + sin "la realidad cambió, el chip no"

    build_hallazgo_slide(
        prs,
        headline_plain="40% de las mujeres es la principal proveedora de su hogar. ",
        headline_italic="40% provee. El rol que se le atribuye al hombre lo ocupa ella.",
        stats=[
            {
                'value': '40%',
                'desc_runs': [
                    make_run('responde que sí es la '),
                    make_run('principal proveedora', bold=True),
                    make_run(' de su hogar. [FW — perfil: jóvenes, jefas de hogar, líderes, emprendedoras. No proyectable a población general.]')
                ]
            },
            {
                'value': '68%',
                'desc_runs': [
                    make_run('tiene '),
                    make_run('trabajo formal', bold=True),
                    make_run('; 13.9% informal o emprendimiento; solo 4.2% se declara ama de casa. La mujer del estudio FW está mayoritariamente vinculada al ingreso.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Las tomo yo sola', bold=True),
                    make_run('" es la frase con mayor acuerdo top-2 en decisiones financieras (μ=3.01). "No las tomo yo, lo hace mi esposo" tiene el mayor desacuerdo (bottom-2 79.6%).')
                ]
            }
        ],
        source_text=SRC_H14
    )

    # Slide H14-B — Cards cualitativas
    # Headline cuali v6: sin "imaginario"
    build_cuali_slide(
        prs,
        headline_plain="El discurso del proveedor no cedió. ",
        headline_italic="Ellas proveen, el rol sigue siendo del hombre en el cuento social.",
        verbatims=[
            {
                'quote': 'No, por tu trabajo tú eres proveedor. Pero se supone que el rol del proveedor siempre ha sido el hombre.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            },
            {
                'quote': 'Yo no resuelvo de nadie. Mi esposo es la madre. [Modismo: ella ocupa el rol de proveedor que culturalmente se le asigna al hombre.]',
                'attribution': 'Familia Masivo Jefas de Hogar'
            }
        ],
        source_text=SRC_H14
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 04 — LA MUJER Y EL DINERO
    # H7 (sueños financieros), H8 (bloqueos), H9 (planificación), H10 (prejuicio)
    # ═══════════════════════════════════════════════════════════════════════════

    build_tension_cover(
        prs,
        tension_num=4,
        tension_title="La mujer y el dinero",
        tension_subtitle="Sueña con un techo propio, planifica mejor que el estereotipo, y carga con un prejuicio que dice lo contrario."
    )

    # ── H07 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "no es consumo — es"

    build_hallazgo_slide(
        prs,
        headline_plain="La meta financiera #1 de la mujer dominicana es un techo propio. ",
        headline_italic="El sueño es no depender del próximo sueldo.",
        stats=[
            {
                'value': '60%',
                'desc_runs': [
                    make_run('elige "'),
                    make_run('Comprar vivienda', bold=True),
                    make_run('" como meta financiera actual (1ª selección); 50.8% "Lograr estabilidad"; 48.3% "Emprender / Invertir".')
                ]
            },
            {
                'value': '48%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Comprar vivienda', bold=True),
                    make_run('" se repite como #1 en la segunda selección (47.9%) y "Salir de deudas" sube a #2 (32.1%).')
                ]
            },
            {
                'value': '14%',
                'desc_runs': [
                    make_run('"'),
                    make_run('Mejorar imagen personal', bold=True),
                    make_run('" queda último o penúltimo en ambas selecciones (13.8% y 2.1%). El prejuicio de la mujer que sueña con consumo estético no aguanta la data.')
                ]
            }
        ],
        source_text=SRC_H07
    )

    build_consumer_voice_slide(
        prs,
        quote='Siempre ha sido mi sueño, siempre ha sido como que pagar mi cosa yo, si yo pudiera pagar todo, o sea, todo lo mío yo.',
        attribution='Familia Preferente Jóvenes',
        source_text=SRC_H07
    )

    # ── H08 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "el problema no es la disciplina, es la plata"

    build_hallazgo_slide(
        prs,
        headline_plain="74.9% siente que algo le bloquea el sueño financiero. ",
        headline_italic="87.1% apunta a lo mismo: falta plata.",
        stats=[
            {
                'value': '75%',
                'desc_runs': [
                    make_run('responde "Sí" a si siente que hay algo que le impide lograr su '),
                    make_run('sueño financiero', bold=True),
                    make_run('.')
                ]
            },
            {
                'value': '87%',
                'desc_runs': [
                    make_run('de las que sienten un bloqueo señala "'),
                    make_run('Falta de ingresos', bold=True),
                    make_run('" como impedimento; muy por encima de "Falta de hábitos" (44.9%) o "Falta de educación financiera" (35.4%).')
                ]
            },
            {
                'value': '60%',
                'desc_runs': [
                    make_run('de todas las mujeres nombra "'),
                    make_run('Falta de ingresos', bold=True),
                    make_run('" como lo que le impide tomar mejores decisiones financieras hoy.')
                ]
            }
        ],
        source_text=SRC_H08
    )

    build_consumer_voice_slide(
        prs,
        quote='Comparado con un año atrás que yo no tenía libertad financiera, tenía que depender de otro y es difícil, porque aunque el otro quiera, a veces uno quiere hacer algo y no puede. Ahora como yo tengo mi negocio, si quiero comprar algo lo compro, si le quiero comprar algo a mis hijos no tengo que esperar que otro me dé.',
        attribution='Familia Masivo Emprendedoras',
        source_text=SRC_H08
    )

    # ── H09 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "no es el estereotipo: es la norma"

    build_hallazgo_slide(
        prs,
        headline_plain="74.2% tiene presupuesto, 74.1% planifica siempre o casi siempre, 78.7% revisa sus cuentas cada semana o cada día. ",
        headline_italic="La administradora es la norma, no la excepción.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('tiene un '),
                    make_run('presupuesto mensual', bold=True),
                    make_run('. Y 74.1% planifica sus gastos "Siempre" (42.1%) o "Casi siempre" (32.1%).')
                ]
            },
            {
                'value': '79%',
                'desc_runs': [
                    make_run('revisa sus cuentas a '),
                    make_run('diario (44.2%) o semanalmente (34.6%)', bold=True),
                    make_run('. Control financiero activo como práctica cotidiana.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('se describe como "'),
                    make_run('Cuidadosa', bold=True),
                    make_run('" con el dinero y 20% como "Práctica"; solo 6.7% se define "Libre". La autopercepción es de control, no de derroche.')
                ]
            }
        ],
        source_text=SRC_H09
    )

    build_consumer_voice_slide(
        prs,
        quote='Hay mujeres que no, que son ocultas, que si se presenta un caso, ellas siempre están ahí tapando la falta del esposo. Hay mujeres que son muy honradas con el dinero.',
        attribution='Familia Masivo Jefas de Hogar',
        source_text=SRC_H09
    )

    # ── H10 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="39.2% siente que se espera algo distinto de ella en dinero por ser mujer. ",
        headline_italic="La esperan derrochadora y a la vez le exigen ser la administradora perfecta del hogar.",
        stats=[
            {
                'value': '39%',
                'desc_runs': [
                    make_run('responde "Sí" a si siente que se espera algo distinto de ella en temas de '),
                    make_run('dinero por ser mujer', bold=True),
                    make_run('; 31.7% dice "No" y 29.2% "No estoy segura". La mayoría no descarta la presión.')
                ]
            }
        ],
        source_text=SRC_H10
    )

    build_cuali_slide(
        prs,
        headline_plain="Dos prejuicios, una mujer: ",
        headline_italic="derrochadora y mala administradora al mismo tiempo.",
        verbatims=[
            {
                'quote': 'Algunas personas tienen la mentalidad que las mujeres son derrochadoras. El otro extremo es que la mujer siempre siempre debe guardar. Ambos extremos son malos.',
                'attribution': 'Respuesta abierta M45b · FW'
            },
            {
                'quote': 'Que al no ser proveedoras no somos buenas administradoras.',
                'attribution': 'Familia Masivo Jefas de Hogar'
            }
        ],
        source_text=SRC_H10
    )

    # ═══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 05 — AHORRO, CONSUMO Y AUTONOMÍA ECONÓMICA
    # H11 (emoción financiera: estrés), H12 (compra impulsiva), H13 (ahorro imposible)
    # ═══════════════════════════════════════════════════════════════════════════

    build_tension_cover(
        prs,
        tension_num=5,
        tension_title="Ahorro, consumo y autonomía económica",
        tension_subtitle="El dinero se vive como amenaza, la compra impulsiva como desahogo y el ahorro como una trampa de la que no se sale."
    )

    # ── H11 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="La emoción que más siente la mujer al pensar en su dinero es estrés (44.3%). ",
        headline_italic="El orgullo solo llega al 5.5%: el dinero se vive como amenaza, no como logro.",
        stats=[
            {
                'value': '44%',
                'desc_runs': [
                    make_run('responde "'),
                    make_run('Estrés', bold=True),
                    make_run('" como la emoción que más siente al pensar en sus finanzas; "Tranquilidad" 29.8%, "Miedo" 12.3%, "Culpa" 8.1%, "Orgullo" 5.5%.')
                ]
            },
            {
                'value': '41%',
                'desc_runs': [
                    make_run('siente '),
                    make_run('ansiedad o estrés financiero', bold=True),
                    make_run(' con frecuencia (top-2, 4+5); el promedio en escala 1-5 es 3.48.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('está satisfecha con su situación financiera actual (top-2 solo '),
                    make_run('22.5%', bold=True),
                    make_run('); 30.8% está insatisfecha (bottom-2). El promedio es 2.75 sobre 5.')
                ]
            }
        ],
        source_text=SRC_H11
    )

    build_consumer_voice_slide(
        prs,
        quote='Tengo miedo a no ser lo que yo espero ser, o a que las cosas no se den como yo quiero. Siento que el tiempo está pasando muy rápido, que no tengo una casa mía o no tengo un auto propio.',
        attribution='Familia Masivo Jóvenes',
        source_text=SRC_H11
    )

    # ── H12 ───────────────────────────────────────────────────────────────────
    # Headline v6: sin "El motivo #1 no es la promoción: es consentirse"

    build_hallazgo_slide(
        prs,
        headline_plain="47.9% hizo una compra impulsiva el último mes. ",
        headline_italic="El motivo #1 es consentirse — el único espacio que el rol le deja solo para ella.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('hizo una '),
                    make_run('compra impulsiva', bold=True),
                    make_run(' en el último mes; y 42.5% admite haber comprado algo recientemente solo para sentirse mejor emocionalmente.')
                ]
            },
            {
                'value': '44%',
                'desc_runs': [
                    make_run('El motivo #1: "'),
                    make_run('Para consentirme o darme un gusto', bold=True),
                    make_run('" (44.2%), por encima de "buena oportunidad" (39.2%) o "promoción" (33.8%). 25.8% compra cuando está estresada.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('de categorías son "'),
                    make_run('Comidas/Bebidas', bold=True),
                    make_run('" y 59.8% "Ropa/Accesorios"; "Electrónica" apenas 5.0%. Desahogo de día a día, no derroche aspiracional.')
                ]
            }
        ],
        source_text=SRC_H12
    )

    build_consumer_voice_slide(
        prs,
        quote='Uno se siente como empoderada y señor me lo merezco, por eso yo trabajo. Tiene como su momento, le llega de vez en vez, no siempre. Déjame comprar ese porchecito, porque se me suba el ánimo.',
        attribution='Familia Masivo Jefas de Hogar',
        source_text=SRC_H12
    )

    # ── H13 ───────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="50% no logra ahorrar con facilidad. La estrategia que más funciona: ",
        headline_italic="esconderse el dinero de sí mismas antes de que llegue a la cuenta.",
        stats=[
            {
                'value': '50%',
                'desc_runs': [
                    make_run('responde que NO le resulta fácil '),
                    make_run('ahorrar', bold=True),
                    make_run('. En escala 1-5, el bottom-2 (35.8%) supera al top-2 (35%), con promedio 2.98.')
                ]
            },
            {
                'value': '24%',
                'desc_runs': [
                    make_run('no destina ningún porcentaje de su dinero mensual al '),
                    make_run('ahorro', bold=True),
                    make_run('; otro 28.3% destina menos del 5%. Más de la mitad ahorra poco o nada.')
                ]
            },
            {
                'value': '66%',
                'desc_runs': [
                    make_run('"'),
                    make_run('No me alcanza el dinero', bold=True),
                    make_run('" es el obstáculo #1 (65.5%), por encima de "Falta de educación financiera" (58.4%). La brecha de ahorro es de ingresos, no de conocimiento.')
                ]
            }
        ],
        source_text=SRC_H13
    )

    build_cuali_slide(
        prs,
        headline_plain="El mecanismo del ahorro: ",
        headline_italic="interceptar el dinero antes de que llegue a la mano.",
        verbatims=[
            {
                'quote': 'Me entré a la cooperativa de la empresa, no me lo van a extraer de la cuenta. Lo hice con ese fin porque una de mis metas al principio del año fue esa: ahorrar, porque crear un fondo de emergencia.',
                'attribution': 'Familia Preferente Jefas de Hogar'
            },
            {
                'quote': 'Me entré en la cooperativa del trabajo, que es la única forma que ese dinero no entre a mi cuenta y salga.',
                'attribution': 'Familia Preferente Jefas de Hogar'
            }
        ],
        source_text=SRC_H13
    )

    return prs


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mujer-v2-deck-flat.pptx")

    print("Construyendo deck Mujer v2...")
    prs = build_deck()

    prs.save(output_path)
    print(f"Deck guardado: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides, 1):
        print(f"  Slide {i:02d}")


if __name__ == "__main__":
    main()
