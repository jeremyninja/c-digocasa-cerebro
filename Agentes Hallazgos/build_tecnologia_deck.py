#!/usr/bin/env python3
"""
build_tecnologia_deck.py
========================
Deck tecnologia-deck-flat.pptx desde cero con python-pptx.
20 slides — set editorial cerrado tecnologia-hallazgos-editados.md (11 hallazgos).

Distribución por hallazgo:
  H01 (cuanti 3 stats + Consumer Voice)          → 2 slides
  H02 (cuanti 3 stats + Consumer Voice)          → 2 slides
  H03 (cuanti+cuali integrado, 2 verbatims)      → 2 slides (cuanti + cuali s/s)
  H04 (cuanti 1 stat  + Consumer Voice)          → 2 slides
  H05 (cuanti 2 stats + Consumer Voice)          → 2 slides
  H06 (cuanti 3 stats + Consumer Voice)          → 2 slides
  H07 (cuanti 1 stat  + Consumer Voice)          → 2 slides
  H08 (cuanti 2 stats + Consumer Voice)          → 2 slides
  H09 (cuanti+cuali integrado, 2 verbatims)      → 2 slides (cuanti + cuali s/s)
  H10 (solo-cuali 2 verbatims cards s/s)         → 1 slide
  H11 (solo-cuali 2 verbatims cards s/s)         → 1 slide

TOTAL: 9×2 + 2×1 = 20 slides

Specs visuales (aprendizajes-montador-cc.md v6):
  Slide: 1920×1080 px = 18,288,000 × 10,287,000 EMU
  Fonts: Instrument Serif (headlines, stats, verbatims CV) + Poppins (cuerpo)
  Fondo: #000000 negro plano (NO masterslide)
  Consumer Voice (Layout 4): verbatim SUELTO sobre fondo negro — SIN card
  Cards cuali (Layout 5): 567×227pt side by side con fill rgba(0,0,0,0.45)
  Stat number: 135pt Instrument Serif italic
  Stat desc: 13pt Poppins
  Headline: 50pt fijo Instrument Serif MAYÚSCULAS
  NO Source al pie de página (eliminado per v5/v6)

Constantes canónicas verificadas (aprendizajes § lección operativa crítica):
  HEADLINE_PT     = 50        (fijo, NO auto-size)
  STAT_DESC_PT    = 13        (13pt fijo)
  STAT_BOX_W      = 378       (10cm fijo)
  STAT_BOX_H      = 113       (3cm fijo)
  CV_VERBATIM_PT  = 50        (v6: era 60)
  CARD_W          = 567       (v6: 15cm)
  CARD_H          = 227       (v6: 6cm)
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

def pt_from_px(px_val):
    return round(px_val * 0.75, 1)

# ── Fuentes canónicas (aprendizajes v6) ─────────────────────────────────────
HEADLINE_PT       = 50.0    # fijo — NO auto-size, NO pick_headline_size_px()
STAT_NUMBER_PT    = pt_from_px(180)  # 135pt
STAT_DESC_PT      = 13.0   # 13pt fijo (NO 16, NO 16.5)
CV_HEADER_PT      = 14.0
CV_VERBATIM_PT    = 50.0   # v6: era 60
CV_ATTRIB_PT      = 20.0
CARD_VERBATIM_PT  = 15.0
CARD_ATTRIB_PT    = 12.0

# ── Dimensiones de stat boxes (v6: 10cm × 3cm = 378×113pt) ─────────────────
STAT_BOX_W = 378   # 10cm @ 96dpi — fijo para 1, 2 y 3 stats
STAT_BOX_H = 113   # 3cm @ 96dpi — fijo
STAT_TOP_PX = 480  # posición vertical de la cifra grande

# ── Dimensiones de cards cuali (v6: 15×6cm = 567×227pt side by side) ──────
CARD_W    = 567    # 15cm @ 96dpi
CARD_H    = 227    # 6cm @ 96dpi
CARD_GAP_2 = 80
CARD_GAP_3 = 36
CARD_TOP   = 429

# ── Consumer Voice ────────────────────────────────────────────────────────
CV_TEXT_W    = 1637
CV_TEXT_LEFT = (SLIDE_W_PX - CV_TEXT_W) // 2  # = 141px


# ── Helpers base ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en el run XML."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """Headline 50pt fijo Instrument Serif MAYUSCULAS centrado."""
    txBox = slide.shapes.add_textbox(
        px(100), px(80), px(1720), px(320)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(HEADLINE_PT)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(HEADLINE_PT)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    """Línea vertical fina #2E2E2E entre stats."""
    line = slide.shapes.add_shape(
        1,
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left_px, top_px, width_px, height_px):
    """Card rounded rectangle: fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass

    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_el.set('val', '45000')  # 45% opacidad

    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w   = STAT_BOX_W
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=STAT_BOX_H,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


# ── Layouts de slides ──────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic, stats):
    """
    Slide de Hallazgo cuanti (Layouts 1/2/3).
    NO source — eliminado per aprendizajes v5/v6.
    stats: lista de dicts con 'value' y 'desc_runs'.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice (Layout 4 — verbatim SUELTO sobre fondo negro, SIN card).
    Header "CONSUMER VOICE" Poppins 14pt gris arriba.
    Verbatim Instrument Serif 50pt blanco, comillas «».
    Atribución Poppins italic 20pt blanco.
    """
    slide = blank_slide(prs)

    # Header
    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Posición vertical dinámica según longitud del verbatim
    char_count = len(quote)
    if char_count > 250:
        verb_top = 240
    elif char_count > 180:
        verb_top = 290
    elif char_count > 120:
        verb_top = 330
    else:
        verb_top = 380

    full_quote = f"«{quote}»"

    verb_box = slide.shapes.add_textbox(
        px(CV_TEXT_LEFT), px(verb_top),
        px(CV_TEXT_W), px(500)
    )
    tf = verb_box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_q = p.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(CV_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Atribución
    attrib_top = verb_top + 340
    add_textbox(
        slide,
        left_px=0, top_px=attrib_top,
        width_px=1920, height_px=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIB_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_cuali_slide_side_by_side(prs, headline_plain, headline_italic, verbatims):
    """
    Slide Card cualitativa (Layout 5 — cards 567×227pt SIDE BY SIDE).
    Headline arriba 50pt Instrument Serif MAYUSCULAS.
    Cards con fill rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt.
    Verbatim Poppins 15pt blanco con «», atribución Poppins italic 12pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)

    if n == 1:
        card_left = (SLIDE_W_PX - CARD_W) // 2
        positions = [(card_left, CARD_TOP)]
    elif n == 2:
        total_w = CARD_W * 2 + CARD_GAP_2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_2, CARD_TOP)
        ]
    else:
        total_w = CARD_W * 3 + CARD_GAP_3 * 2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_3, CARD_TOP),
            (start_x + CARD_W * 2 + CARD_GAP_3 * 2, CARD_TOP)
        ]

    for i, v in enumerate(verbatims):
        left, top = positions[i]
        make_card_rgba_45_with_outline(slide, left, top, CARD_W, CARD_H)

        full_quote = f"«{v['quote']}»"
        pad_h = 30
        pad_v = 25
        text_w = CARD_W - (pad_h * 2)
        text_h = CARD_H - (pad_v * 2)

        txBox = slide.shapes.add_textbox(
            px(left + pad_h), px(top + pad_v),
            px(text_w), px(text_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left   = Emu(0)
        tf.margin_right  = Emu(0)
        tf.margin_top    = Emu(0)
        tf.margin_bottom = Emu(0)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run_q = p1.add_run()
        run_q.text = full_quote
        run_q.font.name   = FONT_BODY
        run_q.font.size   = Pt(CARD_VERBATIM_PT)
        run_q.font.italic = False
        run_q.font.bold   = False
        run_q.font.color.rgb = WHITE
        _set_kern(run_q._r)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run_a = p2.add_run()
        run_a.text = f"— {v['attribution']}"
        run_a.font.name   = FONT_BODY
        run_a.font.size   = Pt(CARD_ATTRIB_PT)
        run_a.font.italic = True
        run_a.font.bold   = False
        run_a.font.color.rgb = WHITE
        _set_kern(run_a._r)

    return slide


# ── BUILD DECK ─────────────────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — El celular manda. Lo demás, lejos.
    # Headline: "92.4% tiene el celular como dispositivo esencial del hogar.
    #   La televisión queda en 48.8%, la computadora en 20.8%.
    #   El teléfono fijo: 6.2%. La consola: 1.8%."
    # Stats: 92.4%, 48.8%/20.8%, 6.2%/1.8%
    # Verbatim: Familia Sin Hijos con Mascota
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 01 — H01 Hallazgo cuanti (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="92.4% TIENE EL CELULAR COMO DISPOSITIVO ESENCIAL DEL HOGAR. LA TELEVISIÓN QUEDA EN 48.8%, LA COMPUTADORA EN 20.8%. ",
        headline_italic="EL TELÉFONO FIJO: 6.2%. LA CONSOLA: 1.8%.",
        stats=[
            {
                'value': '92.4%',
                'desc_runs': [
                    make_run('nombra el '),
                    make_run('celular', bold=True),
                    make_run(' como dispositivo tecnológico esencial del hogar — casi el doble que la segunda opción. P49, Base 500.')
                ]
            },
            {
                'value': '48.8%',
                'desc_runs': [
                    make_run('nombra la '),
                    make_run('televisión', bold=True),
                    make_run('; 20.8% la computadora; 14.4% la tableta. La brecha entre el celular y cualquier otro dispositivo supera los 43 puntos. P49, Base 500.')
                ]
            },
            {
                'value': '6.2%',
                'desc_runs': [
                    make_run('nombra el '),
                    make_run('teléfono fijo', bold=True),
                    make_run(' como esencial; 1.8% la consola de videojuegos. Dos tecnologías estándar en otros mercados, fuera del mapa doméstico dominicano. P49.')
                ]
            }
        ]
    )

    # Slide 02 — H01 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo uso hasta para pedir una pizza. Yo uso tecnología porque yo descargo todas las aplicaciones, pido la compra del supermercado por ahí (...) tengo muchísimas cosas en mi teléfono, de verdad, tengo el teléfono — en yo hago mis sueños.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — La hiperconexión la tienen los que mandan
    # Headline: "52.2% del NSE AB pasa más de 6 horas diarias en pantalla.
    #   En el NSE E, 46.9% pasa menos de 1 hora.
    #   La desconexión es de los que menos tienen."
    # Stats: 52.2% AB (tendencia), 46.9% E (tendencia), 26.2%
    # Verbatim: Familia Biparental con Hijos Pequeños
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 03 — H02 Hallazgo cuanti (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="52.2% DEL NSE AB PASA MÁS DE 6 HORAS DIARIAS EN PANTALLA. EN EL NSE E, 46.9% PASA MENOS DE 1 HORA. ",
        headline_italic="LA DESCONEXIÓN ES DE LOS QUE MENOS TIENEN.",
        stats=[
            {
                'value': '52.2%',
                'desc_runs': [
                    make_run('del NSE AB pasa más de '),
                    make_run('6 horas diarias', bold=True),
                    make_run(' en pantalla — hiperconexión como respuesta top. Tendencia direccional: NSE AB n≈23. P48 × NSE.')
                ]
            },
            {
                'value': '46.9%',
                'desc_runs': [
                    make_run('del NSE E pasa '),
                    make_run('menos de 1 hora', bold=True),
                    make_run(' diaria — único estrato donde la desconexión es la respuesta top. Tendencia direccional: NSE E n≈32. P48 × NSE.')
                ]
            },
            {
                'value': '26.2%',
                'desc_runs': [
                    make_run('del total (n=131) pasa más de '),
                    make_run('6 horas diarias', bold=True),
                    make_run(' frente a pantallas; 19.0% (n=95) pasa menos de 1 hora. La distribución promedio cubre un rango amplio. P48, Base 500.')
                ]
            }
        ]
    )

    # Slide 04 — H02 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Después de la pandemia la tecnología es vital, es vital, es vital. Y según lo que estoy observando aquí en el grupo de nosotros, la mayoría trabajan vía, o sea, no presencial.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — Dicen que la tecnología los une. En la cena, cada uno en su pantalla.
    # Headline: "54.8% califica de muy positivo el impacto de la tecnología en la
    #   comunicación familiar. En la cena, ella con la tablet, él con el celular — pero hablamos."
    # Stats: 54.8% / 16.0%
    # Verbatims (2, cuanti+cuali integrado): Familia Monoparental × 2
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 05 — H03 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="54.8% CALIFICA DE MUY POSITIVO EL IMPACTO DE LA TECNOLOGÍA EN LA COMUNICACIÓN FAMILIAR. ",
        headline_italic="EN LA CENA, ELLA CON LA TABLET, ÉL CON EL CELULAR — PERO HABLAMOS.",
        stats=[
            {
                'value': '54.8%',
                'desc_runs': [
                    make_run('califica con 5 — '),
                    make_run('muy positivo', bold=True),
                    make_run(' — el impacto del uso de la tecnología en la comunicación familiar. Solo 16.0% (n=80) lo califica negativo (1+2). P50, Base 500.')
                ]
            },
            {
                'value': '16.0%',
                'desc_runs': [
                    make_run('califica el impacto como '),
                    make_run('negativo', bold=True),
                    make_run('. Puede ser el porcentaje que ya no cree la versión "pero hablamos". El 54.8% incluye hogares donde pantallas simultáneas y comunicación se perciben compatibles. P50, Base 500.')
                ]
            }
        ]
    )

    # Slide 06 — H03 Cards cualitativas (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="DICEN QUE LA TECNOLOGÍA LOS UNE. ",
        headline_italic="EN LA CENA, CADA UNO EN SU PANTALLA.",
        verbatims=[
            {
                'quote': 'Ella se pone a ver su tablet porque siempre viendo su tablet (...) las cinco — porque los fines de semana también a veces compran pizza o cualquier otra comida rápida y nos la comemos así, o viendo televisión. Mi hermanito siempre vive en su mundo.',
                'attribution': 'Familia Monoparental'
            },
            {
                'quote': 'Y allí, ¿qué hacen? Hablamos, y yo con mi celular, ella con el celular. Ok. Sí, exacto. Pero hablamos.',
                'attribution': 'Familia Monoparental'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — Le doy el celular para que se quede tranquila. Sé que está mal.
    # Headline: "La madre da el celular para poder hacer sus cosas. La hija hace
    #   berrinche cuando no se lo dan. El ciclo tiene nombre: lo llaman error. Y siguen."
    # Stats: 54.8% muy positivo / 13.4% muy negativo (1 stat con ambas cifras en desc)
    # Verbatim: Familia Biparental con Hijos Pequeños
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 07 — H04 Hallazgo cuanti (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="LA MADRE DA EL CELULAR PARA PODER HACER SUS COSAS. LA HIJA HACE BERRINCHE CUANDO NO SE LO DAN. ",
        headline_italic="EL CICLO TIENE NOMBRE: LO LLAMAN ERROR. Y SIGUEN.",
        stats=[
            {
                'value': '13.4%',
                'desc_runs': [
                    make_run('califica el impacto de la tecnología como '),
                    make_run('muy negativo', bold=True),
                    make_run('. El 54.8% que dice "muy positivo" (n=274) puede incluir hogares ya dentro del ciclo error-enviciando que lo aceptaron como normal. P50, Base 500.')
                ]
            }
        ]
    )

    # Slide 08 — H04 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Es un error de nosotros porque como le doy el celular ella ve a los muñequitos y se está más tranquila, y como que para nosotros poder hacer las cosas le damos el celular. Entonces es un error en el que caemos, y realmente yo sé que debemos trabajar en eso porque la estamos como enviciando.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — Los que crían son los más enganchados
    # Headline: "36.5% del rango 25–34 años pasa más de 6 horas diarias en
    #   pantalla — el único grupo etario donde la hiperconexión es la respuesta top.
    #   Los que crían, los más pegados."
    # Stats: 36.5% rango 25-34, resto de rangos
    # Verbatim: Familia Biparental con Hijos Adultos
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 09 — H05 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="36.5% DEL RANGO 25–34 AÑOS PASA MÁS DE 6 HORAS DIARIAS EN PANTALLA — EL ÚNICO GRUPO ETARIO DONDE LA HIPERCONEXIÓN ES LA RESPUESTA TOP. ",
        headline_italic="LOS QUE CRÍAN, LOS MÁS PEGADOS.",
        stats=[
            {
                'value': '36.5%',
                'desc_runs': [
                    make_run('del rango '),
                    make_run('25–34 años', bold=True),
                    make_run(' pasa más de 6 horas diarias en pantalla — único rango etario donde la hiperconexión es la respuesta top. Base 25-34 n≈115. P48 × edad.')
                ]
            },
            {
                'value': '1–3h',
                'desc_runs': [
                    make_run('es la respuesta top en el resto de rangos: '),
                    make_run('18–24 (40.0%), 35–44 (42.0%), 45–54 (47.4%), 55+ (49.6%)', bold=True),
                    make_run('. La conexión baja con la edad. Todos n≥55. P48 × edad.')
                ]
            }
        ]
    )

    # Slide 10 — H05 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo hasta para ir al baño voy con el teléfono. Estoy sentada en mi baño, estoy viendo TikTok o lo que sea, ya eso, en muchas casas ya. Eso es una adicción.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — La supervisión digital es mirar por encima del hombro
    # Headline: "36.4% protege a los niños 'supervisando directamente' su uso de
    #   internet. 19.2% usa controles parentales.
    #   El mecanismo de protección es la mano del adulto, no el software."
    # Stats: 36.4%, 19.2%/7.2%, 10.0%/15.2%/9.6%
    # Verbatim: Familia Monoparental
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 11 — H06 Hallazgo cuanti (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="36.4% PROTEGE A LOS NIÑOS SUPERVISANDO DIRECTAMENTE SU USO DE INTERNET. 19.2% USA CONTROLES PARENTALES. ",
        headline_italic="EL MECANISMO DE PROTECCIÓN ES LA MANO DEL ADULTO, NO EL SOFTWARE.",
        stats=[
            {
                'value': '36.4%',
                'desc_runs': [
                    make_run('cita "'),
                    make_run('supervisión directa', bold=True),
                    make_run('" como mecanismo de protección digital — opción más mencionada de toda la batería P53. Base 500.')
                ]
            },
            {
                'value': '19.2%',
                'desc_runs': [
                    make_run('usa '),
                    make_run('controles parentales', bold=True),
                    make_run(' en dispositivos o aplicaciones; 7.2% (n=36) usa filtros de contenido. La herramienta digital es elegida por menos de 1 de cada 5 hogares. P53, Base 500.')
                ]
            },
            {
                'value': '10.0%',
                'desc_runs': [
                    make_run('revisa '),
                    make_run('historial de navegación', bold=True),
                    make_run('; 15.2% restringe horarios; 9.6% usa cuentas compartidas. La acción analógica —mirar, revisar, restringir— supera siempre a la herramienta. P53.')
                ]
            }
        ]
    )

    # Slide 12 — H06 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Ella tiene Instagram pero bajo, como ella dice, bajo mucho control, bajo mucha supervisión — porque si no se lo superviso yo, se lo supervisa su papá (...) yo puedo coger el celular y puedo revisar, y lo hago.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — Uno de cada cuatro hogares no protege nada
    # Headline: "26.4% de los hogares no usa ningún mecanismo de protección digital
    #   para los niños. La segunda opción más votada de la batería es la rendición."
    # Stats: 26.4% (1 stat)
    # Verbatim: Familia Biparental con Hijos Pequeños
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 13 — H07 Hallazgo cuanti (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="26.4% DE LOS HOGARES NO USA NINGÚN MECANISMO DE PROTECCIÓN DIGITAL PARA LOS NIÑOS. ",
        headline_italic="LA SEGUNDA OPCIÓN MÁS VOTADA DE LA BATERÍA ES LA RENDICIÓN.",
        stats=[
            {
                'value': '26.4%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('no utilizamos ningún mecanismo de protección', bold=True),
                    make_run('" para niños en el entorno digital — segunda opción más votada de P53, solo por debajo de la supervisión directa (36.4%). Base 500.')
                ]
            }
        ]
    )

    # Slide 14 — H07 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo creí que sí lo podía hacer y pude caer en un gancho. Entonces descubrí muchos mensajes de niñas que yo creía y veía de una forma, y en el Instagram con las cosas que hacían (...) le tuvo que quitar el teléfono al niño porque, o sea, eso es algo con lo que yo te hablo muy fuerte.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — La conversación sobre riesgos de internet llega después del susto
    # Headline: "5.8% de los hogares cita 'conversaciones sobre riesgos en internet'
    #   como mecanismo de protección. Es la opción menos votada de la batería.
    #   La acción gana siempre a la palabra."
    # Stats: 5.8%, batería comparativa
    # Verbatim: Familia Monoparental
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 15 — H08 Hallazgo cuanti (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="5.8% DE LOS HOGARES CITA CONVERSACIONES SOBRE RIESGOS EN INTERNET COMO MECANISMO DE PROTECCIÓN. ES LA OPCIÓN MENOS VOTADA DE LA BATERÍA. ",
        headline_italic="LA ACCIÓN GANA SIEMPRE A LA PALABRA.",
        stats=[
            {
                'value': '5.8%',
                'desc_runs': [
                    make_run('cita "'),
                    make_run('conversaciones sobre riesgos en internet', bold=True),
                    make_run('" (acoso, privacidad, redes sociales) como mecanismo de protección digital — opción menos mencionada de las 9 en P53. Base 500.')
                ]
            },
            {
                'value': '7.2%',
                'desc_runs': [
                    make_run('usa '),
                    make_run('filtros de contenido', bold=True),
                    make_run('; 9.6% cuentas compartidas; 10.0% revisa historial. La conversación queda por debajo de vigilar, bloquear y restringir. Siempre antes que hablar. P53, Base 500.')
                ]
            }
        ]
    )

    # Slide 16 — H08 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Los delincuentes no le hacen daño a los niños en la calle — es desde tu casa, ese niño que tú entiendes que es un niño. Pero esas son conversaciones que no son de niños, es un depredador adulto.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — El control parental dominicano se llama Alexa con cámara
    # Headline: "19.2% declara usar controles parentales formales. Lo que usan los
    #   padres clase media-alta: Alexa en cada habitación, también en el carro,
    #   también en el trabajo."
    # Stats: 19.2% (1 stat)
    # Verbatims (2, cuanti+cuali integrado): Biparental Hijos Pequeños + Extendida
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 17 — H09 Hallazgo cuanti (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="19.2% DECLARA USAR CONTROLES PARENTALES FORMALES. LO QUE USAN LOS PADRES CLASE MEDIA-ALTA: ",
        headline_italic="ALEXA EN CADA HABITACIÓN, TAMBIÉN EN EL CARRO, TAMBIÉN EN EL TRABAJO.",
        stats=[
            {
                'value': '19.2%',
                'desc_runs': [
                    make_run('declara usar "'),
                    make_run('controles parentales en dispositivos o aplicaciones', bold=True),
                    make_run('" como mecanismo de protección digital de niños. Lo que el cuestionario llama "control parental" y lo que el hogar construye son categorías distintas. P53, Base 500.')
                ]
            }
        ]
    )

    # Slide 18 — H09 Cards cualitativas (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="EL CONTROL PARENTAL DOMINICANO ",
        headline_italic="SE LLAMA ALEXA CON CÁMARA.",
        verbatims=[
            {
                'quote': 'Yo le puse una Alexa con cámara, entonces yo puedo ver. Yo le digo a Alexa que lo acerque y lo acecho. Y yo le tengo la Alexa puesta así, queda de frente a donde está, en su escritorio, con su computadora, de frente a mí. Y yo veo todo lo que está haciendo.',
                'attribution': 'Familia Biparental con Hijos Pequeños'
            },
            {
                'quote': 'Yo tengo dos Alexa, la televisión la controlo también por ahí (...) por el Samsung yo se la prendo, la veo por la cámara que está viendo y le controlo todo. Tengo una robotina que también la activo desde el trabajo.',
                'attribution': 'Familia Extendida'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — La tienda más cercana del dominicano se llama Marketplace
    # Headline: "La cifra de compras por redes no está tabulada. En los grupos,
    #   el Marketplace y TikTok son el retail principal de amas de casa,
    #   revendedores y emprendedores."
    # Solo-cuali — 2 verbatims en cards side by side
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 19 — H10 Card cualitativa (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="LA TIENDA MÁS CERCANA DEL DOMINICANO ",
        headline_italic="SE LLAMA MARKETPLACE.",
        verbatims=[
            {
                'quote': 'Yo diría que el Facebook — el Marketplace de Facebook — yo siempre lo chequeo. Eso está vendiendo mucho. (...) Yo compré todo el equipo de cocina por esa vía, con el aparato (...) que te pueden costar 145 mil pesos.',
                'attribution': 'Familia Biparental con Hijos Adultos'
            },
            {
                'quote': 'Importo cosas por Alibaba, o sea, algo que hago alterno los sábados y los domingos. Y básicamente el teléfono es esa fuente de comunicación.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — La tecnología es trabajo. El ocio llega después.
    # Headline: "El ranking de usos de la tecnología no pudo tabularse. En los
    #   grupos, la respuesta arranca siempre por trabajo.
    #   El entretenimiento llega al final, como aside."
    # Solo-cuali — 2 verbatims en cards side by side
    # ─────────────────────────────────────────────────────────────────────────

    # Slide 20 — H11 Card cualitativa (2 verbatims side by side)
    build_cuali_slide_side_by_side(
        prs,
        headline_plain="EN LOS GRUPOS, LA RESPUESTA ARRANCA SIEMPRE POR TRABAJO. ",
        headline_italic="EL ENTRETENIMIENTO LLEGA AL FINAL, COMO ASIDE.",
        verbatims=[
            {
                'quote': 'La tecnología es esencial para mí en mi trabajo porque me ayuda a yo poder hacer varias cosas al mismo tiempo ocupando un solo lugar (...). Yo no tengo que andar con papel, con sello, con nada de eso porque yo lo tengo ahí.',
                'attribution': 'Familia Biparental con Hijos Pequeños'
            },
            {
                'quote': 'Yo utilizo mucho chat GPT. Mi trabajo es vital, la tecnología — utilizo mucho el Excel, el paquete Office, redes sociales también.',
                'attribution': 'Familia Homoparental'
            }
        ]
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "tecnologia-deck-flat.pptx"
    )
    prs.save(output_path)

    total = len(prs.slides)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {total}")

    print("\n-- Verificacion specs -------------------------------------------")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"  = {int(prs.slide_width / PX)}px x {int(prs.slide_height / PX)}px")
    ok = prs.slide_width == Emu(SLIDE_W_PX * PX) and prs.slide_height == Emu(SLIDE_H_PX * PX)
    print(f"  Slide 1920x1080 OK: {ok}")
    print(f"HEADLINE_PT:      {HEADLINE_PT}pt (fijo — NO auto-size)")
    print(f"STAT_NUMBER_PT:   {STAT_NUMBER_PT}pt Instrument Serif italic")
    print(f"STAT_DESC_PT:     {STAT_DESC_PT}pt Poppins")
    print(f"STAT_BOX_W:       {STAT_BOX_W}pt (10cm)")
    print(f"STAT_BOX_H:       {STAT_BOX_H}pt (3cm)")
    print(f"CV_VERBATIM_PT:   {CV_VERBATIM_PT}pt (v6: era 60)")
    print(f"CV_ATTRIB_PT:     {CV_ATTRIB_PT}pt Poppins italic")
    print(f"CARD_W:           {CARD_W}pt (15cm v6)")
    print(f"CARD_H:           {CARD_H}pt (6cm v6)")
    print("Cards cuali: rgba(0,0,0,0.45) + outline blanco 1pt + border-radius 15pt")
    print("Consumer Voice: verbatim SUELTO sobre fondo negro (SIN card)")
    print("Comillas espanolas en todos los verbatims")
    print("Kerning 0 en todo el deck")
    print("NO masterslide — fondo negro plano")
    print("NO source en ningun slide (v6)")

    print("\n-- QA FLAGS —-------------------------------------------------")
    print("H01 (slides  1-2):  cuanti 3 stats — 92.4% / 48.8% / 6.2%")
    print("H02 (slides  3-4):  cuanti 3 stats — 52.2% AB (tend.) / 46.9% E (tend.) / 26.2%")
    print("H03 (slides  5-6):  cuanti 2 stats + cuali 2 cards s/s — 54.8% / 16.0%")
    print("H04 (slides  7-8):  cuanti 1 stat  — 13.4% muy negativo (framing vs 54.8%)")
    print("H05 (slides  9-10): cuanti 2 stats — 36.5% rango 25-34 / 1-3h resto rangos")
    print("H06 (slides 11-12): cuanti 3 stats — 36.4% / 19.2% / 10.0%")
    print("H07 (slides 13-14): cuanti 1 stat  — 26.4% ningun mecanismo")
    print("H08 (slides 15-16): cuanti 2 stats — 5.8% / 7.2%")
    print("H09 (slides 17-18): cuanti 1 stat + cuali 2 cards s/s — 19.2% / Alexa x2")
    print("H10 (slide  19):    solo-cuali 2 verbatims side by side — Marketplace")
    print("H11 (slide  20):    solo-cuali 2 verbatims side by side — Tecnologia=Trabajo")
    print(f"Total slides verificados: {total} (esperado: 20)")
    print("Verbatims en blanco (WHITE) — ningun verde lima")
    print("NO lineas decorativas entre headline y stats")
    print("NO conclusiones italic separadas")
    print("H03 y H12 (P50 54.8%) — mismo stat base, frames distintos: OK")
    print("H10 y H11 solo-cuali — cifra pendiente de retabulacion (editorial decision cerrada)")

    return output_path


if __name__ == "__main__":
    build_deck()
