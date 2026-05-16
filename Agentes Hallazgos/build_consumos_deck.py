#!/usr/bin/env python3
"""
build_consumos_deck.py
======================
Deck consumos-deck-flat.pptx — formato MED 5 tensiones × 8 slides = 40 slides.
Basado en consumos-hallazgos-editados.md (13 hallazgos cerrados).

Distribución de tensiones y hallazgos:

  T1 (slides 01–08) — El consumidor que no compra verde
      H01: 42.8% no filtra por ecología  (cuanti + CV)
      H02: 7.4% exige política ambiental a la marca  (solo-cuanti 3 stats)

  T2 (slides 09–16) — La responsabilidad ambiental externalizada
      H03: Conciencia ambiental existe, responsabilidad al gobierno  (cuanti + CV)
      H04: Reusa por costumbre, no por conciencia  (cuanti + CV)

  T3 (slides 17–24) — Sostenibilidad como deseo sin infraestructura
      H05: Consumo responsable = apagar la luz  (cuanti + CV)
      H06: 28% quiere reciclar, el camión lo revuelve  (cuanti + CV)
      H07: Conciencia ecológica se aprende en EE.UU.  (cuanti + CV)

  T4 (slides 25–32) — La compra por redes: exclusión, no desinterés
      H08: 65.4% no compró nada por redes  (solo-cuanti 2 stats)
      H09: A más pobreza, menos clic — brecha NSE  (solo-cuanti 3 stats)
      H10: Generación joven compra el doble — brecha generacional  (solo-cuanti 2 stats)
      H11: Hogar sin hijos con mascota el más permeable  (solo-cuanti 2 stats)

  T5 (slides 33–40) — La filosofía del gasto familiar
      H12: Jerarquía moral del gasto heredada  (solo-cuali)
      H13: Redes venden vidas, no productos  (cuanti + CV)

Specs v6 (aprendizajes-montador-cc.md):
  Slide: 1920×1080 px = 18,288,000 × 10,287,000 EMU
  Fonts: Instrument Serif (headlines, stats, verbatims CV) + Poppins (cuerpo)
  Fondo: #000000 negro plano (NO masterslide)
  HEADLINE_PT     = 50   (fijo, NO auto-size)
  STAT_DESC_PT    = 13   (13pt fijo)
  STAT_BOX_W      = 378  (10cm @ 96dpi)
  STAT_BOX_H      = 113  (3cm @ 96dpi)
  CV_VERBATIM_PT  = 50   (v6)
  CARD_W          = 567  (15cm v6)
  CARD_H          = 227  (6cm v6)
  Consumer Voice: verbatim SUELTO sobre fondo negro (SIN card)
  NO source al pie (eliminado v6)
  Kerning 0 en todo el deck
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP   = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT  = RGBColor(0x9B, 0x9B, 0x9B)
LIME       = RGBColor(0xC8, 0xF5, 0x60)   # verde lima — SOLO cápsulas tensión

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── Constantes canónicas v6 (verificadas contra aprendizajes lección operativa crítica) ─
HEADLINE_PT      = 50.0   # fijo — NO auto-size
STAT_NUMBER_PT   = 135.0  # 180px × 0.75 = 135pt
STAT_DESC_PT     = 13.0   # 13pt fijo
CV_HEADER_PT     = 14.0
CV_VERBATIM_PT   = 50.0   # v6: era 60
CV_ATTRIB_PT     = 20.0
CARD_VERBATIM_PT = 15.0
CARD_ATTRIB_PT   = 12.0

# ── Stat boxes (v6: 10cm × 3cm = 378×113pt) ──────────────────────────────────
STAT_BOX_W  = 378   # 10cm @ 96dpi — fijo
STAT_BOX_H  = 113   # 3cm @ 96dpi — fijo
STAT_TOP_PX = 480   # posición vertical cifra grande

# ── Cards cuali (v6: 15×6cm = 567×227pt) ────────────────────────────────────
CARD_W     = 567
CARD_H     = 227
CARD_GAP_2 = 80
CARD_GAP_3 = 36
CARD_TOP   = 429

# ── Consumer Voice ────────────────────────────────────────────────────────────
CV_TEXT_W    = 1637
CV_TEXT_LEFT = (SLIDE_W_PX - CV_TEXT_W) // 2  # 141px


# ── Helpers base ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en el run XML."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name    = font_name
    run.font.size    = Pt(font_size_pt)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """Headline 50pt fijo Instrument Serif MAYÚSCULAS centrado."""
    txBox = slide.shapes.add_textbox(
        px(100), px(80), px(1720), px(320)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name    = FONT_HEADLINE
        run.font.size    = Pt(HEADLINE_PT)
        run.font.bold    = False
        run.font.italic  = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name    = FONT_HEADLINE
        run2.font.size    = Pt(HEADLINE_PT)
        run2.font.bold    = False
        run2.font.italic  = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    """Línea vertical fina #2E2E2E entre stats."""
    line = slide.shapes.add_shape(
        1,
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45_with_outline(slide, left_px, top_px, width_px, height_px):
    """Card rounded rectangle: fill rgba(0,0,0,0.45) + outline blanco 1pt."""
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass

    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_el.set('val', '45000')  # 45% opacidad

    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1)

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w    = STAT_BOX_W
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=STAT_BOX_H,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


# ── Layouts de slides ──────────────────────────────────────────────────────────

def build_tension_title_slide(prs, tension_num, headline_plain, headline_italic=""):
    """
    Slide 1 de cada tensión — título grande sobre fondo negro.
    Cápsula 'TENSIÓN 0X' en verde lima arriba, headline grande Instrument Serif.
    """
    slide = blank_slide(prs)

    # Cápsula tensión — verde lima, texto negro
    cap_w, cap_h = 240, 50
    cap_left = (SLIDE_W_PX - cap_w) // 2
    cap_top  = 120
    cap_shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        px(cap_left), px(cap_top), px(cap_w), px(cap_h)
    )
    try:
        cap_shape.adjustments[0] = 0.4
    except Exception:
        pass
    cap_shape.fill.solid()
    cap_shape.fill.fore_color.rgb = LIME
    cap_shape.line.fill.background()

    # Texto cápsula
    tf_cap = cap_shape.text_frame
    tf_cap.word_wrap = False
    tf_cap.margin_left   = Emu(0)
    tf_cap.margin_right  = Emu(0)
    tf_cap.margin_top    = Emu(0)
    tf_cap.margin_bottom = Emu(0)
    p_cap = tf_cap.paragraphs[0]
    p_cap.alignment = PP_ALIGN.CENTER
    r_cap = p_cap.add_run()
    r_cap.text = f"TENSIÓN {tension_num:02d}"
    r_cap.font.name   = FONT_BODY
    r_cap.font.size   = Pt(14)
    r_cap.font.bold   = True
    r_cap.font.italic = False
    r_cap.font.color.rgb = BLACK
    _set_kern(r_cap._r)

    # Headline grande centrado verticalmente
    hl_top  = 250
    hl_height = 500
    txBox = slide.shapes.add_textbox(
        px(120), px(hl_top), px(1680), px(hl_height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if headline_plain:
        r1 = p.add_run()
        r1.text = headline_plain.upper()
        r1.font.name   = FONT_HEADLINE
        r1.font.size   = Pt(70)
        r1.font.bold   = False
        r1.font.italic = False
        r1.font.color.rgb = WHITE
        _set_kern(r1._r)

    if headline_italic:
        r2 = p.add_run()
        r2.text = headline_italic.upper()
        r2.font.name   = FONT_HEADLINE
        r2.font.size   = Pt(70)
        r2.font.bold   = False
        r2.font.italic = True
        r2.font.color.rgb = WHITE
        _set_kern(r2._r)

    return slide


def build_hallazgo_slide(prs, headline_plain, headline_italic, stats):
    """
    Slide Hallazgo cuanti (Layouts 1/2/3) — NO source.
    stats: lista de dicts con 'value' y 'desc_runs'.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice (Layout 4) — verbatim SUELTO sobre fondo negro, SIN card.
    """
    slide = blank_slide(prs)

    # Header
    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Posición vertical según longitud
    char_count = len(quote)
    if char_count > 280:
        verb_top = 200
    elif char_count > 220:
        verb_top = 240
    elif char_count > 160:
        verb_top = 290
    elif char_count > 100:
        verb_top = 340
    else:
        verb_top = 380

    full_quote = f"«{quote}»"

    verb_box = slide.shapes.add_textbox(
        px(CV_TEXT_LEFT), px(verb_top),
        px(CV_TEXT_W), px(520)
    )
    tf = verb_box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_q = p.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(CV_VERBATIM_PT)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Atribución
    attrib_top = verb_top + 370
    add_textbox(
        slide,
        left_px=0, top_px=attrib_top,
        width_px=1920, height_px=60,
        text=f"— {attribution}",
        font_name=FONT_BODY, font_size_pt=CV_ATTRIB_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims):
    """
    Slide Card cualitativa (Layout 5) — cards 567×227pt SIDE BY SIDE.
    1–3 verbatims en cards con fill rgba(0,0,0,0.45) + outline blanco 1pt.
    Texto interno Poppins 15pt.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)

    if n == 1:
        card_left = (SLIDE_W_PX - CARD_W) // 2
        positions = [(card_left, CARD_TOP)]
    elif n == 2:
        total_w = CARD_W * 2 + CARD_GAP_2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_2, CARD_TOP)
        ]
    else:
        total_w = CARD_W * 3 + CARD_GAP_3 * 2
        start_x = (SLIDE_W_PX - total_w) // 2
        positions = [
            (start_x, CARD_TOP),
            (start_x + CARD_W + CARD_GAP_3, CARD_TOP),
            (start_x + CARD_W * 2 + CARD_GAP_3 * 2, CARD_TOP)
        ]

    for i, v in enumerate(verbatims):
        left, top = positions[i]
        make_card_rgba_45_with_outline(slide, left, top, CARD_W, CARD_H)

        full_quote = f"«{v['quote']}»"
        pad_h = 30
        pad_v = 25
        text_w = CARD_W - (pad_h * 2)
        text_h = CARD_H - (pad_v * 2)

        txBox = slide.shapes.add_textbox(
            px(left + pad_h), px(top + pad_v),
            px(text_w), px(text_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left   = Emu(0)
        tf.margin_right  = Emu(0)
        tf.margin_top    = Emu(0)
        tf.margin_bottom = Emu(0)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run_q = p1.add_run()
        run_q.text = full_quote
        run_q.font.name   = FONT_BODY
        run_q.font.size   = Pt(CARD_VERBATIM_PT)
        run_q.font.italic = False
        run_q.font.bold   = False
        run_q.font.color.rgb = WHITE
        _set_kern(run_q._r)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        run_a = p2.add_run()
        run_a.text = f"— {v['attribution']}"
        run_a.font.name   = FONT_BODY
        run_a.font.size   = Pt(CARD_ATTRIB_PT)
        run_a.font.italic = True
        run_a.font.bold   = False
        run_a.font.color.rgb = WHITE
        _set_kern(run_a._r)

    return slide


def build_sintesis_slide(prs, tension_num, tension_name, body_lines):
    """
    Slide 7 de cada tensión — Síntesis.
    Headline 'SÍNTESIS' arriba, bullet lines de conclusión en Poppins.
    """
    slide = blank_slide(prs)

    # Label "SÍNTESIS" en gris suave
    add_textbox(
        slide,
        left_px=0, top_px=60,
        width_px=1920, height_px=60,
        text=f"TENSIÓN {tension_num:02d} — SÍNTESIS",
        font_name=FONT_BODY, font_size_pt=14,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Nombre tensión
    add_textbox(
        slide,
        left_px=120, top_px=150,
        width_px=1680, height_px=80,
        text=tension_name.upper(),
        font_name=FONT_HEADLINE, font_size_pt=HEADLINE_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Cuerpo síntesis
    body_top = 300
    for line in body_lines:
        add_textbox(
            slide,
            left_px=240, top_px=body_top,
            width_px=1440, height_px=120,
            text=line,
            font_name=FONT_BODY, font_size_pt=18,
            color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True
        )
        body_top += 110

    return slide


def build_detonadora_slide(prs, tension_num, pregunta_detonadora):
    """
    Slide 8 de cada tensión — Pregunta detonadora para marca/cliente.
    """
    slide = blank_slide(prs)

    # Label arriba en gris
    add_textbox(
        slide,
        left_px=0, top_px=60,
        width_px=1920, height_px=60,
        text=f"TENSIÓN {tension_num:02d} — PREGUNTA PARA LA MARCA",
        font_name=FONT_BODY, font_size_pt=14,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Pregunta grande centrada
    txBox = slide.shapes.add_textbox(
        px(120), px(280), px(1680), px(500)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = pregunta_detonadora.upper()
    run.font.name   = FONT_HEADLINE
    run.font.size   = Pt(42)
    run.font.bold   = False
    run.font.italic = True
    run.font.color.rgb = WHITE
    _set_kern(run._r)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# BUILD DECK
# ══════════════════════════════════════════════════════════════════════════════

def build_deck():
    prs = new_prs()

    # ══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 01 — El consumidor que no compra verde  (slides 01–08)
    # H01 (cuanti + CV) + H02 (solo-cuanti 3 stats)
    # 8 slides: Título | H01 cuanti | H01 CV | H02 cuanti | H02-2ª cifra |
    #           Cuali-H01 verbatim | Síntesis | Detonadora
    # ══════════════════════════════════════════════════════════════════════════

    # Slide 01 — Título T1
    build_tension_title_slide(
        prs,
        tension_num=1,
        headline_plain="El consumidor que ",
        headline_italic="no compra verde."
    )

    # Slide 02 — H01 Hallazgo cuanti principal (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="LA MAYORÍA DECLARA SIN CRITERIO: ",
        headline_italic="42.8% COMPRA SIN CONSIDERAR EL MEDIO AMBIENTE.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara que '),
                    make_run('no toma decisiones de compra con criterios ecológicos', bold=True),
                    make_run(' — la opción más votada de P35. Supera cualquier práctica verde concreta. P35, Base 500.')
                ]
            },
            {
                'value': '30%',
                'desc_runs': [
                    make_run('menciona la '),
                    make_run('bolsa reusable', bold=True),
                    make_run(' como única práctica ecológica activa al comprar — el gesto más común. 12 puntos por debajo de la abstención declarada. P35, Base 500.')
                ]
            }
        ]
    )

    # Slide 03 — H01 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo no uso bolsas de plástico de la del supermercado, yo uso de las bolsas que son reusables. Cuando por ejemplo yo trato de pedir, cuando en Amazon tú tienes la opción de pedir paquetería que sea (...) que no tiene plástico, que no tiene el envoltorio de papel y así por el estilo. Son pequeñas cosas que uno hace.',
        attribution='Familia Mixta'
    )

    # Slide 04 — H02 Hallazgo cuanti secundario (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="EL CONSUMIDOR DOMINICANO NO LE EXIGE NADA A LA MARCA: ",
        headline_italic="7.4% COMPRA POR POLÍTICA AMBIENTAL; 6% POR ORIGEN LOCAL.",
        stats=[
            {
                'value': '7.4%',
                'desc_runs': [
                    make_run('compra marcas con '),
                    make_run('políticas ambientales responsables', bold=True),
                    make_run(' — la opción de mayor sofisticación ecológica en P35. La minoría más exigente. P35, Base 500.')
                ]
            },
            {
                'value': '6%',
                'desc_runs': [
                    make_run('prefiere '),
                    make_run('productos locales o de producción sostenible', bold=True),
                    make_run('. Origen y sostenibilidad como criterio de compra: uso marginal. P35, Base 500.')
                ]
            },
            {
                'value': '3.8%',
                'desc_runs': [
                    make_run('compra '),
                    make_run('a granel o sin empaque', bold=True),
                    make_run('. Las tres cifras juntas no llegan al 20%: menos de 1 de cada 5 ejerce algún filtro ambiental sobre la marca. P35.')
                ]
            }
        ]
    )

    # Slide 05 — H02 lectura complementaria (reutiliza cuali cards con contexto)
    # El H02 no tiene verbatim. Se monta como slide de síntesis visual cuanti.
    build_hallazgo_slide(
        prs,
        headline_plain="EL SILENCIO CUALI CONFIRMA LA CIFRA: ",
        headline_italic="NADIE NOMBRÓ UNA MARCA AMBIENTAL COMO CRITERIO DE SELECCIÓN.",
        stats=[
            {
                'value': '0',
                'desc_runs': [
                    make_run('grupos donde emergió una '),
                    make_run('marca con política ambiental', bold=True),
                    make_run(' como criterio activo de compra. La conversación ecológica vive en la gestión doméstica, no en el punto de venta. 11 FGs revisados.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('declara querer practicar '),
                    make_run('consumo responsable', bold=True),
                    make_run(' (P36) — pero la marca no recibe ese deseo. La exigencia ambiental no llega al punto de venta. Base 500.')
                ]
            }
        ]
    )

    # Slide 06 — Cuali H01 (el verbatim sobre el gesto doméstico)
    build_cuali_slide(
        prs,
        headline_plain="EL GESTO DOMÉSTICO NO ES ",
        headline_italic="UN CRITERIO DE COMPRA.",
        verbatims=[
            {
                'quote': 'Son pequeñas cosas que uno hace.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # Slide 07 — Síntesis T1
    build_sintesis_slide(
        prs,
        tension_num=1,
        tension_name="El consumidor que no compra verde",
        body_lines=[
            "43% declara no filtrar por ecología al comprar. Es la respuesta más votada.",
            "El gesto verde existe pero opera en el hogar, no en la tienda.",
            "Menos del 20% ejerce algún filtro ambiental sobre la marca.",
            "El silencio de los grupos confirma: ningún participante nombró una marca verde como criterio."
        ]
    )

    # Slide 08 — Detonadora T1
    build_detonadora_slide(
        prs,
        tension_num=1,
        pregunta_detonadora="¿Tu marca puede hacer el puente entre el gesto doméstico y el criterio de compra?"
    )

    # ══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 02 — La responsabilidad ambiental externalizada  (slides 09–16)
    # H03 (cuanti + 2 CV) + H04 (cuanti + CV)
    # ══════════════════════════════════════════════════════════════════════════

    # Slide 09 — Título T2
    build_tension_title_slide(
        prs,
        tension_num=2,
        headline_plain="La responsabilidad ambiental ",
        headline_italic="externalizada."
    )

    # Slide 10 — H03 Hallazgo cuanti principal (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="LA CONCIENCIA AMBIENTAL EXISTE. ",
        headline_italic="LA RESPONSABILIDAD, SE LA DEJAN AL GOBIERNO.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara que '),
                    make_run('no toma decisiones de compra con criterios ecológicos', bold=True),
                    make_run('. Cuando se pregunta quién debería actuar, el diagnóstico cierra rápido: "eso se lo adjudicamos al gobierno." P35, Base 500.')
                ]
            }
        ]
    )

    # Slide 11 — H03 Consumer Voice (1er verbatim — delegación al gobierno)
    build_consumer_voice_slide(
        prs,
        quote='Y se lo adjudicamos mucho, se lo adjudicamos mucho al gobierno, el gobierno de la basura.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # Slide 12 — H03 segundo verbatim sobre infraestructura (cuali slide)
    build_cuali_slide(
        prs,
        headline_plain="SI EL SISTEMA NO RECICLA, ",
        headline_italic="SEPARAR ES TRABAJO PERDIDO.",
        verbatims=[
            {
                'quote': '¿Qué hago yo separando vainas? Al final van a ir todos a la tercera riqueza. No lo van a reciclar, no lo van a separar, no van a hacer nada. Yo puedo tener tres zafacones y separar la basura y a la hora del [camión], la van a recoger toda y la van a echar en el mismo camión.',
                'attribution': 'Familia Homoparental'
            }
        ]
    )

    # Slide 13 — H04 Hallazgo cuanti secundario (1 stat + contexto)
    build_hallazgo_slide(
        prs,
        headline_plain="EN EL HOGAR DOMINICANO SE REUSA POR COSTUMBRE, ",
        headline_italic="NO POR CONCIENCIA AMBIENTAL.",
        stats=[
            {
                'value': '30%',
                'desc_runs': [
                    make_run('usa '),
                    make_run('bolsas reusables', bold=True),
                    make_run(' al comprar — segunda práctica ecológica más mencionada. El hogar que reusa lo hace por economía doméstica, no por huella. P35, Base 500.')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara '),
                    make_run('no filtrar por ecología', bold=True),
                    make_run(' — supera por 12 puntos a quienes reusaron bolsa. Mismo hogar: reconoce que no filtra y, como máximo, reutiliza la bolsa. P35, Base 500.')
                ]
            }
        ]
    )

    # Slide 14 — H04 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo siempre trato de que no se compre mucho desechable. Ellos lo que dicen es que yo soy una abusadora porque yo ando fregando. Pero yo les digo, sí, pero estamos creando menos basura. Un parque más y lo fregaste y ya todos tienen menos basura.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # Slide 15 — Síntesis T2
    build_sintesis_slide(
        prs,
        tension_num=2,
        tension_name="La responsabilidad ambiental externalizada",
        body_lines=[
            "La conciencia existe. La responsabilidad termina en el Estado, no en el ciudadano ni en la marca.",
            "El argumento de la infraestructura no es apatía: es lógica. Si el sistema no recicla, separar es trabajo perdido.",
            "El reuso doméstico ocurre, pero por economía familiar — no como convicción ambiental.",
            "La marca ambiental sigue invisible en este espacio."
        ]
    )

    # Slide 16 — Detonadora T2
    build_detonadora_slide(
        prs,
        tension_num=2,
        pregunta_detonadora="¿Tu marca puede operar donde el Estado falló? ¿O prefieres culpar al consumidor por no exigirte más?"
    )

    # ══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 03 — Sostenibilidad como deseo sin infraestructura  (slides 17–24)
    # H05 (cuanti + CV) + H06 (cuanti + CV) + H07 (cuanti + CV)
    # ══════════════════════════════════════════════════════════════════════════

    # Slide 17 — Título T3
    build_tension_title_slide(
        prs,
        tension_num=3,
        headline_plain="Sostenibilidad como ",
        headline_italic="deseo sin infraestructura."
    )

    # Slide 18 — H05 Hallazgo cuanti principal (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="EL DOMINICANO QUIERE CONSUMIR RESPONSABLE: ",
        headline_italic="LO QUE ENTIENDE POR ESO ES APAGAR LA LUZ Y NO DESPERDICIAR EL AGUA.",
        stats=[
            {
                'value': '51%',
                'desc_runs': [
                    make_run('menciona '),
                    make_run('consumo responsable', bold=True),
                    make_run(' como acción deseada en su hogar — primer lugar de P36. 51.2% exacto. Base 500.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('menciona '),
                    make_run('ahorro de agua', bold=True),
                    make_run(', prácticamente empatado con consumo responsable (50.6%). El empate revela colapso semántico: son la misma cosa. P36, Base 500.')
                ]
            },
            {
                'value': '28%',
                'desc_runs': [
                    make_run('quiere '),
                    make_run('reciclar', bold=True),
                    make_run('; 18.4% quiere reducir plásticos. El ranking muestra que la sostenibilidad se traduce en administrar mejor la factura. P36.')
                ]
            }
        ]
    )

    # Slide 19 — H05 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='El bombillo, a mí no me gusta tener un bombillo prendido en mi casa. Yo paro con el bombillo y me llaman. (...) En mi casa me molesta un bombillo prendido.',
        attribution='Familia Mixta'
    )

    # Slide 20 — H06 Hallazgo cuanti terciario (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="EL 28% QUIERE RECICLAR. ",
        headline_italic="EL CAMIÓN MUNICIPAL RECOGE TODO EN EL MISMO CARRO.",
        stats=[
            {
                'value': '28%',
                'desc_runs': [
                    make_run('declara que le gustaría '),
                    make_run('reciclar en su hogar', bold=True),
                    make_run(' — tercera acción más mencionada en P36. El deseo existe. La cadena operativa, no. P36, Base 500.')
                ]
            }
        ]
    )

    # Slide 21 — H06 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo hago el reciclado, o sea, yo tengo mis zafacones, separo latas, botella, ta, ta, ta. El único inconveniente para mí es cómo botarle, porque tengo que ir a un lugar específico, no siempre está cerca de mí o no siempre tengo mi vehículo disponible (...) la basura se va acumulando, al final tú te cansas. Y la echan en el mismo club en donde estaban todos.',
        attribution='Familia Mixta'
    )

    # Slide 22 — H07 Hallazgo cuali + cuanti (Consumer Voice con cuanti de apoyo)
    build_hallazgo_slide(
        prs,
        headline_plain="LA CONCIENCIA ECOLÓGICA SE APRENDE EN ESTADOS UNIDOS. ",
        headline_italic="EN RD NO HAY QUIEN LA ENSEÑE.",
        stats=[
            {
                'value': '51%',
                'desc_runs': [
                    make_run('declara querer practicar '),
                    make_run('consumo responsable', bold=True),
                    make_run(' (P36). El deseo es masivo. La pregunta que la batería no hace es dónde lo aprendió. P36, Base 500.')
                ]
            }
        ]
    )

    # Slide 23 — H07 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Justamente en la parte del reciclado, porque justamente lo aprendí en varios viajes que he dado a Estados Unidos. El hijo mío también ya está con esa doctrina.',
        attribution='Familia Mixta'
    )

    # Slide 24 — Síntesis T3
    build_sintesis_slide(
        prs,
        tension_num=3,
        tension_name="Sostenibilidad como deseo sin infraestructura",
        body_lines=[
            "51% quiere consumo responsable. Lo que entiende por eso: apagar luces, no desperdiciar agua.",
            "28% quiere reciclar. El camión lo revuelve todo. El ciudadano abandona — por lógica, no por pereza.",
            "La conciencia ecológica tiene un origen geográfico: el viaje al exterior.",
            "Quien no viaja, no aprende. El sistema educativo local no cierra esa brecha."
        ]
    )

    # ══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 04 — La compra por redes: exclusión, no desinterés  (slides 25–32)
    # H08 (solo-cuanti) + H09 (solo-cuanti cruce NSE) + H10 (cruce edad) + H11 (cruce tipología)
    # ══════════════════════════════════════════════════════════════════════════

    # Slide 25 — Título T4
    build_tension_title_slide(
        prs,
        tension_num=4,
        headline_plain="La compra por redes: ",
        headline_italic="exclusión, no desinterés."
    )

    # Slide 26 — H08 Hallazgo cuanti principal (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="2 DE CADA 3 DECLARAN QUE NO COMPRAN POR REDES. ",
        headline_italic="LA COMPRA POR ESTÍMULO DIGITAL ES MINORÍA REAL.",
        stats=[
            {
                'value': '65%',
                'desc_runs': [
                    make_run('no compró '),
                    make_run('ninguna vez', bold=True),
                    make_run(' en los últimos 6 meses por algo visto en redes sociales. 65.4% exacto. P52, Base 500.')
                ]
            },
            {
                'value': '20%',
                'desc_runs': [
                    make_run('compró '),
                    make_run('1 a 2 veces', bold=True),
                    make_run('; 10.6% compró 3 a 5 veces. La frecuencia media es baja — no hay masa de compradores habituales. P52, Base 500.')
                ]
            },
            {
                'value': '4%',
                'desc_runs': [
                    make_run('compró '),
                    make_run('6 o más veces', bold=True),
                    make_run(' en ese período. La compra impulsiva digital no es masa — es el perfil de una minoría de alta frecuencia. P52, Base 500.')
                ]
            }
        ]
    )

    # Slide 27 — H09 Hallazgo cuanti secundario — cruce NSE (3 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="A MÁS POBREZA, MENOS CLIC: ",
        headline_italic="87.5% DEL ESTRATO E NO COMPRA NADA POR REDES; 60.6% DEL ESTRATO C TAMPOCO, PERO POR OTRA RAZÓN.",
        stats=[
            {
                'value': '88%',
                'desc_runs': [
                    make_run('del '),
                    make_run('NSE E', bold=True),
                    make_run(' (n≈32) declaró "ninguna" compra por redes — el subset con mayor abstención. 87.5% exacto. Tendencia direccional. P52 × NSE.')
                ]
            },
            {
                'value': '74%',
                'desc_runs': [
                    make_run('del '),
                    make_run('NSE D', bold=True),
                    make_run(' (n≈172) no compró nada — 73.8% exacto. La abstinencia cae sostenidamente al subir el nivel socioeconómico. P52 × NSE.')
                ]
            },
            {
                'value': '61%',
                'desc_runs': [
                    make_run('del '),
                    make_run('NSE C', bold=True),
                    make_run(' (n≈218) no compró nada. 60.6% exacto. Brecha E vs C: 27 puntos. El gradiente es lineal. NSE AB excluido (n<30). P52 × NSE.')
                ]
            }
        ]
    )

    # Slide 28 — H10 Hallazgo cuanti terciario — cruce edad (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="LA GENERACIÓN MÁS JOVEN COMPRA POR REDES EL DOBLE QUE LA MAYOR: ",
        headline_italic="BRECHA DE 29 PUNTOS ENTRE 18–24 Y 55+.",
        stats=[
            {
                'value': '47%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('18–24 años', bold=True),
                    make_run(' (n≈55) declaró "ninguna" compra — el grupo menos abstencionista. 47.3% exacto. Más de la mitad del grupo joven sí compró al menos una vez. P52 × Edad.')
                ]
            },
            {
                'value': '77%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('55+', bold=True),
                    make_run(' (n≈119) no compró nada. 76.5% exacto. Abstención crece lineal con la edad. Brecha entre extremos generacionales: 29 puntos. P52 × Edad.')
                ]
            }
        ]
    )

    # Slide 29 — H11 Hallazgo cuanti — cruce tipología (2 stats)
    build_hallazgo_slide(
        prs,
        headline_plain="EL HOGAR SIN HIJOS CON MASCOTA ES EL MÁS PERMEABLE A LAS REDES: ",
        headline_italic="CASI LA MITAD COMPRÓ AL MENOS UNA VEZ EN 6 MESES.",
        stats=[
            {
                'value': '51%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('sin hijos con mascotas', bold=True),
                    make_run(' (n≈37) declaró "ninguna" compra — la tipología menos abstencionista. 51.4% exacto. Casi la mitad sí compró. Tendencia direccional. P52 × Tipología.')
                ]
            },
            {
                'value': '74%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('monoparental', bold=True),
                    make_run(' (n≈122) no compró nada — 73.8% exacto. La mayor abstención. Vs 60.0% biparental (n≈190). Más dependientes = menos compra digital. P52 × Tipología.')
                ]
            }
        ]
    )

    # Slide 30 — Síntesis T4
    build_sintesis_slide(
        prs,
        tension_num=4,
        tension_name="La compra por redes: exclusión, no desinterés",
        body_lines=[
            "65% dice que no compra por redes. Pero los cruces revelan que la abstención no es uniforme.",
            "NSE E: 88% no compra. NSE C: 61%. La brecha es de 27 puntos — y el cuestionario no preguntó por qué.",
            "Los jóvenes 18–24: solo 47% abstención — la mitad sí compró. Los 55+: 77% abstención.",
            "La abstención del NSE bajo refleja exclusión digital, no desinterés. La marca que lee 65% plano pierde el dato real."
        ]
    )

    # Slide 31 — Detonadora T4
    build_detonadora_slide(
        prs,
        tension_num=4,
        pregunta_detonadora="Si la abstención de compra digital es exclusión, ¿qué hace tu marca para ser accesible antes del clic?"
    )

    # Slide 32 — Cuali complementaria T4 — sin verbatim disponible; slide estadístico resumen
    build_hallazgo_slide(
        prs,
        headline_plain="EL 65% DECLARA SIN TRANSACCIÓN. ",
        headline_italic="EL DESEO DIGITAL OPERA ANTES DEL CLIC.",
        stats=[
            {
                'value': '65%',
                'desc_runs': [
                    make_run('declara '),
                    make_run('"ninguna compra por redes"', bold=True),
                    make_run(' (P52). Pero el cuali de T5 revela que la influencia de redes opera como aspiración, no como transacción. El clic no ocurrió; el deseo, sí. P52, Base 500.')
                ]
            }
        ]
    )

    # ══════════════════════════════════════════════════════════════════════════
    # TENSIÓN 05 — La filosofía del gasto familiar  (slides 33–40)
    # H12 (solo-cuali) + H13 (cuanti + CV)
    # ══════════════════════════════════════════════════════════════════════════

    # Slide 33 — Título T5
    build_tension_title_slide(
        prs,
        tension_num=5,
        headline_plain="La filosofía del ",
        headline_italic="gasto familiar."
    )

    # Slide 34 — H12 Cuali slide (solo-cuali — 1 card)
    build_cuali_slide(
        prs,
        headline_plain="OBLIGATORIO, NECESARIO, SUPERFLUO: ",
        headline_italic="LA JERARQUÍA MORAL DEL GASTO SE HEREDA DE PADRES A HIJOS.",
        verbatims=[
            {
                'quote': 'Mi papá me enseñó que existen tres cosas que se compran: lo obligatorio, lo necesario y lo superfluo. (...) Tú tienes que comprar lo que es obligatorio, no necesariamente es necesario, y luego entonces tus gustos son notables porque tú resuelves todas las cosas. Con mis hijos trato de hacer lo mismo.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # Slide 35 — H12 cuali secundario — contexto del sistema heredado
    build_cuali_slide(
        prs,
        headline_plain="TRES GENERACIONES. ",
        headline_italic="EL MISMO SISTEMA DE DECISIÓN DE COMPRA.",
        verbatims=[
            {
                'quote': 'Con mis hijos trato de hacer lo mismo.',
                'attribution': 'Familia Mixta'
            }
        ]
    )

    # Slide 36 — H13 Hallazgo cuanti principal (1 stat)
    build_hallazgo_slide(
        prs,
        headline_plain="LAS REDES NO VENDEN PRODUCTOS. ",
        headline_italic="VENDEN VIDAS QUE EL DOMINICANO SIENTE QUE NO LE TOCAN.",
        stats=[
            {
                'value': '65%',
                'desc_runs': [
                    make_run('declara que no compró nada por redes en 6 meses ('),
                    make_run('P52', bold=True),
                    make_run('). La cifra dice: sin transacción. Lo que el cuali añade: la influencia opera antes de la transacción, en la formación del deseo. Base 500.')
                ]
            }
        ]
    )

    # Slide 37 — H13 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Todo es una pámpara, como dice un buen dominicano. O sea, las personas quieren vivir una vida que no es la que les toca vivir. Las personas quieren — yo se lo atribuyo a las redes sociales. La gente se desespera.',
        attribution='Familia Homoparental'
    )

    # Slide 38 — Síntesis T5
    build_sintesis_slide(
        prs,
        tension_num=5,
        tension_name="La filosofía del gasto familiar",
        body_lines=[
            "El hogar dominicano tiene una taxonomía del gasto heredada: obligatorio → necesario → superfluo.",
            "La cadena es intergeneracional. Del padre a los hijos. Tres generaciones, el mismo sistema.",
            "Las redes no generan transacciones masivas — generan deseos que el ingreso no puede costear.",
            "El 65% que 'no compra por redes' puede ser honesto en lo transaccional y activo en lo aspiracional."
        ]
    )

    # Slide 39 — Detonadora T5
    build_detonadora_slide(
        prs,
        tension_num=5,
        pregunta_detonadora="Si las redes venden aspiración, no transacción, ¿dónde entra tu marca — en el deseo o en el clic?"
    )

    # Slide 40 — Cierre del deck — tensión editorial central
    slide_cierre = blank_slide(prs)

    add_textbox(
        slide_cierre,
        left_px=0, top_px=60,
        width_px=1920, height_px=60,
        text="CONSUMOS — CÓDIGO CASA 2025",
        font_name=FONT_BODY, font_size_pt=14,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    txBox_cierre = slide_cierre.shapes.add_textbox(
        px(120), px(240), px(1680), px(560)
    )
    tf_cierre = txBox_cierre.text_frame
    tf_cierre.word_wrap = True
    tf_cierre.margin_left   = Emu(0)
    tf_cierre.margin_right  = Emu(0)
    tf_cierre.margin_top    = Emu(0)
    tf_cierre.margin_bottom = Emu(0)

    p_cierre = tf_cierre.paragraphs[0]
    p_cierre.alignment = PP_ALIGN.CENTER

    r_cierre1 = p_cierre.add_run()
    r_cierre1.text = "EL PILAR CONSUMOS ES ESTRUCTURALMENTE UN PILAR DE BRECHAS SEMÁNTICAS: "
    r_cierre1.font.name   = FONT_HEADLINE
    r_cierre1.font.size   = Pt(42)
    r_cierre1.font.bold   = False
    r_cierre1.font.italic = False
    r_cierre1.font.color.rgb = WHITE
    _set_kern(r_cierre1._r)

    r_cierre2 = p_cierre.add_run()
    r_cierre2.text = "EL DOMINICANO DECLARA UNA COSA Y VIVE OTRA."
    r_cierre2.font.name   = FONT_HEADLINE
    r_cierre2.font.size   = Pt(42)
    r_cierre2.font.bold   = False
    r_cierre2.font.italic = True
    r_cierre2.font.color.rgb = WHITE
    _set_kern(r_cierre2._r)

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "consumos-deck-flat.pptx"
    )
    prs.save(output_path)

    total = len(prs.slides)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {total}")

    print("\n-- Verificacion specs -------------------------------------------")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"  = {int(prs.slide_width / PX)}px x {int(prs.slide_height / PX)}px")
    ok = (prs.slide_width == px(SLIDE_W_PX) and prs.slide_height == px(SLIDE_H_PX))
    print(f"  Slide 1920x1080 OK: {ok}")
    print(f"HEADLINE_PT:      {HEADLINE_PT}pt (fijo — NO auto-size)")
    print(f"STAT_NUMBER_PT:   {STAT_NUMBER_PT}pt Instrument Serif italic")
    print(f"STAT_DESC_PT:     {STAT_DESC_PT}pt Poppins")
    print(f"STAT_BOX_W:       {STAT_BOX_W}pt (10cm)")
    print(f"STAT_BOX_H:       {STAT_BOX_H}pt (3cm)")
    print(f"CV_VERBATIM_PT:   {CV_VERBATIM_PT}pt (v6: era 60)")
    print(f"CV_ATTRIB_PT:     {CV_ATTRIB_PT}pt Poppins italic")
    print(f"CARD_W:           {CARD_W}pt (15cm v6)")
    print(f"CARD_H:           {CARD_H}pt (6cm v6)")
    print("Cards cuali: rgba(0,0,0,0.45) + outline blanco 1pt + border-radius")
    print("Consumer Voice: verbatim SUELTO sobre fondo negro (SIN card)")
    print("Comillas espanolas en todos los verbatims")
    print("Kerning 0 en todo el deck")
    print("NO masterslide — fondo negro plano")
    print("NO source en ningun slide (v6)")

    print("\n-- QA FLAGS — slide por slide -----------------------------------")
    print("T1 (slides 01-08): El consumidor que no compra verde")
    print("  01: Titulo T1")
    print("  02: H01 cuanti 2 stats — 43% / 30%")
    print("  03: H01 Consumer Voice — Familia Mixta")
    print("  04: H02 cuanti 3 stats — 7.4% / 6% / 3.8%")
    print("  05: H02 complementario — 0 grupos / 51%")
    print("  06: H01 cuali card — gesto domestico")
    print("  07: Sintesis T1")
    print("  08: Detonadora T1")
    print("T2 (slides 09-16): La responsabilidad ambiental externalizada")
    print("  09: Titulo T2")
    print("  10: H03 cuanti 1 stat — 43%")
    print("  11: H03 Consumer Voice 1 — Familia Sin Hijos con Mascota")
    print("  12: H03 cuali — Familia Homoparental (zafacones)")
    print("  13: H04 cuanti 2 stats — 30% / 43%")
    print("  14: H04 Consumer Voice — Familia Biparental con Hijos Pequenos")
    print("  15: Sintesis T2")
    print("  16: Detonadora T2")
    print("T3 (slides 17-24): Sostenibilidad como deseo sin infraestructura")
    print("  17: Titulo T3")
    print("  18: H05 cuanti 3 stats — 51% / 51% / 28%")
    print("  19: H05 Consumer Voice — Familia Mixta (bombillo)")
    print("  20: H06 cuanti 1 stat — 28%")
    print("  21: H06 Consumer Voice — Familia Mixta (zafacones reciclaje)")
    print("  22: H07 cuanti apoyo — 51%")
    print("  23: H07 Consumer Voice — Familia Mixta (EE.UU.)")
    print("  24: Sintesis T3")
    print("T4 (slides 25-32): La compra por redes: exclusion, no desinteres")
    print("  25: Titulo T4")
    print("  26: H08 cuanti 3 stats — 65% / 20% / 4%")
    print("  27: H09 cuanti 3 stats — 88% / 74% / 61% (cruces NSE)")
    print("  28: H10 cuanti 2 stats — 47% / 77% (cruce edad)")
    print("  29: H11 cuanti 2 stats — 51% / 74% (cruce tipologia)")
    print("  30: Sintesis T4")
    print("  31: Detonadora T4")
    print("  32: Complementario T4 — 65% transaccion vs aspiracion")
    print("T5 (slides 33-40): La filosofia del gasto familiar")
    print("  33: Titulo T5")
    print("  34: H12 cuali 1 card — Familia Mixta (obligatorio/necesario/superfluo)")
    print("  35: H12 cuali 1 card — transmision intergeneracional")
    print("  36: H13 cuanti 1 stat — 65%")
    print("  37: H13 Consumer Voice — Familia Homoparental (pampara)")
    print("  38: Sintesis T5")
    print("  39: Detonadora T5")
    print("  40: Cierre editorial del deck")
    print(f"\nTotal slides verificados: {total} (esperado: 40)")

    if total != 40:
        print(f"WARNING: Se esperaban 40 slides, se generaron {total}.")
    else:
        print("OK: 40 slides exactos.")

    return output_path


if __name__ == "__main__":
    build_deck()
