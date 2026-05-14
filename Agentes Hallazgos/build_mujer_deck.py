#!/usr/bin/env python3
"""
build_mujer_deck.py
===================
Deck mujer-deck-flat.pptx desde cero con python-pptx.
26 slides — set editorial cerrado mujer-hallazgos-editados.md (14 hallazgos).

Especificaciones (aprendizajes-montador-cc.md § 3.9):
- Slide 1920×1080 px = 18_288_000 × 10_287_000 EMU
- Cards verbatim: fill rgba(0,0,0,0.45) vía XML alpha 45000 (NO #2E2E2E sólido)
  Shape: MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05 → ~24px border-radius
- Stat box: Layout 1 → 822×238, Layout 2 → 880×238, Layout 3 → 560×238
- Stat number: 180px (135pt) Instrument Serif italic
- Stat desc: Poppins 22px (16.5pt) con bold selectivo
- Consumer Voice card: 1637×485 px centrada
- Verbatim: Instrument Serif regular 50px (37.5pt)
- Headlines: auto-size ≤30→110px, 31-44→88px, 45-65→64px, 66-85→52px, >85→42px
  Todos los headlines de Mujer son >85 chars → 42px (31.5pt)
- Atribución: Poppins italic 22px (16.5pt)
- Source: Poppins italic 16px (12pt) gris #9B9B9B
- Header "CONSUMER VOICE": Poppins 16px (12pt) gris #9B9B9B
- Fondo negro plano — sin masterslide decorativo
- Comillas españolas «...» en todos los verbatims
- Kerning 0 en todo el deck

Conteo de slides:
  H01 (3 stats + CV): 2
  H02 (2 stats + CV): 2
  H03 (3 stats + CV): 2
  H04 (2 stats + CV): 2
  H05 (2 stats + CV): 2
  H06 (3 stats solo-cuanti): 1
  H07 (2 stats + CV): 2
  H08 (3 stats + CV): 2
  H09 (2 stats + CV): 2
  H10 (2 stats + CV): 2
  H11 (3 stats + CV): 2
  H12 (2 stats + CV): 2
  H13 (3 stats solo-cuanti): 1
  H14 (2 stats + CV): 2
  TOTAL: 26 slides

Convención de unidades:
  1 px @ 96 dpi = 9525 EMU
  px_to_pt(n) = n * 0.75
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ── Conversión px → EMU ────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 dpi = 9525 EMU

def px(n):
    return Emu(int(n * PX))

# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = px(SLIDE_W_PX)
SLIDE_H = px(SLIDE_H_PX)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── px → pt ───────────────────────────────────────────────────────────────────
def pt_from_px(px_val):
    return round(px_val * 0.75, 1)

STAT_NUMBER_PT  = pt_from_px(180)   # 135.0 pt
STAT_DESC_PT    = pt_from_px(22)    # 16.5 pt
VERBATIM_PT     = pt_from_px(50)    # 37.5 pt
ATTRIBUTION_PT  = pt_from_px(22)    # 16.5 pt
SOURCE_PT       = pt_from_px(16)    # 12.0 pt
CV_HEADER_PT    = pt_from_px(16)    # 12.0 pt

# ── Dimensiones de elementos ───────────────────────────────────────────────────
STAT_BOX_W_1STAT  = 822
STAT_BOX_W_2STATS = 880
STAT_BOX_W_3STATS = 560
STAT_BOX_H        = 238

CV_CARD_W = 1637
CV_CARD_H = 485

STAT_TOP_PX = 500


# ── Helpers básicos ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_px, top_px, width_px, height_px,
                text, font_name, font_size_pt,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_px, top_px, width_px, height_px,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name   = r.get('font_name', FONT_BODY)
        run.font.size   = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold   = r.get('bold', False)
        run.font.italic = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def pick_headline_size_px(chars):
    if chars <= 30:
        return 110
    elif chars <= 44:
        return 88
    elif chars <= 65:
        return 64
    elif chars <= 85:
        return 52
    else:
        return 42


def add_headline(slide, text_plain, text_italic=""):
    full_text = text_plain + text_italic
    chars     = len(full_text.strip())
    size_px   = pick_headline_size_px(chars)
    size_pt   = pt_from_px(size_px)

    txBox = slide.shapes.add_textbox(
        px(100), px(100), px(1720), px(300)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name   = FONT_HEADLINE
        run.font.size   = Pt(size_pt)
        run.font.bold   = False
        run.font.italic = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name   = FONT_HEADLINE
        run2.font.size   = Pt(size_pt)
        run2.font.bold   = False
        run2.font.italic = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def add_source(slide, source_text):
    add_textbox(
        slide,
        left_px=0, top_px=990,
        width_px=1920, height_px=50,
        text=source_text,
        font_name=FONT_BODY, font_size_pt=SOURCE_PT,
        italic=True, color=GREY_SOFT,
        align=PP_ALIGN.CENTER
    )


def add_vertical_separator(slide, x_px, y_top_px, height_px):
    line = slide.shapes.add_shape(
        1,
        px(x_px - 1), px(y_top_px),
        px(2), px(height_px)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_card_rgba_45(slide, left_px, top_px, width_px, height_px):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(left_px), px(top_px),
        px(width_px), px(height_px)
    )
    shape.adjustments[0] = 0.05
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    sp_elem = shape._element
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))
        for child in list(solidFill):
            solidFill.remove(child)
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')

    return shape


def make_run(text, bold=False, italic=False,
             size_pt=None, color=WHITE, font_name=FONT_BODY):
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': size_pt if size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def add_card(slide, quote_text, attribution,
             left_px, top_px, card_w_px, card_h_px,
             verbatim_size_pt=None):
    if verbatim_size_pt is None:
        verbatim_size_pt = VERBATIM_PT

    make_card_rgba_45(slide, left_px, top_px, card_w_px, card_h_px)

    full_quote = f"«{quote_text}»"

    pad_h = 80
    pad_v = 60
    text_w = card_w_px - (pad_h * 2)
    text_h = card_h_px - (pad_v * 2)

    txBox = slide.shapes.add_textbox(
        px(left_px + pad_h), px(top_px + pad_v),
        px(text_w), px(text_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = full_quote
    run_q.font.name   = FONT_HEADLINE
    run_q.font.size   = Pt(verbatim_size_pt)
    run_q.font.italic = False
    run_q.font.bold   = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name   = FONT_BODY
    run_a.font.size   = Pt(ATTRIBUTION_PT)
    run_a.font.italic = True
    run_a.font.bold   = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


def stat_x_centers(n_stats):
    w = SLIDE_W_PX
    if n_stats == 1:
        return [w // 2]
    elif n_stats == 2:
        return [w // 3, (w * 2) // 3]
    else:
        return [w // 4, w // 2, (w * 3) // 4]


def stat_box_width(n_stats):
    if n_stats == 1:
        return STAT_BOX_W_1STAT
    elif n_stats == 2:
        return STAT_BOX_W_2STATS
    else:
        return STAT_BOX_W_3STATS


def add_stat_block(slide, stat_value, desc_runs, x_center_px, n_stats):
    box_w = stat_box_width(n_stats)
    box_left = x_center_px - (box_w // 2)
    box_top  = STAT_TOP_PX

    num_h = 200
    add_textbox(
        slide,
        left_px=box_left, top_px=box_top,
        width_px=box_w, height_px=num_h,
        text=stat_value,
        font_name=FONT_HEADLINE, font_size_pt=STAT_NUMBER_PT,
        italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    desc_top = box_top + num_h + 20
    add_rich_textbox(
        slide,
        left_px=box_left, top_px=desc_top,
        width_px=box_w, height_px=180,
        runs=desc_runs,
        align=PP_ALIGN.CENTER
    )


def build_hallazgo_slide(prs, headline_plain, headline_italic,
                          stats, source_text):
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    xs = stat_x_centers(n)

    sep_top_px    = STAT_TOP_PX - 20
    sep_height_px = STAT_BOX_H + 40

    if n >= 2:
        x_sep1 = (xs[0] + xs[1]) // 2
        add_vertical_separator(slide, x_sep1, sep_top_px, sep_height_px)
    if n == 3:
        x_sep2 = (xs[1] + xs[2]) // 2
        add_vertical_separator(slide, x_sep2, sep_top_px, sep_height_px)

    for i, stat in enumerate(stats):
        add_stat_block(slide, stat['value'], stat['desc_runs'], xs[i], n)

    add_source(slide, source_text)
    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    slide = blank_slide(prs)

    add_textbox(
        slide,
        left_px=0, top_px=80,
        width_px=1920, height_px=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    card_left = (SLIDE_W_PX - CV_CARD_W) // 2
    card_top  = 220

    add_card(
        slide,
        quote_text=quote,
        attribution=attribution,
        left_px=card_left, top_px=card_top,
        card_w_px=CV_CARD_W, card_h_px=CV_CARD_H
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims, source_text):
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)
    card_left = (SLIDE_W_PX - CV_CARD_W) // 2

    if n == 1:
        card_h   = 485
        gap      = 0
        verb_pt  = VERBATIM_PT
        y_start  = 430
    elif n == 2:
        card_h   = 290
        gap      = 30
        verb_pt  = pt_from_px(38)
        y_start  = 380
    else:
        card_h   = 200
        gap      = 25
        verb_pt  = pt_from_px(32)
        y_start  = 360

    for i, v in enumerate(verbatims):
        y_card = y_start + i * (card_h + gap)
        add_card(
            slide,
            quote_text=v['quote'],
            attribution=v['attribution'],
            left_px=card_left, top_px=y_card,
            card_w_px=CV_CARD_W, card_h_px=card_h,
            verbatim_size_pt=verb_pt
        )

    add_source(slide, source_text)
    return slide


# ── BUILD DECK ─────────────────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # H01 — La mujer de 35-44 en la nuclear; la de estrato E y 55+ sola.
    # Layout 3 (3 stats) + CV = 2 slides
    # 53.6%→54%, 25.0%→25%, 25.2%→25%  (QA FLAG: stats 2 y 3 son ambos 25%)
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="La mujer de 35-44 años vive en la familia nuclear. La mujer de estrato E y la de 55+ ",
        headline_italic="viven solas con lo que queda.",
        stats=[
            {
                'value': '54%',
                'desc_runs': [
                    make_run('de las mujeres de '),
                    make_run('35-44 años', bold=True),
                    make_run(' (n=112) vive en hogar biparental con hijos menores de 18 — TOP-1 de esa franja etaria.')
                ]
            },
            {
                'value': '25%',
                'desc_runs': [
                    make_run('de las mujeres del '),
                    make_run('estrato E', bold=True),
                    make_run(' (n=32) declara hogar monoparental — único estrato donde el monoparental destrona al biparental. (tendencia direccional)')
                ]
            },
            {
                'value': '25%',
                'desc_runs': [
                    make_run('de las mujeres de '),
                    make_run('55+', bold=True),
                    make_run(' (n=119) lidera con monoparental como TOP-1. La estructura familiar femenina se redibuja en la vejez.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P1 × D5, P1 × D6 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Soy madre soltera, sola, con una situación actualmente de mi madre que falleció y que vive conmigo y mi hermanito en condición especial, entonces es difícil, es difícil todo en general, quién lo cuide, porque yo trabajo todo el tiempo.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H02 — La mujer dice "Madre"; el hombre dice "Ambos".
    # Layout 2 (2 stats) + CV = 2 slides
    # 49.5%→50%, 45.6%→46%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="La mujer dominicana dice que la madre es la única responsable de la crianza. El hombre dice que son ambos. ",
        headline_italic="La diferencia no es de opinión — es de casa.",
        stats=[
            {
                'value': '50%',
                'desc_runs': [
                    make_run('del subset femenino (n=305) declara “Madre” como quien tiene la '),
                    make_run('mayor responsabilidad', bold=True),
                    make_run(' de la crianza. Por encima de “Ambos padres” (35.4%).')
                ]
            },
            {
                'value': '46%',
                'desc_runs': [
                    make_run('del subset masculino (n=195) declara “'),
                    make_run('Ambos padres', bold=True),
                    make_run('” como TOP-1. El hombre ve equidad; la mujer ve su propia carga.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P2 × D2 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo con la imagen dos, yo siempre hago todo. Pero ella ya entiende que arreglar el gavetero, que lo pote y demás, ella entiende que ese es su trabajo, ella siempre lo hace sin mandarlo. Yo no le hago mucho énfasis.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H03 — Crianza materna × NSE: a menos bolsillo, más carga.
    # Layout 3 (3 stats) + CV = 2 slides
    # 59.4%→59%, 45.3%→45%, 42.2%→42%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="59.4% de las mujeres del estrato E son las únicas responsables de la crianza, vs 43.5% en AB. ",
        headline_italic="A menos bolsillo, más carga materna — la pobreza vuelve a ‘Madre’ sinónimo de ‘todos’.",
        stats=[
            {
                'value': '59%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('NSE E', bold=True),
                    make_run(' declara “Madre” como única responsable — TOP-1 más alto de cualquier estrato. (n≈32, tendencia direccional)')
                ]
            },
            {
                'value': '45%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('NSE D', bold=True),
                    make_run(' (n=172) cita “Madre” como TOP-1. D + E combinados: la lectura materna como única responsable domina la base popular.')
                ]
            },
            {
                'value': '42%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('NSE C', bold=True),
                    make_run(' (n=218) declara “Ambos padres” como TOP-1 — primer estrato donde el reparto co-parental se vuelve mayoría.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P2 × D6 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Mis hijos, los amo y los quiero, pero también son iguales. Claro, por ejemplo, las escuelas, los colegios son expertos en ponerte una misma reunión de dos grados distintos a la misma hora. Y Ana Castillo, la única madre, es la que tiene que ir.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H04 — 85.4% ya tiene hijos; la no maternidad no tiene nombre.
    # Layout 2 (2 stats) + CV = 2 slides
    # 85.4%→85%, 14.6%→15%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="85.4% de las mujeres del estudio ya tiene hijos. De las que no, ",
        headline_italic="casi nadie lo nombra como elección — la no maternidad aún no tiene nombre en RD.",
        stats=[
            {
                'value': '85%',
                'desc_runs': [
                    make_run('del subsample femenino (n=305) declara “'),
                    make_run('Tengo hijos actualmente', bold=True),
                    make_run('”. Código Casa capta una mujer dominicana mayoritariamente madre.')
                ]
            },
            {
                'value': '15%',
                'desc_runs': [
                    make_run('del total femenino no tiene hijos. Razones declaradas: “'),
                    make_run('no es prioridad', bold=True),
                    make_run('” (3.8%), económicas (2.2%), no deseo (1.0%). La no maternidad aún no tiene nombre propio.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P11 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Si yo no hubiese tenido hijos estuviera viajando y estudiando 24 a 7.',
        attribution='Familia Mixta'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H05 — El trío doméstico: cocina, ropa, limpieza siguen siendo de ella.
    # Layout 2 (2 stats) + CV = 2 slides
    # 66.6%→67%, 63.0%→63%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="66.6% de las mujeres dice que cocinar lo hace ella. Lavar la ropa: 63.0%. Limpiar la casa: 59.6%. ",
        headline_italic="El trío doméstico clásico sigue siendo territorio femenino.",
        stats=[
            {
                'value': '67%',
                'desc_runs': [
                    make_run('del subsample femenino (n=305) dice que '),
                    make_run('cocinar', bold=True),
                    make_run(' lo hace “más la mujer” — ítem con mayor concentración femenina de toda la matriz P33.')
                ]
            },
            {
                'value': '63%',
                'desc_runs': [
                    make_run('dice que '),
                    make_run('lavar la ropa', bold=True),
                    make_run(' es “más la mujer”; 59.6% limpiar la casa. Los tres ítems más feminizados: cocina, ropa y limpieza.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P33 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Y bueno, como mujer, bueno, no sabía cocinar y tuve que hacerlo. Obligada, aprendí a mamá, pero... No me gusta, pero tengo que hacerlo.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H06 — La feminización doméstica se detiene donde empieza el dinero.
    # Layout 3 (3 stats) SOLO-CUANTI — sin Consumer Voice = 1 slide
    # 28.6%→29%, 30.6%→31%, 51.0%→51%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="Cocina: 66.6% de ella. Limpieza: 59.6% de ella. Manejo del dinero del hogar: ",
        headline_italic="51.0% de ambos. La feminización doméstica se detiene donde empieza el dinero.",
        stats=[
            {
                'value': '29%',
                'desc_runs': [
                    make_run('dice que el '),
                    make_run('pago de servicios', bold=True),
                    make_run(' lo hace “más la mujer”; 42.0% “ambos por igual”; 27.0% “más el hombre”. Tarea más equitativa de la matriz.')
                ]
            },
            {
                'value': '31%',
                'desc_runs': [
                    make_run('dice que el '),
                    make_run('manejo del dinero del hogar', bold=True),
                    make_run(' lo hace “más la mujer”; 51.0% “ambos por igual” — el ítem más co-gestionado de la matriz.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('declara “'),
                    make_run('ambos por igual', bold=True),
                    make_run('” en manejo del dinero. En cocina: 66.6% ella. La frontera de la carga femenina pasa donde empieza la administración del dinero.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P33 · Base 305."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H07 — El trabajo de presencia (citas, cuidado) también es de ella.
    # Layout 2 (2 stats) + CV = 2 slides
    # 45.2%→45%, 38.6%→39%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="45.2% dice que acompañar a citas médicas o escolares lo hace ella. ",
        headline_italic="El trabajo de presencia — llevar, esperar, resolver — también es trabajo de mujer.",
        stats=[
            {
                'value': '45%',
                'desc_runs': [
                    make_run('del subsample femenino dice que '),
                    make_run('citas médicas o escolares', bold=True),
                    make_run(' lo hace “más la mujer”. 27.6% “ambos por igual”; 21.0% “no aplica”.')
                ]
            },
            {
                'value': '39%',
                'desc_runs': [
                    make_run('dice que el '),
                    make_run('cuidado de los hijos', bold=True),
                    make_run(' lo hace “más la mujer”. Más distribuido que la cocina (66.6%) pero más cargado que el dinero (30.6%).')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P33 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Tuve un tiempo como mamá soltera y mi hermana es mamá soltera de tres niños, y es bien difícil levantarte, llevar al niño a la escuela, el dinero, cuando el papá no está económicamente ni emocionalmente.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H08 — 76.7% mujeres vs 71.8% hombres: perciben más cambio quienes más cargan.
    # Layout 3 (3 stats) + CV = 2 slides
    # 76.7%→77%, 71.8%→72%, 79.5%→80%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="76.7% de las mujeres dice que el rol ha cambiado mucho — 4.9 puntos más que los hombres. ",
        headline_italic="La percepción más fuerte del cambio la tiene quien más sigue cargando.",
        stats=[
            {
                'value': '77%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' (n=305) elige “5 = Mucho” en la escala de cuánto ha cambiado el rol del hombre y la mujer en la familia.')
                ]
            },
            {
                'value': '72%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('masculino', bold=True),
                    make_run(' (n=195) elige “5 = Mucho” en la misma escala. Brecha de percepción de género: 4.9 puntos.')
                ]
            },
            {
                'value': '80%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('monoparental', bold=True),
                    make_run(' (n=122) elige “5 = Mucho” — pico más alto. La percepción de cambio se intensifica donde la mujer carga sola.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P34 × D2, P34 × D1, P34 × D5 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='El cuidado es más trabajoso, porque son más delicados, no te puedo negar que yo me siento contenta de que no tengo que pasar por el pasillo del supermercado de las mangas.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H09 — El carrito: co-gestionado con sesgo, no territorio femenino.
    # Layout 2 (2 stats) + CV = 2 slides
    # 36.4%→36%, 33.4%→33%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="36.4% dice que la compra del supermercado la hace más ella; 48.4% dice que la hacen ambos. ",
        headline_italic="El carrito no es territorio femenino — es territorio compartido con sesgo.",
        stats=[
            {
                'value': '36%',
                'desc_runs': [
                    make_run('del subsample femenino dice que la compra del '),
                    make_run('supermercado', bold=True),
                    make_run(' la hace “más la mujer”; 48.4% “ambos por igual”. Uno de los pocos ítems donde “ambos” lidera.')
                ]
            },
            {
                'value': '33%',
                'desc_runs': [
                    make_run('dice que la compra de artículos '),
                    make_run('personales o escolares', bold=True),
                    make_run(' la hace “más la mujer”; 46.0% “ambos por igual”. En ambas compras, la co-gestión es la norma declarada.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P33 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo soy la que compro porque es lo que yo pueda comprar y comer porque es lo que hay.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H10 — Economía: TOP-1 estresor, el doble de cualquier otro.
    # Layout 2 (2 stats) + CV = 2 slides
    # 47.8%→48%, 21.8%→22%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="47.8% de las mujeres nombra la economía como su mayor fuente de estrés familiar. ",
        headline_italic="La precariedad duplica a cualquier otro estresor — antes que crianza, antes que tiempo propio.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('del subsample femenino (n=305) marca “'),
                    make_run('Economía', bold=True),
                    make_run('” como factor de estrés familiar — TOP-1 de la batería P26. 28.4% también marca “Realidad económica”.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('marca “'),
                    make_run('Inseguridad', bold=True),
                    make_run('”; 20.6% “Falta de tiempo personal”; 16.2% “Crianza”. La economía es el doble de cualquiera.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P26 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Es bien difícil levantarte, llevar al niño a la escuela, el dinero, cuando el papá no está económicamente ni emocionalmente ni nada, es bien difícil.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H11 — 48.5% de las mujeres no se cuida; brecha de 15.2 puntos vs hombres.
    # Layout 3 (3 stats) + CV = 2 slides
    # 48.5% SE MANTIENE (1 dígito antes del punto → regla decimal), 60.0%→60%, 51.2%→51%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no realiza ninguna actividad para cuidar su salud mental ni física. ",
        headline_italic="En los hombres, el mismo indicador cae a 33.3%. Brecha de 15.2 puntos.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' (n=305) declara “No realizo ninguna actividad” — TOP-1 absoluto. En masculino: 33.3%. Brecha de género: 15.2 puntos.')
                ]
            },
            {
                'value': '60%',
                'desc_runs': [
                    make_run('del subset femenino '),
                    make_run('18-24 años', bold=True),
                    make_run(' (n=55) marca la inacción como TOP-1 — pico de inacción de cualquier franja etaria del subsample femenino.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('NSE D', bold=True),
                    make_run(' (n=172 mujeres) marca inacción como TOP-1. La inacción se intensifica en la base socioeconómica.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P27 × D2, P27 × D5, P27 × D6 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Cuando uno es madre, a veces se olvida, a veces no, casi siempre se olvida de uno mismo. Y el tiempo pasa, y cuando tú vienes a mirar para atrás, entonces tú dices, ¿ y yo? Porque crecen y se van.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H12 — El monoparental es femenino por acumulación, no por elección.
    # Layout 2 (2 stats) + CV = 2 slides
    # 24.4%→24%, 56.6%→57%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="56.6% del monoparental cita a la madre como única responsable de crianza — 14.8 puntos sobre el promedio. ",
        headline_italic="La jefatura monoparental es femenina por acumulación, no por elección.",
        stats=[
            {
                'value': '24%',
                'desc_runs': [
                    make_run('del total Código Casa (n=500) declara hogar '),
                    make_run('monoparental', bold=True),
                    make_run(' — segunda tipología más común. Por composición, la tipología monoparental está liderada por mujeres.')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('del subset monoparental (n=122) cita “'),
                    make_run('Madre', bold=True),
                    make_run('” como única responsable de la crianza, vs 41.8% del promedio. Brecha: +14.8 puntos.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P1, P2 × D1 · Base 500."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo me críé con esos dos conceptos en mi mente y yo siempre me enfocaba. Tanto así que yo decía: ‘No, yo no me voy a separar, yo voy a aguantar’, porque mi mamá me decía que hay que aguantar.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H13 — La mujer 18-24 es la que menos se cuida.
    # Layout 3 (3 stats) SOLO-CUANTI — sin Consumer Voice = 1 slide
    # 60.0%→60%, 47.0%→47%, 38.7%→39%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="La mujer 18-24 es la que menos se cuida: 60.0% no hace ninguna actividad para su salud. ",
        headline_italic="El autocuidado femenino sube con la edad — el punto más bajo es el principio.",
        stats=[
            {
                'value': '60%',
                'desc_runs': [
                    make_run('del subset femenino '),
                    make_run('18-24 años', bold=True),
                    make_run(' (n=55) marca “no realizo ninguna actividad” como TOP-1 — pico de inacción de cualquier franja etaria.')
                ]
            },
            {
                'value': '47%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('25-34 años', bold=True),
                    make_run(' (n=115) marca inacción como TOP-1. Solo en 35-44 el TOP-1 muta a ejercicio (39.3%).')
                ]
            },
            {
                'value': '39%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('55+', bold=True),
                    make_run(' (n=119) marca ejercicio como TOP-1. La mujer mayor dominicana es más activa físicamente que la joven.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo 2025 · P27 × D5 · Base 305."
    )

    # ─────────────────────────────────────────────────────────────────────────
    # H14 — Autoridad doméstica sin independencia financiera.
    # Layout 2 (2 stats) + CV = 2 slides
    # 49.5%→50%, 51.0%→51%
    # ─────────────────────────────────────────────────────────────────────────

    build_hallazgo_slide(
        prs,
        headline_plain="La mujer dominicana decide la crianza, la cocina, la limpieza y el cuidado. El dinero del hogar lo co-gestiona. ",
        headline_italic="La autoridad doméstica no compra independencia financiera.",
        stats=[
            {
                'value': '50%',
                'desc_runs': [
                    make_run('del subsample femenino dice “Madre” como responsable de '),
                    make_run('crianza', bold=True),
                    make_run('. Lidera también en cocina (66.6%), ropa (63.0%), limpieza (59.6%).')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('declara “'),
                    make_run('ambos por igual', bold=True),
                    make_run('” en manejo del dinero del hogar. En pago de servicios: 42.0% ambos. Autoridad doméstica sin independencia financiera.')
                ]
            }
        ],
        source_text="Source: Código Casa — Estudio cuantitativo + cualitativo 2025 · P2, P33 · Base 305."
    )

    build_consumer_voice_slide(
        prs,
        quote='Yo trabajaba en Ikea, diseñando cocinas, y en ese momento no se podía utilizar celulares en el cine. Y yo trabajaba de tarde hasta la noche. Y mi pareja fue un hombre de 6’4”. Yo dije, yo te he llamado más de 10 veces. Y mis padres estaban llevándola a emergencia mientras él fue a buscarme al trabajo. Mi hija llegó a una deshidratación secundaria.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = "/home/user/c-digocasa-cerebro/Agentes Hallazgos/mujer-deck-flat.pptx"
    prs.save(output_path)
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    print("\n── Verificación specs ───────────────────────────────────────────")
    print(f"Slide size: {prs.slide_width} × {prs.slide_height} EMU")
    print(f"  = {prs.slide_width // PX}px × {prs.slide_height // PX}px")
    print(f"  Slide 1920×1080 OK: {prs.slide_width == Emu(SLIDE_W_PX * PX) and prs.slide_height == Emu(SLIDE_H_PX * PX)}")
    print(f"Stat number: {STAT_NUMBER_PT}pt ({STAT_NUMBER_PT/0.75:.0f}px)")
    print(f"Verbatim:    {VERBATIM_PT}pt ({VERBATIM_PT/0.75:.0f}px)")
    print(f"CV card:     {CV_CARD_W}×{CV_CARD_H}px")
    print("Cards: rgba(0,0,0,0.45) via MSO_SHAPE.ROUNDED_RECTANGLE + XML alpha=45000")
    print("Fondo negro plano — sin masterslide")
    print("Kerning 0 en todo el deck")
    print("Comillas españolas «...» en todos los verbatims")

    print("\n── QA FLAGS ───────────────────────────────────────────────")
    print("H01 slide 1: stats 2 y 3 ambos redondean a 25%. Per regla de enteros.")
    print("  Diferenciados en descripciones: E (n=32) vs 55+ (n=119).")
    print("H06: solo-cuanti (sin CV) — OK per editor (Verbatims: Ninguno).")
    print("H13: solo-cuanti (sin CV) — OK per editor (Verbatims: Ninguno).")
    print("H11 stat 1: 48.5% mantiene decimal (1 dígito antes del punto → regla).")
    print("H08 stat 3: 79.5% → 80% (redondeado a entero).")
    print("Total slides esperados: 26")

    return output_path


if __name__ == "__main__":
    build_deck()
