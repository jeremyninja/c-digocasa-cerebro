#!/usr/bin/env python3
"""
build_bienestar_v5.py
=====================
Deck bienestar-deck-flat-v5.pptx desde cero con python-pptx.
23 slides — set editorial cerrado bienestar-hallazgos-editados.md.

CAMBIOS CRÍTICOS v4 → v5 (specs DEFINITIVAS en pt de aprendizajes-montador-cc.md § 3.9):

1. Convención de unidades:
   - 1 pt Keynote = 9525 EMU (equivalente a 1 px @ 96 dpi en python-pptx)
   - Se define PT_EMU = 9525 y se usa como multiplicador.
   - Todas las dimensiones del archivo de aprendizajes son en pt (sistema Keynote).

2. Slide: 1920×1080 pt = 18_288_000 × 10_287_000 EMU (mismos valores absolutos).

3. Headlines: 50pt FIJO (no auto-size por largo). Todos los headlines van a 50pt.
   v4 usaba 42px (31.5pt) para headlines >85 chars. Ahora 50pt fijo.

4. Cifras grandes: 180pt Instrument Serif italic blanco.
   v4 usaba 135pt (derivado de 180px × 0.75pt/px). Ahora 180pt directos.

5. Caja descripción stat: 360×75pt. Poppins 16pt.
   v4 usaba cajas de 560–880px × 238px con Poppins 22px.

6. Respiración entre stats:
   - 2 stats: gap 120pt mínimo → conjunto total 360+120+360=840pt,
     empieza en (1920-840)/2 = 540pt.
   - 3 stats: gap 80pt mínimo → conjunto total 360*3+80*2=1240pt,
     empieza en (1920-1240)/2 = 340pt.

7. Consumer Voice card: 1637×485pt centrada.
   - Verbatim: 60pt Instrument Serif regular blanco, comillas «...»
   - Atribución: 20pt Poppins italic blanco, centrada debajo
   - Header "CONSUMER VOICE": 14pt Poppins gris #9B9B9B, top 80pt

8. Cards cualitativas: 500×245pt SIDE BY SIDE (no apiladas).
   - 1 card: centrada en slide.
   - 2 cards: side by side con gap 80pt.
   - 3 cards: side by side con gap 50pt.
   - Verbatim: 22-28pt según largo.
   - Atribución: 14pt Poppins italic.
   H10 lleva 2 cards SIDE BY SIDE.
   H12 lleva 1 card centrada.

9. NO SOURCE en ningún slide. Pie de página técnico eliminado.

10. NO cajas "Pregunta P##. Base..." debajo de stats.

11. NO masterslide decorativo. Fondo negro plano #000000.

12. Cards: fill rgba(0,0,0,0.45) vía XML alpha=45000 + ROUNDED_RECTANGLE.
    Border-radius: adjustments[0] = 0.05 (15pt aprox en card de 485pt alto).

13. Kerning 0 en todo el deck.

Posiciones calculadas literalmente (§ 3.9 Layouts):
  - Headline: left 100pt, top 100pt, width 1720pt
  - Layout 1 (1 stat):
    - Cifra: top 420pt, centrada en slide
    - Desc: left 780pt, top 660pt, width 360pt, height 75pt
  - Layout 2 (2 stats):
    - Desc1: left 540pt, top 660pt
    - Desc2: left 1020pt, top 660pt
    - Gap entre cajas: 120pt
  - Layout 3 (3 stats):
    - Desc1: left 340pt, top 660pt
    - Desc2: left 780pt, top 660pt  (gap 80pt después de 340+360=700)
    - Desc3: left 1220pt, top 660pt
  - Cifras centradas sobre sus descripciones, top 420pt
  - Consumer Voice card: left 141.5pt, top 297pt, width 1637pt, height 485pt
  - Cards cuali 2 side by side: left 420, 1000; top centrado
  - Cards cuali 1 sola: left 710, top 417
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ── Conversión pt (Keynote) → EMU ─────────────────────────────────────────────
# Convención: 1 pt Keynote = 9525 EMU (igual que 1 px @ 96 dpi en python-pptx)
PT_EMU = 9525  # EMU por pt Keynote

def pt(n):
    """Convierte pt (Keynote) a EMU."""
    return Emu(int(n * PT_EMU))


# ── Dimensiones del slide ──────────────────────────────────────────────────────
SLIDE_W_PT = 1920
SLIDE_H_PT = 1080
SLIDE_W = pt(SLIDE_W_PT)
SLIDE_H = pt(SLIDE_H_PT)

# ── Colores ────────────────────────────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_SEP  = RGBColor(0x2E, 0x2E, 0x2E)  # separadores verticales entre stats
GREY_SOFT = RGBColor(0x9B, 0x9B, 0x9B)  # header CONSUMER VOICE

# ── Tipografías ────────────────────────────────────────────────────────────────
FONT_HEADLINE = "Instrument Serif"
FONT_BODY     = "Poppins"

# ── Tamaños de fuente en pt (SPECS DEFINITIVAS § 3.9) ─────────────────────────
HEADLINE_PT     = 50    # FIJO para todos los headlines. No auto-size.
STAT_NUMBER_PT  = 180   # Cifras grandes Instrument Serif italic
STAT_DESC_PT    = 16    # Caja descripción Poppins
CV_VERBATIM_PT  = 60    # Verbatim Consumer Voice card (1637×485pt)
CV_ATTR_PT      = 20    # Atribución Consumer Voice
CV_HEADER_PT    = 14    # Header "CONSUMER VOICE"

# Cards cualitativas (500×245pt) — más pequeñas
CUALI_ATTR_PT   = 14    # Atribución en cards cuali

# ── Posiciones canónicas (§ 3.9 Layouts) en pt ────────────────────────────────
HEADLINE_LEFT   = 100
HEADLINE_TOP    = 100
HEADLINE_WIDTH  = 1720   # 1920 - 100*2
HEADLINE_HEIGHT = 300    # suficiente para 3-4 líneas a 50pt

STAT_NUM_TOP    = 420    # top de las cifras grandes
STAT_DESC_TOP   = 660    # top de las cajas de descripción
STAT_DESC_W     = 360    # anchura de cada caja descripción (spec exacta)
STAT_DESC_H     = 75     # altura de cada caja descripción (spec exacta)

# Gaps entre stats (spec § 3.9)
GAP_2_STATS     = 120    # gap horizontal entre cajas con 2 stats
GAP_3_STATS     = 80     # gap horizontal entre cajas con 3 stats

# Posiciones calculadas para 2 stats (conjunto = 360+120+360 = 840pt)
STATS2_START    = (SLIDE_W_PT - (STAT_DESC_W * 2 + GAP_2_STATS)) // 2  # 540pt
STATS2_DESC_X   = [STATS2_START, STATS2_START + STAT_DESC_W + GAP_2_STATS]  # [540, 1020]

# Posiciones calculadas para 3 stats (conjunto = 360*3 + 80*2 = 1240pt)
STATS3_START    = (SLIDE_W_PT - (STAT_DESC_W * 3 + GAP_3_STATS * 2)) // 2  # 340pt
STATS3_DESC_X   = [
    STATS3_START,                              # 340
    STATS3_START + STAT_DESC_W + GAP_3_STATS,  # 780
    STATS3_START + (STAT_DESC_W + GAP_3_STATS) * 2  # 1220
]

# Consumer Voice card (§ 3.9 Layout 4)
CV_CARD_W       = 1637
CV_CARD_H       = 485
CV_CARD_LEFT    = (SLIDE_W_PT - CV_CARD_W) // 2   # 141pt (≈141.5 redondeado)
CV_CARD_TOP     = (SLIDE_H_PT - CV_CARD_H) // 2   # 297pt (centrada vertical)
CV_HEADER_TOP   = 80

# Cards cualitativas (§ 3.9 Layout 5)
CUALI_CARD_W    = 500
CUALI_CARD_H    = 245

# 1 card centrada
CUALI_1_LEFT    = (SLIDE_W_PT - CUALI_CARD_W) // 2   # 710pt
CUALI_1_TOP     = (SLIDE_H_PT - CUALI_CARD_H) // 2   # 417pt (centrado)

# 2 cards side by side, gap 80pt
CUALI_2_GAP     = 80
CUALI_2_TOTAL   = CUALI_CARD_W * 2 + CUALI_2_GAP     # 1080pt
CUALI_2_START   = (SLIDE_W_PT - CUALI_2_TOTAL) // 2  # 420pt
CUALI_2_LEFT    = [CUALI_2_START, CUALI_2_START + CUALI_CARD_W + CUALI_2_GAP]  # [420, 1000]
CUALI_2_TOP     = (SLIDE_H_PT - CUALI_CARD_H) // 2   # 417pt

# 3 cards side by side, gap 50pt
CUALI_3_GAP     = 50
CUALI_3_TOTAL   = CUALI_CARD_W * 3 + CUALI_3_GAP * 2  # 1600pt
CUALI_3_START   = (SLIDE_W_PT - CUALI_3_TOTAL) // 2   # 160pt
CUALI_3_LEFT    = [
    CUALI_3_START,
    CUALI_3_START + CUALI_CARD_W + CUALI_3_GAP,
    CUALI_3_START + (CUALI_CARD_W + CUALI_3_GAP) * 2
]  # [160, 710, 1260]
CUALI_3_TOP     = (SLIDE_H_PT - CUALI_CARD_H) // 2   # 417pt

# Ajuste top de cards cuali cuando hay headline arriba (top 100+300+40=440 aprox)
# Las cards cuali van en vertical 400-750pt aprox cuando hay headline
CUALI_TOP_WITH_HEADLINE = 420   # top de cards con headline arriba

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Slide en blanco con fondo negro plano #000000."""
    layout = prs.slide_layouts[6]   # Blank
    slide  = prs.slides.add_slide(layout)
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def _set_kern(r_elem):
    """Kerning 0 en <a:rPr> del run."""
    rpr = r_elem.find(qn('a:rPr'))
    if rpr is not None:
        rpr.set('kern', '0')
        rpr.set('spc', '0')


def add_textbox(slide, left_pt_val, top_pt_val, width_pt_val, height_pt_val,
                text, font_name, font_size_pt_val,
                bold=False, italic=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox de un solo run. Todas las dimensiones en pt (Keynote)."""
    txBox = slide.shapes.add_textbox(
        pt(left_pt_val), pt(top_pt_val),
        pt(width_pt_val), pt(height_pt_val)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name    = font_name
    run.font.size    = Pt(font_size_pt_val)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    _set_kern(run._r)
    return txBox


def add_rich_textbox(slide, left_pt_val, top_pt_val, width_pt_val, height_pt_val,
                     runs, align=PP_ALIGN.LEFT, word_wrap=True):
    """Textbox con múltiples runs (bold selectivo). Dimensiones en pt."""
    txBox = slide.shapes.add_textbox(
        pt(left_pt_val), pt(top_pt_val),
        pt(width_pt_val), pt(height_pt_val)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name    = r.get('font_name', FONT_BODY)
        run.font.size    = Pt(r.get('font_size_pt', STAT_DESC_PT))
        run.font.bold    = r.get('bold', False)
        run.font.italic  = r.get('italic', False)
        run.font.color.rgb = r.get('color', WHITE)
        _set_kern(run._r)
    return txBox


def add_headline(slide, text_plain, text_italic=""):
    """
    Headline en MAYÚSCULAS, 50pt FIJO, centrado.
    Posición: top 100pt, left 100pt, width 1720pt, height 300pt.
    Split: text_plain en regular, text_italic en italic.
    """
    txBox = slide.shapes.add_textbox(
        pt(HEADLINE_LEFT), pt(HEADLINE_TOP),
        pt(HEADLINE_WIDTH), pt(HEADLINE_HEIGHT)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if text_plain:
        run = p.add_run()
        run.text = text_plain.upper()
        run.font.name    = FONT_HEADLINE
        run.font.size    = Pt(HEADLINE_PT)
        run.font.bold    = False
        run.font.italic  = False
        run.font.color.rgb = WHITE
        _set_kern(run._r)

    if text_italic:
        run2 = p.add_run()
        run2.text = text_italic.upper()
        run2.font.name    = FONT_HEADLINE
        run2.font.size    = Pt(HEADLINE_PT)
        run2.font.bold    = False
        run2.font.italic  = True
        run2.font.color.rgb = WHITE
        _set_kern(run2._r)

    return txBox


def make_card_rgba_45(slide, left_pt_val, top_pt_val, width_pt_val, height_pt_val):
    """
    ROUNDED_RECTANGLE con fill rgba(0,0,0,0.45) vía XML alpha=45000.
    adjustments[0] = 0.05 → border-radius ~15pt en card de 485pt alto.
    Sin borde.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        pt(left_pt_val), pt(top_pt_val),
        pt(width_pt_val), pt(height_pt_val)
    )

    shape.adjustments[0] = 0.05   # ~15pt border-radius
    shape.line.fill.background()  # sin borde

    # Fill negro sólido base
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK

    # Aplicar opacidad 45% vía XML directo
    sp_elem  = shape._element
    spPr     = sp_elem.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(spPr, qn('a:solidFill'))

        for child in list(solidFill):
            solidFill.remove(child)

        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '000000')

        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '45000')   # 45% en notación OOXML (0–100000)

    return shape


def add_vertical_separator(slide, x_pt_val, y_top_pt_val, height_pt_val):
    """
    Línea vertical fina #2E2E2E entre stats.
    Ancho 2pt visual.
    """
    line = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        pt(x_pt_val - 1), pt(y_top_pt_val),
        pt(2), pt(height_pt_val)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREY_SEP
    line.line.fill.background()


def make_run(text, bold=False, italic=False,
             font_size_pt=None, color=WHITE, font_name=FONT_BODY):
    """Helper para construir dict de run en rich textbox."""
    return {
        'text': text,
        'font_name': font_name,
        'font_size_pt': font_size_pt if font_size_pt is not None else STAT_DESC_PT,
        'bold': bold,
        'italic': italic,
        'color': color
    }


def add_card_with_text(slide, quote_text, attribution,
                       left_pt_val, top_pt_val, card_w_pt, card_h_pt,
                       verbatim_pt=None, attr_pt=None,
                       pad_h=80, pad_v=50):
    """
    Card verbatim: ROUNDED_RECTANGLE rgba(0,0,0,0.45) + texto sobre ella.
    quote_text: sin comillas (se agregan como «...»).
    verbatim_pt: tamaño en pt del verbatim.
    attr_pt: tamaño en pt de la atribución.
    pad_h, pad_v: padding interno en pt.
    """
    if verbatim_pt is None:
        verbatim_pt = CV_VERBATIM_PT   # 60pt — Consumer Voice
    if attr_pt is None:
        attr_pt = CV_ATTR_PT            # 20pt — Consumer Voice

    # Fondo de la card
    make_card_rgba_45(slide, left_pt_val, top_pt_val, card_w_pt, card_h_pt)

    # Textbox sobre la card con padding interno
    text_left   = left_pt_val + pad_h
    text_top    = top_pt_val  + pad_v
    text_width  = card_w_pt   - (pad_h * 2)
    text_height = card_h_pt   - (pad_v * 2)

    txBox = slide.shapes.add_textbox(
        pt(text_left), pt(text_top),
        pt(text_width), pt(text_height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    # Párrafo 1: verbatim con comillas españolas
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run_q = p1.add_run()
    run_q.text = f"«{quote_text}»"
    run_q.font.name    = FONT_HEADLINE
    run_q.font.size    = Pt(verbatim_pt)
    run_q.font.italic  = False    # NO italic — spec § 3.9
    run_q.font.bold    = False
    run_q.font.color.rgb = WHITE
    _set_kern(run_q._r)

    # Párrafo 2: atribución
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run_a = p2.add_run()
    run_a.text = f"— {attribution}"
    run_a.font.name    = FONT_BODY
    run_a.font.size    = Pt(attr_pt)
    run_a.font.italic  = True
    run_a.font.bold    = False
    run_a.font.color.rgb = WHITE
    _set_kern(run_a._r)

    return txBox


# ── Stat positions por número de stats ────────────────────────────────────────

def stat_desc_x_positions(n_stats):
    """
    Retorna lista de posiciones left en pt para las cajas de descripción.
    Calculadas literalmente según § 3.9.
    """
    if n_stats == 1:
        # Caja centrada: (1920 - 360) / 2 = 780pt
        return [(SLIDE_W_PT - STAT_DESC_W) // 2]
    elif n_stats == 2:
        return STATS2_DESC_X
    else:
        return STATS3_DESC_X


def stat_num_x_center(desc_left_pt):
    """Centro horizontal de la cifra = centro de la caja de descripción."""
    return desc_left_pt + STAT_DESC_W // 2


# ── Builders de slides ─────────────────────────────────────────────────────────

def build_hallazgo_slide(prs, headline_plain, headline_italic, stats):
    """
    Slide Hallazgo cuanti. NO SOURCE.
    stats: lista de dicts {value, desc_runs} — 1, 2 o 3 elementos.

    Posiciones (§ 3.9):
    - Headline: left 100, top 100, width 1720, height 300
    - Cifras: top 420pt, centradas sobre su caja de descripción
    - Cajas desc: top 660pt, width 360pt, height 75pt
    - Gaps: 120pt (2 stats), 80pt (3 stats)
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(stats)
    desc_xs = stat_desc_x_positions(n)

    # Separadores verticales entre stats (§ 3.9 — finos, #2E2E2E)
    sep_top    = STAT_NUM_TOP - 20         # 400pt
    sep_height = (STAT_DESC_TOP + STAT_DESC_H) - sep_top  # 400→735 = 335pt

    if n >= 2:
        # Separador entre stat 1 y 2: punto medio entre las cajas
        x_sep1 = desc_xs[0] + STAT_DESC_W + GAP_2_STATS // 2
        add_vertical_separator(slide, x_sep1, sep_top, sep_height)
    if n == 3:
        x_sep2 = desc_xs[1] + STAT_DESC_W + GAP_3_STATS // 2
        add_vertical_separator(slide, x_sep2, sep_top, sep_height)

    # Montar cada stat
    for i, stat in enumerate(stats):
        desc_left = desc_xs[i]
        num_center = stat_num_x_center(desc_left)

        # Cifra grande: textbox centrado sobre la caja de descripción
        # La cifra a 180pt puede superar los 300pt de ancho — usamos 600pt de ancho
        # centrado en num_center para que no desborde hacia el stat adyacente
        num_box_w = min(STAT_DESC_W + 120, 480)  # max 480pt para no chocar
        num_box_left = num_center - num_box_w // 2

        add_textbox(
            slide,
            left_pt_val=num_box_left, top_pt_val=STAT_NUM_TOP,
            width_pt_val=num_box_w, height_pt_val=220,
            text=stat['value'],
            font_name=FONT_HEADLINE, font_size_pt_val=STAT_NUMBER_PT,
            italic=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

        # Caja descripción 360×75pt
        add_rich_textbox(
            slide,
            left_pt_val=desc_left, top_pt_val=STAT_DESC_TOP,
            width_pt_val=STAT_DESC_W, height_pt_val=STAT_DESC_H,
            runs=stat['desc_runs'],
            align=PP_ALIGN.CENTER
        )

    return slide


def build_consumer_voice_slide(prs, quote, attribution):
    """
    Slide Consumer Voice — 1 verbatim en card 1637×485pt centrada.
    NO SOURCE.

    Posiciones (§ 3.9 Layout 4):
    - Header "CONSUMER VOICE": top 80pt, centrado
    - Card: left 141pt, top 297pt (centrada vertical)
    - Verbatim: 60pt Instrument Serif regular
    - Atribución: 20pt Poppins italic
    """
    slide = blank_slide(prs)

    # Header "CONSUMER VOICE" — top 80pt, Poppins 14pt gris #9B9B9B
    add_textbox(
        slide,
        left_pt_val=0, top_pt_val=CV_HEADER_TOP,
        width_pt_val=SLIDE_W_PT, height_pt_val=50,
        text="CONSUMER VOICE",
        font_name=FONT_BODY, font_size_pt_val=CV_HEADER_PT,
        color=GREY_SOFT, align=PP_ALIGN.CENTER
    )

    # Card centrada — 1637×485pt
    add_card_with_text(
        slide,
        quote_text=quote,
        attribution=attribution,
        left_pt_val=CV_CARD_LEFT, top_pt_val=CV_CARD_TOP,
        card_w_pt=CV_CARD_W, card_h_pt=CV_CARD_H,
        verbatim_pt=CV_VERBATIM_PT,   # 60pt
        attr_pt=CV_ATTR_PT,            # 20pt
        pad_h=80, pad_v=60
    )

    return slide


def build_cuali_slide(prs, headline_plain, headline_italic, verbatims):
    """
    Slide solo-cuali: headline arriba + 1–3 cards SIDE BY SIDE. NO SOURCE.

    Dimensiones de cards (§ 3.9 Layout 5):
    - 1 card: 500×245pt centrada, verbatim 28pt, attr 14pt
    - 2 cards: side by side gap 80pt, left [420, 1000], top 417pt
    - 3 cards: side by side gap 50pt, left [160, 710, 1260], top 417pt

    Con headline arriba → top de cards ajustado a 420pt aprox.
    """
    slide = blank_slide(prs)
    add_headline(slide, headline_plain, headline_italic)

    n = len(verbatims)

    if n == 1:
        # 1 card centrada — con headline arriba ajustar top
        card_top  = CUALI_TOP_WITH_HEADLINE   # 420pt
        card_left = CUALI_1_LEFT              # 710pt

        # Con headline arriba la card puede necesitar subir un poco
        # para centrar en espacio restante (1080 - (420 + 245) = 415pt abajo)
        # Ajuste: centrar en zona disponible [360, 1080]
        available_start = 370   # después del headline (~100+300+margen)
        available_h     = SLIDE_H_PT - available_start   # ~710pt
        card_top_adj    = available_start + (available_h - CUALI_CARD_H) // 2  # ~607?
        # Pero el spec dice top 417pt para slide sin headline — con headline subir a ~430
        card_top = max(430, (SLIDE_H_PT + 370 - CUALI_CARD_H) // 2)

        verbatim_len = len(verbatims[0]['quote'])
        verb_pt = 28 if verbatim_len < 100 else 24

        add_card_with_text(
            slide,
            quote_text=verbatims[0]['quote'],
            attribution=verbatims[0]['attribution'],
            left_pt_val=card_left, top_pt_val=card_top,
            card_w_pt=CUALI_CARD_W, card_h_pt=CUALI_CARD_H,
            verbatim_pt=verb_pt, attr_pt=CUALI_ATTR_PT,
            pad_h=30, pad_v=30
        )

    elif n == 2:
        # 2 cards side by side, gap 80pt
        card_top = CUALI_TOP_WITH_HEADLINE   # 420pt

        for i, v in enumerate(verbatims):
            verbatim_len = len(v['quote'])
            verb_pt = 22 if verbatim_len > 120 else 26

            add_card_with_text(
                slide,
                quote_text=v['quote'],
                attribution=v['attribution'],
                left_pt_val=CUALI_2_LEFT[i], top_pt_val=card_top,
                card_w_pt=CUALI_CARD_W, card_h_pt=CUALI_CARD_H,
                verbatim_pt=verb_pt, attr_pt=CUALI_ATTR_PT,
                pad_h=30, pad_v=30
            )

    else:   # n == 3
        # 3 cards side by side, gap 50pt
        card_top = CUALI_TOP_WITH_HEADLINE

        for i, v in enumerate(verbatims):
            verbatim_len = len(v['quote'])
            verb_pt = 20 if verbatim_len > 100 else 22

            add_card_with_text(
                slide,
                quote_text=v['quote'],
                attribution=v['attribution'],
                left_pt_val=CUALI_3_LEFT[i], top_pt_val=card_top,
                card_w_pt=CUALI_CARD_W, card_h_pt=CUALI_CARD_H,
                verbatim_pt=verb_pt, attr_pt=CUALI_ATTR_PT,
                pad_h=25, pad_v=25
            )

    return slide


# ── Construcción del deck ──────────────────────────────────────────────────────

def build_deck():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — H01 — T1
    # Hallazgo cuanti — 2 stats — Layout 2
    # Headline 162 chars → 50pt fijo
    # 3.0% → 3%, 9.0% → 9%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="La terapia es la actividad de autocuidado menos elegida. El dominicano no ignora su salud mental — ",
        headline_italic="la terceriza a la fe, la mascota, el ejercicio y la familia.",
        stats=[
            {
                'value': '3%',
                'desc_runs': [
                    make_run('declara "voy a '),
                    make_run('terapia', bold=True),
                    make_run('" como actividad principal de autocuidado. La opción menos votada de toda la batería P27.')
                ]
            },
            {
                'value': '9%',
                'desc_runs': [
                    make_run('elige '),
                    make_run('meditación', bold=True),
                    make_run(' como autocuidado. Terapia + meditación: 12%. El 88% restante canaliza su bienestar por vías no clínicas.')
                ]
            }
        ]
    )

    # SLIDE 2 — H01 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo no iba a hacer lo que fue medicarme, pasar por psiquiatría. (...) Después de que pasó todo lo que pasó con la otra pareja que yo tuve, tuve que ir a la psiquiatría, tuve que medicarme, pero al final valió la pena, rindió sus frutos.',
        attribution='Familia Homoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — H02 — T1
    # Hallazgo cuanti — 1 stat — Layout 1
    # 85.0% → 85%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="85% del dominicano considera a la mascota parte de la familia. ",
        headline_italic="En los hombres, el perro cumple la función emocional que la cultura no les permite buscar en otra parte.",
        stats=[
            {
                'value': '85%',
                'desc_runs': [
                    make_run('considera a las mascotas '),
                    make_run('parte de la familia', bold=True),
                    make_run('. En hogares sin hijos con mascota sube a 94.6%.')
                ]
            }
        ]
    )

    # SLIDE 4 — H02 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Estadísticamente hablando, el perro suele ser un apoyo muy emocional para el hombre, en muchas ocasiones hasta más que la mujer, y no es por un tema de que no sentamos ese amor, sino porque tenemos un aspecto cultural y hasta biológico (...) de que no nos desahogamos como las mujeres lo hacen (...) nosotros simplemente necesitamos compañía, tacto y simplemente presencia.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5 — H03 — T2
    # Hallazgo cuanti — 2 stats — Layout 2
    # 64.0% → 64%, ratio 1:6 se mantiene como dato
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="64% recurre siempre a la religión en momentos de dificultad. ",
        headline_italic="La fe es el primer canal de salud mental: gratis, disponible las 24 horas, sin estigma ni copago.",
        stats=[
            {
                'value': '64%',
                'desc_runs': [
                    make_run('marca "5 = Siempre" recurrir a la '),
                    make_run('religión', bold=True),
                    make_run(' o espiritualidad en momentos de dificultad. Respuesta dominante en todas las tipologías.')
                ]
            },
            {
                'value': '1:6',
                'desc_runs': [
                    make_run('menciones clínicas por cada 6 '),
                    make_run('religiosas', bold=True),
                    make_run(' en 11 grupos (216 menciones fe vs 37 de terapia/psicólogo). Corpus completo de FGs.')
                ]
            }
        ]
    )

    # SLIDE 6 — H03 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='En mi caso nosotros con la biblia (...) si eso me ayudó bastante, bastante, la cercanía con Dios, a nuestra salud mental, de ambos.',
        attribution='Familia Sin Hijos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7 — H04 — T2 — SOLO-CUANTI (1 slide, sin Consumer Voice)
    # Hallazgo cuanti — 2 stats — Layout 2
    # 70.8% → 71%, 72.7% → 73%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="70.8% de las mujeres y 72.7% del estrato D recurren siempre a la fe. ",
        headline_italic="La intensidad religiosa sube donde más aprieta: si eres mujer y el bolsillo está flaco, rezas más.",
        stats=[
            {
                'value': '71%',
                'desc_runs': [
                    make_run('de las '),
                    make_run('mujeres', bold=True),
                    make_run(' marca "Siempre" recurrir a la religión en momentos difíciles, vs 53.3% de los hombres. Brecha de género: 17.5 puntos.')
                ]
            },
            {
                'value': '73%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' marca "Siempre" recurrir a la fe, vs 52.8% del estrato C+ y 61.5% del estrato C.')
                ]
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8 — H05 — T3
    # Hallazgo cuanti — 2 stats — Layout 2
    # 42.6% → 43%, 45.9% → 46%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="42.6% del dominicano no hace ninguna actividad para cuidar su salud, ni física ni mental. ",
        headline_italic="La inacción es la respuesta más común de toda la batería de autocuidado.",
        stats=[
            {
                'value': '43%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('no realizo ninguna actividad', bold=True),
                    make_run('" para cuidar su salud. TOP-1 de la batería P27, por encima de ejercicio (34.8%), oración (27.6%).')
                ]
            },
            {
                'value': '46%',
                'desc_runs': [
                    make_run('de hogares '),
                    make_run('sin hijos con mascota', bold=True),
                    make_run(' lidera la inacción (45.9%). La inacción encabeza 5 de 7 tipologías: biparental (44.2%), extendido (42.9%), monoparental (41.8%).')
                ]
            }
        ]
    )

    # SLIDE 9 — H05 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Bueno yo para lo físico yo no estoy haciendo nada. Yo no hago ejercicio, yo me como todo lo que yo quiera y a la hora que yo quiera. Bueno, para la salud mental a veces uno tiene que hacerse loco en ciertas cosas.',
        attribution='Familia Biparental con Hijos Adultos'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10 — H06 — T3
    # Hallazgo cuanti — 2 stats — Layout 2
    # 34.8% → 35%, 43.1% → 43%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="34.8% del dominicano se ejercita con regularidad — el único autocuidado donde el hombre lidera. ",
        headline_italic="Para él, entrenar no es físico: es la válvula emocional que la cultura sí le permite.",
        stats=[
            {
                'value': '35%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('me ejercito con regularidad', bold=True),
                    make_run('" — segundo lugar de la batería P27, detrás de "ninguna actividad" (42.6%).')
                ]
            },
            {
                'value': '43%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('masculino', bold=True),
                    make_run(' tiene el ejercicio como TOP-1. En el femenino el TOP-1 es "ninguna actividad" (48.5%). Único canal donde el hombre supera a la mujer.')
                ]
            }
        ]
    )

    # SLIDE 11 — H06 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Yo soy una persona súper activa, súper activa (...) La salud mental mía es deporte. Yo me irrito cuando yo duro tres días que no puedo hacer ejercicio (...) Yo entro a la cancha y yo me olvido de todo lo que yo tengo.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 12 — H07 — T4
    # Hallazgo cuanti — 3 stats — Layout 3
    # 47.8% → 48%, 28.4% → 28%, 51.2% → 51%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="47.8% nombra la economía como su mayor fuente de estrés familiar. ",
        headline_italic="Sumando los dos ítems económicos de la batería, el estresor número uno del hogar dominicano alcanza al 76% de los hogares.",
        stats=[
            {
                'value': '48%',
                'desc_runs': [
                    make_run('declara "'),
                    make_run('Economía', bold=True),
                    make_run('" como el factor que más estrés genera en su vida familiar. TOP-1 nacional de la batería P26.')
                ]
            },
            {
                'value': '28%',
                'desc_runs': [
                    make_run('marca también "'),
                    make_run('Realidad económica', bold=True),
                    make_run('" en la misma batería. Los dos ítems económicos combinados: 76% del total.')
                ]
            },
            {
                'value': '51%',
                'desc_runs': [
                    make_run('del '),
                    make_run('estrato D', bold=True),
                    make_run(' declara la economía como estresor principal, vs el promedio nacional de 47.8%.')
                ]
            }
        ]
    )

    # SLIDE 13 — H07 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Tuve un tiempo como mamá soltera y mi hermana es mamá soltera de tres niños, y es bien difícil levantarte, llevar al niño a la escuela, el dinero, cuando el papá no está económicamente ni emocionalmente.',
        attribution='Familia Biparental con Hijos Pequeños'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 14 — H08 — T5
    # Hallazgo cuanti — 3 stats — Layout 3
    # 42.2% → 42%, 22.3% → 22%, 5.5% mantiene decimal (1 dígito antes)
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="42.2% dice que la presión sobre los padres es demasiada. ",
        headline_italic="El agotamiento no aparece de golpe — en la franja 35–44 ya lo declara como estresor el 22.3%, el triple que a los 18–24.",
        stats=[
            {
                'value': '42%',
                'desc_runs': [
                    make_run('declara "5 — Muy de acuerdo" con "'),
                    make_run('hay demasiada presión', bold=True),
                    make_run(' para que los papás lo tengan todo." Ítem con mayor consenso de la batería P8 sobre crianza.')
                ]
            },
            {
                'value': '22%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('35–44 años', bold=True),
                    make_run(' declara la crianza como factor de estrés, vs 15.5% del subset 25–34. Pico en la franja media adulta.')
                ]
            },
            {
                'value': '5.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('18–24 años', bold=True),
                    make_run(' declara la crianza como estresor. La presión parental se construye con los años.')
                ]
            }
        ]
    )

    # SLIDE 15 — H08 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Es un mito que una mujer va a estar (...) toda ella relajada, limpiando una casa, con el marido sentado allá cargando el niño y qué feliz de la vida. Eso es un mito. De toda mujer que tiene hijos siempre agotada, cansada, batida.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 16 — H09 — T6
    # Hallazgo cuanti — 2 stats — Layout 2
    # 24.4% → 24%, 56.6% → 57%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="24.4% de los hogares es monoparental. En más de la mitad, la madre es la única responsable de la crianza. ",
        headline_italic="El mercado diseña para dos. El hogar dominicano no siempre llega de a dos.",
        stats=[
            {
                'value': '24%',
                'desc_runs': [
                    make_run('del dominicano declara que su hogar es '),
                    make_run('monoparental', bold=True),
                    make_run('. Segunda tipología más común después del biparental con hijos menores de 18 (38.0%).')
                ]
            },
            {
                'value': '57%',
                'desc_runs': [
                    make_run('del subset monoparental cita a la '),
                    make_run('madre', bold=True),
                    make_run(' como única responsable de la crianza, vs 41.8% del promedio nacional. Crianza monoparental de facto materna.')
                ]
            }
        ]
    )

    # SLIDE 17 — H09 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Soy madre soltera, sola, con una situación actualmente de mi madre que falleció y que vive conmigo y mi hermanito en condición especial, entonces es difícil, es difícil todo en general, quién lo cuide, porque yo trabajo todo el tiempo.',
        attribution='Familia Monoparental'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 18 — H10 — T7 — SOLO-CUALI
    # 2 cards SIDE BY SIDE (spec § 3.9: Layout 5 con 2 verbatims)
    # Headline 185 chars → 50pt fijo
    # ─────────────────────────────────────────────────────────────────────────
    build_cuali_slide(
        prs,
        headline_plain="El sistema público de salud falla antes de la consulta. ",
        headline_italic="La pregunta dominicana no es qué le pasa al médico — es si el seguro aprueba, si hay cama, si alguien dentro te puede ayudar.",
        verbatims=[
            {
                'quote': 'Vaya con un dolor, aunque sea un dolor de una hebra de cabello, y que te digan "hasta que el seguro no te apruebe, no." O tú ir a comprar un medicamento y que te digan: "ah no, tu límite del medicamento ya se agotó." Exactamente. Tú vas a una clínica y lo primero que te preguntan es qué seguro tú tienes.',
                'attribution': 'Familia Homoparental'
            },
            {
                'quote': 'Si tú no tienes alguien que te puede ayudar en el sistema, tú estás muerto.',
                'attribution': 'Familia Sin Hijos con Mascota'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 19 — H11 — T8
    # Hallazgo cuanti — 3 stats — Layout 3
    # 73.8% → 74%, 24.0% → 24%, 12.8% → 13%
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="73.8% pide pensiones para que los mayores vivan con dignidad. Solo 12.8% es cuidador directo. ",
        headline_italic="Los dominicanos saben lo que hace falta — solo que no están en posición de darlo.",
        stats=[
            {
                'value': '74%',
                'desc_runs': [
                    make_run('marca "'),
                    make_run('Pensiones o apoyos económicos suficientes', bold=True),
                    make_run('" como lo que más falta para que la tercera edad viva con dignidad. TOP-1 absoluto de P79.')
                ]
            },
            {
                'value': '24%',
                'desc_runs': [
                    make_run('menciona "'),
                    make_run('Apoyo y capacitación para cuidadores familiares', bold=True),
                    make_run('" en la misma batería. El cuidador familiar pide ayuda — el mercado no la ofrece como categoría.')
                ]
            },
            {
                'value': '13%',
                'desc_runs': [
                    make_run('declara ser el '),
                    make_run('cuidador principal', bold=True),
                    make_run(' de un adulto mayor. 72.8% declara que no vive con ninguna persona mayor en su hogar.')
                ]
            }
        ]
    )

    # SLIDE 20 — H11 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='La frustración en cuanto a salud y bienestar para mí es el tema del presupuesto que requiere. (...) Porque tengo un choque entre lo que yo digo que tengo que aceptar y las costumbres pasadas. Mi mamá era de la mujer que me enseñó a mí a que si algo no es necesario, no gastes dinero en eso, economiza.',
        attribution='Familia Sin Hijos con Mascota'
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 21 — H12 — T9 — SOLO-CUALI
    # 1 card centrada (spec § 3.9: Layout 5 con 1 verbatim)
    # Headline 189 chars → 50pt fijo
    # ─────────────────────────────────────────────────────────────────────────
    build_cuali_slide(
        prs,
        headline_plain='"Yo no me voy a separar, yo voy a aguantar — porque mi mamá me decía que hay que aguantar." ',
        headline_italic="El aguante dominicano no es una decisión. Es una instrucción que llega de la generación anterior.",
        verbatims=[
            {
                'quote': 'Yo me crié con esos dos conceptos en mi mente y yo siempre me enfocaba. Tanto así que yo decía: "No, yo no me voy a separar, yo voy a aguantar", porque mi mamá me decía que hay que aguantar.',
                'attribution': 'Familia Monoparental'
            }
        ]
    )

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 22 — H13 — T10
    # Hallazgo cuanti — 1 stat — Layout 1
    # 48.5% mantiene decimal (1 dígito antes del punto)
    # ─────────────────────────────────────────────────────────────────────────
    build_hallazgo_slide(
        prs,
        headline_plain="48.5% de las mujeres no hace ninguna actividad para cuidarse, vs 33.3% de los hombres. ",
        headline_italic="Brecha de género en autocuidado: 15.2 puntos. La mujer se cuida menos — y lo sabe.",
        stats=[
            {
                'value': '48.5%',
                'desc_runs': [
                    make_run('del subset '),
                    make_run('femenino', bold=True),
                    make_run(' declara "no realizo ninguna actividad para cuidar mi salud" como TOP-1 de P27, vs 33.3% del subset masculino. Brecha de género en inacción: 15.2 puntos.')
                ]
            }
        ]
    )

    # SLIDE 23 — H13 Consumer Voice
    build_consumer_voice_slide(
        prs,
        quote='Cuando uno es madre, a veces se olvida, a veces no... casi siempre se olvida de uno mismo. (...) Y cuando tú vienes a mirar para atrás, entonces tú dices: "¿y yo?" Porque crecen y se van.',
        attribution='Familia Monoparental'
    )

    # ── Guardar ────────────────────────────────────────────────────────────────
    output_path = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Agentes Hallazgos/bienestar-deck-flat-v5.pptx"
    prs.save(output_path)

    # ── Reporte de specs v5 ────────────────────────────────────────────────────
    print(f"GUARDADO: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print()
    print("── CHECKLIST SPECS V5 ──────────────────────────────────────────────────")
    print(f"Slide size: {prs.slide_width} × {prs.slide_height} EMU")
    print(f"  = {prs.slide_width / PT_EMU:.0f}pt × {prs.slide_height / PT_EMU:.0f}pt")
    ok_w = prs.slide_width == pt(1920)
    ok_h = prs.slide_height == pt(1080)
    print(f"  Slide 1920×1080pt: {'OK' if (ok_w and ok_h) else 'ERROR'}")
    print()
    print(f"Headlines: {HEADLINE_PT}pt fijo (NO auto-size)")
    print(f"Cifras grandes: {STAT_NUMBER_PT}pt Instrument Serif italic")
    print(f"Cajas descripción: {STAT_DESC_W}×{STAT_DESC_H}pt Poppins {STAT_DESC_PT}pt")
    print()
    print("Respiración entre stats:")
    print(f"  2 stats: gap {GAP_2_STATS}pt (cajas desc: [{STATS2_DESC_X[0]}, {STATS2_DESC_X[1]}]pt)")
    print(f"  3 stats: gap {GAP_3_STATS}pt (cajas desc: [{STATS3_DESC_X[0]}, {STATS3_DESC_X[1]}, {STATS3_DESC_X[2]}]pt)")
    print()
    print(f"Consumer Voice card: {CV_CARD_W}×{CV_CARD_H}pt (left {CV_CARD_LEFT}pt, top {CV_CARD_TOP}pt)")
    print(f"  Verbatim: {CV_VERBATIM_PT}pt Instrument Serif regular")
    print(f"  Atribución: {CV_ATTR_PT}pt Poppins italic")
    print(f"  Header: {CV_HEADER_PT}pt Poppins gris #9B9B9B top {CV_HEADER_TOP}pt")
    print()
    print("Cards cualitativas:")
    print(f"  Dimensiones: {CUALI_CARD_W}×{CUALI_CARD_H}pt")
    print(f"  1 card: left {CUALI_1_LEFT}pt (centrada)")
    print(f"  2 cards side by side: gap {CUALI_2_GAP}pt, left {CUALI_2_LEFT}pt")
    print(f"  3 cards side by side: gap {CUALI_3_GAP}pt, left {CUALI_3_LEFT}pt")
    print()
    print("Fill cards: rgba(0,0,0,0.45) vía ROUNDED_RECTANGLE + XML alpha=45000")
    print("Border-radius: adjustments[0]=0.05 (~15pt)")
    print("Kerning: 0 en todo el deck")
    print("Source: NO en ningún slide (eliminado)")
    print("Cajas P##/Base: NO en ningún slide (eliminado)")
    print("Masterslide: NO (fondo negro plano — Jeremy aplica masterslide después)")
    print()
    print("Slide H10 (solo-cuali): 2 cards SIDE BY SIDE ✓")
    print("Slide H12 (solo-cuali): 1 card centrada ✓")
    print("Slide H04 (solo-cuanti): 1 slide, sin Consumer Voice ✓")

    return output_path


if __name__ == "__main__":
    build_deck()
