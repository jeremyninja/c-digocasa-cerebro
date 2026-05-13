#!/usr/bin/env python3
"""
generate_pdf.py — Genera PDF de QA del deck usando reportlab.
Crea una página por slide con el texto de cada uno para revisión.
"""
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from pptx import Presentation
from pptx.util import Pt
import re

PPTX_PATH = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v2-2026-05-07/bienestar-humo-v2-2026-05-07.pptx"
PDF_PATH  = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v2-2026-05-07/bienestar-humo-v2-2026-05-07-QA.pdf"

PAGE_W, PAGE_H = landscape(A4)
MARGIN = 1.8 * cm

# Colores
C_BLACK  = colors.HexColor('#000000')
C_WHITE  = colors.HexColor('#FFFFFF')
C_GRAY   = colors.HexColor('#2E2E2E')
C_LIME   = colors.HexColor('#B8FF4D')
C_SUBTEXT = colors.HexColor('#AAAAAA')
C_BG     = colors.HexColor('#111111')


def extract_slide_content(slide):
    """Extrae texto de cada shape, agrupado."""
    shapes_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        lines = []
        for para in shape.text_frame.paragraphs:
            line = "".join(r.text for r in para.runs).strip()
            if line:
                lines.append(line)
        if lines:
            shapes_text.append("\n".join(lines))
    return shapes_text


def build_pdf():
    prs = Presentation(PPTX_PATH)

    doc = SimpleDocTemplate(
        PDF_PATH,
        pagesize=landscape(A4),
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN,
    )

    styles = getSampleStyleSheet()

    # Estilos custom
    s_title = ParagraphStyle(
        'SlideTitle',
        fontName='Helvetica-Bold',
        fontSize=11,
        textColor=C_LIME,
        alignment=TA_LEFT,
        spaceAfter=4,
    )
    s_headline = ParagraphStyle(
        'Headline',
        fontName='Helvetica-Bold',
        fontSize=13,
        textColor=C_WHITE,
        alignment=TA_CENTER,
        spaceAfter=8,
        spaceBefore=4,
    )
    s_body = ParagraphStyle(
        'Body',
        fontName='Helvetica',
        fontSize=8.5,
        textColor=C_WHITE,
        alignment=TA_LEFT,
        spaceAfter=4,
        leading=12,
    )
    s_insight = ParagraphStyle(
        'Insight',
        fontName='Helvetica-Oblique',
        fontSize=9,
        textColor=C_WHITE,
        alignment=TA_CENTER,
        spaceAfter=4,
        spaceBefore=4,
    )
    s_source = ParagraphStyle(
        'Source',
        fontName='Helvetica-Oblique',
        fontSize=7.5,
        textColor=C_SUBTEXT,
        alignment=TA_LEFT,
        spaceBefore=4,
    )
    s_verbatim = ParagraphStyle(
        'Verbatim',
        fontName='Helvetica-Oblique',
        fontSize=8.5,
        textColor=C_WHITE,
        alignment=TA_LEFT,
        spaceAfter=2,
    )
    s_tag = ParagraphStyle(
        'Tag',
        fontName='Helvetica-Bold',
        fontSize=8,
        textColor=C_SUBTEXT,
        alignment=TA_RIGHT,
    )

    # Slides data (parallel a build_deck.py)
    SLIDE_META = [
        {"tag": "H01", "label": "Estrés Económico — Factor #1"},
        {"tag": "H02", "label": "Peso Relativo Bloque Financiero"},
        {"tag": "H03", "label": "Autocuidado — La inacción domina"},
        {"tag": "H04", "label": "Autocuidado con género y NSE"},
        {"tag": "H05", "label": "Sistema de Salud Pública — Convergencia cualitativa"},
    ]

    story = []

    for i, (slide, meta) in enumerate(zip(prs.slides, SLIDE_META), 1):
        # Fondo negro simulado con tabla
        shapes_text = extract_slide_content(slide)

        # Header del slide en QA
        story.append(Paragraph(
            f"SLIDE {i} · {meta['tag']} · {meta['label']}",
            s_title,
        ))
        story.append(HRFlowable(width="100%", thickness=0.5, color=C_LIME))
        story.append(Spacer(1, 6))

        # Texto completo del slide
        for j, block in enumerate(shapes_text):
            lines = block.split('\n')
            for line in lines:
                if not line.strip():
                    continue
                # Clasificar el bloque
                if j == 0:  # Tag H0X
                    story.append(Paragraph(line, s_tag))
                elif j == 1:  # Headline (parte plain)
                    story.append(Paragraph(f"<b>{line}</b>", s_headline))
                elif j == 2:  # Headline (parte italic)
                    story.append(Paragraph(f"<i>{line}</i>", s_headline))
                elif line.startswith("Source:"):
                    story.append(Paragraph(line, s_source))
                elif line.startswith('"') or line.startswith('“'):
                    story.append(Paragraph(line, s_verbatim))
                elif any(kw in line for kw in ["El estrés", "En la escala", "Antes que", "La brecha", "El diagnóstico"]):
                    story.append(Paragraph(f"<i>{line}</i>", s_insight))
                elif line.startswith("CONVERGENCIA"):
                    story.append(Paragraph(f"<b><i>{line}</i></b>", s_source))
                else:
                    story.append(Paragraph(line, s_body))

        story.append(Spacer(1, 12))
        story.append(HRFlowable(width="100%", thickness=0.3, color=C_GRAY))
        story.append(Spacer(1, 16))

        # Page break entre slides excepto el último
        if i < len(SLIDE_META):
            from reportlab.platypus import PageBreak
            story.append(PageBreak())

    # Build con fondo negro
    def on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(C_BG)
        canvas.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)
        canvas.restoreState()

    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    print(f"PDF QA generado: {PDF_PATH}")


if __name__ == "__main__":
    build_pdf()
