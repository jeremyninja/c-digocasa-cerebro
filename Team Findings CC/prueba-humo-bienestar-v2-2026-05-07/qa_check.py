#!/usr/bin/env python3
"""
qa_check.py — QA de texto del deck generado
Verifica: número de slides, presencia de textos clave, colores.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor

DECK = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v2-2026-05-07/bienestar-humo-v2-2026-05-07.pptx"

GREEN_LIME = "B8FF4D"

def check_deck():
    prs = Presentation(DECK)
    slides = list(prs.slides)
    print(f"\n{'='*60}")
    print(f"QA DECK — Prueba de Humo Bienestar · {len(slides)} slides")
    print(f"{'='*60}\n")

    if len(slides) != 5:
        print(f"  ERROR: se esperan 5 slides, hay {len(slides)}")
    else:
        print(f"  OK: 5 slides exactos\n")

    lime_issues = []

    for i, slide in enumerate(slides, 1):
        all_text = []
        print(f"  SLIDE {i} ─────────────────────────────────")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text.strip()
                    if txt:
                        all_text.append(txt)
                    # Check verde lima en texto
                    try:
                        clr = run.font.color.rgb
                        if str(clr).upper() == GREEN_LIME:
                            lime_issues.append(f"Slide {i}: texto verde lima encontrado: '{txt[:40]}'")
                    except Exception:
                        pass
                    # Check kerning
                    rPr = run._r.get_or_add_rPr()
                    spc = rPr.get('spc')
                    if spc and spc != '0':
                        print(f"    WARN kerning spc={spc} en: '{txt[:30]}'")

        # Mostrar primeras líneas de texto del slide
        preview = " | ".join(all_text[:6])
        print(f"    Texto: {preview[:120]}")
        # Verificar que no hay 'placeholder' o texto de relleno
        full = " ".join(all_text).lower()
        if "placeholder" in full or "aquí va" in full or "[pilar]" in full:
            print(f"    ERROR: placeholder sin llenar detectado")
        else:
            print(f"    OK: sin placeholders")
        print()

    if lime_issues:
        print(f"  ERROR verde lima en texto:")
        for issue in lime_issues:
            print(f"    {issue}")
    else:
        print(f"  OK: sin verde lima en texto")

    # Verificar sources
    print(f"\n  SOURCES detectados:")
    for i, slide in enumerate(slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs)
                    if line.strip().startswith("Source:"):
                        print(f"    Slide {i}: {line.strip()}")

    print(f"\n{'='*60}")
    print("QA completado.")
    print(f"{'='*60}\n")

if __name__ == "__main__":
    check_deck()
