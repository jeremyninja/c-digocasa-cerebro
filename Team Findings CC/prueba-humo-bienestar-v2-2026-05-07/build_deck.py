#!/usr/bin/env python3
"""
build_deck.py — Prueba de Humo Bienestar · Código Casa MED
============================================================
Genera 5 slides flat (uno por hallazgo) con python-pptx.
Cada slide sigue el layout:
  - Headline: Instrument Serif MAYÚSCULAS, auto-tamaño, italic en el giro
  - Caja cuanti / convergencia cualitativa
  - Caja de insight (italic)
  - Caja de verbatim (blanco, nunca verde)
  - Source al pie

Reglas aplicadas:
  - Kerning 0 en todo el deck
  - Verbatims en blanco #FFFFFF, nunca verde lima
  - Headlines en MAYÚSCULAS centrados
  - Auto-tamaño según largo (tabla del knowledge pack)
  - Decimales en descripciones inline: se mantienen
  - Stat grande: redondeado a entero
  - Source: formato estándar Código Casa
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────
# PALETA
# ─────────────────────────────────────────────
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x2E, 0x2E, 0x2E)
LIME      = RGBColor(0xB8, 0xFF, 0x4D)   # SOLO cápsula, NUNCA texto
BLUE_DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_MED  = RGBColor(0x66, 0x66, 0x66)

# ─────────────────────────────────────────────
# DIMENSIONES (widescreen 13.33" × 7.5")
# ─────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.7)
MARGIN_Y = Inches(0.55)
CONTENT_W = SLIDE_W - 2 * MARGIN_X

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def auto_headline_size(text: str) -> int:
    """Retorna tamaño en puntos según largo del headline (tabla knowledge pack)."""
    n = len(text.strip())
    if n < 30:
        return 70
    elif n < 45:
        return 55
    elif n < 66:
        return 42
    elif n < 86:
        return 36
    else:
        return 30


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    return txBox


def set_kerning_zero(run):
    """Kerning a 0: sin spacing entre letras."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', '0')


def set_font_base(run, font_name, font_size, bold=False, italic=False, color=WHITE):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    set_kerning_zero(run)


def add_rect(slide, left, top, width, height, fill_color, corner_radius=None):
    """Agrega un rectángulo con color de fondo."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # sin borde
    return shape


def add_headline(slide, plain_text: str, italic_text: str, top_y: float):
    """
    Agrega headline split plain + italic en Instrument Serif MAYÚSCULAS.
    plain_text: parte en regular
    italic_text: parte en italic (el giro)
    Ambas van en la misma caja, centradas.
    """
    full_text = (plain_text + italic_text).strip()
    size = auto_headline_size(full_text)
    # Altura dinámica según tamaño
    if size >= 55:
        box_h = Inches(1.2)
    elif size >= 42:
        box_h = Inches(1.5)
    elif size >= 36:
        box_h = Inches(1.7)
    else:
        box_h = Inches(1.9)

    txBox = add_textbox(slide, MARGIN_X, top_y, CONTENT_W, box_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    if plain_text:
        run1 = p.add_run()
        run1.text = plain_text.upper()
        set_font_base(run1, "Instrument Serif", size, italic=False, color=WHITE)

    if italic_text:
        run2 = p.add_run()
        run2.text = italic_text.upper()
        set_font_base(run2, "Instrument Serif", size, italic=True, color=WHITE)

    return box_h


def add_divider(slide, y):
    """Línea horizontal separadora en gris oscuro."""
    line = slide.shapes.add_shape(1, MARGIN_X, y, CONTENT_W, Pt(0.75))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY_DARK
    line.line.fill.background()


def build_stat_block(slide, stats: list, y_start: float, box_h: float):
    """
    Construye el bloque de stats cuanti.
    stats: lista de dicts con keys: 'big', 'desc'
    'big': texto del stat grande (ya redondeado si aplica)
    'desc': descripción con posible bold inline (lista de tuplas (texto, bold))
    """
    col_w = CONTENT_W / len(stats)
    sep_width = Pt(0.75)

    for i, stat in enumerate(stats):
        col_x = MARGIN_X + col_w * i

        # Separador vertical (entre stats)
        if i > 0:
            sep_x = MARGIN_X + col_w * i - Pt(0.5)
            sep = slide.shapes.add_shape(1, sep_x, y_start, sep_width, box_h)
            sep.fill.solid()
            sep.fill.fore_color.rgb = GRAY_DARK
            sep.line.fill.background()

        # Stat grande (96pt)
        big_box = add_textbox(slide, col_x + Inches(0.1), y_start, col_w - Inches(0.2), Inches(1.2))
        big_tf = big_box.text_frame
        big_p = big_tf.paragraphs[0]
        big_p.alignment = PP_ALIGN.CENTER
        big_run = big_p.add_run()
        big_run.text = stat['big']
        set_font_base(big_run, "Instrument Serif", 72, italic=True, color=WHITE)
        # 96pt a veces desborda, usamos 72pt que es más seguro en columnas

        # Descripción (Poppins 11pt, partes con bold)
        desc_y = y_start + Inches(1.15)
        desc_box = add_textbox(slide, col_x + Inches(0.1), desc_y, col_w - Inches(0.2), Inches(1.4))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_p = desc_tf.paragraphs[0]
        desc_p.alignment = PP_ALIGN.CENTER

        for (txt, is_bold) in stat['desc']:
            r = desc_p.add_run()
            r.text = txt
            set_font_base(r, "Poppins", 10, bold=is_bold, italic=False, color=WHITE)


def build_insight_box(slide, text: str, y: float):
    """Conclusión italic centrada — Poppins italic 12pt."""
    box_h = Inches(0.65)
    txBox = add_textbox(slide, MARGIN_X, y, CONTENT_W, box_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font_base(run, "Poppins", 11, italic=True, color=WHITE)


def build_verbatim_box(slide, quote: str, attribution: str, y: float):
    """
    Caja de verbatim: fondo gris oscuro, comillas, atribución.
    Color del texto: BLANCO siempre.
    """
    box_h = Inches(1.2)
    rect = add_rect(slide, MARGIN_X, y, CONTENT_W, box_h, GRAY_DARK)

    # Texto del quote
    txBox = add_textbox(slide, MARGIN_X + Inches(0.25), y + Inches(0.15),
                        CONTENT_W - Inches(0.5), Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"“{quote}”"
    set_font_base(run, "Poppins", 9.5, italic=False, color=WHITE)

    # Atribución
    attr_box = add_textbox(slide, MARGIN_X + Inches(0.25), y + Inches(0.85),
                           CONTENT_W - Inches(0.5), Inches(0.25))
    attr_tf = attr_box.text_frame
    attr_p = attr_tf.paragraphs[0]
    attr_p.alignment = PP_ALIGN.RIGHT
    attr_run = attr_p.add_run()
    attr_run.text = f"— {attribution}"
    set_font_base(attr_run, "Poppins", 8.5, italic=True, color=WHITE)


def build_convergencia_box(slide, profiles: list, y: float):
    """
    H05: bloque convergencia cualitativa con 3 perfiles en lugar de caja cuanti.
    profiles: lista de dicts {quote, attribution}
    """
    # Label de caveat
    caveat_box = add_textbox(slide, MARGIN_X, y, CONTENT_W, Inches(0.3))
    caveat_tf = caveat_box.text_frame
    caveat_p = caveat_tf.paragraphs[0]
    caveat_p.alignment = PP_ALIGN.CENTER
    caveat_run = caveat_p.add_run()
    caveat_run.text = "CONVERGENCIA CUALITATIVA — Datos cuantitativos pendientes de re-tabulación (P28)"
    set_font_base(caveat_run, "Poppins", 8, italic=True, color=RGBColor(0xAA, 0xAA, 0xAA))

    col_w = CONTENT_W / 3
    for i, prof in enumerate(profiles):
        col_x = MARGIN_X + col_w * i
        box_y = y + Inches(0.35)
        box_h = Inches(1.55)

        rect = add_rect(slide, col_x + Inches(0.05), box_y,
                        col_w - Inches(0.1), box_h, GRAY_DARK)

        txBox = add_textbox(slide, col_x + Inches(0.15), box_y + Inches(0.12),
                            col_w - Inches(0.3), Inches(1.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"“{prof['quote']}”"
        set_font_base(run, "Poppins", 8.5, italic=False, color=WHITE)

        attr_box = add_textbox(slide, col_x + Inches(0.15), box_y + Inches(1.2),
                               col_w - Inches(0.3), Inches(0.25))
        attr_tf = attr_box.text_frame
        attr_p = attr_tf.paragraphs[0]
        attr_p.alignment = PP_ALIGN.RIGHT
        attr_run = attr_p.add_run()
        attr_run.text = f"— {prof['attribution']}"
        set_font_base(attr_run, "Poppins", 8, italic=True, color=WHITE)


def build_source_line(slide, source_text: str):
    """Source al pie de página — Poppins italic 9pt."""
    src_y = SLIDE_H - Inches(0.38)
    txBox = add_textbox(slide, MARGIN_X, src_y, CONTENT_W, Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = source_text
    set_font_base(run, "Poppins", 9, italic=True, color=RGBColor(0x99, 0x99, 0x99))


def add_slide_number_tag(slide, tag: str):
    """Tag discreto en la esquina superior derecha (H01, H02, etc.)."""
    txBox = add_textbox(slide, SLIDE_W - Inches(1.1), Inches(0.2), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = tag
    set_font_base(run, "Poppins", 9, italic=False, color=RGBColor(0x55, 0x55, 0x55))


def new_slide(prs):
    """Agrega un slide en blanco (layout 6 = blank) con fondo negro."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    return slide


# ─────────────────────────────────────────────────────────────────
# DATOS DE LOS 5 HALLAZGOS
# ─────────────────────────────────────────────────────────────────
# DECISIONES EDITORIALES RESUELTAS:
#   H01: verbatim 1 (más concisa) — según recomendación del editor
#   H02: slide independiente (profundización distinta del headline)
#   H03: top 1 vs top 2 + datos de contraste (ejercicio 34.8%,
#         terapia 3.0%, meditación 9.0%) — legibilidad en 3 stats
#   H04: comparativa él/ella en stat 1, brecha NSE/edad en stat 2
#         (según knowledge pack: layout columnar cuanti estándar)
#   H05: caja cuanti → bloque convergencia cualitativa con 3 perfiles

SLIDES_DATA = [
    # ─── H01 ─────────────────────────────────────────────────────
    {
        "tag": "H01",
        "headline_plain": "Dicen que el dinero no lo es todo, pero ",
        "headline_italic": "47.8% nombra la economía como su principal fuente de estrés — doblando a la inseguridad y a la crianza.",
        "stats": [
            {
                "big": "48%",
                "desc": [
                    ("nombra ", False),
                    ('"Economía"', True),
                    (" como factor #1 de estrés en la vida familiar, por encima de todos los demás factores evaluados. · P26 · Base 500", False),
                ]
            },
            {
                "big": "28.4%",
                "desc": [
                    ('"Realidad económica"', True),
                    (" adicional — el bloque financiero casi ", False),
                    ("duplica al resto", True),
                    (" (Inseguridad 21.8%, Crianza 16.2%). · P26 · Base 500", False),
                ]
            },
        ],
        "insight": "El estrés del hogar dominicano tiene dirección: casi la mitad del país señala la economía antes que la violencia, la crianza o el tiempo.",
        "verbatim": {
            "quote": "Como lo dijo una compañera acá ahorita: el dinero. Tú no podés resolver tal vez algo que se te presente. Eso es frustra, eso no es fácil.",
            "attribution": "Grupo Monoparental",
        },
        "source": "Source: Código Casa — Estudio cuantitativo 2025 · P26 · Base 500.",
        "type": "cuanti",
    },

    # ─── H02 ─────────────────────────────────────────────────────
    {
        "tag": "H02",
        "headline_plain": "El bloque financiero pesa más que todos los demás factores de estrés juntos. ",
        "headline_italic": "76.2 puntos vs 74.8 repartidos entre inseguridad, tiempo, crianza y hogar.",
        "stats": [
            {
                "big": "76.2",
                "desc": [
                    ("puntos de peso relativo: Economía (", False),
                    ("47.8%", True),
                    (") + Realidad económica (28.4%). · P26 · Base 500", False),
                ]
            },
            {
                "big": "74.8",
                "desc": [
                    ("puntos distribuidos entre: Inseguridad ", False),
                    ("21.8%", True),
                    (", Falta de tiempo 20.6%, Crianza 16.2%, Responsabilidades del hogar 16.2%. · P26 · Base 500", False),
                ]
            },
        ],
        "insight": "En la escala de presiones del hogar dominicano, lo económico no compite con los demás factores — los supera en conjunto.",
        "verbatim": {
            "quote": "No tengo paciencia cuando no tengo dinero. Cuando no tengo dinero, que no sé de qué voy a hacer con los compromisos y todo eso, entonces yo como que pienso, le doy vuelta a la cosa, un estrés.",
            "attribution": "Grupo Monoparental",
        },
        "source": "Source: Código Casa — Estudio cuantitativo 2025 · P26 · Base 500.",
        "note": "Nota técnica: 76.2 y 74.8 son pesos relativos de dimensión, no porcentajes de personas.",
        "type": "cuanti",
    },

    # ─── H03 ─────────────────────────────────────────────────────
    {
        "tag": "H03",
        "headline_plain": "El autocuidado existe como aspiración. ",
        "headline_italic": "La primera respuesta del país es no hacer nada.",
        "stats": [
            {
                "big": "43%",
                "desc": [
                    ("responde ", False),
                    ('"No realizo ninguna actividad"', True),
                    (" para cuidar su salud mental y física — la opción más mencionada de todas. · P27 · Base 500", False),
                ]
            },
            {
                "big": "35%",
                "desc": [
                    ("hace ", False),
                    ("ejercicio regular", True),
                    (", la segunda respuesta. Terapia llega a 3.0% y meditación a 9.0% — el autocuidado activo es la excepción. · P27 · Base 500", False),
                ]
            },
        ],
        "insight": "Antes que el ejercicio o la terapia, la inacción es el hábito de salud más extendido en el hogar dominicano.",
        "verbatim": {
            "quote": "Estar un poquito más estable en todos los sentidos. He estado con mucho estrés. Sí. Para entonces poder ir a la gimnasia y comprar alimentos que debo de comer.",
            "attribution": "Grupo Monoparental",
        },
        "source": "Source: Código Casa — Estudio cuantitativo 2025 · P27 · Base 500.",
        "type": "cuanti",
    },

    # ─── H04 ─────────────────────────────────────────────────────
    {
        "tag": "H04",
        "headline_plain": "El autocuidado tiene género en RD: él se ejercita, ",
        "headline_italic": "ella no hace nada.",
        "stats": [
            {
                "big": "48.5%",
                "desc": [
                    ("de las ", False),
                    ("mujeres", True),
                    (' dice "No realizo ninguna actividad". Entre los hombres, la respuesta dominante es ejercicio regular (43.1%). · P27×D2 · Base 500', False),
                ]
            },
            {
                "big": "60%",
                "desc": [
                    ("de inacción en el grupo ", False),
                    ("18–24 años", True),
                    (". Estrato D: 51.2%. En NSE AB y C el ejercicio regular domina (39.1% y 38.5%). · P27×D7×D3 · Base 500", False),
                ]
            },
        ],
        "insight": "La brecha no es solo de género — el autocuidado activo es privilegio de estrato: a menor NSE, mayor inacción.",
        "verbatim": {
            "quote": "Mi esposo va al gimnasio, yo camino diario menos los domingos hago pilates, tenemos una… tratamos de tener una buena.",
            "attribution": "Grupo Biparental Hijos Pequeños",
        },
        "source": "Source: Código Casa — Estudio cuantitativo 2025 · P27, D2, D7, D3 · Base 500.",
        "type": "cuanti",
    },

    # ─── H05 ─────────────────────────────────────────────────────
    {
        "tag": "H05",
        "headline_plain": "El sistema de salud pública se critica por calidad, costo y acceso. ",
        "headline_italic": "Los tres problemas llegan juntos, no en orden.",
        "convergencia": [
            {
                "quote": "El sistema de salud pública no sirve. Es un caos. En mi caso no me ha tocado, gracias a Dios; siempre tuve un empleo que me ha tenido un pago seguro y puedo ir a cualquier clínica tranquilamente. Pero la persona que no tiene un seguro médico hace mucho trabajo, mucho trabajo.",
                "attribution": "Grupo Extendido",
            },
            {
                "quote": "De farmacia… si esto te cuesta 500 pesos y te cura, te ponen [medicamento] que vale 5 mil porque se alió a tu farmacia, a tu laboratorio.",
                "attribution": "Grupo Biparental Hijos Adultos",
            },
            {
                "quote": "Yo fui a visitar a alguien ahí en el Marcelo Rivas… eso parece un hotel y las atenciones son enormes, lo que pasa es que la demanda de personas es mucha.",
                "attribution": "Grupo Extendido",
            },
        ],
        "insight": "El diagnóstico del sistema de salud pública no aparece como lista ordenada — aparece como queja total, donde calidad, costo y acceso se refuerzan entre sí.",
        "source": "Source: Código Casa — Estudio cualitativo 2025 · P28 · Evidencia cualitativa convergente — cifra cuantitativa pendiente de re-tabulación.",
        "type": "cualitativo",
    },
]


# ─────────────────────────────────────────────────────────────────
# BUILD PRINCIPAL
# ─────────────────────────────────────────────────────────────────

def build_deck(output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for data in SLIDES_DATA:
        slide = new_slide(prs)

        # Tag (H01, etc.) en esquina
        add_slide_number_tag(slide, data["tag"])

        # ── Headline ──────────────────────────────────────────────
        headline_h = add_headline(
            slide,
            data["headline_plain"],
            data["headline_italic"],
            top_y=MARGIN_Y,
        )

        current_y = MARGIN_Y + headline_h + Inches(0.15)

        # ── Separador ─────────────────────────────────────────────
        add_divider(slide, current_y)
        current_y += Pt(4)

        # ── Stats o Convergencia cualitativa ──────────────────────
        if data["type"] == "cuanti":
            stat_block_h = Inches(2.6)
            build_stat_block(slide, data["stats"], current_y, stat_block_h)
            current_y += stat_block_h + Inches(0.1)

            # Nota técnica si existe
            if data.get("note"):
                note_box = add_textbox(slide, MARGIN_X, current_y, CONTENT_W, Inches(0.25))
                note_tf = note_box.text_frame
                note_p = note_tf.paragraphs[0]
                note_p.alignment = PP_ALIGN.LEFT
                note_run = note_p.add_run()
                note_run.text = data["note"]
                set_font_base(note_run, "Poppins", 8, italic=True, color=RGBColor(0x88, 0x88, 0x88))
                current_y += Inches(0.3)

        else:  # cualitativo — convergencia H05
            convergencia_h = Inches(2.0)
            build_convergencia_box(slide, data["convergencia"], current_y)
            current_y += convergencia_h + Inches(0.1)

        # ── Separador ─────────────────────────────────────────────
        add_divider(slide, current_y)
        current_y += Pt(6)

        # ── Insight ───────────────────────────────────────────────
        build_insight_box(slide, data["insight"], current_y)
        current_y += Inches(0.7)

        # ── Verbatim (H01–H04) ────────────────────────────────────
        if data.get("verbatim"):
            build_verbatim_box(
                slide,
                data["verbatim"]["quote"],
                data["verbatim"]["attribution"],
                current_y,
            )

        # ── Source al pie ─────────────────────────────────────────
        build_source_line(slide, data["source"])

    prs.save(output_path)
    print(f"PPTX guardado: {output_path}")


if __name__ == "__main__":
    OUT = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v2-2026-05-07/bienestar-humo-v2-2026-05-07.pptx"
    build_deck(OUT)
