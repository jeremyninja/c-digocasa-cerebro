"""
Prueba de humo - Bienestar, Codigo Casa MED
5 slides de hallazgo (+ portada) para validar voz visual.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

LDQ = "“"
RDQ = "”"

BASE_DIR_REAL = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-2026-05-07"
LOGO_PATH = "/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Logo/CC - White.png"
OUT_PPTX  = os.path.join(BASE_DIR_REAL, "bienestar-prueba-humo-2026-05-07.pptx")

BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2E, 0x2E, 0x2E)
LIME      = RGBColor(0xB8, 0xFF, 0x4D)
MID_GRAY  = RGBColor(0x88, 0x88, 0x88)
RED_FLAG  = RGBColor(0xFF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L  = Inches(0.55)
MARGIN_R  = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

HL_TOP = Inches(0.55)
HL_H   = Inches(1.85)

BOX_TOP = Inches(2.5)
BOX_H   = Inches(3.6)
BOX_PAD = Inches(0.18)

COL_GAP = Inches(0.18)
COL_W   = (CONTENT_W - 2 * COL_GAP) / 3

SOURCE_TOP = Inches(6.9)
SOURCE_H   = Inches(0.45)

def col_left(n):
    return MARGIN_L + n * (COL_W + COL_GAP)

def black_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide

def add_logo(slide, path, right_margin=Inches(0.3), top_margin=Inches(0.25), height=Inches(0.45)):
    if os.path.exists(path):
        slide.shapes.add_picture(path, SLIDE_W - height * 2 - right_margin, top_margin, height=height)

def horiz_line(slide, left, top, width, color=None):
    if color is None:
        color = DARK_GRAY
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(0.75)

def headline_font_size(n):
    if n < 30: return 52
    if n < 45: return 44
    if n < 66: return 36
    if n < 86: return 30
    return 26

def add_headline_box(slide, text_plain, text_italic, left, top, width, height):
    fsize = headline_font_size(len(text_plain) + len(text_italic))
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if text_plain:
        r1 = p.add_run()
        r1.text = text_plain
        r1.font.name = "Instrument Serif"
        r1.font.size = Pt(fsize)
        r1.font.italic = False
        r1.font.color.rgb = WHITE
    if text_italic:
        r2 = p.add_run()
        r2.text = text_italic
        r2.font.name = "Instrument Serif"
        r2.font.size = Pt(fsize)
        r2.font.italic = True
        r2.font.color.rgb = WHITE

def add_source_line(slide, source):
    tb = slide.shapes.add_textbox(MARGIN_L, SOURCE_TOP, CONTENT_W, SOURCE_H)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = source
    r.font.name = "Poppins"
    r.font.size = Pt(8)
    r.font.italic = True
    r.font.color.rgb = MID_GRAY

def add_textrun(p, text, font_name="Poppins", font_size=9.5,
                bold=False, italic=False, color=None):
    if color is None:
        color = WHITE
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color

def build_stats_block(slide, stats):
    stats_w = COL_W * 2 + COL_GAP
    stats_left = col_left(0)
    bg = slide.shapes.add_textbox(stats_left, BOX_TOP, stats_w, BOX_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GRAY
    n = len(stats)
    stat_block_h = BOX_H / n
    for i, stat in enumerate(stats):
        s_top = BOX_TOP + i * stat_block_h + BOX_PAD
        tb_big = slide.shapes.add_textbox(
            stats_left + BOX_PAD, s_top, stats_w - 2 * BOX_PAD, Inches(0.7))
        tf_big = tb_big.text_frame
        tf_big.word_wrap = False
        pb = tf_big.paragraphs[0]
        pb.alignment = PP_ALIGN.LEFT
        rb = pb.add_run()
        rb.text = stat["big"]
        rb.font.name = "Instrument Serif"
        rb.font.size = Pt(38)
        rb.font.italic = True
        rb.font.color.rgb = LIME
        if stat.get("unit"):
            ru = pb.add_run()
            ru.text = "  " + stat["unit"]
            ru.font.name = "Poppins"
            ru.font.size = Pt(10)
            ru.font.color.rgb = MID_GRAY
        desc_top = s_top + Inches(0.72)
        desc_h = stat_block_h - Inches(0.72) - BOX_PAD * 2
        tb_desc = slide.shapes.add_textbox(
            stats_left + BOX_PAD, desc_top, stats_w - 2 * BOX_PAD, desc_h)
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        pd = tf_desc.paragraphs[0]
        pd.alignment = PP_ALIGN.LEFT
        for rdef in stat.get("desc_runs", []):
            add_textrun(pd, rdef["text"],
                        bold=rdef.get("bold", False),
                        italic=rdef.get("italic", False),
                        color=rdef.get("color", WHITE))
        if i < n - 1:
            sep_y = BOX_TOP + (i + 1) * stat_block_h - Inches(0.05)
            horiz_line(slide, stats_left + Inches(0.15), sep_y,
                       stats_w - Inches(0.3), color=RGBColor(0x55, 0x55, 0x55))

def build_right_column(slide, cuali, verbatim_text, verbatim_attr, nota=None):
    right_left = col_left(2)
    right_w = COL_W
    cuali_h = Inches(1.5)
    verb_h = BOX_H - cuali_h - Inches(0.15)
    tb_c = slide.shapes.add_textbox(right_left, BOX_TOP, right_w, cuali_h)
    tb_c.fill.solid()
    tb_c.fill.fore_color.rgb = DARK_GRAY
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    pc = tf_c.paragraphs[0]
    pc.alignment = PP_ALIGN.LEFT
    rc = pc.add_run()
    rc.text = cuali
    rc.font.name = "Poppins"
    rc.font.size = Pt(10)
    rc.font.italic = True
    rc.font.color.rgb = WHITE
    verb_top = BOX_TOP + cuali_h + Inches(0.15)
    tb_v = slide.shapes.add_textbox(right_left, verb_top, right_w, verb_h)
    tb_v.fill.solid()
    tb_v.fill.fore_color.rgb = BLACK
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    pv = tf_v.paragraphs[0]
    pv.alignment = PP_ALIGN.LEFT
    rv = pv.add_run()
    rv.text = LDQ + verbatim_text + RDQ
    rv.font.name = "Instrument Serif"
    rv.font.size = Pt(9.5)
    rv.font.italic = True
    rv.font.color.rgb = WHITE
    pa = tf_v.add_paragraph()
    pa.alignment = PP_ALIGN.LEFT
    pa.space_before = Pt(6)
    ra = pa.add_run()
    ra.text = "— " + verbatim_attr
    ra.font.name = "Poppins"
    ra.font.size = Pt(8.5)
    ra.font.color.rgb = MID_GRAY
    if nota:
        pn = tf_v.add_paragraph()
        pn.alignment = PP_ALIGN.LEFT
        pn.space_before = Pt(5)
        rn = pn.add_run()
        rn.text = nota
        rn.font.name = "Poppins"
        rn.font.size = Pt(7.5)
        rn.font.italic = True
        rn.font.color.rgb = MID_GRAY

def add_slide_number(slide, num):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(0.7), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "0" + str(num)
    r.font.name = "Poppins"
    r.font.size = Pt(9)
    r.font.color.rgb = LIME

def build_cover(prs):
    slide = black_slide(prs)
    add_logo(slide, LOGO_PATH, right_margin=Inches(0.5), top_margin=Inches(0.3), height=Inches(0.5))
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "CODIGO CASA - NINJA"
    r.font.name = "Poppins"
    r.font.size = Pt(10)
    r.font.color.rgb = LIME
    add_headline_box(slide, "PILAR BIENESTAR: ", "PRUEBA DE HUMO",
                     Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.0))
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "5 hallazgos - Estudio cuantitativo y cualitativo 2025"
    r2.font.name = "Poppins"
    r2.font.size = Pt(12)
    r2.font.italic = True
    r2.font.color.rgb = MID_GRAY
    tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Mayo 2026"
    r3.font.name = "Poppins"
    r3.font.size = Pt(10)
    r3.font.color.rgb = MID_GRAY

def build_hallazgo_slide(prs, slide_num, headline_plain, headline_italic,
                          stats, cuali, verbatim_text, verbatim_attr,
                          source, nota=None, flag_text=None):
    slide = black_slide(prs)
    add_logo(slide, LOGO_PATH, right_margin=Inches(0.4), top_margin=Inches(0.22), height=Inches(0.38))
    add_slide_number(slide, slide_num)
    add_headline_box(slide, headline_plain, headline_italic,
                     MARGIN_L, HL_TOP, CONTENT_W, HL_H)
    horiz_line(slide, MARGIN_L, BOX_TOP - Inches(0.12), CONTENT_W)
    build_stats_block(slide, stats)
    build_right_column(slide, cuali, verbatim_text, verbatim_attr, nota)
    add_source_line(slide, source)
    if flag_text:
        tb_f = slide.shapes.add_textbox(MARGIN_L, SOURCE_TOP - Inches(0.32), CONTENT_W, Inches(0.3))
        pf = tb_f.text_frame.paragraphs[0]
        pf.alignment = PP_ALIGN.LEFT
        rf = pf.add_run()
        rf.text = flag_text
        rf.font.name = "Poppins"
        rf.font.size = Pt(8)
        rf.font.bold = True
        rf.font.color.rgb = RED_FLAG

def build_h04_slide(prs):
    slide = black_slide(prs)
    add_logo(slide, LOGO_PATH, right_margin=Inches(0.4), top_margin=Inches(0.22), height=Inches(0.38))
    add_slide_number(slide, 4)
    add_headline_box(slide,
                     "EL SISTEMA DE SALUD PUBLICA FALLA EN TODO. ",
                     "LA QUEJA NO VARIA POR NIVEL SOCIAL, TIPO DE SEGURO NI PUNTO DE ENTRADA.",
                     MARGIN_L, HL_TOP, CONTENT_W, HL_H)
    horiz_line(slide, MARGIN_L, BOX_TOP - Inches(0.12), CONTENT_W)
    stats_w = COL_W * 2 + COL_GAP
    stats_left = col_left(0)
    tb_conv = slide.shapes.add_textbox(stats_left, BOX_TOP, stats_w, BOX_H)
    tb_conv.fill.solid()
    tb_conv.fill.fore_color.rgb = DARK_GRAY
    tf_conv = tb_conv.text_frame
    tf_conv.word_wrap = True
    ph = tf_conv.paragraphs[0]
    ph.alignment = PP_ALIGN.LEFT
    rh = ph.add_run()
    rh.text = "CONVERGENCIA CUALITATIVA"
    rh.font.name = "Poppins"
    rh.font.size = Pt(8)
    rh.font.bold = True
    rh.font.color.rgb = LIME
    perfiles = [
        (
            "Trabajadora del sistema publico",
            ("Pesimo. Yo trabajo en un hospital, pesimo. Todo esta mal: "
             "no hay cama, no hay insumos, los hospitales no les pagan "
             "a los servidores. Todo esta mal."),
            "Grupo Biparental Hijos Pequenos"
        ),
        (
            "Usuaria de seguro privado",
            ("Depende del seguro, y el trato se une: si tu no tienes seguro, "
             "tienes que hacer un avance de 30 mil pesos para atras antes de "
             "que ni siquiera te atiendan en una clinica privada."),
            "Grupo Homoparental"
        ),
        (
            "Sobreviviente de cancer",
            ("Aqui el sistema de salud es horrible, no importa si tienes "
             "seguro medico o no. Aqui el medico que le interesaba el "
             "paciente se murio. No existe."),
            "Grupo Mixto"
        ),
    ]
    for label, quote_text, grupo in perfiles:
        pl = tf_conv.add_paragraph()
        pl.alignment = PP_ALIGN.LEFT
        pl.space_before = Pt(10)
        rl = pl.add_run()
        rl.text = label.upper()
        rl.font.name = "Poppins"
        rl.font.size = Pt(8)
        rl.font.bold = True
        rl.font.color.rgb = WHITE
        pq = tf_conv.add_paragraph()
        pq.alignment = PP_ALIGN.LEFT
        rq = pq.add_run()
        rq.text = LDQ + quote_text + RDQ
        rq.font.name = "Instrument Serif"
        rq.font.size = Pt(8.5)
        rq.font.italic = True
        rq.font.color.rgb = WHITE
        pa_attr = tf_conv.add_paragraph()
        pa_attr.alignment = PP_ALIGN.LEFT
        ra_attr = pa_attr.add_run()
        ra_attr.text = "— " + grupo
        ra_attr.font.name = "Poppins"
        ra_attr.font.size = Pt(7.5)
        ra_attr.font.color.rgb = MID_GRAY
    right_left = col_left(2)
    right_w = COL_W
    tb_cuali = slide.shapes.add_textbox(right_left, BOX_TOP, right_w, Inches(2.5))
    tb_cuali.fill.solid()
    tb_cuali.fill.fore_color.rgb = DARK_GRAY
    tf_cuali = tb_cuali.text_frame
    tf_cuali.word_wrap = True
    pc = tf_cuali.paragraphs[0]
    pc.alignment = PP_ALIGN.LEFT
    rc = pc.add_run()
    rc.text = ("La salud publica falla desde adentro y desde afuera: "
               "lo dice quien trabaja en el sistema, quien paga seguro privado, "
               "y quien sobrevivio un proceso de cancer dentro de el.")
    rc.font.name = "Poppins"
    rc.font.size = Pt(10)
    rc.font.italic = True
    rc.font.color.rgb = WHITE
    tb_flag = slide.shapes.add_textbox(MARGIN_L, SOURCE_TOP - Inches(0.32), CONTENT_W, Inches(0.3))
    pf = tb_flag.text_frame.paragraphs[0]
    pf.alignment = PP_ALIGN.LEFT
    rf = pf.add_run()
    rf.text = "! Evidencia cualitativa convergente - Datos cuantitativos pendientes de re-tabulacion (P28)"
    rf.font.name = "Poppins"
    rf.font.size = Pt(8)
    rf.font.bold = True
    rf.font.color.rgb = RED_FLAG
    add_source_line(slide,
                    "Source: Codigo Casa - Estudio cualitativo 2025 · P28 · "
                    "Evidencia cualitativa convergente - cifra cuantitativa pendiente de re-tabulacion.")

SLIDES_DATA = [
    dict(
        slide_num=1,
        headline_plain="EL DINERO ES LA PRESION DOMINANTE EN EL HOGAR DOMINICANO. ",
        headline_italic="47.8% LO NOMBRA PRIMERO, CASI EL DOBLE DE LO QUE PESA LA INSEGURIDAD.",
        stats=[
            dict(big="48%", unit="",
                 desc_runs=[
                     dict(text="nombra la economia como factor #1 de estres familiar"),
                     dict(text=" — por encima de inseguridad (21.8%) y crianza (16.2%).", bold=True),
                     dict(text=" · P26 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="28%", unit="",
                 desc_runs=[
                     dict(text="adicional menciona "),
                     dict(text="realidad economica", bold=True),
                     dict(text=" como segunda causa financiera. Sumados ambos campos financieros, ninguna otra causa se acerca."),
                     dict(text=" · P26 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="22%", unit="",
                 desc_runs=[
                     dict(text="identifica la "),
                     dict(text="inseguridad", bold=True),
                     dict(text=" como segunda fuente de estres — menos de la mitad del peso que carga la economia."),
                     dict(text=" · P26 · Base 500", color=MID_GRAY),
                 ]),
        ],
        cuali="El estres financiero no compite con las demas presiones del hogar: las duplica.",
        verbatim_text=("Porque la mayor parte del problema hoy en dia es el dinero, "
                       "todo lo resuelve el dinero aunque la gente no lo quiera admitir, "
                       "eso es, por ahi es que la gente se jode en su salud mental, el dinero."),
        verbatim_attr="Grupo Monoparental",
        source="Source: Codigo Casa - Estudio cuantitativo 2025 · P26 · Base 500.",
    ),
    dict(
        slide_num=2,
        headline_plain="4 DE CADA 10 DOMINICANOS NO HACE NADA POR SU SALUD MENTAL NI FISICA. ",
        headline_italic="LA INACCION ES LA RESPUESTA MAS FRECUENTE DEL PAIS.",
        stats=[
            dict(big="43%", unit="",
                 desc_runs=[
                     dict(text="declara que "),
                     dict(text="no realiza ninguna actividad", bold=True),
                     dict(text=(" para cuidar su salud mental y fisica — primera opcion del ranking, "
                                "por encima de ejercitarse (34.8%).")),
                     dict(text=" · P27 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="3%", unit="",
                 desc_runs=[
                     dict(text="reporta ir a "),
                     dict(text="terapia.", bold=True),
                     dict(text=(" 9% practica meditacion. "
                                "El apoyo profesional de salud mental es marginal en el pais.")),
                     dict(text=" · P27 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="35%", unit="",
                 desc_runs=[
                     dict(text="se "),
                     dict(text="ejercita", bold=True),
                     dict(text=(" con regularidad — segunda opcion del ranking, "
                                "pero a 8 puntos de la inaccion.")),
                     dict(text=" · P27 · Base 500", color=MID_GRAY),
                 ]),
        ],
        cuali="El primer lugar del ranking no es una practica de cuidado — es su ausencia.",
        verbatim_text=("Yo para lo fisico no estoy haciendo nada, no hago ejercicio, "
                       "me como todo lo que quiera y a la hora que quiera. "
                       "Bueno, para la salud mental, a veces uno tiene que hacerse loco "
                       "en ciertas cosas, porque si uno le da mucha mente a la cosa, "
                       "uno se va a volver loco."),
        verbatim_attr="Grupo Biparental Hijos Adultos",
        source="Source: Codigo Casa - Estudio cuantitativo 2025 · P27 · Base 500.",
    ),
    dict(
        slide_num=3,
        headline_plain=("CASI LA MITAD DE LAS MUJERES NO HACE NADA POR SU SALUD. "
                        "LOS HOMBRES VAN AL GIMNASIO. "),
        headline_italic="CUIDAN COSAS DISTINTAS — O NO CUIDAN LO MISMO.",
        stats=[
            dict(big="49%", unit="mujeres",
                 desc_runs=[
                     dict(text="declara que "),
                     dict(text="no realiza ninguna actividad", bold=True),
                     dict(text=(" para cuidar su salud. "
                                "Respuesta mas frecuente del subgrupo femenino.")),
                     dict(text=" · P27 x D2 · Subset fem. n=305", color=MID_GRAY),
                 ]),
            dict(big="43%", unit="hombres",
                 desc_runs=[
                     dict(text="reporta "),
                     dict(text="ejercitarse con regularidad", bold=True),
                     dict(text=(" como principal practica de autocuidado. "
                                "Primer lugar del subgrupo masculino.")),
                     dict(text=" · P27 x D2 · Subset masc. n=195", color=MID_GRAY),
                 ]),
            dict(big="6 pt", unit="",
                 desc_runs=[
                     dict(text="de diferencia entre la "),
                     dict(text="inaccion femenina (48.5%)", bold=True),
                     dict(text=(" y la masculina (42.5%). "
                                "La brecha existe, y el tipo de practica diverge aun mas.")),
                     dict(text=" · P27 x D2", color=MID_GRAY),
                 ]),
        ],
        cuali="Los hombres mueven el cuerpo. Las mujeres, en mayoria, no tienen espacio ni para eso.",
        verbatim_text=("Bueno, la ansiedad, como te dije, a mi me da mucho, "
                       "el tema de la ansiedad, y aumente muchisimo de peso. "
                       "Me hizo una barriada, pero yo me frustro porque hice tantas cosas, "
                       "gaste tantos cuartos, estoy volviendo a guardar porque "
                       "sigo mi vida igualita de desorden."),
        verbatim_attr="Grupo Biparental Hijos Pequenos (mujer)",
        source="Source: Codigo Casa - Estudio cuantitativo 2025 · P27 x D2 · Base 500.",
        nota=('* "barriada" = impacto fisico severo. '
              '"cuartos" = dinero. '
              '"volviendo a guardar" = recuperando peso.'),
    ),
    dict(
        slide_num=5,
        headline_plain="SOLO 11.6% RECONOCE QUE ALGUN HIJO SUYO VIVIO BULLYING. ",
        headline_italic=("EN LOS GRUPOS, EL BULLYING APARECE COMO UN MIEDO "
                         "PARENTAL DE FONDO PERMANENTE."),
        stats=[
            dict(big="12%", unit="",
                 desc_runs=[
                     dict(text="de la base reconoce que algun hijo ha vivido "),
                     dict(text="acoso o bullying", bold=True),
                     dict(text=", escolar o digital. Base efectiva n=58."),
                     dict(text=" · P75 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="75%", unit="",
                 desc_runs=[
                     dict(text="responde "),
                     dict(text="No", bold=True),
                     dict(text=(" — cifra que, a la luz del cualitativo, refleja subreporte "
                                "mas que ausencia real del fenomeno.")),
                     dict(text=" · P75 · Base 500", color=MID_GRAY),
                 ]),
            dict(big="13%", unit="",
                 desc_runs=[
                     dict(text="no tiene hijos (NS/NR). El "),
                     dict(text="universo parental efectivo", bold=True),
                     dict(text=(" hace que el 11.6% represente "
                                "una base de 58 respondientes.")),
                     dict(text=" · P75 · Base 500", color=MID_GRAY),
                 ]),
        ],
        cuali=("La cifra de reconocimiento es baja. "
               "El cualitativo muestra que el bullying no es un evento puntual "
               "sino una alerta que los padres ya llevan encendida."),
        verbatim_text=("Bueno, tuve al principio del ano escolar un momento de tension "
                       "porque le estaban haciendo bullying a las ninas, y las ninas, "
                       "que les gustaba su colegio, ya no querian ir. "
                       "Eso fue una senal de alerta para mi."),
        verbatim_attr="Grupo Extendido",
        source=("Source: Codigo Casa - Estudio cuantitativo 2025 · P75 · Base 500. "
                "Cruce NSE AB (n=23) excluido por base debil; disponible en nota tecnica."),
    ),
]

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    build_cover(prs)
    for data in SLIDES_DATA[:3]:
        build_hallazgo_slide(prs, **data)
    build_h04_slide(prs)
    build_hallazgo_slide(prs, **SLIDES_DATA[3])
    prs.save(OUT_PPTX)
    print("SAVED:", OUT_PPTX)

if __name__ == "__main__":
    main()
