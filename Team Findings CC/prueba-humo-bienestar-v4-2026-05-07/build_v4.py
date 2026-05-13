#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_v4.py
============
Genera prueba-humo-bienestar-v4-2026-05-07.pptx desde cero.

Estructura nueva (mayo 2026 - CSS Jeremy):
  - 4 hallazgos cuanti x 2 slides (Slide A: Hallazgo + Slide B: Consumer Voice)
  - 1 hallazgo cualitativo x 1 slide (Cards cualitativas)
  Total: 9 slides

Reglas no negociables:
  - NO lineas separadoras horizontales.
  - Cajas de stats: 7.5 cm x 3.5 cm exactos.
  - Cards cualitativas: rounded edges, fill negro 45% opacidad.
  - Verbatims SIEMPRE en blanco.
  - NO redondear ninguna cifra.
  - NO insight italic debajo de stats.
  - Kerning 0 en todo el deck.
  - Headlines MAYUSCULAS, Instrument Serif, auto-tamano 30-70pt.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
OUT_DIR  = Path("/Users/jeremyrodriguez/Documents/Cerebro/Codigo Casa/Team Findings CC/prueba-humo-bienestar-v4-2026-05-07")
OUT_DIR2 = Path("/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v4-2026-05-07")

# Intentar la ruta con tilde primero, luego sin tilde
import os
if os.path.exists("/Users/jeremyrodriguez/Documents/Cerebro/Código Casa"):
    OUT_DIR = Path("/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v4-2026-05-07")
else:
    OUT_DIR = Path("/Users/jeremyrodriguez/Documents/Cerebro/Codigo Casa/Team Findings CC/prueba-humo-bienestar-v4-2026-05-07")

PPTX_OUT = OUT_DIR / "bienestar-humo-v4-2026-05-07.pptx"

# ---------------------------------------------------------------------------
# PALETA
# ---------------------------------------------------------------------------
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x2E, 0x2E, 0x2E)
LIME  = RGBColor(0xB8, 0xFF, 0x4D)
NAVY  = RGBColor(0x1A, 0x1A, 0x2E)
SILVER = RGBColor(0xAA, 0xAA, 0xAA)
DIM   = RGBColor(0x77, 0x77, 0x77)
DIMMER = RGBColor(0x55, 0x55, 0x55)
FAINT  = RGBColor(0x88, 0x88, 0x88)

# Comillas curvas como constantes unicode
LQUOTE = "“"  # "
RQUOTE = "”"  # "

# ---------------------------------------------------------------------------
# DIMENSIONES DE SLIDE (widescreen 33.87 cm x 19.05 cm)
# ---------------------------------------------------------------------------
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)


# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def new_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background
    fill   = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def set_run(run, text, font_name="Poppins", size_pt=11,
            bold=False, italic=False, color=None):
    run.text       = text
    run.font.name  = font_name
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "0")


def add_textbox(slide, x, y, w, h, text, font_name="Poppins", font_size_pt=11,
                bold=False, italic=False, color=None,
                align=PP_ALIGN.LEFT, word_wrap=True):
    if color is None:
        color = WHITE
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap      = word_wrap
    tf.margin_left    = 0
    tf.margin_right   = 0
    tf.margin_top     = 0
    tf.margin_bottom  = 0
    para              = tf.paragraphs[0]
    para.alignment    = align
    run               = para.add_run()
    set_run(run, text, font_name=font_name, size_pt=font_size_pt,
            bold=bold, italic=italic, color=color)
    return txb


def headline_size(text):
    n = len(text.strip())
    if   n < 30: return 70
    elif n < 45: return 55
    elif n < 66: return 42
    elif n < 86: return 36
    else:        return 30


def add_headline_split(slide, plain, italic_part, y_cm, h_cm):
    full = plain + italic_part
    size = headline_size(full)
    margin = Cm(1.0)
    txb = slide.shapes.add_textbox(margin, Cm(y_cm),
                                   SLIDE_W - 2 * margin, Cm(h_cm))
    tf  = txb.text_frame
    tf.word_wrap     = True
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    if plain:
        r = para.add_run()
        set_run(r, plain.upper(), font_name="Instrument Serif",
                size_pt=size, italic=False, color=WHITE)
    if italic_part:
        r = para.add_run()
        set_run(r, italic_part.upper(), font_name="Instrument Serif",
                size_pt=size, italic=True, color=WHITE)
    return txb


def add_stat_box(slide, x_cm, y_cm, stat_value, desc_runs, ref_text):
    """
    Caja de stat: 7.5 cm x 3.5 cm (CSS Jeremy).
    NO lineas separadoras.
    """
    BOX_W = Cm(7.5)
    x = Cm(x_cm)
    y = Cm(y_cm)

    n = len(stat_value.strip())
    if   n <= 3: stat_pt = 72
    elif n <= 5: stat_pt = 60
    else:        stat_pt = 48

    # Cifra grande
    stat_h   = Cm(1.8)
    stat_txb = slide.shapes.add_textbox(x, y, BOX_W, stat_h)
    stat_tf  = stat_txb.text_frame
    stat_tf.word_wrap     = False
    stat_tf.margin_left   = 0
    stat_tf.margin_right  = 0
    stat_tf.margin_top    = 0
    stat_tf.margin_bottom = 0
    sp = stat_tf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    r = sp.add_run()
    set_run(r, stat_value, font_name="Instrument Serif",
            size_pt=stat_pt, italic=True, color=WHITE)

    # Descripcion multi-run
    desc_y   = y + stat_h + Cm(0.1)
    desc_h   = Cm(1.3)
    desc_txb = slide.shapes.add_textbox(x, desc_y, BOX_W, desc_h)
    desc_tf  = desc_txb.text_frame
    desc_tf.word_wrap     = True
    desc_tf.margin_left   = 0
    desc_tf.margin_right  = 0
    desc_tf.margin_top    = 0
    desc_tf.margin_bottom = 0
    dp = desc_tf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    for spec in desc_runs:
        r = dp.add_run()
        set_run(r, spec["text"],
                font_name=spec.get("font_name", "Poppins"),
                size_pt=spec.get("font_size_pt", 11),
                bold=spec.get("bold", False),
                italic=spec.get("italic", False),
                color=spec.get("color", WHITE))

    # Referencia tecnica
    ref_y   = desc_y + desc_h
    ref_h   = Cm(0.6)
    ref_txb = slide.shapes.add_textbox(x, ref_y, BOX_W, ref_h)
    ref_tf  = ref_txb.text_frame
    ref_tf.word_wrap     = True
    ref_tf.margin_left   = 0
    ref_tf.margin_right  = 0
    ref_tf.margin_top    = 0
    ref_tf.margin_bottom = 0
    rp = ref_tf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    r = rp.add_run()
    set_run(r, ref_text, font_name="Poppins",
            size_pt=9, italic=True, color=SILVER)
    return stat_txb


def add_source(slide, source_text, y_cm=17.6):
    add_textbox(slide, x=Cm(1.0), y=Cm(y_cm),
                w=SLIDE_W - Cm(2.0), h=Cm(0.7),
                text=source_text, font_name="Poppins", font_size_pt=9,
                italic=True, color=SILVER, align=PP_ALIGN.LEFT)


def add_pilar_label(slide, pilar="BIENESTAR"):
    add_textbox(slide, x=Cm(1.0), y=Cm(0.4),
                w=Cm(10), h=Cm(0.6),
                text="CÓDIGO CASA — " + pilar,
                font_name="Poppins", font_size_pt=8,
                color=DIM, align=PP_ALIGN.LEFT)


def add_slide_label(slide, label_text):
    add_textbox(slide, x=Cm(31.5), y=Cm(0.4),
                w=Cm(2.0), h=Cm(0.6),
                text=label_text, font_name="Poppins", font_size_pt=8,
                color=DIMMER, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------

def build_slide_hallazgo(prs, headline_plain, headline_italic,
                         stats, source, slide_label=""):
    """
    Slide A: Hallazgo cuanti.
    SIN verbatim. SIN conclusion italic. SIN lineas separadoras.
    Cajas de stats 7.5 cm x 3.5 cm.
    """
    slide = blank_slide(prs)
    add_pilar_label(slide)
    if slide_label:
        add_slide_label(slide, slide_label)

    add_headline_split(slide, headline_plain, headline_italic,
                       y_cm=1.2, h_cm=5.0)

    n       = len(stats)
    box_w   = 7.5
    gap_cm  = (33.87 - 2.0 - n * box_w) / (n + 1)
    y_stats = 7.5

    for i, st in enumerate(stats):
        x_cm = 1.0 + gap_cm + i * (box_w + gap_cm)
        add_stat_box(slide, x_cm=x_cm, y_cm=y_stats,
                     stat_value=st["value"],
                     desc_runs=st["desc_runs"],
                     ref_text=st["ref"])

    add_source(slide, source)


def build_slide_consumer_voice(prs, verbatim, atribucion, slide_label=""):
    """
    Slide B: Consumer Voice.
    Header CONSUMER VOICE arriba. Quote centralizada. Atribucion abajo.
    Verbatim en blanco. NUNCA verde lima.
    """
    slide = blank_slide(prs)
    add_pilar_label(slide)
    if slide_label:
        add_slide_label(slide, slide_label)

    add_textbox(slide, x=Cm(1.0), y=Cm(1.2),
                w=SLIDE_W - Cm(2.0), h=Cm(0.8),
                text="CONSUMER VOICE",
                font_name="Poppins", font_size_pt=10,
                color=SILVER, align=PP_ALIGN.CENTER)

    quote_text = LQUOTE + verbatim + RQUOTE
    vlen = len(verbatim)
    if   vlen < 80:  v_pt = 22
    elif vlen < 160: v_pt = 18
    elif vlen < 240: v_pt = 15
    else:            v_pt = 13

    add_textbox(slide, x=Cm(2.5), y=Cm(3.5),
                w=SLIDE_W - Cm(5.0), h=Cm(11.0),
                text=quote_text,
                font_name="Instrument Serif", font_size_pt=v_pt,
                color=WHITE, align=PP_ALIGN.CENTER, word_wrap=True)

    add_textbox(slide, x=Cm(2.5), y=Cm(15.5),
                w=SLIDE_W - Cm(5.0), h=Cm(0.8),
                text=atribucion,
                font_name="Poppins", font_size_pt=11,
                color=SILVER, align=PP_ALIGN.CENTER)


def build_slide_cards_cualitativas(prs, headline_plain, headline_italic,
                                   cards, caveat, source, slide_label=""):
    """
    Slide unico de cards cualitativas.
    3 cards con rounded edges, fill negro 45% opacidad.
    Verbatims en blanco.
    """
    slide = blank_slide(prs)
    add_pilar_label(slide)
    if slide_label:
        add_slide_label(slide, slide_label)

    add_headline_split(slide, headline_plain, headline_italic,
                       y_cm=1.2, h_cm=4.0)

    n_cards    = len(cards)
    card_w_emu = Cm(9.5)
    card_h_emu = Cm(10.0)
    gap_emu    = (SLIDE_W - Cm(2.0) - n_cards * card_w_emu) / (n_cards + 1)
    y_card     = Cm(5.5)

    for i, card in enumerate(cards):
        x_card = Cm(1.0) + gap_emu + i * (card_w_emu + gap_emu)

        # Shape con rounded corners y fill negro 45% opacidad
        sp = slide.shapes.add_shape(
            1,  # Rectangle
            x_card, y_card, card_w_emu, card_h_emu
        )

        # Rounded corners via XML
        sp_el   = sp._element
        spPr_el = sp_el.find(qn("p:spPr"))
        if spPr_el is not None:
            prstGeom = spPr_el.find(qn("a:prstGeom"))
            if prstGeom is not None:
                prstGeom.set("prst", "roundRect")
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, qn("a:avLst"))
                gd = etree.SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", "val 8000")

        # Fill negro solido
        fill = sp.fill
        fill.solid()
        fill.fore_color.rgb = BLACK

        # Aplicar alpha 45% via XML (45000 de 100000)
        solidFill_el = spPr_el.find(qn("a:solidFill"))
        if solidFill_el is not None:
            srgbClr = solidFill_el.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha_el.set("val", "45000")

        # Borde sutil gris
        sp.line.color.rgb = GRAY
        sp.line.width     = Pt(0.5)

        # Quote dentro de la card
        quote_text = LQUOTE + card["quote"] + RQUOTE
        q_len = len(card["quote"])
        q_pt  = 11 if q_len < 200 else 10

        add_textbox(slide,
                    x=x_card + Cm(0.35),
                    y=y_card  + Cm(0.5),
                    w=card_w_emu - Cm(0.7),
                    h=card_h_emu - Cm(1.6),
                    text=quote_text,
                    font_name="Poppins", font_size_pt=q_pt,
                    color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True)

        # Atribucion dentro de la card
        add_textbox(slide,
                    x=x_card + Cm(0.35),
                    y=y_card  + card_h_emu - Cm(1.2),
                    w=card_w_emu - Cm(0.7),
                    h=Cm(1.0),
                    text=card["atribucion"],
                    font_name="Poppins", font_size_pt=9,
                    italic=True, color=SILVER, align=PP_ALIGN.LEFT)

    if caveat:
        add_textbox(slide, x=Cm(1.0), y=Cm(16.2),
                    w=SLIDE_W - Cm(2.0), h=Cm(0.7),
                    text=caveat, font_name="Poppins", font_size_pt=8,
                    italic=True, color=FAINT, align=PP_ALIGN.LEFT)

    add_source(slide, source)


# ---------------------------------------------------------------------------
# DATA v4 -- BIENESTAR
# Comillas curvas usando constantes LQUOTE / RQUOTE para evitar conflicto
# con delimitadores de string Python.
# ---------------------------------------------------------------------------

# Alias cortos para construir cadenas con comillas tipograficas
LQ = LQUOTE
RQ = RQUOTE

HALLAZGOS = [
    # ------------------------------------------------------------------
    # H01
    # ------------------------------------------------------------------
    {
        "type":     "cuanti",
        "label":    "H01",
        "headline_plain":  ("EL ESTRÉS DEL HOGAR DOMINICANO TIENE DIRECCIÓN:"
                            " CASI LA MITAD DEL PAÍS SEÑALA LA ECONOMÍA "),
        "headline_italic": ("ANTES QUE LA VIOLENCIA, LA CRIANZA O EL TIEMPO."),
        "stats": [
            {
                "value": "47.8%",
                "desc_runs": [
                    {"text": "nombra "},
                    {"text": LQ + "Economía" + RQ,
                     "bold": True},
                    {"text": " como factor #1 de estrés en la vida familiar."},
                ],
                "ref": "P26 multi-select · Base 500",
            },
            {
                "value": "28.4%",
                "desc_runs": [
                    {"text": LQ + "Realidad económica" + RQ + " suma "},
                    {"text": "28.4% adicional;", "bold": True},
                    {"text": (" " + LQ + "Inseguridad" + RQ
                              + " llega a 21.8% y "
                              + LQ + "Crianza" + RQ + " a 16.2%.")},
                ],
                "ref": "P26 multi-select · Base 500",
            },
        ],
        "source":     ("Source: Código Casa — Estudio cuantitativo 2025"
                       " · P26 · Base 500."),
        "verbatim":   ("Como lo dijo una compañera acá ahorita: el dinero."
                       " Tú no podés resolver tal vez algo que se te presente."
                       " Eso es frustra, eso no es fácil."),
        "atribucion": "— Grupo Monoparental",
    },

    # ------------------------------------------------------------------
    # H02
    # ------------------------------------------------------------------
    {
        "type":     "cuanti",
        "label":    "H02",
        "headline_plain":  ("LO ECONÓMICO NO COMPITE CON LOS DEMÁS"
                            " ESTRESORES DEL HOGAR DOMINICANO. "),
        "headline_italic": "LOS SUPERA EN CONJUNTO.",
        "stats": [
            {
                "value": "76.2",
                "desc_runs": [
                    {"text": "puntos de peso relativo suman "},
                    {"text": (LQ + "Economía" + RQ
                               + " + " + LQ + "Realidad económica" + RQ),
                     "bold": True},
                    {"text": (" (47.8% + 28.4%) vs 74.8 pts distribuidos entre"
                              " los cuatro estresores restantes.")},
                ],
                "ref": ("P26 multi-select · Base 500"
                        " · Nota: peso relativo, no % de personas"),
            },
        ],
        "source":     ("Source: Código Casa — Estudio cuantitativo 2025"
                       " · P26 · Base 500."),
        "verbatim":   ("No tengo paciencia cuando no tengo dinero. Cuando no tengo"
                       " dinero, que no sé de qué voy a hacer con los"
                       " compromisos y todo eso, entonces yo como que pienso,"
                       " le doy vuelta a la cosa, un estrés."),
        "atribucion": "— Grupo Monoparental",
    },

    # ------------------------------------------------------------------
    # H03
    # ------------------------------------------------------------------
    {
        "type":     "cuanti",
        "label":    "H03",
        "headline_plain":  "EL AUTOCUIDADO EXISTE COMO ASPIRACIÓN. ",
        "headline_italic": "LA PRIMERA RESPUESTA DEL PAÍS ES NO HACER NADA.",
        "stats": [
            {
                "value": "42.6%",
                "desc_runs": [
                    {"text": "responde "},
                    {"text": (LQ + "No realizo ninguna actividad para cuidar"
                               " mi salud mental y física" + RQ),
                     "bold": True},
                    {"text": " — la opción más mencionada."},
                ],
                "ref": "P27 single-select · Base 500",
            },
            {
                "value": "34.8%",
                "desc_runs": [
                    {"text": "El ejercicio regular llega a "},
                    {"text": "34.8%, segunda respuesta;", "bold": True},
                    {"text": (" la terapia queda en 3.0%"
                               " y la meditación en 9.0%.")},
                ],
                "ref": "P27 single-select · Base 500",
            },
        ],
        "source":     ("Source: Código Casa — Estudio cuantitativo 2025"
                       " · P27 · Base 500."),
        "verbatim":   ("Estar un poquito más estable en todos los sentidos."
                       " He estado con mucho estrés. Sí. Para entonces"
                       " poder ir a la gimnasia y comprar alimentos que debo de comer."),
        "atribucion": "— Grupo Monoparental",
    },

    # ------------------------------------------------------------------
    # H04
    # ------------------------------------------------------------------
    {
        "type":     "cuanti",
        "label":    "H04",
        "headline_plain":  ("EL AUTOCUIDADO TIENE GÉNERO EN RD:"
                            " ÉL SE EJERCITA, "),
        "headline_italic": "ELLA NO HACE NADA. Y A MENOR NSE, MÁS INACCIÓN.",
        "stats": [
            {
                "value": "48.5%",
                "desc_runs": [
                    {"text": "Entre las mujeres, la respuesta más frecuente es "},
                    {"text": (LQ + "No realizo ninguna actividad" + RQ
                               + " (48.5%)."),
                     "bold": True},
                    {"text": (" Entre los hombres, es "
                               + LQ + "Me ejercito con regularidad" + RQ
                               + " (43.1%).")},
                ],
                "ref": "P27 × D2 sexo · Base 500",
            },
            {
                "value": "60.0%",
                "desc_runs": [
                    {"text": "La inacción se concentra en "},
                    {"text": ("estrato D (51.2%) y en el grupo"
                               " 18–24 años (60.0%)."),
                     "bold": True},
                    {"text": (" En NSE AB y C, la respuesta dominante es"
                               " el ejercicio regular (39.1% y 38.5%).")},
                ],
                "ref": "P27 × D7 NSE × D3 edad · Base 500",
            },
        ],
        "source":     ("Source: Código Casa — Estudio cuantitativo 2025"
                       " · P27, D2, D7, D3 · Base 500."),
        "verbatim":   ("Mi esposo va al gimnasio, yo camino diario menos los"
                       " domingos hago pilates, tenemos una… tratamos de"
                       " tener una buena."),
        "atribucion": "— Grupo Biparental Hijos Pequeños",
    },

    # ------------------------------------------------------------------
    # H05 -- cualitativo
    # ------------------------------------------------------------------
    {
        "type":     "cualitativo",
        "label":    "H05",
        "headline_plain":  ("EL DIAGNÓSTICO DEL SISTEMA DE SALUD PÚBLICA"
                            " NO LLEGA EN LISTA. "),
        "headline_italic": ("LLEGA COMO QUEJA TOTAL: CALIDAD, COSTO Y ACCESO"
                            " SE REFUERZAN ENTRE SÍ."),
        "cards": [
            {
                "quote":      ("El sistema de salud pública no sirve. Es un caos."
                               " En mi caso no me ha tocado, gracias a Dios; siempre"
                               " tuve un empleo que me ha tenido un pago seguro y puedo"
                               " ir a cualquier clínica tranquilamente. Pero la persona"
                               " que no tiene un seguro médico hace mucho trabajo,"
                               " mucho trabajo."),
                "atribucion": "— Grupo Extendido",
            },
            {
                "quote":      ("De farmacia… si esto te cuesta 500 pesos y te cura,"
                               " te ponen [medicamento] que vale 5 mil porque se alió"
                               " a tu farmacia, a tu laboratorio."),
                "atribucion": "— Grupo Biparental Hijos Adultos",
            },
            {
                "quote":      ("Yo fui a visitar a alguien ahí en el Marcelo Rivas…"
                               " eso parece un hotel y las atenciones son enormes, lo que"
                               " pasa es que la demanda de personas es mucha."),
                "atribucion": "— Grupo Extendido",
            },
        ],
        "caveat": ("Convergencia cualitativa — datos cuantitativos pendientes"
                   " de re-tabulación (P28)."),
        "source": ("Source: Código Casa — Estudio cualitativo 2025"
                   " · P28 · Evidencia cualitativa convergente."),
    },
]


# ---------------------------------------------------------------------------
# BUILD
# ---------------------------------------------------------------------------

def build():
    prs         = new_presentation()
    slide_count = 0

    for h in HALLAZGOS:
        label = h["label"]

        if h["type"] == "cuanti":
            slide_count += 1
            build_slide_hallazgo(
                prs,
                headline_plain  = h["headline_plain"],
                headline_italic = h["headline_italic"],
                stats           = h["stats"],
                source          = h["source"],
                slide_label     = label + "A",
            )
            print("  OK Slide %d -- %s Hallazgo cuanti" % (slide_count, label))

            slide_count += 1
            build_slide_consumer_voice(
                prs,
                verbatim    = h["verbatim"],
                atribucion  = h["atribucion"],
                slide_label = label + "B",
            )
            print("  OK Slide %d -- %s Consumer Voice" % (slide_count, label))

        elif h["type"] == "cualitativo":
            slide_count += 1
            build_slide_cards_cualitativas(
                prs,
                headline_plain  = h["headline_plain"],
                headline_italic = h["headline_italic"],
                cards           = h["cards"],
                caveat          = h.get("caveat", ""),
                source          = h["source"],
                slide_label     = label,
            )
            print("  OK Slide %d -- %s Cards cualitativas" % (slide_count, label))

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    print("\nGuardado: %s" % PPTX_OUT)
    print("Total slides: %d" % slide_count)
    return slide_count


if __name__ == "__main__":
    print("=== BUILD bienestar-humo-v4 ===\n")
    n = build()
    print("\n=== COMPLETADO: %d slides ===" % n)
