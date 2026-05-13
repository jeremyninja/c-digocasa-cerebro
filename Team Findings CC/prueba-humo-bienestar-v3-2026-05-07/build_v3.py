#!/usr/bin/env python3
"""
build_v3.py
============
Corrige las cifras grandes del deck Bienestar prueba-humo v2 → v3.

Cambios de regla: NO redondear ninguna cifra, incluyendo el stat grande.
Mantener el decimal exacto que entregó el cazador.

Cambios por slide:
  Slide 1 (H01): stat grande "48%" → "47.8%"
  Slide 2 (H02): sin cambio (76.2 y 74.8 ya están exactos)
  Slide 3 (H03): stat grande "43%" → "42.6%", "35%" → "34.8%"
  Slide 4 (H04): stat grande "60%" → "60.0%"  (48.5% ya estaba correcto)
  Slide 5 (H05): sin cifra cuanti grande, sin cambio

El texto de las descripciones inline ya tenía los decimales correctos en v2,
se mantiene sin modificar.
"""

import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

SRC = Path("/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v2-2026-05-07/bienestar-humo-v2-2026-05-07.pptx")
DST = Path("/Users/jeremyrodriguez/Documents/Cerebro/Código Casa/Team Findings CC/prueba-humo-bienestar-v3-2026-05-07/bienestar-humo-v3-2026-05-07.pptx")

# Mapa de correcciones: {(slide_idx_0based, shape_idx): (old_text, new_text)}
# Solo los stats grandes (las cifras prominentes)
CORRECTIONS = {
    # Slide 1 (H01): shape 3 = stat "48%" → "47.8%"
    (0, 3): ("48%", "47.8%"),
    # Slide 3 (H03): shape 3 = "43%" → "42.6%", shape 6 = "35%" → "34.8%"
    (2, 3): ("43%", "42.6%"),
    (2, 6): ("35%", "34.8%"),
    # Slide 4 (H04): shape 6 = "60%" → "60.0%"
    (3, 6): ("60%", "60.0%"),
}

def apply_corrections():
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)

    change_log = []

    for (slide_idx, shape_idx), (old_text, new_text) in CORRECTIONS.items():
        slide = prs.slides[slide_idx]
        shape = slide.shapes[shape_idx]

        if not shape.has_text_frame:
            print(f"WARNING: Slide {slide_idx+1}, Shape {shape_idx} no tiene text frame")
            continue

        actual_text = shape.text_frame.text.strip()
        if actual_text != old_text:
            print(f"WARNING: Slide {slide_idx+1}, Shape {shape_idx}: esperaba '{old_text}', encontré '{actual_text}' — revisando runs...")

        # Reemplazar en cada run del text frame
        replaced = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == old_text:
                    run.text = new_text
                    replaced = True
                    size_pt = round(run.font.size.pt, 1) if run.font.size else "inherit"
                    change_log.append({
                        "slide": slide_idx + 1,
                        "shape": shape_idx,
                        "old": old_text,
                        "new": new_text,
                        "font_size_pt": size_pt,
                        "font_size_adjusted": False,  # 72pt cabe para 5 chars en slot 5.76"
                    })
                    print(f"  OK Slide {slide_idx+1} Shape {shape_idx}: '{old_text}' → '{new_text}' @ {size_pt}pt")

        if not replaced:
            # Intentar reemplazar run por run incluyendo whitespace
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old_text in run.text:
                        size_pt = round(run.font.size.pt, 1) if run.font.size else "inherit"
                        run.text = run.text.replace(old_text, new_text)
                        replaced = True
                        change_log.append({
                            "slide": slide_idx + 1,
                            "shape": shape_idx,
                            "old": old_text,
                            "new": new_text,
                            "font_size_pt": size_pt,
                            "font_size_adjusted": False,
                        })
                        print(f"  OK (partial) Slide {slide_idx+1} Shape {shape_idx}: '{old_text}' → '{new_text}' @ {size_pt}pt")
                        break

            if not replaced:
                print(f"  ERROR: No se pudo reemplazar '{old_text}' en Slide {slide_idx+1}, Shape {shape_idx}")

    prs.save(DST)
    print(f"\nGuardado: {DST}")
    return change_log


def verify_result(change_log):
    prs = Presentation(DST)
    print("\n--- VERIFICACIÓN ---")
    for (slide_idx, shape_idx), (old_text, new_text) in CORRECTIONS.items():
        slide = prs.slides[slide_idx]
        shape = slide.shapes[shape_idx]
        actual = shape.text_frame.text.strip()
        status = "OK" if actual == new_text else f"FAIL (encontré: '{actual}')"
        print(f"  Slide {slide_idx+1} Shape {shape_idx}: '{new_text}' — {status}")

    print("\n--- RESUMEN DE CAMBIOS ---")
    for c in change_log:
        adj_note = " [tamaño de fuente bajado]" if c["font_size_adjusted"] else f" [tamaño mantenido a {c['font_size_pt']}pt — cabe en slot 5.76\"]"
        print(f"  Slide {c['slide']} H0{c['slide']}: '{c['old']}' → '{c['new']}'{adj_note}")


if __name__ == "__main__":
    print(f"Origen: {SRC}")
    print(f"Destino: {DST}")
    print()
    log = apply_corrections()
    verify_result(log)
